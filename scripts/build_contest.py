"""인프라그룹 AI경진대회 과제 출품 자료 빌더.

메인 양식 `templates/사내양식/260511_인프라그룹 AI경진대회 과제_양식.pptx`
의 Slide 1·2 placeholder 를 FA 기술혁신 Agent 팀 콘텐츠로 채우고,
회색조 첨부 슬라이드 3장(팀원+SOP, 데모 산출물 4종, ROI+다음단계)을 뒤에 추가한다.

출력: outputs/AI경진대회_FA기술혁신Agent팀_2026-05-13.pptx
"""
from __future__ import annotations

import glob
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ─── 양식 색상 (메인 슬라이드용, 양식과 동일) ─────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x66, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ─── 첨부 색상 (회색조, 단조) ────────────────────────────────────
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

TEMPLATE = glob.glob("templates/사내양식/260511*.pptx")
if not TEMPLATE:
    print("ERROR: template not found", file=sys.stderr)
    sys.exit(1)
TEMPLATE_PATH = TEMPLATE[0]
OUT_PATH = Path("outputs/AI경진대회_FA기술혁신Agent팀_2026-05-13.pptx")


# ─── 헬퍼 ────────────────────────────────────────────────────────
def replace_textframe(tf, lines, *, font=FONT_BODY, size=12, bold=False,
                      color=BLACK, align=None):
    """text_frame 전체를 lines (list[str]) 로 교체, 폰트/크기/색 강제."""
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def replace_textframe_mixed(tf, segments, *, align=None):
    """text_frame 을 [{text, font, size, bold, color, newline}] 시퀀스로 교체.

    newline=True 이면 새 paragraph 시작.
    """
    tf.clear()
    para = tf.paragraphs[0]
    if align is not None:
        para.alignment = align
    first_run_added = False
    for seg in segments:
        if seg.get("newline") and first_run_added:
            para = tf.add_paragraph()
            if align is not None:
                para.alignment = align
        run = para.add_run()
        run.text = seg["text"]
        run.font.name = seg.get("font", FONT_BODY)
        run.font.size = Pt(seg.get("size", 12))
        run.font.bold = seg.get("bold", False)
        run.font.color.rgb = seg.get("color", BLACK)
        first_run_added = True


def replace_cell_text(cell, text, *, font=FONT_BODY, size=10, bold=False,
                      color=BLACK):
    """표 셀 텍스트 교체 (단일 paragraph)."""
    replace_textframe(cell.text_frame, [text], font=font, size=size,
                      bold=bold, color=color, align=PP_ALIGN.CENTER)


# ─── 첨부 슬라이드용 카드/박스 헬퍼 (회색조) ──────────────────────
def add_band(slide, x, y, w, h, fill=DIM_GRAY):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def add_text(slide, x, y, w, h, text, *, size=10, bold=False,
             color=CHARCOAL, font=FONT_BODY, align=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    box.text_frame.word_wrap = True
    if isinstance(text, list):
        replace_textframe(box.text_frame, text, font=font, size=size,
                          bold=bold, color=color, align=align)
    else:
        replace_textframe(box.text_frame, [text], font=font, size=size,
                          bold=bold, color=color, align=align)
    return box


def add_gray_card(slide, x, y, w, h, *, header, body_lines,
                  header_size=11, body_size=9):
    """회색조 카드: 짙은 회색 헤더 띠 + 라이트 그레이 본문."""
    # 카드 외곽
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(x), Cm(y), Cm(w), Cm(h))
    base.fill.solid()
    base.fill.fore_color.rgb = LIGHT_GRAY
    base.line.color.rgb = LINE_GRAY
    base.line.width = Pt(0.5)

    # 헤더 띠
    head_h = 0.7
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(x), Cm(y), Cm(w), Cm(head_h))
    head.fill.solid()
    head.fill.fore_color.rgb = DIM_GRAY
    head.line.fill.background()
    head_tf = head.text_frame
    head_tf.margin_left = Cm(0.2)
    head_tf.margin_top = Cm(0.05)
    head_tf.margin_bottom = Cm(0.05)
    replace_textframe(head_tf, [header], font=FONT_BODY,
                      size=header_size, bold=True, color=WHITE)

    # 본문
    body = slide.shapes.add_textbox(Cm(x + 0.15), Cm(y + head_h + 0.1),
                                    Cm(w - 0.3), Cm(h - head_h - 0.2))
    body.text_frame.word_wrap = True
    replace_textframe(body.text_frame, body_lines, font=FONT_BODY,
                      size=body_size, color=CHARCOAL)


def add_appendix_titlebar(slide, slide_w, title, no_label):
    """첨부 슬라이드 상단 타이틀 바."""
    add_band(slide, 0, 0, slide_w, 1.4, fill=CHARCOAL)
    add_text(slide, 0.5, 0.25, slide_w - 1.0, 1.0, title,
             size=15, bold=True, color=WHITE)
    add_text(slide, slide_w - 3.5, 0.35, 3.0, 0.7,
             no_label, size=11, bold=False, color=LIGHT_GRAY,
             font=FONT_EMPH, align=PP_ALIGN.RIGHT)


def add_appendix_footer(slide, slide_w, slide_h, footer):
    add_band(slide, 0, slide_h - 0.5, slide_w, 0.5, fill=SOFT_GRAY)
    add_text(slide, 0.5, slide_h - 0.45, slide_w - 1.0, 0.4, footer,
             size=8, color=DIM_GRAY, font=FONT_EMPH)


# ─── 메인 콘텐츠 (Slide 1 / Slide 2) ─────────────────────────────
TITLE = "FA 기술혁신 Agent 팀"
SUBTITLE = (
    "5인 AI 에이전트 + 3 도메인 스킬 + 2 Hook 으로 주간보고·회의록·"
    "기술과제 발굴을 자동화 — 보고 업무 시간 단위 → 분 단위 단축"
)
RESPONSIBLE = "FA기술혁신파트 강모원책임"

ISSUE_BULLETS = [
    " • 주간/월간/회의록 작성에 보고서당 수시간 소요 (자료 수집·정리·서식 정렬)",
    " • 출처 누락·서로 다른 양식·과거 자료 재사용 어려움 — 보고 신뢰도 저하",
    " • 회의록·녹취 정리 지연 (작성자 1인 부담), 액션아이템 누락",
    " • 외부 기술 동향·경쟁사 정보가 보고서·로드맵에 늦게 반영 — 사외 시야 부족",
    " • 보고 표준화 부재로 영역 간 (생산·품질·구매) 자산 재사용 불가",
]

IMPROVE_BULLETS = [
    " • 5인 AI 에이전트 팀 (사내 3 · 사외 1 · 작성 2) 우편함 통신·병렬 develop",
    " • 도메인 스킬 3개: 회의록 자동화 · 녹취+OneNote 통합 · FA 6분야 동향 발굴",
    " • 결정론 PPT 파이프라인 (extract → fill) — 양식 도형/색 무결성 100% 보존",
    " • Hook 2개: SessionStart 자동 인덱싱 + TaskCompleted 출처 자동 검증",
    " • Agent Teams (CLAUDE_CODE_EXPERIMENTAL_AGENT_TEAMS=1) + Claude Code 만으로 구현",
]

EFFECT_BULLETS = [
    " • 작성 시간 시간 단위 → 분 단위 (주간보고 PR 요약 등 검증된 산출물 다수)",
    " • 출처 100% 자동 검증 (미기입 시 TaskCompleted hook 이 작업 차단)",
    " • 외부 동향·경쟁사 자료가 references/ 에 누적 → 다음 보고서가 더 강해짐",
    " • 사람 워크플로우와 동일 패턴 → 즉시 그룹 내 다른 영역(생산·품질·구매) 확장 가능",
    " • 단순 LLM 호출이 아닌 팀 기반 협업 — 신뢰성·재현성·감사 가능성 확보",
]

PARTICIPANT = ("FA기술혁신파트", "강모원",
               "FA 기술혁신 Agent 팀 설계·구축·운영 (5인 에이전트 + 3 스킬 + 2 Hook)")


# ─── 빌드 ────────────────────────────────────────────────────────
def build_main(prs):
    """Slide 1, Slide 2 의 placeholder 를 채운다."""
    # ── Slide 1 ──
    s1 = prs.slides[0]
    # [5] 제목 9 — '과제 명 : ㅇㅇㅇ'
    title_shape = s1.shapes[5]
    replace_textframe_mixed(title_shape.text_frame, [
        {"text": "과제 명 ", "font": FONT_BODY, "size": 16, "bold": True, "color": BLACK},
        {"text": ": ",       "font": FONT_BODY, "size": 16, "bold": True, "color": BLACK},
        {"text": TITLE,      "font": FONT_BODY, "size": 16, "bold": True, "color": GREEN},
    ])
    # [6] 직사각형 1 — 부제 라인
    replace_textframe(s1.shapes[6].text_frame, [SUBTITLE],
                      font=FONT_EMPH, size=13, bold=False, color=GREEN)
    # [7] TextBox 15 — 이슈 및 배경
    replace_textframe(s1.shapes[7].text_frame, ISSUE_BULLETS,
                      font=FONT_BODY, size=11, color=BLACK)
    # [8] TextBox 21 — 우상단 책임자
    replace_textframe(s1.shapes[8].text_frame, [RESPONSIBLE],
                      font=FONT_EMPH, size=12, bold=True, color=GREEN,
                      align=PP_ALIGN.RIGHT)
    # [10] TextBox 3 — 기대 효과
    replace_textframe(s1.shapes[10].text_frame, EFFECT_BULLETS,
                      font=FONT_BODY, size=11, color=BLACK)
    # [12] TextBox 8 — 개선 사항
    replace_textframe(s1.shapes[12].text_frame, IMPROVE_BULLETS,
                      font=FONT_BODY, size=11, color=BLACK)
    # [14] 표 6×4 — 참여 인원 (1행만 채우고 r2~r5 는 빈 칸 유지)
    tbl = s1.shapes[14].table
    replace_cell_text(tbl.cell(1, 0), "1", size=10)
    replace_cell_text(tbl.cell(1, 1), PARTICIPANT[0], size=10)
    replace_cell_text(tbl.cell(1, 2), PARTICIPANT[1], size=10)
    # 기여 내용은 좌측 정렬이 더 자연스러움
    replace_textframe(tbl.cell(1, 3).text_frame, [PARTICIPANT[2]],
                      font=FONT_BODY, size=10, color=BLACK,
                      align=PP_ALIGN.LEFT)
    # 빈 행 r2~r5 의 c0 번호 유지 (이미 양식에 있음), c1~c3 비워두기 → 양식 그대로

    # ── Slide 2 (보조 — 시스템 구성) ──
    s2 = prs.slides[1]
    # [0] 제목 — '보조.' → '보조 ① — 시스템 구성'
    replace_textframe(s2.shapes[0].text_frame,
                      ["보조 ① — 시스템 구성 (5 에이전트 + 3 스킬 + 2 Hook)"],
                      font=FONT_BODY, size=16, bold=True, color=BLACK)
    # [1] 우상단 책임자
    replace_textframe(s2.shapes[1].text_frame, [RESPONSIBLE],
                      font=FONT_EMPH, size=12, bold=True, color=GREEN,
                      align=PP_ALIGN.RIGHT)
    # 본문 다이어그램 — 양식 빈 공간에 도형 그룹 추가
    _build_arch_diagram(s2)


def _build_arch_diagram(slide):
    """Slide 2 본문에 5 에이전트 + 3 스킬 + 2 Hook 다이어그램 추가."""
    # ── 사용자 → 팀장 흐름 (상단 띠)
    add_text(slide, 1.0, 1.5, 25.5, 0.7,
             "사용자 요청 → 팀장 (메인 세션, 분해·우편함 라우팅)",
             size=12, bold=True, color=BLACK)

    # ── 5 에이전트 (1행 5열)
    agents = [
        ("data-teammate",         "KPI·Capex·MRM·DST", "사내 정량"),
        ("ops-teammate",          "협력사·HR·결재·문서", "사내 운영"),
        ("tech-research-teammate","외부 동향·경쟁사·평가", "사외 리서치"),
        ("document-writer",       "Markdown 양식 채움", "작성"),
        ("ppt-writer",            "PPT 양식 채움 (결정론)", "작성"),
    ]
    card_w, gap = 5.1, 0.15
    start_x = (27.52 - (card_w * 5 + gap * 4)) / 2
    for i, (name, role, tag) in enumerate(agents):
        x = start_x + i * (card_w + gap)
        # 카드 박스
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(x), Cm(2.5), Cm(card_w), Cm(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GREEN
        box.line.width = Pt(0.75)
        # 헤더
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Cm(x), Cm(2.5), Cm(card_w), Cm(0.65))
        head.fill.solid()
        head.fill.fore_color.rgb = GREEN
        head.line.fill.background()
        replace_textframe(head.text_frame, [name], font=FONT_BODY,
                          size=10, bold=True, color=WHITE,
                          align=PP_ALIGN.CENTER)
        # 태그
        add_text(slide, x, 3.25, card_w, 0.5, tag,
                 size=9, bold=True, color=GREEN, font=FONT_EMPH,
                 align=PP_ALIGN.CENTER)
        # 역할
        add_text(slide, x + 0.15, 3.85, card_w - 0.3, 2.9, role,
                 size=9, color=BLACK, align=PP_ALIGN.CENTER)

    # ── 3 스킬 (1행 3열, 카드형, 그린 라인)
    skills = [
        ("meeting-minutes",   "녹취 → 회의록 PPTX"),
        ("meeting-merger",    "녹취 + OneNote → 통합 Word"),
        ("fa-task-discovery", "FA 6분야 동향 → 과제 docx + 카톡"),
    ]
    sk_w, sk_gap = 8.5, 0.25
    sk_start = (27.52 - (sk_w * 3 + sk_gap * 2)) / 2
    add_text(slide, 1.0, 7.4, 25.5, 0.5,
             "도메인 스킬 3개 — 자동 발동 워크플로우",
             size=11, bold=True, color=BLACK)
    for i, (name, desc) in enumerate(skills):
        x = sk_start + i * (sk_w + sk_gap)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(x), Cm(8.0), Cm(sk_w), Cm(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GREEN
        box.line.width = Pt(0.5)
        # 헤더
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Cm(x), Cm(8.0), Cm(sk_w), Cm(0.65))
        head.fill.solid()
        head.fill.fore_color.rgb = GREEN
        head.line.fill.background()
        replace_textframe(head.text_frame, [name], font=FONT_BODY,
                          size=11, bold=True, color=WHITE,
                          align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.2, 8.75, sk_w - 0.4, 1.3, desc,
                 size=10, color=BLACK, align=PP_ALIGN.CENTER)

    # ── 2 Hook (1행 2열, 그린 띠)
    add_text(slide, 1.0, 10.6, 25.5, 0.5,
             "Hook 2개 — Claude Code 자동 검증",
             size=11, bold=True, color=BLACK)
    hooks = [
        ("Hook 1 — SessionStart",
         "references/ 자동 인덱싱  ·  tech-research 가 저장한 신규 자료 자동 등록"),
        ("Hook 2 — TaskCompleted",
         "outputs/*.md 출처 인용 자동 검증  ·  미기입 시 작업 차단 (exit code)"),
    ]
    h_w, h_gap = 12.8, 0.4
    h_start = (27.52 - (h_w * 2 + h_gap)) / 2
    for i, (name, desc) in enumerate(hooks):
        x = h_start + i * (h_w + h_gap)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Cm(x), Cm(11.15), Cm(h_w), Cm(2.0))
        band.fill.solid()
        band.fill.fore_color.rgb = WHITE
        band.line.color.rgb = BLACK
        band.line.width = Pt(0.75)
        # 좌측 그린 띠 (강조)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Cm(x), Cm(11.15), Cm(0.25), Cm(2.0))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GREEN
        accent.line.fill.background()
        add_text(slide, x + 0.5, 11.35, h_w - 0.7, 0.7, name,
                 size=11, bold=True, color=BLACK, font=FONT_EMPH)
        add_text(slide, x + 0.5, 12.0, h_w - 0.7, 1.2, desc,
                 size=10, color=BLACK)

    # ── 아래쪽 협업 패턴 요약
    add_text(slide, 1.0, 13.6, 25.5, 0.7,
             "SOP: 분해 → 병렬 수집 (data·ops·tech-research) → 위임 작성 (writer) → 품질 게이트 (Hook 자동 검증)",
             size=11, bold=True, color=GREEN, font=FONT_EMPH,
             align=PP_ALIGN.CENTER)

    # 인용
    add_text(slide, 1.0, 14.5, 25.5, 0.5,
             "출처: .claude/agents/*.md, .claude/skills/*/SKILL.md, CLAUDE.md, scripts/hooks/check_output_citations.py",
             size=8, color=DIM_GRAY, font=FONT_EMPH)


# ─── 첨부 슬라이드 빌드 ───────────────────────────────────────────
def build_appendix(prs):
    slide_w = prs.slide_width / 360000  # EMU → cm
    slide_h = prs.slide_height / 360000

    # 양식은 layout 이 1개뿐 ('제목 슬라이드') — 그대로 사용 (도형은 모두 직접 추가)
    blank_layout = prs.slide_layouts[0]

    # ── 첨부 1 — 팀원 5명 + SOP 4단계 + Hook 2개 ──
    s = prs.slides.add_slide(blank_layout)
    add_appendix_titlebar(s, slide_w, "첨부 ① — 팀원 5명 · SOP · Hook",
                          "Appendix 1 / 3")

    # 5명 카드 (1행 5열, 상단)
    members = [
        ("data-teammate",
         ["[Read·Grep·Glob·Bash]",
          "사내 정량:",
          "KPI · Capex · MRM(양산) · DST(R&D)",
          "추측 금지, 항상 출처 인용",
          "외부 동향 요청은 tech-research 위임"]),
        ("ops-teammate",
         ["[Read·Grep·Glob·Bash]",
          "사내 운영:",
          "협력사 · HR · 결재 큐 · 사내 PDF/PPT",
          "사람 승인 필요한 작업은 제안만",
          "KPI·외부 요청은 위임"]),
        ("tech-research-teammate",
         ["[+ WebSearch · WebFetch + MCP]",
          "사외 리서치:",
          "외부 동향 · 경쟁사 · 기술 평가",
          "AMR/협동로봇/디지털트윈/AI/물류",
          "references/ 자동 아카이빙"]),
        ("document-writer",
         ["[Read·Write·Edit·Glob·Grep]",
          "Markdown 작성:",
          "templates/*.md placeholder 채움",
          "다른 팀원에 우편함 직접 요청",
          "plan-mode 로 outputs/ 저장"]),
        ("ppt-writer",
         ["[+ Bash, 결정론 파이프라인]",
          "PPT 작성:",
          "ppt_extract.py + ppt_fill.py",
          "5단계 placeholder 추론",
          "양식 도형/색상 무결성 보존"]),
    ]
    card_w, gap = 5.1, 0.15
    start_x = (slide_w - (card_w * 5 + gap * 4)) / 2
    for i, (name, lines) in enumerate(members):
        x = start_x + i * (card_w + gap)
        add_gray_card(s, x, 1.7, card_w, 6.0,
                      header=name, body_lines=lines,
                      header_size=10, body_size=8)

    # SOP 4 단계 (하단 좌측 절반)
    add_text(s, 0.8, 8.1, slide_w / 2 - 1.0, 0.55,
             "표준 작업 절차 (SOP) — Agent Teams 운영 패턴",
             size=12, bold=True, color=CHARCOAL)
    sop_steps = [
        ("Step 1 — 분해",
         "리더가 요청을 도메인별 작업으로 쪼갬 (사내 수치 ↔ 사내 운영 ↔ 사외 리서치, 작성 ↔ 검토)"),
        ("Step 2 — 병렬 수집",
         "필요한 수집 팀원(data·ops·tech-research) 동시 spawn, 출처와 함께 수집 (리더 대화 미상속)"),
        ("Step 3 — 위임 작성",
         "결과를 모아 document-writer(.md) / ppt-writer(.pptx) 가 placeholder 만 치환, 부족 데이터는 우편함 라운드"),
        ("Step 4 — 품질 게이트",
         "TaskCompleted hook 이 outputs/*.md 출처 인용 자동 검증, 미기입 시 작업 차단"),
    ]
    step_y = 8.75
    step_h = 1.1
    for i, (title_, desc) in enumerate(sop_steps):
        y = step_y + i * (step_h + 0.05)
        # 좌측 라벨
        lab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(0.8), Cm(y), Cm(3.0), Cm(step_h))
        lab.fill.solid()
        lab.fill.fore_color.rgb = DIM_GRAY
        lab.line.fill.background()
        replace_textframe(lab.text_frame, [title_],
                          font=FONT_BODY, size=10, bold=True, color=WHITE,
                          align=PP_ALIGN.CENTER)
        # 본문
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(3.85), Cm(y), Cm(slide_w / 2 - 4.65), Cm(step_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = SOFT_GRAY
        bg.line.color.rgb = LINE_GRAY
        bg.line.width = Pt(0.4)
        add_text(s, 4.05, y + 0.2, slide_w / 2 - 5.05, step_h - 0.4, desc,
                 size=9, color=CHARCOAL)

    # Hook 2개 (하단 우측 절반)
    add_text(s, slide_w / 2 + 0.2, 8.1, slide_w / 2 - 1.0, 0.55,
             "Hook 2개 — Claude Code 자동 검증",
             size=12, bold=True, color=CHARCOAL)
    hooks = [
        ("SessionStart Hook",
         "scripts/build_reference_index.py",
         ["세션 시작 시 references/ 폴더 전체 스캔",
          "frontmatter (title·category·date·summary·tags) 추출",
          "→ references/INDEX.md 자동 생성·갱신",
          "tech-research 가 저장한 신규 자료 다음 세션부터 자동 인덱싱"]),
        ("TaskCompleted Hook",
         "scripts/hooks/check_output_citations.py",
         ["outputs/*.md 생성 후 자동 실행",
          "각 섹션에 (출처:…) / [N행] / 자료 미확인 중 1개 이상 필수",
          "미기입 시 exit non-zero → 작업 차단",
          "URL 인용도 통과 ((출처: <URL>, accessed YYYY-MM-DD))"]),
    ]
    hook_x = slide_w / 2 + 0.2
    hook_w = slide_w / 2 - 0.4
    hook_h = 2.3
    for i, (title_, script, lines) in enumerate(hooks):
        y = 8.75 + i * (hook_h + 0.15)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(hook_x), Cm(y), Cm(hook_w), Cm(hook_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = LINE_GRAY
        bg.line.width = Pt(0.5)
        # 좌측 강조 띠
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Cm(hook_x), Cm(y), Cm(0.2), Cm(hook_h))
        accent.fill.solid()
        accent.fill.fore_color.rgb = DIM_GRAY
        accent.line.fill.background()
        add_text(s, hook_x + 0.4, y + 0.15, hook_w - 0.6, 0.55, title_,
                 size=11, bold=True, color=CHARCOAL)
        add_text(s, hook_x + 0.4, y + 0.7, hook_w - 0.6, 0.4, script,
                 size=8, color=DIM_GRAY, font=FONT_EMPH)
        add_text(s, hook_x + 0.4, y + 1.15, hook_w - 0.6, hook_h - 1.3, lines,
                 size=8, color=CHARCOAL)

    add_appendix_footer(s, slide_w, slide_h,
                        "출처: .claude/agents/*.md, .claude/settings.json, scripts/hooks/check_output_citations.py, scripts/build_reference_index.py")

    # ── 첨부 2 — 데모 산출물 4종 격자 ──
    s = prs.slides.add_slide(blank_layout)
    add_appendix_titlebar(s, slide_w, "첨부 ② — 검증된 산출물 (실제 outputs/)",
                          "Appendix 2 / 3")
    add_text(s, 0.8, 1.5, slide_w - 1.6, 0.55,
             "본 시스템으로 이미 생성된 산출물 — 모두 (출처:…) 자동 검증 통과",
             size=11, color=CHARCOAL, font=FONT_EMPH)

    demos = [
        ("① 주간보고 (출처 자동 검증 모범)",
         "outputs/주간보고_2026-05-12.md",
         ["• MD 본문 모든 KPI/Capex 값에 (출처: sample_data/…)·[N행] 인용 자동 부착",
          "• data-teammate · ops-teammate 가 같은 라운드에 병렬 수집",
          "• document-writer 가 templates/주간보고.md 의 17개 placeholder 만 치환",
          "• 같은 양식으로 매주 재현 가능 — 사람·날짜에 의존 X"]),
        ("② 회의록 (녹취 → 표준 PPTX, meeting-minutes 스킬)",
         "outputs/회의록_신규R&D과제착수_2026-05-12.md",
         ["• 녹취록 txt 또는 텍스트 붙여넣기 → LGES 표준 회의록 PPTX 자동 생성",
          "• 참석자·안건·결정사항·액션아이템 인터랙티브 UI 로 확정",
          "• 부가: meeting-merger 스킬 — 녹취 + OneNote → 통합 Word 1건 출력",
          "• 회의 종료 직후 즉시 정리, 작성자 부담 0"]),
        ("③ 로드맵 v2 — 4축 KPI 슬라이드 (ppt-writer)",
         "outputs/roadmap_v2.pptx",
         ["• 중장기 로드맵 4축 (양산준비·구매·R&D·표준화) KPI 자동 시각화",
          "• templates/사내양식/생산기술혁신센터 MRM 과제 요약서 양식_AX 과제 포함.pptx",
          "• ppt_extract.py 로 placeholder 추론 → ppt_fill.py 로 결정론 채움",
          "• 양식 도형/색상/폰트 100% 보존 — 회사 표준 그대로 유지"]),
        ("④ FA 기술과제 후보 docx + 카톡 (fa-task-discovery 스킬)",
         "outputs/FA_기술과제후보_20260512.docx",
         ["• AMR/AGV · 협동로봇 · 디지털트윈 · 스마트물류 · AI · 피지컬AI 6 분야",
          "• 한/영 외부 동향 자동 수집 → LGES 적용 가능성 평가 → 과제 후보 도출",
          "• docx 보고서 + 카카오톡 요약 동시 전송",
          "• tech-research-teammate 가 references/ 에 결과 자동 아카이빙"]),
    ]
    cell_w = (slide_w - 1.6 - 0.3) / 2  # 좌우 카드 폭
    cell_h = (slide_h - 1.4 - 2.5 - 0.6) / 2  # 상하 카드 높이
    for i, (title_, path, lines) in enumerate(demos):
        col = i % 2
        row = i // 2
        x = 0.8 + col * (cell_w + 0.3)
        y = 2.3 + row * (cell_h + 0.25)
        # 카드 외곽
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(x), Cm(y), Cm(cell_w), Cm(cell_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.color.rgb = LINE_GRAY
        bg.line.width = Pt(0.5)
        # 상단 라벨 띠
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(x), Cm(y), Cm(cell_w), Cm(0.65))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = DIM_GRAY
        hdr.line.fill.background()
        replace_textframe(hdr.text_frame, [title_],
                          font=FONT_BODY, size=10, bold=True, color=WHITE)
        # 경로 (소형)
        add_text(s, x + 0.2, y + 0.75, cell_w - 0.4, 0.4, path,
                 size=8, color=DIM_GRAY, font=FONT_EMPH)
        # 본문 라인
        add_text(s, x + 0.2, y + 1.2, cell_w - 0.4, cell_h - 1.4, lines,
                 size=9, color=CHARCOAL)

    add_appendix_footer(s, slide_w, slide_h,
                        "출처: outputs/주간보고_2026-05-12.md, outputs/회의록_신규R&D과제착수_2026-05-12.md, outputs/roadmap_v2.pptx, outputs/FA_기술과제후보_20260512.docx")

    # ── 첨부 3 — ROI + 다음 단계 + 확장 시나리오 ──
    s = prs.slides.add_slide(blank_layout)
    add_appendix_titlebar(s, slide_w, "첨부 ③ — ROI · 다음 단계 · 확장 시나리오",
                          "Appendix 3 / 3")

    # 좌측 — ROI 정량 지표
    add_text(s, 0.8, 1.7, slide_w / 2 - 1.0, 0.55,
             "ROI — 본 시스템 도입 효과",
             size=12, bold=True, color=CHARCOAL)
    roi_items = [
        ("작성 시간",
         "보고서당 수시간 → 분 단위 (PR #18 4 커밋 1일 산출 17 슬라이드 + MD)"),
        ("출처 신뢰도",
         "100% 자동 검증 (TaskCompleted hook 이 미기입 차단)"),
        ("자료 누적",
         "외부 동향·경쟁사가 references/ 에 누적 → 보고서 품질 시간에 비례 증가"),
        ("재현성",
         "결정론 PPT 파이프라인 — 같은 입력 = 같은 출력, 양식 무결성 100%"),
        ("표준화",
         "5인 팀 + 4단계 SOP — 영역·작성자에 무관하게 동일 품질"),
        ("그룹 확장성",
         "사람 워크플로우와 동일 — 생산·품질·구매 영역에 즉시 이식 가능"),
    ]
    roi_y = 2.4
    for i, (label, desc) in enumerate(roi_items):
        y = roi_y + i * 1.45
        lab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(0.8), Cm(y), Cm(2.8), Cm(1.25))
        lab.fill.solid()
        lab.fill.fore_color.rgb = DIM_GRAY
        lab.line.fill.background()
        replace_textframe(lab.text_frame, [label],
                          font=FONT_BODY, size=10, bold=True, color=WHITE,
                          align=PP_ALIGN.CENTER)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(3.65), Cm(y), Cm(slide_w / 2 - 4.45), Cm(1.25))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY
        bg.line.color.rgb = LINE_GRAY
        bg.line.width = Pt(0.4)
        add_text(s, 3.85, y + 0.2, slide_w / 2 - 4.85, 0.95, desc,
                 size=9, color=CHARCOAL)

    # 우측 — 다음 단계 + 확장 시나리오
    add_text(s, slide_w / 2 + 0.2, 1.7, slide_w / 2 - 1.0, 0.55,
             "다음 단계 (Q3~Q4 2026)",
             size=12, bold=True, color=CHARCOAL)
    next_items = [
        "① 그룹 내 다른 영역 (생산·품질·구매) 의 보고 업무로 본 시스템 이식",
        "② 외부 MCP 추가 (kakao calendar, gmail, github, opendart 등) — 자동 알림·일정",
        "③ 사내 SharePoint·OneNote API 연결 — 사내 문서 검색 정확도 강화",
        "④ 사내 강모원 책임 외 추가 사용자 1~2명 인계 — 운영 노하우 전파",
        "⑤ 정량 ROI 1개월 측정 (작성 시간·출처 인용률·재사용률) → 그룹 보고",
    ]
    next_x = slide_w / 2 + 0.2
    next_w = slide_w / 2 - 0.4
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Cm(next_x), Cm(2.4), Cm(next_w), Cm(4.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.color.rgb = LINE_GRAY
    bg.line.width = Pt(0.5)
    add_text(s, next_x + 0.3, 2.6, next_w - 0.6, 4.1, next_items,
             size=10, color=CHARCOAL)

    # 우측 하단 — 확장 시나리오
    add_text(s, next_x, 7.1, next_w, 0.55,
             "확장 시나리오 — 본 시스템 이식 사례 (예상)",
             size=12, bold=True, color=CHARCOAL)
    scenarios = [
        ("생산기술담당",
         "양산 가동률·CT·Yield KPI 일/주간 자동 보고, 라인별 이상 알림"),
        ("품질담당",
         "QC 데이터·반품 분석 주간 보고, 협력사 품질 자동 평가서"),
        ("구매담당",
         "벤더 평가·발주 현황 보고, 시장 가격·환율 외부 데이터 자동 누적"),
        ("R&D",
         "특허·논문 동향 + 사내 DST 진척 통합, 과제 후보 자동 발굴"),
    ]
    sc_y = 7.75
    sc_h = (slide_h - 0.5 - sc_y) / 4 - 0.05
    for i, (org, desc) in enumerate(scenarios):
        y = sc_y + i * (sc_h + 0.05)
        lab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(next_x), Cm(y), Cm(2.8), Cm(sc_h))
        lab.fill.solid()
        lab.fill.fore_color.rgb = MID_GRAY
        lab.line.fill.background()
        replace_textframe(lab.text_frame, [org],
                          font=FONT_BODY, size=10, bold=True, color=WHITE,
                          align=PP_ALIGN.CENTER)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Cm(next_x + 2.85), Cm(y), Cm(next_w - 2.85), Cm(sc_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = SOFT_GRAY
        bg.line.color.rgb = LINE_GRAY
        bg.line.width = Pt(0.3)
        add_text(s, next_x + 3.05, y + 0.15, next_w - 3.25, sc_h - 0.3, desc,
                 size=9, color=CHARCOAL)

    add_appendix_footer(s, slide_w, slide_h,
                        "출처: outputs/주간보고_2026-05-12_PR요약.md, CLAUDE.md, .claude/agents/*.md, references/roadmap/2026_FA기술담당_중장기로드맵_v2.md")


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(TEMPLATE_PATH)
    print(f"slide_size: {prs.slide_width/360000:.2f}cm × {prs.slide_height/360000:.2f}cm")
    print(f"입력 슬라이드 수: {len(prs.slides)}")
    build_main(prs)
    build_appendix(prs)
    print(f"최종 슬라이드 수: {len(prs.slides)}")
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
