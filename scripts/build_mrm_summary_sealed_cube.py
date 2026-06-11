"""MRM 과제 요약서 양식에 '밀폐형 큐브 물류' 내용을 채우는 빌더 (v2).

양식 ``templates/사내양식/생산기술혁신센터 MRM 과제 요약서 양식_AX 과제 포함.pptx``
의 폰트(Arial Narrow), 사이즈(20/14/11/10/9/8/7pt), 정렬, margins, anchor
를 모두 그대로 유지하면서 placeholder 텍스트만 치환한다.

핵심 원칙:
- 폰트: Arial Narrow (양식 통일)
- 사이즈: 양식 원본 그대로
- 정렬·anchor·margins: 양식 그대로
- As-is/To-be 본문: 빈 AUTO_SHAPE 박스에 직접 입력 (textbox 오버레이 금지)
- 표 셀: 셀별 양식의 alignment 유지

출력: outputs/MRM과제요약_밀폐형큐브물류_2026-05-20.pptx
"""
from __future__ import annotations

import os
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

TEMPLATE_PATH = None
for root, dirs, files in os.walk("templates/사내양식"):
    for f in files:
        if "MRM" in f and f.endswith(".pptx"):
            TEMPLATE_PATH = os.path.join(root, f)
assert TEMPLATE_PATH, "MRM 양식 파일을 찾지 못했습니다"

OUT_PATH = Path("outputs/MRM과제요약_밀폐형큐브물류_2026-05-20.pptx")

# 색상 — 양식 default 는 건드리지 않고, 강조에만 사용
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x00, 0xFF)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)

FONT_LATIN = "Arial Narrow"      # 영문/숫자 (양식 원본 폰트)
FONT_EA = "LG스마트체 Regular"  # 한글 (사내 표준)


def _set_run_fonts(run):
    """run 에 한·영 분리 폰트 적용. 한글=LG스마트체, 영문=Arial Narrow."""
    run.font.name = FONT_LATIN
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", FONT_EA), ("cs", FONT_EA), ("latin", FONT_LATIN)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def cm_emu(v):
    return (v or 0) / 360000


def find_shape_at(slide, target_l, target_t, tol=0.15, skip_groups=True):
    for shp in slide.shapes:
        if skip_groups and shp.shape_type and "GROUP" in str(shp.shape_type):
            continue
        l = cm_emu(shp.left)
        t = cm_emu(shp.top)
        if abs(l - target_l) < tol and abs(t - target_t) < tol:
            return shp
    return None


def find_table_at(slide, target_l, target_t, tol=0.2):
    for shp in slide.shapes:
        if not shp.has_table:
            continue
        l = cm_emu(shp.left)
        t = cm_emu(shp.top)
        if abs(l - target_l) < tol and abs(t - target_t) < tol:
            return shp.table
    return None


def write_runs(text_frame, paragraphs):
    """paragraphs: list of paragraph specs, each = list of (text, opts).

    opts: size, bold, color, align (옵션). 폰트는 항상 Arial Narrow.
    """
    text_frame.clear()
    for pi, runs in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if pi == 0 else text_frame.add_paragraph()
        # paragraph-level alignment (첫 run 의 opts.align 우선)
        if runs:
            first_opts = runs[0][1] if isinstance(runs[0], tuple) else {}
            if "align" in first_opts:
                p.alignment = first_opts["align"]
        for item in runs:
            if isinstance(item, tuple):
                text, opts = item
            else:
                text, opts = item, {}
            run = p.add_run()
            run.text = text
            if "size" in opts:
                run.font.size = Pt(opts["size"])
            if opts.get("bold") is not None:
                run.font.bold = opts["bold"]
            if "color" in opts:
                run.font.color.rgb = opts["color"]
            _set_run_fonts(run)


def set_cell(cell, text, *, size=9, bold=False, color=None,
             align=None, anchor=None):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    _set_run_fonts(run)


def build():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(TEMPLATE_PATH, OUT_PATH)
    prs = Presentation(str(OUT_PATH))
    slide = prs.slides[0]

    # ── 좌상단 과제명 (20pt Bold LEFT) ──────────────────
    shp = find_shape_at(slide, 0.37, 0.24)
    write_runs(shp.text_frame, [
        [("과제명 ", {"size": 20, "bold": True, "color": BLACK,
                      "align": PP_ALIGN.LEFT}),
         (": 밀폐형 큐브 물류 설비",
          {"size": 20, "bold": True, "color": BLACK})],
    ])

    # ── 우상단 수행 부서 (9pt RIGHT) ────────────────────
    shp = find_shape_at(slide, 23.13, 0.08)
    write_runs(shp.text_frame, [
        [("수행 부서 및 담당자 ", {"size": 9, "align": PP_ALIGN.RIGHT}),
         (": FA기술혁신파트 강모원책임", {"size": 9})],
    ])

    # ── 과제 목표 띠 (14pt Bold LEFT, AUTO_SHAPE) ───────
    shp = find_shape_at(slide, 0.58, 1.44)
    write_runs(shp.text_frame, [
        [("과제 목표 ", {"size": 14, "bold": True, "color": BLACK,
                         "align": PP_ALIGN.LEFT}),
         (": 공정·Stock 0.3% DR 자체 유지 + AMR 단기 밀폐 Bridge로 DR Less 표준화 (2028년 양산 적용)",
          {"size": 14, "bold": True, "color": BLACK})],
    ])

    # ── ① 기술 개발 배경 본문 (10pt, T=4.49) ────────────
    shp = find_shape_at(slide, 0.97, 4.49)
    write_runs(shp.text_frame, [
        [("- FOUP (Front Opening Unified Pod): 반도체 표준 밀폐 캐리어",
          {"size": 10, "bold": True, "color": BLACK,
           "align": PP_ALIGN.LEFT})],
        [("  외부 파티클·수분·금속 오염 차단 + AMHS/OHT 자동화 호환 (SEMI 표준)",
          {"size": 10, "color": BLACK})],
        [("- 배터리 0.3% DR 공정 적용 시 FOUP 단독 밀폐는 미충족",
          {"size": 10, "bold": True, "color": BLACK})],
        [("  ▶ Stock(SMC 일체형) + AMR Bridge(1~3h 밀폐) + 공정 국소 DR 통합 컨셉 필요",
          {"size": 10, "color": BLUE})],
    ])

    # ── As-is 본문 (AUTO_SHAPE #7, L=1.66 T=6.82 W=11.30 H=4.31) ──
    shp = find_shape_at(slide, 1.66, 6.82)
    write_runs(shp.text_frame, [
        [("[현행 운영]", {"size": 10, "bold": True, "color": BLACK,
                         "align": PP_ALIGN.LEFT})],
        [("- 대형 DR 룸 + Open Rack 구조 (전체 라인 DR 유지)",
          {"size": 10, "color": BLACK})],
        [("- ESWA 1동 DR 면적 4,050평 (0.5% 공조 대응 공조실 +330평 증축)",
          {"size": 10, "color": BLACK})],
        [("- Rack 분리·재투입 공수 3,234 MD / Line",
          {"size": 10, "color": BLACK})],
        [("- 셀 표면 산화·결함 리스크 + 매몰 Loss 高",
          {"size": 10, "color": BLACK})],
        [("- DR 면적·공조·공수가 CAPEX 의 주요 비중",
          {"size": 10, "color": BLACK})],
    ])

    # ── To-be 본문 (AUTO_SHAPE #8, L=1.66 T=11.25 W=11.30 H=4.30) ──
    shp = find_shape_at(slide, 1.66, 11.25)
    write_runs(shp.text_frame, [
        [("[밀폐형 Cube 물류 통합 컨셉]",
          {"size": 10, "bold": True, "color": BLUE,
           "align": PP_ALIGN.LEFT})],
        [("- Stock = SMC 일체형 밀폐 Cube Rack (분리·재투입 不요)",
          {"size": 10, "bold": True, "color": BLACK})],
        [("- 이송 = AMR + FOUP-like 단기 밀폐 캐리어 (1~3h 0.3% 유지)",
          {"size": 10, "bold": True, "color": BLACK})],
        [("- 공정 = 설비 단위 국소 DR (0.3% 자체 유지) → 전체 DR 룸 不요",
          {"size": 10, "bold": True, "color": BLACK})],
        [("- SMC 효과: 면적 +25% / 공수 -33% / 투자비 -15%",
          {"size": 10, "color": BLACK})],
        [("- Phasing 투자 + 철거 후 타 Line 재활용 (재활용률 50% '28 목표)",
          {"size": 10, "color": BLACK})],
    ])

    # ── 예상 기대 효과 본문 (10pt, T=16.79) ──────────────
    shp = find_shape_at(slide, 0.96, 16.79)
    write_runs(shp.text_frame, [
        [("- DR 면적 대폭 ↓ + 매몰 Loss 절감 + CAPEX -20% / 면적 -20% (2028 4축 KPI)",
          {"size": 10, "bold": True, "color": BLACK,
           "align": PP_ALIGN.LEFT})],
        [("- 재활용률 50% 목표, 신규 Site 적용 시 모듈 표준화 가능",
          {"size": 10, "color": BLACK})],
    ])

    # ── 우측 ② 검증 항목 표 (3×5, T=3.89) ───────────────
    tbl = find_table_at(slide, 13.96, 3.89)
    # 헤더 row0 그대로
    # row1
    set_cell(tbl.cell(1, 0), "1", size=9, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(1, 1), "0.3% DR 유지 시간 (이송)", size=9,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(1, 2), "≥ 1~3h (단기 밀폐)", size=9,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(1, 3), "Pilot 측정 예정", size=9,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(1, 4), "FOUP 사양서", size=8,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # row2
    set_cell(tbl.cell(2, 0), "2", size=9, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(2, 1), "SMC Cube Rack 효과", size=9,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(2, 2), "면적 +25% / 공수 -33%", size=9,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(2, 3), "Demo 검증 (11.72억)", size=9,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(2, 4), "로드맵 v2 §3.1", size=8,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # ── 우측 ③ 세부 진행 사항 표 (3×3, T=7.59) ──────────
    tbl = find_table_at(slide, 13.98, 7.59)
    # 헤더 (양식이 비어있음 — 추가)
    set_cell(tbl.cell(0, 0), "구분", size=10, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(0, 1), "As-is (현행)", size=10, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(0, 2), "To-be (목표)", size=10, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # row1
    set_cell(tbl.cell(1, 0), "DR 면적", size=9, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(1, 1), "4,050평 (ESWA 1동)", size=9,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(1, 2), "국소 DR 化 → 大폭 ↓", size=9, bold=True,
             color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # row2
    set_cell(tbl.cell(2, 0), "투자비 / Line", size=9, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(2, 1), "239.8억", size=9,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    set_cell(tbl.cell(2, 2), "203.3억 (-15%)", size=9, bold=True, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 우측 ④ "1. 마스터플랜" 텍스트 (10pt Bold) ────────
    shp = find_shape_at(slide, 14.16, 13.44)
    write_runs(shp.text_frame, [
        [("1. ", {"size": 10, "bold": True, "color": BLACK,
                   "align": PP_ALIGN.LEFT}),
         ("마스터플랜 (분기별 추진 일정)",
          {"size": 10, "bold": True, "color": BLACK})],
    ])

    # ── 일정 표 (2×6, T=14.72) — 양식 RIGHT alignment 유지 ──
    tbl = find_table_at(slide, 14.05, 14.72)
    # 헤더 row0 그대로
    # row1 — 양식 alignment=RIGHT 유지하나 가독성 위해 CENTER 로 변경
    schedule = [
        "FOUP/SMC\n컨셉 검토",
        "SMC Demo\n착수 (11.72억)",
        "Demo 진행\n(ESGM2)",
        "M3 Gate\n('26.10)",
        "Hybrid PoC\n착수",
        "표준 양산\n적용 결정",
    ]
    for c, txt in enumerate(schedule):
        is_milestone = c in (3, 5)
        set_cell(
            tbl.cell(1, c), txt,
            size=8,
            bold=is_milestone,
            color=BLUE if is_milestone else None,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ── M1~M5 마일스톤 옆 텍스트박스 ────────────────────
    milestones = [
        (14.64, 16.86, "M1: 컨셉"),
        (17.20, 16.88, "M2: PoC"),
        (18.51, 16.94, "(중간 평가)"),
        (19.51, 16.96, "M3: Demo"),
        (21.48, 16.98, "M4: 양산 결정"),
        (24.35, 17.01, "M5: 표준 양산"),
    ]
    for target_l, target_t, label in milestones:
        shp = find_shape_at(slide, target_l, target_t, tol=0.25)
        if shp and shp.has_text_frame:
            write_runs(shp.text_frame, [
                [(label, {"size": 8, "bold": True, "color": BLACK,
                          "align": PP_ALIGN.CENTER})],
            ])

    # ── 풋노트 (8pt, MOLD_GREEN) ────────────────────────
    shp = find_shape_at(slide, 14.02, 17.83)
    # 풋노트 box width 0.83cm 으로 너무 좁음 → width 확장
    shp.width = Cm(13.0)
    write_runs(shp.text_frame, [
        [("* 출처: references/기술자료/FOUP(510x515)_3S.md │ references/roadmap/2026_FA기술담당_중장기로드맵_v2.md [L62-70] │ references/26FA KPI.md [L354-376] │ 이송 1~3h·국소 DR 컨셉은 보고 작성자 정의",
          {"size": 7, "bold": False, "color": MOLD_GREEN,
           "align": PP_ALIGN.LEFT})],
    ])

    # ── 2번째 슬라이드 (유첨) 삭제 ──────────────────────
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    rid_attr = ("{http://schemas.openxmlformats.org/officeDocument/2006/"
                "relationships}id")
    for s_id in slide_ids[1:]:
        rId = s_id.attrib[rid_attr]
        prs.part.drop_rel(rId)
        xml_slides.remove(s_id)

    prs.save(str(OUT_PATH))
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
