"""FA-Agent 온보딩 가이드 → 전 구성원 발표용 PPTX 빌더 (교육용 자체 양식).

온보딩_가이드.md (리포 루트) 의 내용을 발표·설명에 맞춘 교육용 슬라이드로
변환한다. GitHub·Claude·코드를 모르는 전 구성원 대상.

※ 본 빌더는 사용자 지시("설명 가이드이므로 LG 보고 양식을 따르지 말 것")에
   따라 LGES 분류형 보고 레이아웃(_ppt_layout)을 적용하지 않는 예외 산출물.
   16:9 · 한 장에 한 메시지 · 큰 글씨 · 카드/말풍선 중심의 자체 교육 양식.
   폰트는 사내 PC 어디에나 있는 '맑은 고딕'으로 통일(폰트 미설치 깨짐 방지).

구성 (15장):
  표지 → 아젠다 → [1.개념] 등장인물 3·AI 팀원 5·할 수 있는 일 →
  [2.준비] 계정 2개·안심 포인트 → [3.사용법] 시작 4단계·첫 명령·
  화면 대응·치트시트·잘 시키는 요령 → [4.규칙] 4가지·FAQ → 오늘 시작하기

데이터 출처(추측 없음): 온보딩_가이드.md · CLAUDE.md · README.md
[106~113, 149~158, 186~196행].
출력: outputs/FA-Agent_온보딩_발표자료_2026-06-12.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

OUT_PATH = Path("outputs/FA-Agent_온보딩_발표자료_2026-06-12.pptx")

# ─── 자체 팔레트 (교육용: 잉크 + 파랑 1색 + 회색) ───
INK = RGBColor(0x1F, 0x29, 0x37)        # 본문·제목 (짙은 잉크)
SUB = RGBColor(0x6B, 0x72, 0x80)        # 보조 텍스트
ACCENT = RGBColor(0x25, 0x63, 0xEB)     # 포인트 파랑
ACCENT_SOFT = RGBColor(0xEF, 0xF4, 0xFF)  # 연파랑 배경
CARD_BG = RGBColor(0xF7, 0xF8, 0xFA)    # 카드 배경
LINE = RGBColor(0xD8, 0xDC, 0xE2)       # 카드 테두리
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OK_GREEN = RGBColor(0x0B, 0x7A, 0x3E)   # 체크/안심 포인트 전용

FONT = "맑은 고딕"

# 슬라이드 16:9 (발표 표준)
SLIDE_W = 33.867
SLIDE_H = 19.05
MARGIN = 1.5
USABLE_W = SLIDE_W - 2 * MARGIN


def set_run(run, text, *, size=14, bold=False, color=INK, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        old = rPr.find(qn(f"a:{tag}"))
        if old is not None:
            rPr.remove(old)
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", font)


def write_lines(tf, lines, *, size=14, bold=False, color=INK,
                align=None, anchor=None, space_after=4):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, opts = item if isinstance(item, tuple) else (item, {})
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = opts.get("align", align)
        para.space_before = Pt(opts.get("space_before", 0))
        para.space_after = Pt(opts.get("space_after", space_after))
        run = para.add_run()
        set_run(run, text, size=opts.get("size", size),
                bold=opts.get("bold", bold), color=opts.get("color", color))


def add_text(slide, l, t, w, h, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    write_lines(tf, lines, **kw)
    return tb


def add_shape(slide, shape, l, t, w, h, *, fill=CARD_BG, line=LINE,
              line_w=1.0, radius=None):
    shp = slide.shapes.add_shape(shape, Cm(l), Cm(t), Cm(w), Cm(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if radius is not None and shp.adjustments:
        shp.adjustments[0] = radius
    return shp


def card(slide, l, t, w, h, title, body, *, bg=CARD_BG, border=LINE,
         title_color=INK, title_size=15, body_size=12.5, anchor=None):
    """둥근 카드: 상단 제목 + 본문."""
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h,
                    fill=bg, line=border, radius=0.06)
    tf = shp.text_frame
    tf.margin_left = Cm(0.45)
    tf.margin_right = Cm(0.45)
    tf.margin_top = Cm(0.30)
    tf.margin_bottom = Cm(0.25)
    lines = [(title, {"size": title_size, "bold": True,
                      "color": title_color, "space_after": 6})]
    lines += [it if isinstance(it, tuple) else
              (it, {"size": body_size, "color": INK}) for it in body]
    write_lines(tf, lines, anchor=anchor)
    return shp


def num_badge(slide, l, t, n, *, d=1.05, fill=ACCENT, color=WHITE, size=16):
    shp = add_shape(slide, MSO_SHAPE.OVAL, l, t, d, d, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    write_lines(tf, [(str(n), {"size": size, "bold": True, "color": color,
                               "align": PP_ALIGN.CENTER})],
                anchor=MSO_ANCHOR.MIDDLE)
    return shp


def bubble(slide, l, t, w, h, quote, note=None):
    """말풍선 카드 — '이렇게 말하세요' 인용."""
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h,
                    fill=ACCENT_SOFT, line=ACCENT, line_w=1.2, radius=0.16)
    tf = shp.text_frame
    tf.margin_left = Cm(0.5)
    tf.margin_right = Cm(0.5)
    tf.margin_top = Cm(0.2)
    tf.margin_bottom = Cm(0.2)
    lines = [(f"“{quote}”", {"size": 14, "bold": True, "color": INK,
                             "space_after": 2})]
    if note:
        lines.append((note, {"size": 11, "color": SUB}))
    write_lines(tf, lines, anchor=MSO_ANCHOR.MIDDLE)
    return shp


def title_bar(slide, title, *, section=None, page=None, subtitle=None):
    """상단: 섹션 칩 + 큰 제목 + 구분선 + 하단 푸터."""
    y_title = 0.85
    if section:
        chip = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, 0.72,
                         4.6, 0.85, fill=ACCENT_SOFT, line=None, radius=0.5)
        tf = chip.text_frame
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        write_lines(tf, [(section, {"size": 12.5, "bold": True,
                                    "color": ACCENT,
                                    "align": PP_ALIGN.CENTER})],
                    anchor=MSO_ANCHOR.MIDDLE)
        y_title = 1.75
    add_text(slide, MARGIN, y_title, USABLE_W, 1.5,
             [(title, {"size": 25, "bold": True, "color": INK})])
    if subtitle:
        add_text(slide, MARGIN, y_title + 1.25, USABLE_W, 0.7,
                 [(subtitle, {"size": 13.5, "color": SUB})])
    ln = slide.shapes.add_connector(
        1, Cm(MARGIN), Cm(3.70), Cm(SLIDE_W - MARGIN), Cm(3.70))
    ln.line.color.rgb = LINE
    ln.line.width = Pt(1.0)
    if page is not None:
        add_text(slide, SLIDE_W - 9.0, 18.30, 7.5, 0.5,
                 [(f"FA-Agent 온보딩 (온보딩_가이드.md 기반)  ·  {page}",
                   {"size": 8, "color": SUB, "align": PP_ALIGN.RIGHT})])


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ════════════════════════ 1. 표지 ════════════════════════
def s_cover(prs):
    slide = new_slide(prs)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.45,
              fill=ACCENT, line=None)
    add_text(slide, MARGIN, 4.6, USABLE_W, 1.2,
             [("전 구성원 온보딩 · 30분이면 충분합니다",
               {"size": 15, "bold": True, "color": ACCENT})])
    add_text(slide, MARGIN, 5.8, USABLE_W, 3.6, [
        ("FA-Agent 시작하기", {"size": 44, "bold": True, "color": INK,
                               "space_after": 14}),
        ("한국어로 말하면, AI 팀원 5명이 보고서를 만들어 드립니다",
         {"size": 20, "color": INK}),
    ])
    chips = ["코딩 몰라도 OK", "설치 없이 브라우저로 시작", "준비물은 계정 2개"]
    x = MARGIN
    for c in chips:
        w = 0.85 * len(c) * 0.55 + 1.6
        shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 10.6, w, 1.05,
                        fill=CARD_BG, line=LINE, radius=0.5)
        tf = shp.text_frame
        write_lines(tf, [(f"✓ {c}", {"size": 13, "bold": True, "color": INK,
                                     "align": PP_ALIGN.CENTER})],
                    anchor=MSO_ANCHOR.MIDDLE)
        x += w + 0.5
    add_text(slide, MARGIN, 16.7, USABLE_W, 0.8,
             [("FA기술혁신 Part  ·  2026.06.12  ·  상세 문서: 온보딩_가이드.md",
               {"size": 12, "color": SUB})])


# ════════════════════════ 2. 아젠다 ════════════════════════
def s_agenda(prs):
    slide = new_slide(prs)
    title_bar(slide, "오늘 다룰 내용 — 4가지만 기억하세요", page=2)
    items = [
        ("개념", "FA-Agent가 뭔가요?", "등장인물 3개 · AI 팀원 5명 · 할 수 있는 일"),
        ("준비", "처음 1회, 계정 2개", "GitHub + Claude · 약 30분 · 막히면 관리자에게"),
        ("사용법", "시작부터 결과 받기까지", "웹으로 4단계 · 첫 명령 · 요청 문장 치트시트"),
        ("규칙", "안전하게 쓰기", "추측 금지 · 출처 확인 · 최종 검토는 사람"),
    ]
    w = (USABLE_W - 3 * 0.7) / 4
    for i, (tag, t, b) in enumerate(items):
        x = MARGIN + i * (w + 0.7)
        num_badge(slide, x + 0.25, 4.5, i + 1)
        card(slide, x, 5.9, w, 8.2, "", [], bg=CARD_BG, border=LINE)
        add_text(slide, x + 0.35, 6.4, w - 0.7, 7.2, [
            (tag, {"size": 13, "bold": True, "color": ACCENT,
                   "space_after": 6}),
            (t, {"size": 16.5, "bold": True, "color": INK, "space_after": 10}),
            (b, {"size": 12, "color": SUB}),
        ])


# ════════════ 3. [개념] 등장인물 3개 ════════════
def s_actors(prs):
    slide = new_slide(prs)
    title_bar(slide, "등장인물은 딱 3개입니다", section="1. 개념", page=3,
              subtitle="어려운 용어는 비유로만 기억하면 됩니다")
    actors = [
        ("Claude  (클로드)", "말귀를 알아듣는 AI 직원",
         "한국어를 이해하고 글·문서·분석을\n해 주는 인공지능"),
        ("Claude Code  (클로드 코드)", "그 AI 직원의 사무실",
         "우리 팀 폴더(양식·데이터·자료)를 직접\n열어 보고, 파일을 만들어 주는 작업 환경"),
        ("GitHub  (깃허브)", "우리 팀의 공유 캐비닛",
         "양식·데이터·결과물 보관소.\n누가 뭘 바꿨는지 자동 기록 → 항상 복구 가능"),
    ]
    w = (USABLE_W - 2 * 0.8) / 3
    for i, (name, meta, desc) in enumerate(actors):
        x = MARGIN + i * (w + 0.8)
        shp = card(slide, x, 4.3, w, 7.6, "", [])
        tf = shp.text_frame
        write_lines(tf, [
            (name, {"size": 16, "bold": True, "color": INK, "space_after": 4}),
            (f"= {meta}", {"size": 14, "bold": True, "color": ACCENT,
                           "space_after": 12}),
            *[(ln, {"size": 12.5, "color": INK, "space_after": 3})
              for ln in desc.split("\n")],
        ])
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, 12.9, USABLE_W, 2.6,
              fill=ACCENT_SOFT, line=None, radius=0.12)
    add_text(slide, MARGIN + 0.6, 13.45, USABLE_W - 1.2, 1.6, [
        ("정리하면 :  캐비닛(GitHub) 속 우리 팀 자료로, 사무실(Claude Code)에서, "
         "AI 직원(Claude)이 일합니다.",
         {"size": 15, "bold": True, "color": INK, "space_after": 3}),
        ("여러분이 할 일은 그 직원에게 한국어로 업무 지시를 하는 것뿐입니다.",
         {"size": 15, "bold": True, "color": ACCENT}),
    ])


# ════════════ 4. [개념] AI 팀원 5명 ════════════
def s_team(prs):
    slide = new_slide(prs)
    title_bar(slide, "우리 팀에는 AI 팀원이 5명 있습니다", section="1. 개념",
              page=4, subtitle="누구를 부를지 고를 필요 없음 — 팀장 AI가 알아서 적임자에게 배분합니다")
    members = [
        ("데이터 담당", "KPI 진척 · 투자비(Capex) ·\n양산 일정(MRM) · R&D(DST) 조회"),
        ("운영 담당", "협력사 현황 · HR ·\n결재 대기 · 사내문서 검색"),
        ("리서치 담당", "외부 기술 동향 · 경쟁사 ·\n신규 과제 발굴"),
        ("문서 작성 담당", "회의록 · 주간/월간보고 등\n문서 양식 채움"),
        ("PPT 담당", "사내 PPT 양식 분석 ·\n빈칸 채워 완성본 생성"),
    ]
    w = (USABLE_W - 4 * 0.55) / 5
    for i, (role, desc) in enumerate(members):
        x = MARGIN + i * (w + 0.55)
        shp = card(slide, x, 4.5, w, 6.4, "", [])
        tf = shp.text_frame
        write_lines(tf, [
            (role, {"size": 14, "bold": True, "color": ACCENT,
                    "space_after": 8}),
            *[(ln, {"size": 11.5, "color": INK, "space_after": 3})
              for ln in desc.split("\n")],
        ])
    add_text(slide, MARGIN, 11.8, USABLE_W, 3.0, [
        ("폴더는 4개만 기억하세요", {"size": 15, "bold": True, "color": INK,
                                     "space_after": 8}),
        ("templates/ 양식 보관함 (원본 불변)      sample_data/ 숫자 데이터      "
         "references/ 참고 자료 (넣으면 자동 인식)      outputs/ 결과물 수령함",
         {"size": 13, "color": INK, "space_after": 6}),
        ("→ AI가 만든 완성본은 전부 outputs/ 폴더에 생깁니다",
         {"size": 13, "bold": True, "color": ACCENT}),
    ])


# ════════════ 5. [개념] 할 수 있는 일 ════════════
def s_usecases(prs):
    slide = new_slide(prs)
    title_bar(slide, "이런 일을 시킬 수 있습니다", section="1. 개념", page=5)
    cases = [
        ("회의록", "녹취록만 주면 회의록\nPPT/Word 자동 작성"),
        ("주간·월간 보고", "AI 팀원들이 숫자를 모아\n사내 양식대로 채움"),
        ("데이터 조회", "KPI·투자비·양산 일정 —\n답변에 출처(파일·행)가 함께"),
        ("기술 동향 조사", "외부 자료 조사 →\n신규 과제 후보 보고서"),
        ("파일 작업", "엑셀 정리 · PDF 표 추출 ·\nWord 문서 작성"),
        ("임원 보고 PPT", "사내 표준 색·폰트 적용된\n보고용 슬라이드 생성"),
    ]
    w = (USABLE_W - 2 * 0.8) / 3
    h = 4.6
    for i, (t, d) in enumerate(cases):
        r, c = divmod(i, 3)
        x = MARGIN + c * (w + 0.8)
        y = 4.3 + r * (h + 0.7)
        card(slide, x, y, w, h, t,
             [(ln, {"size": 12.5, "color": INK, "space_after": 3})
              for ln in d.split("\n")],
             title_size=15.5)
    add_text(slide, MARGIN, 14.6, USABLE_W, 1.2,
             [("이 모든 것을 명령어가 아니라 평소 말투의 한국어 한 문장으로 시킵니다.",
               {"size": 14.5, "bold": True, "color": ACCENT})])


# ════════════ 6. [준비] 계정 2개 ════════════
def s_setup(prs):
    slide = new_slide(prs)
    title_bar(slide, "준비물은 계정 2개뿐입니다 (처음 1회, 약 30분)",
              section="2. 준비", page=6,
              subtitle="이미 계정이 있으면 건너뛰면 됩니다 · 막히면 혼자 고민하지 말고 관리자에게")
    w = (USABLE_W - 0.9) / 2
    steps_a = [
        "① github.com 접속 → Sign up (회사 이메일)",
        "② 관리자에게 내 사용자 이름 전달, 초대 요청",
        "③ 초대 메일에서 Accept invitation 클릭",
        "④ github.com/mowonkang/FA-Agent-common 이 열리면 성공",
    ]
    steps_b = [
        "① claude.ai 접속 → 회사 이메일로 가입",
        "② Claude Code 권한이 포함된 시트(라이선스) 배정 요청",
        "    · 회사 계약 → 관리자에게  /  개인 → Pro 이상",
        "③ claude.ai 로그인이 되면 성공",
    ]
    card(slide, MARGIN, 4.4, w, 7.4, "계정 1.  GitHub  — 공유 캐비닛 출입증",
         [(s, {"size": 13, "color": INK, "space_after": 7}) for s in steps_a],
         title_size=16)
    card(slide, MARGIN + w + 0.9, 4.4, w, 7.4,
         "계정 2.  Claude  — AI 직원 이용권",
         [(s, {"size": 13, "color": INK, "space_after": 7}) for s in steps_b],
         title_size=16)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, 12.6, USABLE_W, 2.5,
              fill=CARD_BG, line=LINE, radius=0.12)
    add_text(slide, MARGIN + 0.6, 13.0, USABLE_W - 1.2, 1.8, [
        ("준비 완료 체크 :   □ github.com 로그인   □ 팀 저장소 페이지 열림   "
         "□ claude.ai 로그인   □ Claude Code 시트 확인",
         {"size": 13.5, "bold": True, "color": INK, "space_after": 4}),
        ("회사 라이선스 신청 절차·담당자 연락처는 가이드 문서의 기입란에 정리 예정 (자료 미확인)",
         {"size": 11, "color": SUB}),
    ])


# ════════════ 7. [준비] 안심 포인트 ════════════
def s_safety(prs):
    slide = new_slide(prs)
    title_bar(slide, "걱정 마세요 — 망가뜨릴 수 없게 되어 있습니다",
              section="2. 준비", page=7)
    points = [
        ("실수해도 복구됩니다", "GitHub는 모든 변경 이력을 자동\n보관합니다.",
         "잘못 눌러도 언제든 되돌릴 수 있습니다."),
        ("마음대로 바꾸지 않습니다", "AI는 파일을 만들거나 바꾸기 전에\n반드시 허락을 구합니다.",
         "이상하면 거부하면 그만입니다."),
        ("원본은 항상 그대로", "양식(templates/)과 자료(references/)\n원본은 수정되지 않습니다.",
         "결과물은 outputs/에만 생깁니다."),
    ]
    w = (USABLE_W - 2 * 0.8) / 3
    for i, (t, d1, d2) in enumerate(points):
        x = MARGIN + i * (w + 0.8)
        shp = card(slide, x, 4.4, w, 7.8, "", [])
        write_lines(shp.text_frame, [
            (f"✓ {t}", {"size": 15.5, "bold": True, "color": OK_GREEN,
                        "space_after": 10}),
            *[(ln, {"size": 12.5, "color": INK, "space_after": 3})
              for ln in d1.split("\n")],
            (d2, {"size": 12.5, "bold": True, "color": INK,
                  "space_before": 8}),
        ])
    add_text(slide, MARGIN, 13.3, USABLE_W, 2.4, [
        ("단, 보안은 사람의 몫 :  사외비 문서·개인정보를 붙여넣기 전에는 회사 AI "
         "사용 지침을 확인하세요.", {"size": 13.5, "bold": True, "color": INK,
                                     "space_after": 4}),
        ("현재 sample_data/ 는 시연용 샘플입니다 — 실데이터 반입 전 보안/IT 검토 필수.",
         {"size": 12, "color": SUB}),
    ])


# ════════════ 8. [사용법] 시작 4단계 ════════════
def s_start(prs):
    slide = new_slide(prs)
    title_bar(slide, "시작은 웹에서 — 설치 프로그램 0개", section="3. 사용법",
              page=8, subtitle="PC·노트북·휴대폰(Claude 앱) 어디서나 같은 방법입니다")
    steps = [
        ("claude.ai/code 접속", "Claude 계정으로 로그인"),
        ("GitHub 연결 승인", "최초 1회 — “접근 허용” 화면은\n정상 절차, Authorize 클릭"),
        ("저장소 선택", "FA-Agent-common 선택"),
        ("한국어로 지시", "채팅창에 평소 말투로\n업무를 말하면 끝"),
    ]
    w = (USABLE_W - 3 * 1.5) / 4
    y = 5.2
    for i, (t, d) in enumerate(steps):
        x = MARGIN + i * (w + 1.5)
        num_badge(slide, x + w / 2 - 0.55, y, i + 1, d=1.1)
        card(slide, x, y + 1.5, w, 5.6, "", [])
        add_text(slide, x + 0.3, y + 1.9, w - 0.6, 4.8, [
            (t, {"size": 14.5, "bold": True, "color": INK, "space_after": 7,
                 "align": PP_ALIGN.CENTER}),
            *[(ln, {"size": 11.5, "color": SUB, "space_after": 2,
                    "align": PP_ALIGN.CENTER}) for ln in d.split("\n")],
        ])
        if i < 3:
            add_text(slide, x + w + 0.15, y + 3.4, 1.2, 1.0,
                     [("→", {"size": 22, "bold": True, "color": ACCENT,
                             "align": PP_ALIGN.CENTER})])
    add_text(slide, MARGIN, 13.7, USABLE_W, 2.2, [
        ("매일 쓰는 파워유저는 PC 설치형도 가능 — 설치 명령 1줄, 절차는 "
         "온보딩_가이드.md 3장 참고 (공식 문서: code.claude.com/docs)",
         {"size": 12, "color": SUB}),
    ])


# ════════════ 9. [사용법] 첫 명령 ════════════
def s_first(prs):
    slide = new_slide(prs)
    title_bar(slide, "첫 명령 — 지금 그대로 따라해 보세요", section="3. 사용법",
              page=9, subtitle="팀 세팅이 정상인지 확인하는 공식 테스트 명령입니다")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + 2.5, 4.6,
              USABLE_W - 5.0, 3.4, fill=ACCENT_SOFT, line=ACCENT,
              line_w=1.5, radius=0.1)
    add_text(slide, MARGIN + 3.3, 5.15, USABLE_W - 6.6, 2.4, [
        ("“sample_data/kpi_샘플.md의 Q3 KPI를 회의록 양식으로 정리해줘.",
         {"size": 17, "bold": True, "color": INK, "space_after": 3}),
        ("  일시는 오늘, 참석자는 테스트 사용자.”",
         {"size": 17, "bold": True, "color": INK}),
    ])
    add_text(slide, MARGIN, 8.6, USABLE_W, 1.1,
             [("▼  잠시 후", {"size": 14, "bold": True, "color": ACCENT,
                              "align": PP_ALIGN.CENTER})])
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + 2.5, 9.8,
              USABLE_W - 5.0, 2.3, fill=CARD_BG, line=LINE, radius=0.1)
    add_text(slide, MARGIN + 3.3, 10.35, USABLE_W - 6.6, 1.4, [
        ("outputs/회의록_<오늘날짜>.md  파일이 생기면 성공 — 모든 것이 정상 동작하는 것입니다",
         {"size": 14.5, "bold": True, "color": OK_GREEN}),
    ])
    add_text(slide, MARGIN, 13.2, USABLE_W, 2.8, [
        ("결과물 받기", {"size": 14, "bold": True, "color": INK,
                         "space_after": 5}),
        ("웹 :  “결과 파일 보내줘” → 다운로드  /  “GitHub에 저장해줘” → 팀 보관      "
         "PC :  탐색기에서 outputs/ 폴더 열기",
         {"size": 13, "color": INK, "space_after": 4}),
        ("같은 양식을 같은 날 다시 만들면 파일명에 _v2, _v3 가 자동으로 붙습니다.",
         {"size": 11.5, "color": SUB}),
    ])


# ════════════ 10. [사용법] 화면 대응 ════════════
def s_screens(prs):
    slide = new_slide(prs)
    title_bar(slide, "이런 화면이 떠도 당황하지 마세요", section="3. 사용법",
              page=10)
    rows = [
        ("“~를 실행해도 될까요?”", "파일 생성·실행 전에 허락을 구하는 것",
         "내용 읽고 승인 — 이상하면 거부 후 “왜 필요한지 설명해줘”"),
        ("계획을 보여주며 승인 요청", "“이렇게 채우겠습니다” 하고 미리 검사받는 것",
         "방향이 맞으면 승인 — 아니면 “○○는 빼고 △△ 강조해줘”"),
        ("“자료 미확인” 표기", "근거를 못 찾았다는 정직한 보고 (지어내지 않음)",
         "자료 위치를 알려주거나, 그 칸은 사람이 채우면 됩니다"),
        ("AI 팀원들이 일하는 표시", "팀장이 팀원들에게 업무를 나눠준 상태",
         "기다리면 됩니다 (보통 수 분 이내)"),
    ]
    for x, w, h in ((MARGIN + 0.5, 9.6, "이런 화면이 뜨면"),
                    (MARGIN + 10.4, 9.7, "이런 뜻입니다"),
                    (MARGIN + 20.3, USABLE_W - 20.8, "이렇게 하세요")):
        add_text(slide, x, 4.00, w, 0.6,
                 [(h, {"size": 11.5, "bold": True, "color": SUB})])
    y = 4.75
    for q, mean, act in rows:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, y, USABLE_W,
                  2.55, fill=CARD_BG, line=LINE, radius=0.10)
        add_text(slide, MARGIN + 0.5, y + 0.35, 9.6, 1.9,
                 [(q, {"size": 13.5, "bold": True, "color": INK})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN + 10.4, y + 0.35, 9.7, 1.9,
                 [(mean, {"size": 12, "color": SUB})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN + 20.3, y + 0.35, USABLE_W - 20.8, 1.9,
                 [(f"→ {act}", {"size": 12, "bold": True, "color": ACCENT})],
                 anchor=MSO_ANCHOR.MIDDLE)
        y += 2.95


# ════════════ 11. [사용법] 치트시트 ════════════
def s_cheatsheet(prs):
    slide = new_slide(prs)
    title_bar(slide, "복사해서 쓰는 요청 문장 — 날짜·이름만 바꾸세요",
              section="3. 사용법", page=11)
    quotes = [
        ("이 녹취록으로 회의록 PPT 만들어줘", "녹취록 파일은 채팅창에 끌어다 놓기"),
        ("이번 주 주간보고 작성해줘", "월간은 “5월 월간보고 작성해줘. KPI 중심으로”"),
        ("KPI 진척 현황 알려줘", "투자비·양산 일정도 같은 식 — 출처가 함께 나옵니다"),
        ("FA 기술 동향 조사해서 신규 과제 후보 뽑아줘", "“AMR 동향 조사해줘”처럼 분야 지정도 가능"),
        ("이 PDF에서 표만 뽑아서 엑셀로 만들어줘", "엑셀 정리·Word 변환·사내 공지 초안도 한 문장"),
        ("임원 보고용 PPT 만들어줘. 메인 1장 + 보조 2장으로", "사내 표준 색·폰트가 자동 적용됩니다"),
    ]
    w = (USABLE_W - 0.8) / 2
    h = 2.85
    for i, (q, note) in enumerate(quotes):
        r, c = divmod(i, 2)
        bubble(slide, MARGIN + c * (w + 0.8), 4.1 + r * (h + 0.55), w, h,
               q, note)
    add_text(slide, MARGIN, 14.6, USABLE_W, 1.6, [
        ("“다시” · “더 짧게” · “어조를 공손하게” — 몇 번을 고쳐 시켜도 AI는 지치지 않습니다.   "
         "막히면 “네가 뭘 할 수 있는지 알려줘”라고 물어보세요.",
         {"size": 13.5, "bold": True, "color": INK}),
    ])


# ════════════ 12. [사용법] 잘 시키는 요령 ════════════
def s_tips(prs):
    slide = new_slide(prs)
    title_bar(slide, "잘 시키는 요령 — 신입사원에게 지시하듯 구체적으로",
              section="3. 사용법", page=12)
    pairs = [
        ("보고서 써줘", "다음주 월요일 그룹장 보고용 주간보고 써줘", "목적을 말하기"),
        ("KPI 알려줘", "Q3 KPI 진척률을 목표 대비로 알려줘", "기준·범위 주기"),
        ("정리해줘", "표 1개 + 요약 3줄로 정리해줘", "형식을 지정하기"),
    ]
    y = 4.4
    for bad, good, label in pairs:
        add_text(slide, MARGIN, y + 0.1, 4.4, 1.6,
                 [(label, {"size": 13, "bold": True, "color": ACCENT})],
                 anchor=MSO_ANCHOR.MIDDLE)
        shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + 4.8, y,
                        10.4, 1.85, fill=CARD_BG, line=LINE, radius=0.14)
        write_lines(shp.text_frame,
                    [(f"“{bad}”", {"size": 13, "color": SUB,
                                   "align": PP_ALIGN.CENTER})],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN + 15.4, y + 0.35, 1.4, 1.2,
                 [("→", {"size": 20, "bold": True, "color": ACCENT,
                         "align": PP_ALIGN.CENTER})])
        shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN + 17.0, y,
                        USABLE_W - 17.0, 1.85, fill=ACCENT_SOFT, line=ACCENT,
                        line_w=1.2, radius=0.14)
        write_lines(shp.text_frame,
                    [(f"“{good}”", {"size": 13.5, "bold": True, "color": INK,
                                    "align": PP_ALIGN.CENTER})],
                    anchor=MSO_ANCHOR.MIDDLE)
        y += 2.45
    add_text(slide, MARGIN, y + 0.5, USABLE_W, 3.2, [
        ("그 밖의 요령", {"size": 14, "bold": True, "color": INK,
                          "space_after": 6}),
        ("· 아는 값은 직접 주기 — “작성자는 김OO, 일시는 6/12 14시로 넣어줘”",
         {"size": 13, "color": INK, "space_after": 4}),
        ("· 한 번에 한 묶음 — 회의록 먼저, 끝나면 다음 요청",
         {"size": 13, "color": INK, "space_after": 4}),
        ("· 파일 주기 — 웹은 끌어다 놓기(드래그&드롭), PC는 폴더에 넣고 경로 말하기",
         {"size": 13, "color": INK}),
    ])


# ════════════ 13. [규칙] 4가지 ════════════
def s_rules(prs):
    slide = new_slide(prs)
    title_bar(slide, "꼭 지킬 것 4가지 — 시스템도 자동으로 지켜줍니다",
              section="4. 규칙", page=13)
    rules = [
        ("추측 금지", "AI는 자료가 없으면 지어내는 대신\n“자료 미확인”이라고 적습니다.",
         "사람도 그 칸을 임의로 채우지 않기"),
        ("출처 확인", "답변에는 (출처: 파일명 [행])이 붙고,\n출처 없는 보고서는 자동 차단됩니다.",
         "보고 전 원본 숫자와 대조하기"),
        ("원본 보존", "양식·참고자료 원본은 불변,\n결과물은 outputs/에만 생깁니다.",
         "원본을 직접 고치지 않기"),
        ("최종 책임은 사람", "AI 산출물은 어디까지나 초안입니다.",
         "결재·보고 전 검토는 작성자의 몫"),
    ]
    w = (USABLE_W - 3 * 0.7) / 4
    for i, (t, d, a) in enumerate(rules):
        x = MARGIN + i * (w + 0.7)
        num_badge(slide, x + 0.25, 4.3, i + 1)
        shp = card(slide, x, 5.7, w, 8.0, "", [])
        write_lines(shp.text_frame, [
            (t, {"size": 15.5, "bold": True, "color": INK, "space_after": 9}),
            *[(ln, {"size": 11.5, "color": INK, "space_after": 3})
              for ln in d.split("\n")],
            (f"→ {a}", {"size": 11.5, "bold": True, "color": ACCENT,
                        "space_before": 9}),
        ])
    add_text(slide, MARGIN, 14.3, USABLE_W, 1.4, [
        ("보안 :  사외비·개인정보 입력 전 회사 AI 지침 확인 · 외부 서비스 연동은 보안/IT 검토 후에만",
         {"size": 12.5, "bold": True, "color": INK}),
    ])


# ════════════ 14. [규칙] FAQ ════════════
def s_faq(prs):
    slide = new_slide(prs)
    title_bar(slide, "자주 묻는 질문 — 대부분 한 문장이면 해결됩니다",
              section="4. 규칙", page=14)
    rows = [
        ("영어로 써야 하나요? 명령어를 외워야 하나요?",
         "아니요. 전부 한국어 평소 말투로 — 외울 명령어는 없습니다."),
        ("“작업 완료가 차단되었습니다(missing_citations)”라고 떠요.",
         "출처 인용이 빠졌다는 뜻 → “누락된 출처를 추가해줘” 한마디면 해결."),
        ("어제 세션을 이어갔더니 AI 팀원이 사라졌어요.",
         "알려진 제한 → “팀원을 다시 생성해줘” 하면 같은 구성으로 복원."),
        ("PPT 일부 칸이 빈 채로 나와요.",
         "근거 없는 칸은 추측 대신 비워둔 것 → 값을 알려주고 다시 시키면 됩니다."),
        ("응답이 멈췄거나 이상해요.",
         "Esc(정지) 후 다시 지시 → 새 세션에서 “아까 하던 ○○ 이어서 해줘”."),
    ]
    y = 4.4
    for q, a in rows:
        add_text(slide, MARGIN, y, 14.6, 1.8,
                 [(f"Q.  {q}", {"size": 13, "bold": True, "color": INK})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, MARGIN + 15.2, y, USABLE_W - 15.2, 1.8,
                 [(f"A.  {a}", {"size": 13, "color": INK})],
                 anchor=MSO_ANCHOR.MIDDLE)
        ln = slide.shapes.add_connector(
            1, Cm(MARGIN), Cm(y + 2.0), Cm(SLIDE_W - MARGIN), Cm(y + 2.0))
        ln.line.color.rgb = LINE
        ln.line.width = Pt(0.75)
        y += 2.35
    add_text(slide, MARGIN, y + 0.2, USABLE_W, 1.0,
             [("그 밖의 모든 문제 → 관리자에게 (연락처는 가이드 기입란 — 자료 미확인) · "
               "공식 문서(한국어) code.claude.com/docs",
               {"size": 12, "color": SUB})])


# ════════════ 15. 마무리 ════════════
def s_closing(prs):
    slide = new_slide(prs)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.45,
              fill=ACCENT, line=None)
    add_text(slide, MARGIN, 3.6, USABLE_W, 2.6, [
        ("오늘 바로 시작해 보세요", {"size": 30, "bold": True, "color": INK,
                                     "space_after": 8}),
        ("이 3단계면 오늘 첫 보고서를 받을 수 있습니다",
         {"size": 15, "color": SUB}),
    ])
    steps = [
        ("claude.ai/code 로그인", "저장소 FA-Agent-common 선택"),
        ("첫 명령 복사·붙여넣기", "9쪽의 동작 확인 명령 그대로"),
        ("outputs/ 에서 결과 확인", "이후엔 치트시트(11쪽)로 실전 적용"),
    ]
    w = (USABLE_W - 2 * 1.6) / 3
    for i, (t, d) in enumerate(steps):
        x = MARGIN + i * (w + 1.6)
        num_badge(slide, x + w / 2 - 0.6, 7.2, i + 1, d=1.2)
        card(slide, x, 8.8, w, 3.6, "", [])
        add_text(slide, x + 0.3, 9.3, w - 0.6, 2.8, [
            (t, {"size": 15, "bold": True, "color": INK, "space_after": 6,
                 "align": PP_ALIGN.CENTER}),
            (d, {"size": 12, "color": SUB, "align": PP_ALIGN.CENTER}),
        ])
        if i < 2:
            add_text(slide, x + w + 0.25, 10.0, 1.2, 1.0,
                     [("→", {"size": 22, "bold": True, "color": ACCENT,
                             "align": PP_ALIGN.CENTER})])
    add_text(slide, MARGIN, 14.0, USABLE_W, 2.8, [
        ("더 자세히 :  저장소의 「온보딩_가이드.md」 (배포용 Word 사본 "
         "outputs/FA-Agent_온보딩_가이드_*.docx)",
         {"size": 13, "bold": True, "color": INK, "space_after": 5}),
        ("문의 :  관리자 [담당자 연락처 기입란 — 자료 미확인]   ·   공식 문서(한국어) "
         "code.claude.com/docs",
         {"size": 13, "color": INK, "space_after": 8}),
        ("본 자료 출처 : 온보딩_가이드.md · CLAUDE.md · README.md [106~113·149~158·"
         "186~196행] — 2026.06.12 기준",
         {"size": 9.5, "color": SUB}),
    ])


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    for fn in (s_cover, s_agenda, s_actors, s_team, s_usecases, s_setup,
               s_safety, s_start, s_first, s_screens, s_cheatsheet, s_tips,
               s_rules, s_faq, s_closing):
        fn(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}  (slides: {len(prs.slides)})")


if __name__ == "__main__":
    build()
