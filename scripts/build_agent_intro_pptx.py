"""5개 에이전트 소개 PPT 빌더.

슬라이드 구성:
  1 (메인, 분류형) : 5개 에이전트 개요
  2 (보조)        : data-teammate / ops-teammate 상세
  3 (보조)        : tech-research-teammate / document-writer / ppt-writer 상세
  4 (보조)        : 에이전트 협업 패턴

출력: outputs/에이전트_소개_2026-05-29.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_layout import (  # noqa: E402
    render_classified_slide,
    add_label_band, add_rect, add_text, fill_table, title_block, add_footnote,
    SLIDE_W_CM, SLIDE_H_CM,
    BLACK, BLUE, MID_GRAY, DIM_GRAY, GRAY, WHITE, CHARCOAL, LIGHT_GRAY,
    SZ_HEAD, SZ_TITLE, SZ_BAND, SZ_SECTION, SZ_BODY, SZ_SUB, SZ_FOOT,
    FONT_BODY, write_lines,
)

OUT_PATH = Path("outputs/에이전트_소개_2026-05-29.pptx")
BYLINE = ["FA 기술담당", "(2026.05.29)"]
FOOT_MAIN = "* 본 자료는 레포지토리 .claude/agents/ 정의 기반으로 작성되었습니다."


# ─── 슬라이드 1: 메인 (분류형) ─────────────────────────────────────────────
def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="FA 자동화 에이전트 팀 소개",
        byline=BYLINE,
        head_message=[
            ("5개 전문 에이전트가 데이터 수집·운영 검색·기술 리서치·문서 작성을 분업하고 "
             "협업하여 정형 보고서와 PPT를 자동 생성합니다.",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "data\n-teammate",
                "session_message": "KPI·Capex·MRM·DST 숫자와 일정 데이터 조회",
                "content": [
                    (" • 역할 : 사내 정량 데이터(KPI 진척, 투자비, 양산 로드맵, R&D 과제 진척) 조회",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사용 시점 : KPI 진척 확인, Capex 검토, 양산·R&D 일정 조회 시",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 도구 : Read · Grep · Glob · Bash",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "ops\n-teammate",
                "session_message": "협력사·HR·결재 큐·사내문서 검색",
                "content": [
                    (" • 역할 : 협력사 현황, HR(출장/교육/포상), 승인 대기, 사내 PDF/PPT 검색",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사용 시점 : 협력사 현황 파악, HR 활동 조회, 사내문서 검색 시",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 도구 : Read · Grep · Glob · Bash",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "tech\n-research\n-teammate",
                "session_message": "외부 기술 동향·경쟁사·기술 평가 (FA 6분야)",
                "content": [
                    (" • 역할 : AMR/협동로봇/디지털트윈/스마트물류/AI/피지컬AI 동향·경쟁사 리서치",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사용 시점 : 신규 과제 발굴, KPI 재설정, 로드맵 검토 시",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 도구 : Read · Grep · Glob · Bash · Write · WebSearch · WebFetch",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "document\n-writer",
                "session_message": "markdown(.md) 양식 빈칸 채움 → 완성 보고서 생성",
                "content": [
                    (" • 역할 : templates/*.md 양식의 placeholder를 데이터로 채워 보고서 완성",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사용 시점 : 회의록·주간보고·월간보고 작성 요청 시",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 도구 : Read · Write · Edit · Glob · Grep",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "ppt\n-writer",
                "session_message": ".pptx 양식 분석 → 빈칸 추론 → 데이터 수집 → PPTX 생성",
                "content": [
                    (" • 역할 : templates/*.pptx 양식 구조 추출 후 결정론적 스크립트로 채움",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사용 시점 : PPT 작성·양식 채움 요청 시",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 도구 : Read · Write · Edit · Glob · Grep · Bash",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
        ],
        footnote=FOOT_MAIN,
    )


# ─── 슬라이드 2: 보조 — data / ops 상세 ────────────────────────────────────
def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="데이터·운영 에이전트 상세",
        byline=BYLINE,
        head_message=[
            ("data-teammate는 사내 정량 데이터를, ops-teammate는 협력사·HR·사내문서를 전담합니다.",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "data\n-teammate\n상세",
                "session_message": "KPI·Roadmap / Capex·면적·구매 / MRM 양산 로드맵 / DST R&D 미래기술",
                "content": [
                    (" • KPI·Roadmap — 분기별 KPI 진척률, 목표 대비 실적, 달성률 계산",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • Capex·면적·구매 — 연간 투자비 현황, 설비 면적, 구매 계획 조회",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • MRM 양산 로드맵 — 양산 일정·가동률·수율 데이터 조회",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • DST R&D 미래기술 — R&D 과제 진척, TRL 단계, 마일스톤 확인",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["담당 영역", "주요 데이터 소스", "활용 보고서"],
                    "rows": [
                        ("KPI·Roadmap", "sample_data/kpi_tracker.csv", "주간/월간 보고"),
                        ("Capex·구매", "sample_data/capex_plan.csv", "투자비 현황"),
                        ("MRM 양산", "sample_data/mrm_schedule.md", "양산 로드맵"),
                        ("DST R&D", "sample_data/dst_progress.md", "R&D 진척 보고"),
                    ],
                    "col_w": [4.5, 8.5, 10.8],
                    "row_h": 0.85,
                },
            },
            {
                "label": "ops\n-teammate\n상세",
                "session_message": "협력사 풀/벤더 / HR 출장·교육·포상 / 결재 큐 / 사내문서 검색",
                "content": [
                    (" • 협력사 풀/벤더 — 협력사 현황, 계약 상태, 벤더 등록 정보 검색",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • HR 출장·교육·포상 — 출장 신청 현황, 교육 이수, 포상 대상 조회",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 결재 큐 — 승인 대기 문서 목록 및 결재선 확인",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사내문서 검색 — references/ 내 PDF·PPT·회의록 전문 검색",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["담당 영역", "주요 데이터 소스", "활용 보고서"],
                    "rows": [
                        ("협력사·벤더", "sample_data/vendor_list.csv", "협력사 현황"),
                        ("HR 활동", "sample_data/hr_activities.md", "HR 보고"),
                        ("결재 큐", "sample_data/approval_queue.md", "결재 현황"),
                        ("사내문서", "references/**/*.md, *.pdf", "회의록·정책"),
                    ],
                    "col_w": [4.5, 8.5, 10.8],
                    "row_h": 0.85,
                },
            },
        ],
        footnote=FOOT_MAIN,
    )


# ─── 슬라이드 3: 보조 — tech-research / document-writer / ppt-writer 상세 ──
def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="리서치·작성 에이전트 상세",
        byline=BYLINE,
        head_message=[
            ("tech-research-teammate는 외부 기술·경쟁사 정보를, document-writer·ppt-writer는 "
             "양식 기반 문서를 생성합니다.",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "tech\n-research\n-teammate\n상세",
                "session_message": "글로벌 기술 동향 / 경쟁사(CATL·삼성SDI·BYD·Tesla) / TRL 평가",
                "content": [
                    (" • FA 6분야 : AMR/AGV · 협동로봇 · 디지털트윈 · 스마트물류 · AI · 피지컬AI",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 글로벌 기술 동향 — 최신 논문·특허·산업 동향 WebSearch/Fetch로 수집",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 경쟁사 분석 — CATL·삼성SDI·BYD·Tesla 전략·기술 비교",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • TRL·적용가능성·리스크 평가 — 기술 성숙도 등급 및 도입 리스크 산정",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" ※ 수집 자료는 references/기술자료/ · references/경쟁사/ 에 자동 아카이빙",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "document\n-writer\n상세",
                "session_message": "회의록 / 주간보고 / 월간보고 (templates/*.md 양식 기반)",
                "content": [
                    (" • templates/*.md 양식의 {{ placeholder }} 를 수집 데이터로 채워 완성",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 녹취록 입력 감지 시 → meeting-minutes 스킬로 자동 위임",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • outputs/*.md 생성 후 출처 인용 자동 검증 (TaskCompleted hook)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 부족 데이터는 data-teammate / ops-teammate 에 직접 요청",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "ppt\n-writer\n상세",
                "session_message": "주간/월간 보고 PPT / 회의록 PPT (templates/*.pptx 양식 기반)",
                "content": [
                    (" • ppt_extract.py 로 양식 구조 추출 → {{ key }} · 빈 도형 · 표 셀 추론",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • values.json 작성 후 ppt_fill.py 로 텍스트 치환 (바이너리 직접 수정 금지)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 이미지·SmartArt·차트 데이터는 unfilled 처리 후 사용자에게 안내",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 결과물은 outputs/<이름>_<YYYY-MM-DD>.pptx 로 저장",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
        ],
        footnote=FOOT_MAIN,
    )


# ─── 슬라이드 4: 보조 — 협업 패턴 ─────────────────────────────────────────
def build_slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="에이전트 협업 패턴",
        byline=BYLINE,
        head_message=[
            ("유스케이스별 에이전트 조합으로 3~5명 팀을 구성하여 병렬 수집 → 위임 작성 패턴으로 "
             "보고서를 생성합니다.",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "정형\n보고서",
                "session_message": "주간·월간·회의록 보고 — 3명 구성",
                "content": [
                    (" • 구성 : data-teammate + ops-teammate + (document-writer 또는 ppt-writer)",
                     {"size": SZ_BODY, "bold": True, "color": BLACK}),
                    (" • 흐름 : 리더가 data·ops를 병렬 스폰 → 수집 결과를 writer에 전달 → 양식 채움",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 산출물 : outputs/주간보고_YYYY-MM-DD.md / .pptx",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "신규 과제\n발굴",
                "session_message": "FA 기술 동향 조사 → 과제 후보 도출 — 4명 구성",
                "content": [
                    (" • 구성 : tech-research-teammate + data-teammate + ops-teammate + writer",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 흐름 : tech-research가 6분야 동향 수집 → data·ops가 내부 현황 조회 → writer 통합",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 산출물 : outputs/FA과제후보_YYYY-MM-DD.docx (fa-task-discovery 스킬)",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "경쟁사\n분석",
                "session_message": "경쟁사 분석 → KPI 재설정 — 3~4명 구성",
                "content": [
                    (" • 구성 : tech-research-teammate + data-teammate (+ ops-teammate 선택) + writer",
                     {"size": SZ_BODY, "bold": True, "color": BLACK}),
                    (" • 흐름 : tech-research가 경쟁사(CATL·BYD·Tesla) 정보 수집 → data가 내부 KPI 비교",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 산출물 : outputs/경쟁사분석_YYYY-MM-DD.pptx",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "팀 운영\n원칙",
                "session_message": "권장 규모·분해·병렬 수집·위임 작성·품질 게이트",
                "table": {
                    "headers": ["원칙", "내용"],
                    "rows": [
                        ("팀 규모", "3~5명 / 팀원당 5~6개 작업"),
                        ("병렬 수집", "data·ops·tech-research 동시 스폰"),
                        ("Plan-mode", "writer는 outputs/ 쓰기 전 계획 승인"),
                        ("품질 게이트", "출처 인용 자동 검증 (TaskCompleted hook)"),
                        ("데이터 소유", "수치→data / 운영→ops / 리서치→tech-research"),
                    ],
                    "col_w": [4.5, 19.3],
                    "row_h": 0.85,
                },
            },
        ],
        footnote=FOOT_MAIN,
    )


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(str(OUT_PATH))
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
