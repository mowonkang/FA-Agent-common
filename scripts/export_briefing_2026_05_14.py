"""Convert outputs/FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.md to .docx and .pptx."""
from pathlib import Path
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PPt, Cm as PCm
from pptx.dml.color import RGBColor as PColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
MD = ROOT / "outputs" / "FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.md"
DOCX = ROOT / "outputs" / "FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.docx"
PPTX = ROOT / "outputs" / "FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.pptx"

LGES_RED = RGBColor(0xA5, 0x00, 0x34)
LGES_GRAY = RGBColor(0x55, 0x55, 0x55)
PLGES_RED = PColor(0xA5, 0x00, 0x34)
PLGES_DARK = PColor(0x22, 0x22, 0x22)


def build_docx(text: str) -> None:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "맑은 고딕"
    style.font.size = Pt(10.5)

    lines = text.splitlines()
    i = 0
    in_table = False
    table_rows: list[list[str]] = []

    def flush_table():
        nonlocal table_rows
        if not table_rows:
            return
        cols = len(table_rows[0])
        tbl = doc.add_table(rows=len(table_rows), cols=cols)
        tbl.style = "Light Grid Accent 1"
        for r, row in enumerate(table_rows):
            for c, cell_text in enumerate(row):
                cell = tbl.rows[r].cells[c]
                cell.text = ""
                p = cell.paragraphs[0]
                run = p.add_run(cell_text)
                if r == 0:
                    run.bold = True
                run.font.size = Pt(9)
        doc.add_paragraph("")
        table_rows = []

    while i < len(lines):
        line = lines[i].rstrip()

        if line.startswith("|") and line.endswith("|") and "---" not in line:
            cells = [c.strip() for c in line.strip("|").split("|")]
            table_rows.append(cells)
            in_table = True
            i += 1
            continue
        elif in_table and line.startswith("|"):
            i += 1
            continue
        else:
            if in_table:
                flush_table()
                in_table = False

        if line.startswith("# "):
            h = doc.add_heading(line[2:].strip(), level=0)
            for run in h.runs:
                run.font.color.rgb = LGES_RED
        elif line.startswith("## "):
            h = doc.add_heading(line[3:].strip(), level=1)
            for run in h.runs:
                run.font.color.rgb = LGES_RED
        elif line.startswith("### "):
            h = doc.add_heading(line[4:].strip(), level=2)
            for run in h.runs:
                run.font.color.rgb = LGES_GRAY
        elif line.startswith("- "):
            p = doc.add_paragraph(line[2:].strip(), style="List Bullet")
        elif line.startswith("---"):
            doc.add_paragraph("―" * 30).alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif line.strip() == "":
            doc.add_paragraph("")
        else:
            doc.add_paragraph(line)
        i += 1

    if in_table:
        flush_table()

    doc.save(DOCX)


def add_pptx_slide(prs, title, bullets, footer="LGES FA 기술혁신 Agent팀 | 2026-05-14"):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title bar
    from pptx.shapes.autoshape import Shape
    left, top, width, height = Inches(0), Inches(0), prs.slide_width, Inches(0.7)
    bar = slide.shapes.add_shape(1, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PLGES_RED
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PPt(22)
    p.font.bold = True
    p.font.color.rgb = PColor(0xFF, 0xFF, 0xFF)
    p.font.name = "맑은 고딕"

    # Body bullets
    body = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), prs.slide_width - Inches(1.0), Inches(6.0))
    btf = body.text_frame
    btf.word_wrap = True
    for idx, b in enumerate(bullets):
        if idx == 0:
            para = btf.paragraphs[0]
        else:
            para = btf.add_paragraph()
        # Support nested levels via leading spaces
        level = 0
        text = b
        while text.startswith("  "):
            level += 1
            text = text[2:]
        para.text = ("• " if level == 0 else "– ") + text
        para.level = level
        para.font.size = PPt(16 if level == 0 else 13)
        para.font.name = "맑은 고딕"
        para.font.color.rgb = PLGES_DARK
        para.space_after = PPt(6)

    # Footer
    foot = slide.shapes.add_textbox(Inches(0.3), prs.slide_height - Inches(0.4), prs.slide_width - Inches(0.6), Inches(0.3))
    ftf = foot.text_frame
    fp = ftf.paragraphs[0]
    fp.text = footer
    fp.font.size = PPt(9)
    fp.font.color.rgb = PColor(0x88, 0x88, 0x88)
    fp.font.name = "맑은 고딕"
    fp.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Red bar full
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.5), prs.slide_width, Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PLGES_RED
    bar.line.fill.background()

    title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), prs.slide_width - Inches(1.0), Inches(1.5))
    tt = title_tb.text_frame
    tt.word_wrap = True
    p = tt.paragraphs[0]
    p.text = "FA / AI / AX / Physical AI / 경쟁사"
    p.font.size = PPt(36)
    p.font.bold = True
    p.font.color.rgb = PColor(0xFF, 0xFF, 0xFF)
    p.font.name = "맑은 고딕"
    p2 = tt.add_paragraph()
    p2.text = "통합 동향 보고서"
    p2.font.size = PPt(30)
    p2.font.bold = True
    p2.font.color.rgb = PColor(0xFF, 0xFF, 0xFF)
    p2.font.name = "맑은 고딕"

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), prs.slide_width - Inches(1.0), Inches(1.5))
    st = sub.text_frame
    sp = st.paragraphs[0]
    sp.text = "기준일: 2026-05-14"
    sp.font.size = PPt(18)
    sp.font.color.rgb = PLGES_DARK
    sp.font.name = "맑은 고딕"
    sp2 = st.add_paragraph()
    sp2.text = "작성: FA 기술혁신 Agent팀 (tech-research × 4 병렬 수집)"
    sp2.font.size = PPt(14)
    sp2.font.color.rgb = PColor(0x55, 0x55, 0x55)
    sp2.font.name = "맑은 고딕"


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_pptx_slide(prs, "0. 핵심 요약 — 2026 상반기 4축 동시 변곡점", [
        "FA 4대 분야: AMR 멀티벤더 오케스트레이션 표준화, 협동로봇 ISO 10218:2025 개정, Siemens × NVIDIA Omniverse DSX GA",
        "AI / AX: Copilot → Agentic AI 전환, Industrial AI OS 표준화, ViT 기반 EV 배터리 검사 80m/min · 9.5ms 추론",
        "Physical AI: Figure × BMW, Tesla Optimus(Fremont 7~8월 양산), Apptronik/Agility 자동차·물류 라인 진입",
        "  Foundation Model 3축: GR00T N1.7 + π0.5 + Cosmos WFM",
        "  배터리 라인 직접 도입 사례: 자료 미확인 → LGES 선제 PoC 차별화 기회",
        "경쟁사: CATL Lighthouse(생산성 +17%, 결함 -99%), BYD 시안 자율도 97%, 삼성SDI 디지털트윈 무투자 capa 증대",
        "  SK온 × NVIDIA AI Factory(5만+ GPU), Panasonic × Capgemini Agentic AI, Lyten의 Northvolt 인수 후 2H26 재가동",
    ])

    add_pptx_slide(prs, "1. FA 4대 분야 (AMR / 협동로봇 / 디지털트윈 / 스마트물류)", [
        "AMR/AGV — 멀티벤더 오케스트레이션 표준화, VDA 5050 의무화 (시장 USD 64.5B, CAGR 13.6%)",
        "  LogiMAT 2026: Allient KinetiMax 휠허브 일체형 드라이브 — 차세대 경량 AMR 변화",
        "협동로봇 — ISO 10218-1/-2:2025 전면 개정, 'Cobot' → 'Collaborative Application'",
        "  Class I/II 분류 도입, 사이버보안 요건 신설, ANSI/A3 R15.06 동조 / 2025 글로벌 매출 $2.2B",
        "디지털 트윈 — Siemens × NVIDIA Omniverse DSX Blueprint GA",
        "  Erlangen 전자공장 첫 AI 적응형 제조 블루프린트, FREYR 기가팩토리 화학 시뮬레이션",
        "스마트 물류 — WMS가 '조정 레이어'로 진화, AI 에이전트 + AMR 실시간 오케스트레이션",
        "LGES — RIST 공동연구로 오창공장 안전·정비 디지털화 (승인 1~2일 → 수시간)",
    ])

    add_pptx_slide(prs, "2. AI / AX — 제조 AI · Agentic · Copilot", [
        "Siemens Industrial Copilot 9개 SW 확장 + NVIDIA 'Industrial AI OS' 공동 개발 (CES 2026)",
        "NVIDIA Omniverse, Physical AI OS로 확장 — Samsung, SK하이닉스, HD현대, FANUC, TSMC 가속 (GTC 2026)",
        "Microsoft Azure Foundry + Copilot의 MES/SCADA 통합 — Eyelit MCP 서버 시연 (2026-05-04)",
        "Agentic Manufacturing 본격화 — Accenture·Avanade 공동 솔루션, '2026 Inflection Point'",
        "ABB Robotics × Omniverse Industrial Physical AI — Sim2Real 학습 통합 플랫폼",
        "Schaeffler-Microsoft Factory Operations Agent — LLM + Vision 결합 운영 Agent",
        "EV 배터리 Vision Transformer — 1.5m 폭 호일 80m/min · 9.5ms 추론 · 정밀도 88%/재현율 84%",
    ])

    add_pptx_slide(prs, "3. Physical AI — 휴머노이드 · Foundation Model · World Model", [
        "Figure 02 × BMW Spartanburg — X3 30,000대 생산 기여, 배치 정확도 99%/shift, 2026-02 Leipzig 유럽 첫 배치",
        "Tesla Optimus — Fremont 7~8월 양산(연 100만대), Texas Gen2 연 1,000만대 목표",
        "Apptronik Apollo — Mercedes·GXO·Jabil PoC, $520M Series A (밸류 $5.5B)",
        "Agility Digit × Toyota — Canada Woodstock RAV4 라인 RaaS 배치",
        "Foundation Model 3축: NVIDIA GR00T N1.7 (오픈 VLA, +40% 성능) / Physical Intelligence π0.5 / Cosmos WFM (다운로드 200만+)",
        "Unitree — 2026년 humanoid 출하 20,000대 목표(전년 5,500대 대비 4배), G1 $16K",
        "한국 — Rainbow(시총 13조), HD현대(조선 용접), 두산(AI/humanoid R&D 강화). 배터리 적용은 공백",
        "1X NEO — Hayward 공장 가동(연 1만대), $20K + $499/월 구독 / 가정→산업 모델 전이",
    ])

    add_pptx_slide(prs, "4. 경쟁사 동향 — CATL · BYD · 삼성SDI · SK온", [
        "CATL — Lighthouse 등대공장 생산성 +17%, 결함률 -99% / 헝가리 73억유로 + Stellantis 합작 + Ford 라이선스",
        "  2026년부터 '사람 없는 공장' 단계적 도입 발표",
        "BYD — 시안 공장 자율도 ~97% (AI 로보틱스 + AGV + 지능형 창고)",
        "  칭하이 신공장 결함 -40%, 수명 +20% / Lead Intelligent와 전고체 라인당 5~8 GWh",
        "삼성SDI — 헝가리 디지털트윈 화성공정 가상공장으로 무투자 capa 증대",
        "  2026 하반기 가동률 70%+ 목표, 3공장 착공 임박, LFP 전환 병행",
        "SK온 — 현대 조지아 합작 가동, 안성 배터리캠퍼스 2026년말 완공 (AI 예측 + 자동 검증)",
        "  NVIDIA × SK그룹 AI Factory 5만+ GPU, 1단계 2027말",
    ])

    add_pptx_slide(prs, "4-2. 경쟁사 동향 — Panasonic · Lyten/Northvolt", [
        "Panasonic Energy — De Soto(캔자스) 32 GWh 가동, 노동절감 라인으로 네바다 대비 생산성 +20%",
        "  지열 HVAC + AI 에너지 관리로 CO2 -40%",
        "  Capgemini 협업 Agentic AI/AIP 플랫폼 — 미국 IRA 권역 선례",
        "Lyten — Northvolt 인수 완료(2026-02-27), Skellefteå/Västerås 즉시 재가동 추진",
        "  2H26 상용 셀 출하 목표 (NMC + 리튬-황 R&D 병행)",
        "  Lyten Industrial Hub: 배터리·데이터센터 코로케이션 (제조-AI 융합 모델)",
        "위협 평가: 유럽 capa 공백 일부 회복 → LGES 유럽 시장점유율 압박 재개 가능성 (강도: 중)",
    ])

    add_pptx_slide(prs, "5. LGES 통합 액션 매트릭스 (10건)", [
        "[단기 Q3'26] MES/PLM Copilot PoC — Recipe 작성 시간 단축",
        "[단기 Q3'26] 전극 검사 모델 CNN → ViT 업그레이드 (80m/min · 9.5ms)",
        "[단기] AMR 신규 도입 시 VDA 5050 의무 사양화",
        "[단기] 협동로봇 RA·구매 사양서 ISO 10218:2025 기준 재작성",
        "[중기 Q4'26+] Omniverse DSX 기반 신규 라인 가상 시운전 표준화",
        "[중기 2H'26] Rainbow/두산과 셀·모듈 핸들링 humanoid PoC MOU",
        "[중기] Isaac SIM + Cosmos 합성데이터 파이프라인 PoC",
        "[단기] BYD 자율도 97% / 삼성SDI 디지털트윈을 26FA KPI 상한 벤치마크화",
        "[중기] 미국 라인 Agentic AI (Panasonic-Capgemini 사례) 검토",
        "[상시] Lyten/Northvolt 재가동 모니터링 체계 운영",
    ])

    add_pptx_slide(prs, "6. 리스크 · 자료 미확인 사항", [
        "AMR VDA 5050 표준 KPI 매핑 자료 미확인",
        "CATL '사람 없는 공장' 단계별 일정 — 공식 발표 자료 미확인",
        "SK온 단독 자동화 KPI(수율·자동화율) — 2026 1Q~5월 공개 자료 미확인 (그룹·파트너 자료로 대체)",
        "배터리 라인 humanoid 직접 도입 사례 — 자동차/물류 우회 벤치마크 필요",
        "ISO 10218:2025 전환 비용 — 협동로봇 사양·RA 문서 전면 재작성 부담",
        "Omniverse DSX 종속성 — NVIDIA GPU·OpenUSD 종속, CAPEX·인력 재교육 부담",
    ])

    add_pptx_slide(prs, "7. 후속 권장 작업", [
        "data-teammate — 사내 협동로봇 KPI vs 글로벌 $2.2B 시장의 LGES 점유율 추정 매핑",
        "ops-teammate — Universal Robots / ABB GoFa / FANUC CRX / Rainbow / 두산 협력사 풀 등록 여부",
        "tech-research-teammate — VDA 5050, Cosmos WFM, GR00T N1.7 라이선스 조건 deep-dive",
        "CATL·BYD 등대공장 KPI 정량 비교 차트 작성 (차후 보고서)",
        "출처: outputs/FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.md (64+ 외부 출처)",
    ])

    prs.save(PPTX)


def main() -> None:
    text = MD.read_text(encoding="utf-8")
    build_docx(text)
    build_pptx()
    print(f"DOCX: {DOCX}")
    print(f"PPTX: {PPTX}")


if __name__ == "__main__":
    main()
