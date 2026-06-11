"""채울 값 JSON 을 받아 원본 PPTX 에 적용해 새 PPTX 로 저장한다.

사용법:
    python scripts/ppt_fill.py \
        --template templates/<이름>.pptx \
        --values   outputs/.cache/<이름>.values.json \
        --out      outputs/<이름>_<YYYY-MM-DD>.pptx \
        [--strict]

원본 파일은 절대 덮어쓰지 않는다. 슬라이드 추가/삭제·도형 추가/이동·폰트 변경 금지.
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write(
        "python-pptx 가 설치되어 있지 않습니다. 먼저 다음을 실행하세요:\n"
        "    pip install -r requirements.txt\n"
    )
    sys.exit(2)

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_common import (  # noqa: E402
    PLACEHOLDER_PATTERN,
    is_blank,
    iter_shapes_recursive,
    parse_shape_id,
)

VALID_MODES = {
    "replace_placeholder",
    "replace_text",
    "replace_empty",
    "set_cell_text",
    "append_paragraph",
}


def _find_shape(prs, shape_id_str: str):
    slide_idx, shape_idx = parse_shape_id(shape_id_str)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    shapes = list(slide.shapes)
    if shape_idx >= len(shapes):
        return None
    return shapes[shape_idx]


def _replace_text_keep_first_run(text_frame, new_text: str) -> None:
    """첫 paragraph 의 첫 run 에 새 텍스트를 몰아넣고 나머지 run/문단을 비운다.

    스타일은 첫 run 기준으로 유지된다.
    """
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        text_frame.text = new_text
        return
    first_p = paragraphs[0]
    first_runs = list(first_p.runs)
    if not first_runs:
        first_p.text = new_text
    else:
        first_runs[0].text = new_text
        for r in first_runs[1:]:
            r.text = ""
    for p in paragraphs[1:]:
        for r in p.runs:
            r.text = ""


def _replace_placeholder_in_text_frame(text_frame, key: str, value: str) -> bool:
    """텍스트 프레임 안에서 `{{ key }}` 만 정확히 치환. run 경계가 갈라져도 처리."""
    token_pattern = PLACEHOLDER_PATTERN
    changed = False
    for paragraph in text_frame.paragraphs:
        full_text = "".join(r.text for r in paragraph.runs)
        if not full_text:
            continue

        def _sub(m, _key=key, _value=value):
            nonlocal changed
            if m.group(1) == _key:
                changed = True
                return _value
            return m.group(0)

        new_text = token_pattern.sub(_sub, full_text)
        if new_text != full_text:
            runs = list(paragraph.runs)
            if runs:
                runs[0].text = new_text
                for r in runs[1:]:
                    r.text = ""
    return changed


def _apply_shape_text(shape, fill: dict[str, Any], warnings: list[str]) -> bool:
    if not shape.has_text_frame:
        warnings.append(f"{fill['fill_id']}: shape 에 text_frame 없음")
        return False
    tf = shape.text_frame
    mode = fill["mode"]
    value = str(fill["value"])
    if mode == "replace_placeholder":
        key = fill.get("explicit_placeholder")
        if not key:
            warnings.append(f"{fill['fill_id']}: replace_placeholder 모드에 explicit_placeholder 없음")
            return False
        return _replace_placeholder_in_text_frame(tf, key, value)
    if mode == "replace_text":
        _replace_text_keep_first_run(tf, value)
        return True
    if mode == "replace_empty":
        if not is_blank(tf.text):
            warnings.append(f"{fill['fill_id']}: replace_empty skip (도형이 비어있지 않음)")
            return False
        _replace_text_keep_first_run(tf, value)
        return True
    if mode == "append_paragraph":
        p = tf.add_paragraph()
        p.text = value
        return True
    warnings.append(f"{fill['fill_id']}: unknown mode {mode!r}")
    return False


def _apply_table_cell(shape, fill: dict[str, Any], warnings: list[str]) -> bool:
    if not getattr(shape, "has_table", False):
        warnings.append(f"{fill['fill_id']}: shape 가 표가 아님")
        return False
    target = fill["target"]
    row = int(target["row"])
    col = int(target["col"])
    table = shape.table
    if row >= len(list(table.rows)) or col >= len(list(table.columns)):
        warnings.append(f"{fill['fill_id']}: row/col 범위 초과")
        return False
    cell = table.cell(row, col)
    mode = fill["mode"]
    value = str(fill["value"])
    if mode == "set_cell_text":
        _replace_text_keep_first_run(cell.text_frame, value)
        return True
    if mode == "replace_placeholder":
        key = fill.get("explicit_placeholder")
        if not key:
            warnings.append(f"{fill['fill_id']}: replace_placeholder 모드에 explicit_placeholder 없음")
            return False
        return _replace_placeholder_in_text_frame(cell.text_frame, key, value)
    warnings.append(f"{fill['fill_id']}: 표 셀에 지원되지 않는 mode {mode!r}")
    return False


def apply_fills(prs, values: dict[str, Any]) -> tuple[int, int, list[str]]:
    fills = values.get("fills", [])
    applied = 0
    skipped = 0
    warnings: list[str] = []
    for fill in fills:
        fid = fill.get("fill_id", "?")
        mode = fill.get("mode")
        if mode not in VALID_MODES:
            warnings.append(f"{fid}: unknown mode {mode!r}")
            skipped += 1
            continue
        target = fill["target"]
        kind = target["kind"]
        shape = _find_shape(prs, target["shape_id"])
        if shape is None:
            warnings.append(f"{fid}: shape_id {target['shape_id']} 못 찾음")
            skipped += 1
            continue
        if kind == "shape_text":
            ok = _apply_shape_text(shape, fill, warnings)
        elif kind == "table_cell":
            ok = _apply_table_cell(shape, fill, warnings)
        else:
            warnings.append(f"{fid}: unknown target.kind {kind!r}")
            ok = False
        if ok:
            applied += 1
        else:
            skipped += 1
    return applied, skipped, warnings


def _unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    parent = path.parent
    n = 2
    while True:
        candidate = parent / f"{stem}_v{n}{suffix}"
        if not candidate.exists():
            return candidate
        n += 1


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Fill PPTX template with values JSON.")
    parser.add_argument("--template", required=True)
    parser.add_argument("--values", required=True)
    parser.add_argument("--out", required=True)
    parser.add_argument(
        "--strict",
        action="store_true",
        help="채움 실패가 하나라도 있으면 exit 3",
    )
    args = parser.parse_args(argv)

    template_path = Path(args.template)
    values_path = Path(args.values)
    out_path = Path(args.out)

    if not template_path.exists():
        sys.stderr.write(f"양식 파일을 찾을 수 없습니다: {template_path}\n")
        return 1
    if not values_path.exists():
        sys.stderr.write(f"값 JSON 을 찾을 수 없습니다: {values_path}\n")
        return 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path = _unique_path(out_path)

    values = json.loads(values_path.read_text(encoding="utf-8"))

    try:
        prs = Presentation(str(template_path))
    except Exception as e:
        sys.stderr.write(f"PPTX 로드 실패: {e}\n")
        return 2

    applied, skipped, warnings = apply_fills(prs, values)

    try:
        prs.save(str(out_path))
    except Exception as e:
        sys.stderr.write(f"PPTX 저장 실패: {e}\n")
        return 2

    summary = {
        "saved": str(out_path),
        "applied": applied,
        "skipped": skipped,
        "unfilled_in_input": len(values.get("unfilled", [])),
        "warnings": warnings,
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))

    if args.strict and (skipped or warnings):
        return 3
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
