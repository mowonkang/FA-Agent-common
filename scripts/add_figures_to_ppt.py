#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Swiftfab 비교 PPT에 설명용 다이어그램 4종을 시각 슬라이드로 삽입(메인 다음 위치).
원본 deck(ppt-writer 생성) + outputs/figures/*.png → 그림 슬라이드 추가 후 재정렬.
재실행 시 중복 추가되므로 원본 deck 기준 1회 실행."""
import copy
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DECK = "outputs/일본Swiftfab_vs_SMC-MMF_비교분석_2026-06-04.pptx"
FIGS = [
    "outputs/figures/fig1_개념비교.png",
    "outputs/figures/fig2_분모차이.png",
    "outputs/figures/fig5_일정_면적목표.png",
    "outputs/figures/fig4_정량효과.png",
    "outputs/figures/fig3_타임라인.png",
]
GREEN = RGBColor(0x00, 0x66, 0x00)
GRAY = RGBColor(0x80, 0x80, 0x80)

prs = Presentation(DECK)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]  # blank

# 이미지 비율 1240:720 = 1.7222
ratio = 1240 / 720
margin = Cm(0.7)
img_w = SW - 2 * margin
img_h = int(img_w / ratio)
if img_h > SH - Cm(1.4):
    img_h = SH - Cm(1.4)
    img_w = int(img_h * ratio)
left = int((SW - img_w) / 2)
top = int((SH - img_h) / 2)

new_ids = []
for fig in FIGS:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(fig, left, top, width=img_w, height=img_h)
    # 슬라이드 식별용(재정렬에 사용)
    new_ids.append(s.slide_id)

# 재정렬: 메인(첫 슬라이드) 다음에 그림 4장이 오도록 sldIdLst 조작
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
# 원본 슬라이드 수 = len(ids) - len(FIGS)
orig_n = len(ids) - len(FIGS)
fig_elems = ids[orig_n:]            # 마지막 N개 = 방금 추가한 그림
for fe in fig_elems:
    sldIdLst.remove(fe)
# 메인(index 0) 다음(index 1)부터 삽입
insert_at = 1
for offset, fe in enumerate(fig_elems):
    sldIdLst.insert(insert_at + offset, fe)

prs.save(DECK)
print(f"OK — {len(FIGS)} figure slides inserted after main. total slides:", len(list(sldIdLst)))
