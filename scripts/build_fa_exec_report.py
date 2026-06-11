"""FA 자동화 임원 보고 예시 빌더 (sample_data 기반 — 실데이터는 references/).

본 스크립트는 templates/사내양식/(보고양식)YYMMDD_보고제목_v1.0.pptx 의
12종 보고 블록을 어떻게 조합해 '바로 쓸 수 있는 임원 보고서'를 만드는지
보여주는 **예시 산출물**이다. 양식·톤·색·사이즈는 참고 가이드이며,
가독성·스토리를 위해 자율 조정하되 텍스트는 박스를 벗어나지 않는다.

스토리 (6단계): 현황 → 문제·리스크 → 방안 → 투자 → 일정 → 기대효과

슬라이드 (메인 1 + 보조 4):
  1 (메인)  4사분면 임원 요약   [블록 A 보고프레임]
  2 (보조)  핵심 KPI 5종 표     [블록 C KPI]
  3 (보조)  투자비 As-is→To-be + 유연율 추이  [블록 C·G]
  4 (보조)  2026~2028 중장기 로드맵            [블록 B Master Plan]
  5 (보조)  경쟁사 벤치마킹 + 향후 검증        [블록 E·H]

모든 수치는 references/ 실데이터 인용 (각 슬라이드 풋노트에 출처 명시).
출력: outputs/FA자동화_임원보고_예시_2026-05-18.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 (회색·흰색 기조 / 파랑 강조 / 풋노트 진초록) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x00, 0xFF)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

SZ_TITLE = 16
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

OUT_PATH = Path("outputs/FA자동화_임원보고_예시_2026-05-18.pptx")
REPORTER = "FA기술담당 / 작성: 강모원책임 (예시)"


def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH
    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                color=BLACK, align=None, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = opts.get("ls", 1.0)
        para.space_after = Pt(opts.get("sa", 2))
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        run = para.add_run()
        set_run(
            run, text,
            font=opts.get("font", font),
            size=opts.get("size", size),
            bold=opts.get("bold", bold),
            color=opts.get("color", color),
        )


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.1)
    tb.text_frame.margin_right = Cm(0.1)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, fill=DIM_GRAY,
                   color=WHITE, height=0.7):
    shp = add_rect(slide, left, top, width, height, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=SZ_BAND, bold=True, color=color)
    return shp


def add_header(slide, title, *, reporter_color=BLUE):
    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [(title, {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.7, 0.30, 6.5, 0.9,
        [(REPORTER, {"size": SZ_SUB, "bold": True, "color": reporter_color,
                     "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)


def add_footnote(slide, text):
    add_text(
        slide, 0.6, 18.32, 26.3, 0.55,
        [(text, {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )


def add_step_band(slide, step_text):
    """스토리 6단계 표시 띠 (상단, 옅은 파랑)."""
    add_rect(slide, 0.6, 1.52, 26.3, 0.62, fill=LIGHT_BLUE,
             line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.8, 1.54, 26.0, 0.58,
        [(step_text, {"size": SZ_SUB, "bold": True, "color": BLUE})],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def fill_table(tbl, headers, rows, *, col_widths=None,
                head_size=SZ_BODY, body_size=SZ_BODY,
                left_cols=(), emph_last_row=False, row_h=None):
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Cm(w)
    if row_h:
        for r in range(len(rows) + 1):
            tbl.rows[r].height = Cm(row_h)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        cell.margin_left = Cm(0.1)
        cell.margin_right = Cm(0.1)
        write_lines(cell.text_frame, [h], size=head_size, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    n = len(rows)
    for ri, row in enumerate(rows, start=1):
        is_total = emph_last_row and ri == n
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT_GRAY if is_total else WHITE
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            color = BLACK
            bold = (ci == 0) or is_total
            align = PP_ALIGN.LEFT if ci in left_cols else PP_ALIGN.CENTER
            write_lines(cell.text_frame, [str(val)], size=body_size,
                        bold=bold, color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════
#  슬라이드 1 — 메인 (4사분면 임원 요약)
# ════════════════════════════════════════════════════════════
def build_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "FA 자동화 물류 원가·유연성 진척 보고 (’26 1Q, 예시)")

    # 핵심 메시지 박스
    add_rect(slide, 0.6, 1.50, 26.3, 2.34,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.56, 26.0, 2.26,
        [
            ("[핵심 메시지]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE, "sa": 2}),
            (" • ’26 원가 절감 21.3%(△806억)로 목표 16.2% 초과 — 단, 랜싱 "
             "증액 미반영 시 15.8%로 목표 미달 리스크",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 2}),
            (" • 유연 물류(SMC·SMA·HSC) + AX 자동화로 면적·원가 동시 개선, "
             "유연율 36% → 86%(’28)",
             {"size": SZ_BODY, "color": BLACK, "sa": 2}),
            (" • 랜싱 증액분 사업계획 반영 여부 5/21 CEO 보고 — 의사결정 요청",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ],
    )

    Q_TOP, Q_MID = 3.98, 11.12
    Q_LEFT, Q_RIGHT = 0.6, 14.0
    Q_W = 13.05
    Q_H_TOP, Q_H_BOT = 7.04, 6.90

    # ── Q1 현황 ──
    add_label_band(slide, Q_LEFT, Q_TOP, Q_W, "① 현황 (’26 1Q 진척)")
    add_rect(slide, Q_LEFT, Q_TOP + 0.7, Q_W, Q_H_TOP - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_LEFT + 0.3, Q_TOP + 0.92, Q_W - 0.6, Q_H_TOP - 1.05,
        [
            (" • 원가 절감 실적 514.4억원 (기술 305.9 + 구매 141.0), "
             "절감율 30.2%",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • ESMI Lansing 1차 22.1%(329.3억) · 2차 27.9%(100.9억) · "
             "ESAZ 2170 1차 30.5%(78.6억)",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • ESS 물류 Capa : 목표 4만 셀/Line → ESMI2 42,000 Cell/Line "
             "달성, 타 Site 확산 中",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • ESS Link 면적효율 18.0% 개선 설계 완료 "
             "(374 → 305 평/GWh)",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • SMC 투자 심의 완료·PoC 제작 예정, AI 레이아웃 엔진 인프라 "
             "구축 完",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── Q2 문제·리스크 ──
    add_label_band(slide, Q_RIGHT, Q_TOP, Q_W, "② 문제 · 리스크")
    add_rect(slide, Q_RIGHT, Q_TOP + 0.7, Q_W, Q_H_TOP - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.3, Q_TOP + 0.92, Q_W - 0.6, Q_H_TOP - 1.05,
        [
            (" • 랜싱(ESMI L) 사양 변경(Union 설치비·프로토콜 등)으로 "
             "투자비 증액 예정",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 3}),
            (" • 증액분 사업계획 미반영 시 절감율 21.3% → 15.8% 로 "
             "목표(16.2%) 미달",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 3}),
            (" • 진척 평가 ‘판단 보류’ — 증액분 정합화 진행 中",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • 원통형 HSC(ESWA 1동) M3 Gate 5월·각형 SMA M2 6월 등 "
             "Gate 일정 집중 → 검증 지연 리스크",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • 신규 폼팩터·Capa 변동 시 매몰 Loss — 유연 설비 전환 시급",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── Q3 방안 ──
    add_label_band(slide, Q_LEFT, Q_MID, Q_W, "③ 방안 (유연 물류 + AX)")
    add_rect(slide, Q_LEFT, Q_MID + 0.7, Q_W, Q_H_BOT - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_LEFT + 0.3, Q_MID + 0.92, Q_W - 0.6, Q_H_BOT - 1.05,
        [
            (" • SMC(Smart Modular Cube) : Cell 단위 모듈·재활용, "
             "전극보관창고 투자비 10% 절감",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • SMA 설치 효율화 : Compact Lift·Panelizing Rack 으로 "
             "설치 공수·일정 단축",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • HSC(원통형) : 활성화 면적 20% 절감 PoC 검증",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • AX 레이아웃 자동 설계 : 면적 10%↓ · 설계기간 40%↓",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • 4축 KPI(투자비·면적·유연성·재활용성) 로 관리 체계 전환",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── Q4 임원 결정 포인트 ──
    add_label_band(slide, Q_RIGHT, Q_MID, Q_W, "④ 임원 결정 포인트")
    add_rect(slide, Q_RIGHT, Q_MID + 0.7, Q_W, Q_H_BOT - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.3, Q_MID + 0.92, Q_W - 0.6, 2.55,
        [
            (" • 기대효과 : 유연율 ’25 36% → ’28 86% (+50%p), "
             "각형 투자비 88.4 → 70.7 억/GWh (−20%)",
             {"size": SZ_BODY, "color": BLACK, "sa": 3}),
            (" • 벤치마킹 : LGES 90억/GWh 는 글로벌 중상위·합리적 수준 검증",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )
    add_rect(slide, Q_RIGHT + 0.28, Q_MID + 3.30, Q_W - 0.56, 2.62,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.45, Q_MID + 3.40, Q_W - 0.9, 2.45,
        [
            ("[의사결정 요청]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE, "sa": 2}),
            (" • 랜싱 증액분 사업계획 반영 → 절감 목표 16.2% 달성 (5/21 "
             "CEO 보고 안건)",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 2}),
            (" • SMC Demo Line 투자(11.72억, ’26 계획 12.08억 내) 집행 승인",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    add_footnote(
        slide,
        "* 출처 : references/26FA KPI.md [18·20·39·40행, Slide7 259~289행] │ "
        "references/용어집/flexibility_indicator_summary.md [22~29행]")


# ════════════════════════════════════════════════════════════
#  슬라이드 2 — 보조 : 핵심 KPI 5종
# ════════════════════════════════════════════════════════════
def build_kpi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "[보조 1/4] ’26 FA기술담당 핵심 KPI 진척", reporter_color=BLACK)
    add_step_band(slide, "스토리 ② 현황 — 핵심 전략과제 KPI 5종 목표 대비 1Q 진척")

    add_label_band(slide, 0.6, 2.40, 26.3, "핵심 KPI 진척 현황 (’26 1Q)")
    tbl = slide.shapes.add_table(
        6, 5, Cm(0.6), Cm(3.15), Cm(26.3), Cm(11.6)
    ).table
    fill_table(
        tbl,
        ["KPI", "목표 (D-day)", "1Q 진척 (실적/예상)", "산식·기준", "평가"],
        [
            ("물류 설비원가\n경쟁력 강화",
             "사업계획 대비\n16.2% 절감 (’26.12)",
             "21.3% 절감 예상 (514.4억)\n증액 미반영 시 15.8%",
             "절감액 / 사업계획\n투자비", "판단 보류"),
            ("ESS 물류\nCapa 향상",
             "라인당 4만 셀\n(’26.10)",
             "ESMI2 42,000 Cell/Line\n달성 → 타 Site 확산",
             "반송 Cell / Line", "양호"),
            ("각형 물류 경쟁력\n(SMA)",
             "활성화 면적\n16.7% 개선 (’26.10)",
             "개선 설비 개발 完, 제작 中\n6월 M2 Gate 예정",
             "활성화 면적\n개선율", "진행"),
            ("원통형 물류 경쟁력\n(HSC)",
             "활성화 면적\n20% 절감 (’26.6)",
             "PoC 제작·검증 中\n5월 M3 Gate 예정",
             "활성화 면적\n절감율", "진행"),
            ("ESS Link\n면적 효율화",
             "면적 효율\n18.0% 개선 (’26.12)",
             "374 → 305 평/GWh\n설계 完",
             "평 / GWh", "양호"),
        ],
        col_widths=[5.0, 5.0, 7.6, 4.7, 4.0],
        body_size=SZ_SUB, left_cols=(2, 3), row_h=1.92,
    )

    add_footnote(
        slide,
        "* 출처 : references/26FA KPI.md [Slide1 18~20행, Slide2 38~40행, "
        "Slide7 259~277행] (’26 1그룹 KPI 상반기 예상)")


# ════════════════════════════════════════════════════════════
#  슬라이드 3 — 보조 : 투자비 As-is→To-be + 유연율 추이
# ════════════════════════════════════════════════════════════
def build_capex(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "[보조 2/4] 투자비 절감 실적 & 유연율 추이",
               reporter_color=BLACK)
    add_step_band(slide, "스토리 ④ 투자 — 절감 실적(억원) + 유연율 As-Is→To-Be")

    # 좌 : 투자비 표
    add_label_band(slide, 0.6, 2.40, 13.05,
                   "전체 투자비 절감 (FA, 단위 억원)")
    tbl = slide.shapes.add_table(
        5, 4, Cm(0.6), Cm(3.15), Cm(13.05), Cm(7.0)
    ).table
    fill_table(
        tbl,
        ["구분", "사업계획", "실적(예상)", "증감"],
        [
            ("전체 투자비", "3,791.4", "2,985.3", "△806.1 (21.3%)"),
            ("증액 미반영 시", "—", "—", "△600.7 (15.8%)"),
            ("절감 목표", "—", "—", "△613.7 (16.2%)"),
            ("상반기 투자비", "1,478.0", "1,031.2", "△446.8 (30.2%)"),
        ],
        col_widths=[4.05, 3.0, 3.0, 3.0],
        body_size=SZ_SUB, row_h=1.15,
    )
    add_text(
        slide, 0.6, 10.4, 13.05, 4.0,
        [
            ("[해석]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK, "sa": 2}),
            (" • 목표 16.2% 대비 +5.1%p 절감(21.3%) 예상",
             {"size": SZ_BODY, "color": BLACK, "sa": 2}),
            (" • 랜싱 증액분 미반영 시 15.8% → 목표 미달",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 2}),
            (" • MI Lansing 투자비 : 설비비 1,015.7 + 셋업비 306.6 "
             "= 1,322.3억",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # 우 : 유연율 추이 표 + 정의
    add_label_band(slide, 14.0, 2.40, 13.05,
                   "유연율 추이 (MI Lansing, As-Is → To-Be)")
    tbl2 = slide.shapes.add_table(
        5, 2, Cm(14.0), Cm(3.15), Cm(13.05), Cm(7.0)
    ).table
    fill_table(
        tbl2,
        ["연도", "유연율 (유연 설비 투자비 / 전체)"],
        [
            ("’25 (As-Is)", "36%"),
            ("’26", "59%  (+23%p)"),
            ("’27", "85%  (+26%p)"),
            ("’28 (To-Be)", "86%  (누계 +50%p)"),
        ],
        col_widths=[4.05, 9.0],
        body_size=SZ_SUB, left_cols=(1,), row_h=1.15,
    )
    add_text(
        slide, 14.0, 10.4, 13.05, 4.0,
        [
            ("[유연율 정의]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK, "sa": 2}),
            (" • 유연율 = 유연 설비 투자비 / 전체 설비 투자비",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 2}),
            (" • 분자·분모 모두 (설비비 + 셋업비) 합계 기준",
             {"size": SZ_BODY, "color": BLACK, "sa": 2}),
            (" • 유연 = 셋업비 비중 30% 이하 (AMR·OHT·SMC·대차류 등)",
             {"size": SZ_SUB, "color": MID_GRAY, "sa": 2}),
            (" • 고정 = 셋업비 30% 초과 + 전용설비(포장·세척·포트 등)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    add_footnote(
        slide,
        "* 출처 : references/26FA KPI.md [Slide7 259~333행] │ "
        "references/용어집/flexibility_indicator_summary.md [14~29·76~80행]")


# ════════════════════════════════════════════════════════════
#  슬라이드 4 — 보조 : 2026~2028 중장기 로드맵
# ════════════════════════════════════════════════════════════
def build_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "[보조 3/4] 중장기 로드맵 (2026 → 2028)",
               reporter_color=BLACK)
    add_step_band(slide, "스토리 ⑤ 일정 — 트랙별 연도 마일스톤 & 4축 KPI 목표")

    add_label_band(slide, 0.6, 2.40, 26.3,
                   "트랙별 로드맵 (모듈러 · 셋업 자동화 · AX × Physical AI)")
    tbl = slide.shapes.add_table(
        4, 4, Cm(0.6), Cm(3.15), Cm(26.3), Cm(6.6)
    ).table
    fill_table(
        tbl,
        ["트랙", "2026", "2027", "2028"],
        [
            ("모듈러 설비", "SMC (Smart Modular Cube)\n— Cell 단위 모듈·재활용",
             "MMF (Mobile Micro\nFactory) — 셋업 자동화",
             "밀폐형 Cube 물류 설비\n— 표준 Cube 구조"),
            ("셋업 자동화", "SMA 설치 효율화\n— 모듈 사전 조립",
             "SMS 자동 관리\n— 3D 디지털 트윈",
             "SMS 셋업 자동화\n— 이상 사전 감지"),
            ("AX × Physical AI", "AI 레이아웃 자동 설계\n(면적 10%↓·기간 40%↓)",
             "제품 연동 시뮬레이션\n(면적 5% 추가↓)",
             "Omniverse 디지털 트윈\n통합 검증"),
        ],
        col_widths=[4.3, 7.33, 7.33, 7.34],
        body_size=SZ_SUB, left_cols=(1, 2, 3), row_h=1.65,
    )

    add_label_band(slide, 0.6, 10.15, 26.3,
                   "4축 KPI 목표 (CAPEX · 면적 · 유연성 · 재활용성)")
    tbl2 = slide.shapes.add_table(
        5, 4, Cm(0.6), Cm(10.9), Cm(26.3), Cm(6.4)
    ).table
    fill_table(
        tbl2,
        ["지표", "2026", "2027", "2028"],
        [
            ("각형 투자비 (억/GWh)", "88.4", "79.6 (−10%)", "70.7 (−20%)"),
            ("각형 면적 (평/GWh)", "1,129", "1,073 (−15%)", "1,019 (−20%)"),
            ("Phasing 투자 비율", "Pilot 도입", "50% 적용", "80% 적용"),
            ("설비·Rack 재활용률", "베이스라인 측정", "30%", "50%"),
        ],
        col_widths=[6.3, 6.66, 6.67, 6.67],
        body_size=SZ_SUB, left_cols=(0,), row_h=1.18,
    )

    add_footnote(
        slide,
        "* 출처 : references/roadmap/2026_FA기술담당_중장기로드맵_v2.md "
        "[38~45행, 62~64행, 76~78행, 90~92행]")


# ════════════════════════════════════════════════════════════
#  슬라이드 5 — 보조 : 경쟁사 벤치마킹 + 향후 검증
# ════════════════════════════════════════════════════════════
def build_benchmark(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "[보조 4/4] 경쟁사 벤치마킹 & 향후 검증 계획",
               reporter_color=BLACK)
    add_step_band(slide, "스토리 ⑥ 기대효과 — 글로벌 위치 검증 + 다음 단계")

    add_label_band(slide, 0.6, 2.40, 13.05,
                   "① 글로벌 벤치마킹 결론 (각형 라인)")
    add_rect(slide, 0.6, 3.10, 13.05, 11.4,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, 3.35, 12.45, 11.0,
        [
            (" • LGES 90억원/GWh, 생산동 1,800~2,500 평/GWh, 자동화율 "
             "85~90% → 글로벌 중상위·합리적 수준",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 4}),
            (" • CATL 50~70억원/GWh : 자동화 강도 차이가 아닌 "
             "인건비·EPC·국산설비 자가화 구조적 격차",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" • 삼성SDI 90~110억원/GWh : LGES 와 ±10~15% 동등 수준",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" • 사 간 격차(±15%) < 지역 간 격차(±30~50%)",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" • 동일 북미 IRA 환산 시 3사 70~110억원/GWh 밴드로 수렴",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" → 당사 투자 단가는 경쟁력 있는 수준, 유연성·재활용성 "
             "강화가 차별화 포인트",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    add_label_band(slide, 14.0, 2.40, 13.05, "② 향후 검증 계획")
    add_rect(slide, 14.0, 3.10, 13.05, 11.4,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 14.3, 3.35, 12.45, 8.0,
        [
            (" • 5/21 CEO 보고 : 랜싱 증액분 사업계획 반영 의사결정",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 4}),
            (" • SMC Demo Line PoC 제작(6월) → 3Q 검증 완료",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" • 원통형 HSC 5월 M3 Gate · 각형 SMA 6월 M2 Gate",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" • AI 레이아웃 엔진 8월 개발 완료 → 사용자 Test 정합화",
             {"size": SZ_BODY, "color": BLACK, "sa": 4}),
            (" • 유연성·재활용성 베이스라인 측정 → ’27 목표 산정",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )
    add_rect(slide, 14.28, 11.7, 12.49, 2.8,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 14.45, 11.82, 12.15, 2.6,
        [
            ("[기대효과 요약]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE, "sa": 2}),
            (" • 유연율 +50%p (’25 36% → ’28 86%)",
             {"size": SZ_BODY, "bold": True, "color": BLACK, "sa": 2}),
            (" • 각형 투자비 −20%, 면적 −20% (’28)",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    add_footnote(
        slide,
        "* 출처 : references/roadmap/LGES_FA_Prismatic_Benchmarking_"
        "Integrated_v1-v6.md [36~39·178행] │ references/26FA KPI.md [20행]")


def build():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)
    build_main(prs)
    build_kpi(prs)
    build_capex(prs)
    build_roadmap(prs)
    build_benchmark(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
