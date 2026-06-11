"""유연성(Flexibility) 지표 요약 1장 슬라이드 빌더.

references/용어집/flexibility_indicator_summary.md 의 정의와 ESMI Lansing 사례를
FA기술담당 임원 보고용 1장 요약 슬라이드로 정리.

현재 FA기술담당 진행 중인 과제(SMC, SMA, Hybrid Stacker Crane, AMR 등)와 연계
하여 유연성 지표가 KPI/로드맵에 어떻게 작용하는지 시사한다.

출력: outputs/유연성지표_요약_2026-05-14.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 팔레트 (회색·흰색 기조 / 파랑 강조 / 풋노트 진초록) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)  # 본문 기본색

BLUE = RGBColor(0x00, 0x00, 0xFF)        # 강조색 (R0 G0 B255 순수 파랑)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)  # 하이라이트 배경 (옅은 파랑)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)  # 풋노트 전용 초록곰팡이 (G102)

# 하위 호환 (기존 GREEN 참조 → BLUE 로 매핑)
GREEN = BLUE
LIGHT_GREEN = LIGHT_BLUE

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

# ─── 폰트 사이즈 스케일 (일관성 유지) ───────────────────
SZ_TITLE = 16     # 슬라이드 타이틀
SZ_BAND = 12      # 사분면 라벨 띠 (헤더)
SZ_SECTION = 11   # 카드 내 섹션 헤더 ([핵심 메시지], [분자] 등)
SZ_BODY = 10      # 본문 bullet, 표 셀, 책임자 표기
SZ_SUB = 9        # 보조·예시·시뮬레이션
SZ_FOOT = 8       # 풋노트

OUT_PATH = Path("outputs/유연성지표_요약_2026-05-14.pptx")


def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    """run 폰트 설정.

    font 파라미터는 한글(ea) 폰트를 지정하고, 영문/숫자(latin/cs)는
    항상 Arial Narrow 로 자동 분리 적용한다. (LGES 양식 표준)
    호출자가 font=FONT_EMPH 로 지정해 영문 위주로 표현하고 싶을 때만
    ea/latin 모두 동일 폰트로 강제한다.
    """
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    # 한·영 폰트 분리: 한글=LG스마트체, 영문/숫자=Arial Narrow
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH

    run.font.name = ea_font  # 기본(라틴) 노출 폰트
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                color=BLACK, align=None, anchor=None):
    """lines: list of (text, opts) or str. opts overrides defaults."""
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        run = para.add_run()
        set_run(
            run,
            text,
            font=opts.get("font", font),
            size=opts.get("size", size),
            bold=opts.get("bold", bold),
            color=opts.get("color", color),
        )


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.1)
    tb.text_frame.margin_right = Cm(0.1)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, fill=DIM_GRAY, color=WHITE):
    """그린 라벨 띠 (헤더용)."""
    shp = add_rect(slide, left, top, width, 0.7, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=SZ_BAND, bold=True, color=color)
    return shp


def build_slide():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── 상단 제목 ────────────────────────────────────────
    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [
            ("유연성(Flexibility) 지표 정의 및 FA기술담당 현황",
             {"size": SZ_TITLE, "bold": True, "color": BLACK}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # 우상단 책임자
    add_text(
        slide, 20.8, 0.35, 6.4, 0.8,
        [("FA기술혁신파트 강모원책임",
          {"size": SZ_BODY, "bold": True, "color": BLUE,
           "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # 제목 하단 구분선
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)

    # ── 상단 핵심 메시지 박스 ─────────────────────────────
    add_rect(slide, 0.6, 1.55, 26.3, 2.05,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.60, 26.0, 1.95,
        [
            ("[핵심 메시지]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 산식 : 유연 설비 투자비 / 전체 설비 투자비 (분자·분모 모두 설비비+셋업비 합계)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 유연 설비 선정 = 셋업비 비중 30% 이하, 전용 설비·유니언·개조 case 는 예외 (고정형)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • MI Lansing 기준 '25년 36% → '26년 59% → '27년 85% → '28년 86% (+50%p)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • SMC·SMA·Hybrid Stacker·AMR 등 '26 FA KPI 핵심 과제 = 유연성 지표 직접 견인",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── 4사분면 좌표 ─────────────────────────────────────
    Q_TOP = 3.75
    Q_MID = 11.10
    Q_LEFT = 0.6
    Q_RIGHT = 14.0
    Q_W = 13.05
    Q_H_TOP = 7.20
    Q_H_BOT = 6.90

    # ── Q1 (좌상): 정의·산식 ─────────────────────────────
    add_label_band(slide, Q_LEFT, Q_TOP, Q_W, "① 정의 및 산식")
    add_rect(slide, Q_LEFT, Q_TOP + 0.7, Q_W, Q_H_TOP - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    # 산식 박스
    add_rect(slide, Q_LEFT + 0.3, Q_TOP + 1.0, Q_W - 0.6, 1.85,
             fill=SOFT_GRAY, line=LINE_GRAY)
    add_text(
        slide, Q_LEFT + 0.4, Q_TOP + 1.05, Q_W - 0.8, 1.75,
        [
            ("Flexibility  =  유연 설비 투자비  /  전체 설비 투자비",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            ("", {"size": 6}),
            (" • 분자·분모 모두 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("(설비비 + 셋업비) 합계 ",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("기준",
             {"size": SZ_BODY, "color": BLACK}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # 정의 bullet
    add_text(
        slide, Q_LEFT + 0.35, Q_TOP + 3.05, Q_W - 0.7, Q_H_TOP - 3.20,
        [
            ("[유연 설비 선정 기준 (셋업비 30%)]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 셋업비 비중 30% 이하 → 유연 설비 (분자 포함)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 예시 : AMR, OHT, SMC, Skid · 대차류, 완제품 제작품",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 예외 사항(고정형 분류) 은 우상 ② 참조",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 4}),
            ("[산정 방식 (확정)]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 분자·분모 모두 (설비비 + 셋업비) 합계 기준",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 설비비 단독 산정안은 부결 (지표 과대 평가 우려)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 기존 Reusable(재활용성) 산식 명칭 변경에서 출발",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── Q2 (우상): 분류 기준 ─────────────────────────────
    add_label_band(slide, Q_RIGHT, Q_TOP, Q_W,
                   "② 분류 기준 (셋업비 30%) 및 예외 사항")
    add_rect(slide, Q_RIGHT, Q_TOP + 0.7, Q_W, Q_H_TOP - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    # 분류 표
    tbl_left = Q_RIGHT + 0.3
    tbl_top = Q_TOP + 1.05
    tbl_w = Q_W - 0.6
    tbl_h = 3.85
    tbl_shape = slide.shapes.add_table(
        5, 4, Cm(tbl_left), Cm(tbl_top), Cm(tbl_w), Cm(tbl_h)
    )
    tbl = tbl_shape.table
    # 컬럼 폭
    col_widths = [2.4, 3.5, 2.0, 4.55]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Cm(w)
    headers = ["구분", "Sub-Group", "제작:설치", "주요 설비"]
    rows = [
        ("고정형", "구조물", "50 : 50", "Rack, Stage 등"),
        ("고정형", "기계장치", "75 : 25", "Crane, Conveyor, Lift"),
        ("유연형", "AGV·AMR", "90 : 10", "AGV, OHT, AMR"),
        ("유연형", "제작품", "100 : 0", "Skid, 대차"),
    ]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            color = BLUE if (ci == 0 and row[0] == "유연형") else BLACK
            bold = (ci == 0)
            align = PP_ALIGN.CENTER if ci < 3 else PP_ALIGN.LEFT
            write_lines(cell.text_frame, [val], size=SZ_BODY, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)
    # 예외 사항
    add_text(
        slide, Q_RIGHT + 0.35, Q_TOP + 5.20, Q_W - 0.7, Q_H_TOP - 5.35,
        [
            ("[예외 처리]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 전용 설비·유니언 사용으로 셋업비 30% 초과 → 고정형 분류",
             {"size": SZ_BODY, "color": BLACK}),
            ("    예) 포장기, 세척기, 포트, MSS 구조체, 기존 설비 개조",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 개조에 의한 셋업비 비중 高 Case 는 산정 제외",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── Q3 (좌하): MI Lansing 연도별 유연성 지표 ─────────
    add_label_band(slide, Q_LEFT, Q_MID, Q_W,
                   "③ MI Lansing 기준 연도별 유연성 지표")
    add_rect(slide, Q_LEFT, Q_MID + 0.7, Q_W, Q_H_BOT - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    # 연도별 지표 표
    yr_left = Q_LEFT + 0.3
    yr_top = Q_MID + 1.0
    yr_w = Q_W - 0.6
    yr_h = 2.2
    yr_shape = slide.shapes.add_table(
        2, 5, Cm(yr_left), Cm(yr_top), Cm(yr_w), Cm(yr_h)
    )
    yr = yr_shape.table
    col_wy = [3.05, 2.35, 2.35, 2.35, 2.35]
    for i, w in enumerate(col_wy):
        yr.columns[i].width = Cm(w)
    yr_headers = ["구분", "'25년", "'26년", "'27년", "'28년"]
    yr_rows = [("유연성 지표", "36%", "59%", "85%", "86%")]
    for j, h in enumerate(yr_headers):
        cell = yr.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(yr_rows, start=1):
        for ci, val in enumerate(row):
            cell = yr.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            bold = True
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            # '26년 이후 수치 파랑 강조 (의사결정 포인트)
            color = BLUE if ci >= 2 else BLACK
            size = SZ_SECTION if ci >= 1 else SZ_BODY
            write_lines(cell.text_frame, [val], size=size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    # 추이 해석
    add_text(
        slide, Q_LEFT + 0.35, Q_MID + 3.45, Q_W - 0.7, Q_H_BOT - 3.60,
        [
            ("[지표 추이 해석]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • '25 → '28 사이 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("    +50%p 개선 (36% → 86%)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • '25→'26 +23%p, '26→'27 +26%p — SMC·SMA·Hybrid Stacker 효과로 큰폭 상승",
             {"size": SZ_BODY, "color": BLACK}),
            (" • '27→'28 +1%p — 유연 설비 확산에 따른 안정화 구간",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            (" * MI Lansing 투자비 : 설비비 1,015.7억 + 셋업비 306.6억 = 1,322.3억",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── Q4 (우하): FA기술담당 현황 연계 ───────────────────
    add_label_band(slide, Q_RIGHT, Q_MID, Q_W,
                   "④ FA기술담당 현황 연계 ('26 KPI · 로드맵)")
    add_rect(slide, Q_RIGHT, Q_MID + 0.7, Q_W, Q_H_BOT - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)

    # 매핑 표
    mp_left = Q_RIGHT + 0.3
    mp_top = Q_MID + 1.0
    mp_w = Q_W - 0.6
    mp_h = 4.4
    mp_shape = slide.shapes.add_table(
        6, 3, Cm(mp_left), Cm(mp_top), Cm(mp_w), Cm(mp_h)
    )
    mp = mp_shape.table
    col_w3 = [3.8, 4.7, 3.95]
    for i, w in enumerate(col_w3):
        mp.columns[i].width = Cm(w)
    mp_headers = ["진행 과제", "유연성 기여 포인트", "대상 Site"]
    mp_rows = [
        ("SMC (Smart Modular Cube)",
         "Cube Rack + High-pick AMR → 철거·재배치 용이",
         "ESGM2 (M3 '26.10)"),
        ("SMA (Small Modular Aging)",
         "Compact Lift·Panelizing Rack 규격화",
         "ESMI Lansing"),
        ("Hybrid Stacker Crane",
         "활성화 면적 20% 절감, 유연 적재",
         "ESWA 1동"),
        ("AMR / OHT 표준화",
         "유연 설비 분자 직접 기여, 90:10 구조",
         "ESAZ, ESOT, ESHD 等"),
        ("3D 자동 레이아웃",
         "유연 설비 최적 배치 + 면적 효율 ↑",
         "신규 Site"),
    ]
    for j, h in enumerate(mp_headers):
        cell = mp.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(mp_rows, start=1):
        for ci, val in enumerate(row):
            cell = mp.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            bold = (ci == 0)
            color = CHARCOAL if ci == 0 else BLACK
            align = PP_ALIGN.LEFT if ci != 2 else PP_ALIGN.CENTER
            size = SZ_SUB if ci == 1 else SZ_BODY
            write_lines(cell.text_frame, [val], size=size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    # 임원 시사점
    add_rect(slide, Q_RIGHT + 0.3, Q_MID + 5.55, Q_W - 0.6, 1.30,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.45, Q_MID + 5.55, Q_W - 0.9, 1.30,
        [
            ("[다음 단계]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 산식·MI Lansing 지표 확정 → 全 Site KPI 로 확산 (ESAZ·ESWA·ESGM2)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • '26 KPI 모니터링 체계 구축 (분기별 유연성 지표 트래킹)",
             {"size": SZ_BODY, "color": BLACK}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ── 하단 각주 ─────────────────────────────────────
    add_text(
        slide, 0.6, 18.30, 26.3, 0.55,
        [("* 출처 : references/용어집/flexibility_indicator_summary.md (2026-05-14, 강모원)  │  "
          "references/26FA KPI.md  │  references/roadmap/2026_FA기술담당_중장기로드맵.md",
          {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build_slide()
