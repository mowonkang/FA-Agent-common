"""Swiftfab vs LGES SMC·MMF 비교분석 CEO 보고 PPT 빌더.

슬라이드 구성 (메인 1 + 보조 3):
  1 (메인)  분류형 4섹션 — 일본 Swiftfab 팩트 / 당사 SMC·MMF / 연관성 / 시사점·대응
  2 (보조)  Swiftfab 팩트 상세표 (참여 9개사·역할 / 사업총액·일정 / 컨셉·효과)
  3 (보조)  9개 관점 비교표 (주체/모듈단위/적용범위/비용절감/기간/유연성/디지털트윈/차세대/시간축)
  4 (보조)  시사점·대응 3블록 + SMC·MMF 정량효과 As-is→To-be

출처:
  references/경쟁사/2026-06-04_일본_Swiftfab_전지설비연합체_vs_SMC-MMF.md
  references/roadmap/2026_FA기술담당_중장기로드맵_v2.md

출력: outputs/일본Swiftfab_vs_SMC-MMF_비교분석_2026-06-04.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_layout import (  # noqa: E402
    render_classified_slide,
    add_label_band, add_rect, add_text, fill_table,
    title_block, add_footnote, write_lines,
    SLIDE_W_CM, SLIDE_H_CM,
    BLACK, BLUE, WHITE, GRAY, SOFT_GRAY, DIM_GRAY, MID_GRAY, CHARCOAL, LIGHT_BLUE,
    SZ_TITLE, SZ_HEAD, SZ_BAND, SZ_SECTION, SZ_BODY, SZ_SUB, SZ_FOOT,
)

OUT_PATH = Path("outputs/일본Swiftfab_vs_SMC-MMF_비교분석_2026-06-04.pptx")
BYLINE = ["FA기술담당 (2026.06.04)"]

FOOT_MAIN = (
    "* 출처: 일본경제신문(닛케이) 2026-06-02 / 히타치제작소 보도자료(PR TIMES, 2026) / "
    "Nikkei Asia / references/경쟁사/2026-06-04_일본_Swiftfab_전지설비연합체_vs_SMC-MMF.md / "
    "references/roadmap/2026_FA기술담당_중장기로드맵_v2.md"
)
FOOT_NOTE_SCOPE = (
    "** 절감 대상 설비+설치는 공통. 일본 70%는 건물 포함분 / 당사 -15%는 설비+설치 기준(건물은 건설담당 별도 영역, 26→18개월 단축) — '건물 포함 여부' 감안 비교"
)


# ─── 슬라이드 1 : 메인 (분류형 4섹션) ───────────────────────────────────────

def build_main(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="일본 Swiftfab(블록형 모듈 전지공장) 연합체 vs 당사 SMC·MMF 모듈러 기술 비교 분석",
        byline=BYLINE,
        head_message=[
            (
                "일본 9개 설비사 결성 Swiftfab '블록형 모듈 공장'은 당사 SMC·MMF 모듈러 방향과 본질적으로 동일하며,"
                " 당사는 2026년 실증 선행 및 유연성·재활용성 정량지표로 차별화 우위를 확보할 수 있다.",
                {"size": SZ_HEAD, "bold": True, "color": BLACK},
            ),
        ],
        sections=[
            # ── 섹션 1 : Swiftfab 팩트 ──
            {
                "label": "1.\nSwiftfab\n팩트",
                "session_message": "일본 9개사 2026.4 결성 — 블록형 모듈 전지 공장 연합체",
                "content": [
                    (" • 참여사: 히타치·리코·제이텍트·고마츠NTC·서부기연·도신·豊電子工業·히라타기공·대기사 (BASC 9개사 / 사장 K.Kida)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 사업총액 약 180억 엔(약 1,700억 원), 정부보조 포함 / 모델 완성 2028년, 첫 공장 가동 2030년 목표",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 전지 생산 공정을 규격화된 컨테이너 블록으로 분할 — 약 1,000모듈 = 연 5만 대분",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 총비용(건물+설비+설치) 약 70% 절감**, 공장 건설기간 4~6년 → 2~3년",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 전략: 중국 저가 공세 대항 + 전고체 배터리 시장 반격 목적",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            # ── 섹션 2 : 당사 SMC·MMF ──
            {
                "label": "2.\n당사\nSMC·MMF",
                "session_message": "FA기술담당 자가 개발·자가 적용 — 4축 KPI 중심 모듈러 로드맵",
                "content": [
                    (" • SMC(2026): Cube Rack + High-pick AMR → 면적 +25%, 설치공수 -33%, 투자비 -15%",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • MMF(2027): 이동식 Rack 자동 조립 로봇 → 설치공수 -46%, 투자비 -15%",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 밀폐형 Cube(2028): DR Less 공정 대응, 전고체 폼팩터 전환 표준 Cube 구조",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • KPI(→2028): 면적 1,254→1,019(-20%) / 건설 26→18개월 / 셋업LT -50% / 재활용률 50%",
                     {"size": SZ_BODY, "bold": True, "color": BLACK}),
                    (" • 디지털트윈: CNS Physical Works·Omniverse·NVIDIA Isaac SIM 가상 사전 검증",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            # ── 섹션 3 : 연관성 ──
            {
                "label": "3.\n연관성\n·차이",
                "session_message": "방향성 동일, 범위·비즈니스 모델·성숙도에서 핵심 차이",
                "content": [
                    (" [동일] 설비 규격화·모듈화 → 셋업·건설 공수 절감 + 폼팩터 변화 유연 대응 (산업 정답 외부 검증)",
                     {"size": SZ_BODY, "bold": True, "color": BLACK}),
                    (" [차이①-범위] 절감 대상 설비+설치 공통, '건물 포함 여부'가 차이 — 일본 70%(건물 포함) vs 당사 -15%(건물 제외)**",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" [차이②-모델] 일본=표준 설비 타사 판매(연합 공급망) vs 당사=자가 내재화 경쟁력(self-use)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" [차이③-성숙도] 당사 SMC 2026 Demo 실증 진행 중 vs 일본 모델 2028·가동 2030 — 당사 실증 선행",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" [공통] 디지털트윈 사전 검증: 히타치 HMAX Industry vs CNS Physical Works·Isaac SIM",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            # ── 섹션 4 : 시사점·대응 ──
            {
                "label": "4.\n시사점\n·대응",
                "session_message": "[위협] 설비 범용화 가속·전고체 선점 의도  /  [기회] 실증 선행·지표 차별화",
                "content": [
                    (" [위협] 설비 표준 외판 → 후발 배터리 메이커 진입장벽 저하, 전고체 라인 표준 선점 경쟁",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" [기회①] 당사 SMC/MMF 모듈러 방향이 산업 정답으로 외부 검증 — 그룹장 지표 신설 지시와 정합",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" [기회②] 재활용률 50%·Phasing 80% — 일본 미명시 항목, 당사 고유 정량 차별화 포인트",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" [제언①] 모듈 인터페이스 표준·공통조달 방식 벤치마킹 → 당사 Cube/Rack 표준화 참고",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" [제언②] 디지털트윈 검증 조기화(2027→앞당김) + 전고체 모듈러 대응 명문화 추진",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
        ],
        footnote=FOOT_MAIN + "  |  " + FOOT_NOTE_SCOPE,
        content_top=3.10,
        content_bottom=18.10,
    )


# ─── 슬라이드 2 : 보조 1 — Swiftfab 팩트 상세표 ─────────────────────────────

def build_sub1_swiftfab(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide,
                "[보조 1/4] Swiftfab Energy Systems — 팩트 상세",
                byline=BYLINE)

    # 참여 9개사 표
    add_label_band(slide, 0.6, 1.55, 26.3,
                   "① 참여 9개사 및 역할")

    # 헤더 1행 + 데이터 9행 = 10행, 컬럼 2개
    rows1 = [
        ("히타치제작소", "디지털트윈 시뮬레이션(HMAX Industry), AI·현장 데이터 설계→생산 최적화, 준비 단계 주도"),
        ("리코(리코엘레멕스)", "정밀 가공·소형화 설비(정밀 위치결정·축제어·화상인식)"),
        ("제이텍트(JTEKT, 도요타 계열)", "원천공정(源泉工程) 설비, 심플·슬림·콤팩트 설계"),
        ("서부기연(西部技研)", "데시칸트 제습·전열교환, 드라이룸 환경관리"),
        ("고마츠NTC(Komatsu NTC)", "탭 성형·검사(120m/min), 레이저 용접·반송·모듈 공정 자동화"),
        ("도신(東伸, Toshin)", "전극판 슬리팅·리와인딩 전문"),
        ("豊電子工業(Yutaka Electronics)", "전장·자동화 설비 (아이치현 가리야시)"),
        ("히라타기공(平田機工, 6258)", "양산 조립 라인·엔지니어링, 설비 풋프린트 30% 절감"),
        ("대기사(大気社, Taikisha)", "국소 클린룸·유틸리티 일체형 모듈 공법, NMP 회수"),
    ]
    tbl1 = slide.shapes.add_table(
        len(rows1) + 1, 2, Cm(0.6), Cm(2.35), Cm(26.3), Cm(7.2))
    fill_table(
        tbl1.table,
        headers=["참여사", "역할 (공개 범위)"],
        rows=rows1,
        col_w=[6.5, 19.8],
        body_size=SZ_SUB,
        head_size=SZ_BODY,
    )
    for r in range(len(rows1) + 1):
        tbl1.table.rows[r].height = Cm(0.72)

    # 사업 개요·정량효과 표
    add_label_band(slide, 0.6, 9.75, 26.3,
                   "② 결성 개요 및 정량 효과")

    # 헤더 1행 + 데이터 7행 = 8행, 컬럼 3개
    rows2 = [
        ("사업체명 / CEO", "Swiftfab Energy Systems(Swift+Fabrication) / 사장 Keisuke Kida(키다 케이스케)", "히타치 보도자료 / Nikkei"),
        ("설립·참여", "2025-12-18 BASC 합의 → 2026.4 법인(도쿄 미나토구), BASC 9개사 공동 출자", "히타치 보도자료 / Nikkei Asia"),
        ("사업총액", "약 180억 엔(약 1,700억 원) — 정부 보조금 포함 (원문 約180億円)", "일본경제신문 2026-06-02"),
        ("일정", "모델 설비 2028년경 완성 → 전지 메이커 판매 / 첫 공장 가동 2030년 말", "닛케이 / Nikkei Asia"),
        ("산업 위상", "中 글로벌 전지 ~70% / 日 설비메이커 글로벌 장비시장 9% vs 中 25%", "Automotive World 인용 Nikkei"),
        ("총비용 절감**", "건물+설비+설치 약 70% 절감 (당사 절감 대상은 설비+설치, 건물 제외)", "일본경제신문 2026-06-02"),
        ("건설기간(건물)", "기존 4~6년 → 2~3년(≈24~36개월) / 당사 건설 26→18개월(-31%) 대응", "닛케이 / FA기술담당"),
    ]
    tbl2 = slide.shapes.add_table(
        len(rows2) + 1, 3, Cm(0.6), Cm(10.55), Cm(26.3), Cm(5.6))
    fill_table(
        tbl2.table,
        headers=["항목", "내용", "출처"],
        rows=rows2,
        col_w=[4.0, 16.0, 6.3],
        body_size=SZ_SUB,
        head_size=SZ_BODY,
    )
    for r in range(len(rows2) + 1):
        tbl2.table.rows[r].height = Cm(0.7)

    add_footnote(slide,
                 "* 출처: 히타치제작소 보도자료(PR TIMES) / 일본경제신문 2026-06-02 / Nikkei Asia / Research Report.md  "
                 "|  ** 70%는 건물+설비+설치 기준 / 당사 절감 대상은 설비+설치(건물은 건설담당 별도, 26→18개월)",
                 top=17.90)


# ─── 슬라이드 3 : 보조 2 — 9개 관점 비교표 ──────────────────────────────────

def build_sub2_compare(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide,
                "[보조 2/4] 일본 Swiftfab vs LGES SMC·MMF — 9개 관점 비교",
                byline=BYLINE)

    add_label_band(slide, 0.6, 1.55, 26.3,
                   "방향성 동일, 범위·비즈니스 모델·성숙도에서 핵심 차이 — 당사 실증 시점 선행이 핵심 우위")

    # 헤더 1행 + 데이터 9행 = 10행, 컬럼 3개
    rows_cmp = [
        ("주체",
         "9개 설비사 연합체 (공동 출자, 외판 비즈니스)",
         "FA기술담당 자가 개발·자가 적용 (self-use)"),
        ("모듈 단위",
         "공정 전체를 컨테이너 블록 (약 1,000모듈 = 연 5만 대분)",
         "Cube Rack / 이동식 조립 셀 (물류·셋업 영역 중심)"),
        ("절감 대상(분모)**",
         "건물 + 설비 + 설치 (공장 전체)",
         "설비 + 설치 (건물은 건설담당 별도 영역)"),
        ("비용 절감(주장)**",
         "총비용 70% (건물 포함)",
         "투자비 -15%, 설치공수 -33~46% (설비+설치 기준)"),
        ("건설 기간(건물)",
         "공장 건설 4~6년 → 2~3년 (≈24~36개월)",
         "26개월 → 18개월 (2028, -31%) *건설담당"),
        ("유연성·재활용",
         "모듈 재사용 가능(블록 구조) — 정량 지표 미명시",
         "명시적 지표화: Phasing 투자 80%, 재활용률 50%"),
        ("디지털트윈",
         "히타치 HMAX Industry (소프트 디버깅·AI 최적화)",
         "CNS Physical Works·Omniverse·NVIDIA Isaac SIM"),
        ("차세대 타깃",
         "전고체 배터리 시장 반격 (명시)",
         "폼팩터 변화 대응 표준 Cube (전고체 포함)"),
        ("시간축(성숙도)",
         "모델 2028 / 첫 공장 가동 2030 — 결성 단계",
         "SMC 2026 Demo 실증 진행 / MMF 2027 / 밀폐형 2028"),
    ]
    tbl = slide.shapes.add_table(
        len(rows_cmp) + 1, 3, Cm(0.6), Cm(2.35), Cm(26.3), Cm(14.5))
    fill_table(
        tbl.table,
        headers=["관점", "일본 Swiftfab", "LGES SMC·MMF"],
        rows=rows_cmp,
        col_w=[4.0, 11.15, 11.15],
        body_size=SZ_SUB,
        head_size=SZ_BODY,
    )
    for r in range(len(rows_cmp) + 1):
        tbl.table.rows[r].height = Cm(1.40)

    add_footnote(slide,
                 "* 출처: references/경쟁사/2026-06-04_일본_Swiftfab_전지설비연합체_vs_SMC-MMF.md §3.2 / "
                 "references/roadmap/2026_FA기술담당_중장기로드맵_v2.md  "
                 "|  ** 절감 대상 설비+설치는 공통 / 일본 70%는 건물 포함분, 당사는 건물 제외(건설담당 26→18개월) — '건물 포함 여부' 감안 비교",
                 top=17.90)


# ─── 슬라이드 4 : 보조 3 — 시사점·대응 + 당사 정량효과 ───────────────────────

def build_sub3_implications(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide,
                "[보조 3/4] 시사점·대응 방향 및 당사 정량효과",
                byline=BYLINE)

    # ── 위협·기회·제언 3블록 ──
    add_label_band(slide, 0.6, 1.55, 26.3,
                   "① 위협 / 기회 / 제언")

    BLOCK_W = 8.5
    block_data = [
        ("[위협]", BLOCK_W, 1.75,
         [
             ("공급망 표준화 압력",
              {"size": SZ_SECTION, "bold": True, "color": BLACK}),
             ("설비 표준 외판 → 후발 배터리 메이커 진입장벽 저하, 당사 설비 내재화 경쟁력 상대적 희석 우려",
              {"size": SZ_BODY, "color": BLACK}),
             ("전고체 선점 의도",
              {"size": SZ_SECTION, "bold": True, "color": BLACK}),
             ("일본이 차세대 전고체 라인 표준 선점 시 폼팩터 전환기 주도권 경쟁 불리",
              {"size": SZ_BODY, "color": BLACK}),
         ]),
        ("[기회]", BLOCK_W, 10.25,
         [
             ("방향성 외부 검증",
              {"size": SZ_SECTION, "bold": True, "color": BLUE}),
             ("일본 9개사 연합이 당사 SMC/MMF 모듈러 방향의 산업 정답 검증 — 그룹장 지표 신설 지시와 정합",
              {"size": SZ_BODY, "color": BLACK}),
             ("실증 시점 선행",
              {"size": SZ_SECTION, "bold": True, "color": BLUE}),
             ("SMC 2026 Demo·MMF 1차 설계 진행 중 — 일본(2028 모델) 대비 실증 시점 선행",
              {"size": SZ_BODY, "color": BLACK}),
             ("차별화 축",
              {"size": SZ_SECTION, "bold": True, "color": BLUE}),
             ("일본 미명시 재활용률 50%·Phasing 80% 정량지표 — 당사 고유 차별화 포인트",
              {"size": SZ_BODY, "color": BLACK}),
         ]),
        ("[제언]", BLOCK_W, 18.75,
         [
             ("① 모듈 인터페이스 표준·공통조달 벤치마킹",
              {"size": SZ_BODY, "bold": True, "color": BLACK}),
             ("Swiftfab 컨테이너 블록 규격화+공통조달 방식을 당사 Cube/Rack 표준화에 참고",
              {"size": SZ_BODY, "color": BLACK}),
             ("② 디지털트윈 검증 가속",
              {"size": SZ_BODY, "bold": True, "color": BLACK}),
             ("CNS Physical Works·Isaac SIM 가상 검증 로드맵(2027~2028) 조기화 검토",
              {"size": SZ_BODY, "color": BLACK}),
             ("③ 전고체 모듈러 대응 명문화",
              {"size": SZ_BODY, "bold": True, "color": BLACK}),
             ("2028 밀폐형 Cube·표준 Cube 구조를 전고체 폼팩터 전환 대응으로 명시 연결",
              {"size": SZ_BODY, "color": BLACK}),
         ]),
    ]

    for label, bw, bx, lines in block_data:
        box = add_rect(slide, bx, 2.35, bw, 7.30,
                       fill=None, line=BLACK, line_w=0.75)
        box.text_frame.word_wrap = True
        # label header
        add_label_band(slide, bx, 2.35, bw, label)
        add_text(slide, bx + 0.1, 3.10, bw - 0.2, 6.40, lines)

    # ── 당사 정량효과 As-is→To-be ──
    add_label_band(slide, 0.6, 10.00, 26.3,
                   "② 당사 SMC·MMF 정량효과 (As-is → To-be)")

    # 헤더 1행 + 데이터 5행 = 6행, 컬럼 5개
    rows_kpi = [
        ("SMC (2026)", "면적 효율 (개/라인)", "908", "1,149", "+25%"),
        ("SMC (2026)", "설치 공수 (MD)", "3,234", "2,100", "-33%"),
        ("SMC (2026)", "투자비 (억 원)", "239.8", "203.3", "-15%"),
        ("MMF (2027)", "설치 공수 (MD)", "2,280", "1,232", "-46%"),
        ("MMF (2027)", "투자비 (억 원)", "296.0", "250.6", "-15%"),
    ]
    tbl = slide.shapes.add_table(
        len(rows_kpi) + 1, 5, Cm(0.6), Cm(10.80), Cm(26.3), Cm(6.0))
    fill_table(
        tbl.table,
        headers=["과제", "지표", "As-is", "To-be", "효과"],
        rows=rows_kpi,
        col_w=[4.5, 6.0, 4.8, 4.8, 6.2],
        body_size=SZ_BODY,
        head_size=SZ_BODY,
    )
    for r in range(len(rows_kpi) + 1):
        tbl.table.rows[r].height = Cm(0.95)

    add_footnote(slide,
                 "* 출처: references/roadmap/2026_FA기술담당_중장기로드맵_v2.md [62~69행, 286행, 291행] / "
                 "references/경쟁사/2026-06-04_일본_Swiftfab_전지설비연합체_vs_SMC-MMF.md §4",
                 top=17.90)


# ─── 슬라이드 5 : 보조 4 — 판단 전환 기준 + 심층 리서치 추가 팩트 ──────────────

def build_sub4_watchlist(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide,
                "[보조 4/4] 판단 전환 기준 및 심층 리서치 추가 팩트",
                byline=BYLINE)

    # ① 판단 전환 기준
    add_label_band(slide, 0.6, 1.55, 26.3,
                   "① 판단 전환 기준 (위협 격상 임계값) — 3개 중 2개 이상 충족 시 '실질 위협'으로 격상")
    box = add_rect(slide, 0.6, 2.35, 26.3, 3.55, fill=None, line=BLACK, line_w=0.75)
    box.text_frame.word_wrap = True
    add_text(slide, 0.8, 2.55, 25.9, 3.2, [
        ("Swiftfab은 현재 컨셉·로드맵 단계 — 2028~2030 설비 모듈화·표준화 경쟁의 신호탄. 아래 모니터링 지표 추적.",
         {"size": SZ_BODY, "color": BLACK}),
        ("① 실제 수주처(비중국 전지 메이커·OEM) 확보 발표",
         {"size": SZ_SECTION, "bold": True, "color": BLUE}),
        ("② 총비용 70% 절감의 제3자 검증",
         {"size": SZ_SECTION, "bold": True, "color": BLUE}),
        ("③ 첫 공장 양산 수율 ≥ 90% 달성",
         {"size": SZ_SECTION, "bold": True, "color": BLUE}),
    ])

    # ② 심층 리서치 추가 팩트 (3블록)
    add_label_band(slide, 0.6, 6.25, 26.3,
                   "② 심층 리서치 추가 팩트 (출처: references/Research Report.md)")
    BLOCK_W = 8.5
    blocks = [
        ("전고체 OEM 일정 (정합)", 0.6, [
            ("• 도요타 2027~2028 양산 timeframe", {"size": SZ_BODY, "color": BLACK}),
            ("• 닛산 FY2028 상용화 (ASSB 파일럿 가동)", {"size": SZ_BODY, "color": BLACK}),
            ("• 혼다 2020년대 후반", {"size": SZ_BODY, "color": BLACK}),
            ("→ Swiftfab 차세대 라인(모델 2028) 시점 부합", {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
        ("컨테이너 모듈 기술 한계 (FA)", 9.1, [
            ("• 셀조립·주액 노점 −40℃↓, 권취/조립 ≤RH10%", {"size": SZ_BODY, "color": BLACK}),
            ("• 클린룸 ISO 7~8, 미크론급 정밀도·저진동", {"size": SZ_BODY, "color": BLACK}),
            ("→ 모듈 분할 시 노점·차압·기밀·제진 난제", {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("→ 대기사·서부기연 역할이 성패 관건 (LGES 비교점)", {"size": SZ_BODY, "color": BLACK}),
        ]),
        ("정부정책·산업 위상", 17.6, [
            ("• 경산성 축전지산업전략 7대 과제 中", {"size": SZ_BODY, "color": BLACK}),
            ("  ②전지설비산업의 구조변혁 직접 대응", {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("• 경제안보법 보조(설비 1/3·기술개발 1/2)", {"size": SZ_BODY, "color": BLACK}),
            ("• 中 전지 ~70% / 日 장비 9% vs 中 25%", {"size": SZ_BODY, "color": BLACK}),
        ]),
    ]
    for label, bx, lines in blocks:
        b = add_rect(slide, bx, 7.05, BLOCK_W, 8.6, fill=None, line=BLACK, line_w=0.75)
        b.text_frame.word_wrap = True
        add_label_band(slide, bx, 7.05, BLOCK_W, label)
        add_text(slide, bx + 0.15, 7.85, BLOCK_W - 0.3, 7.5, lines)

    add_footnote(slide,
                 "* 단일 1차 출처(2026.6.2 닛케이) 의존 항목 포함 — 현시점 '계획·목표'이며 실현 성과 아님(수주처 미확보). "
                 "출처: references/Research Report.md, references/경쟁사/2026-06-02_닛케이_블록형공장_원문아카이브.md",
                 top=16.30)


# ─── 메인 ─────────────────────────────────────────────────────────────────────

def build() -> None:
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)

    build_main(prs)
    build_sub1_swiftfab(prs)
    build_sub2_compare(prs)
    build_sub3_implications(prs)
    build_sub4_watchlist(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    # suffix 처리 (동일 날짜 파일 존재 시 _v2, _v3 ...)
    out = OUT_PATH
    if out.exists():
        stem = out.stem
        suffix = out.suffix
        v = 2
        while out.exists():
            out = OUT_PATH.with_name(f"{stem}_v{v}{suffix}")
            v += 1
    prs.save(out)
    print(f"saved: {out}")


if __name__ == "__main__":
    build()
