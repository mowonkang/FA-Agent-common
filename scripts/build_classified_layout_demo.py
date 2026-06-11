"""분류형 레이아웃 데모 PPT 빌더 (사내 보고양식 슬라이드 1 기반).

`scripts/_ppt_layout.py` 의 ``render_classified_slide`` 를 사용해 사내 양식
``(보고양식)YYMMDD_보고제목_v1.0.pptx`` 슬라이드 1 레이아웃(제목 → 배경 없는
검은 Head message → 좌측 분류 레일 → 우측 세션메시지+내용)을 시연한다.

기존 2×2 사분면(quadrant) 레이아웃이 "눈에 안 들어온다"는 피드백에 따라
위 → 아래로 읽히는 좌측 분류 레일 + 우측 내용 구조로 전환한 견본.

- 슬라이드 1 : 분류 1~4 (텍스트 / 표 / 이미지 박스 혼합)
- 슬라이드 2 : 3단(개요·내용·결론) — 가운데 "내용"을 길게 두어 높이 자동 가변 시연

모든 텍스트는 플레이스홀더(〈예시〉)이며 실데이터가 아니다.

출력: outputs/분류형_레이아웃_데모.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_layout import (  # noqa: E402
    render_classified_slide, SLIDE_W_CM, SLIDE_H_CM,
    BLACK, BLUE, MID_GRAY, SZ_HEAD, SZ_BODY, SZ_SUB,
)

OUT_PATH = Path("outputs/분류형_레이아웃_데모.pptx")
FOOT = ("* 샘플 양식 데모 — 실데이터 아님. 레이아웃 기준 : "
        "templates/사내양식/(보고양식)YYMMDD_보고제목_v1.0.pptx 슬라이드 1")


def build_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="분류형 레이아웃 데모 (4분류)",
        byline="FA기술혁신파트",
        head_message=[
            ("좌측 분류 박스 높이는 우측 내용에 맞춰 자동 조절되며, 분류 수·구성은 "
             "내용에 따라 달라집니다",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "분류 1.\n개요",
                "session_message": "〈세션 메시지〉 검토 배경과 목적을 한 줄로 요약",
                "content": [
                    (" • 〈예시〉 추진 배경 — 현행 방식의 한계와 개선 필요성을 기술",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 〈예시〉 검토 범위 — 대상 공정·설비·기간을 명시",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "분류 2.\n현황",
                "session_message": "〈세션 메시지〉 핵심 현황을 표로 정리",
                "content": [
                    (" • 〈예시〉 항목별 As-is 수치와 목표치를 비교",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["구분", "As-is", "To-be", "비고"],
                    "rows": [
                        ("항목 A", "〈예시〉 100", "〈예시〉 130", "+30%"),
                        ("항목 B", "〈예시〉 12", "〈예시〉 8", "△33%"),
                    ],
                    "col_w": [4.0, 6.0, 6.0, 7.8],
                    "row_h": 0.95,
                },
            },
            {
                "label": "분류 3.\n이슈",
                "session_message": "〈세션 메시지〉 주요 이슈와 이미지 첨부 영역",
                "content": [
                    (" • 〈예시〉 이슈 1 — 원인과 영향 범위를 기술",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 〈예시〉 이슈 2 — 대응 우선순위를 기술",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "image_box": 2.4,
            },
            {
                "label": "분류 4.\n다음 단계",
                "session_message": "〈세션 메시지〉 액션 플랜과 일정",
                "content": [
                    (" ① 〈예시〉 단기 과제 — 담당/기한 명시",
                     {"size": SZ_BODY, "bold": True, "color": BLACK}),
                    (" ② 〈예시〉 중기 과제 — 검증·확산 계획",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" · 〈예시〉 세부 일정은 보조 슬라이드 참조",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
        ],
        footnote=FOOT,
    )


def build_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="3단 구성 데모 (개요 · 내용 · 결론)",
        byline="FA기술혁신파트",
        head_message=[
            ("개요·내용·결론 3단 구성 — 가운데 '내용'이 길어 박스 높이가 자동으로 커집니다",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "1.\n개요",
                "session_message": "〈세션 메시지〉 한 줄 배경",
                "content": [
                    (" • 〈예시〉 검토 배경과 목적을 간결히 정리",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "2.\n내용",
                "session_message": "〈세션 메시지〉 상세 검토 내용 (분량이 많은 본문)",
                "content": [
                    (" • 〈예시〉 세부 항목 1 — 분석 내용을 구체적으로 기술",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 〈예시〉 세부 항목 2 — 비교·대안 검토 결과",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 〈예시〉 세부 항목 3 — 정량 효과(원가/시간/품질) 추정",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 〈예시〉 세부 항목 4 — 리스크와 제약 사항",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 〈예시〉 세부 항목 5 — 타 부문/협력사 협의 필요 사항",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 〈예시〉 세부 항목 6 — 검증 방법과 판정 기준",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["대안", "장점", "단점", "결론"],
                    "rows": [
                        ("Alt 1", "〈예시〉 저비용", "〈예시〉 확장성 제약", "보류"),
                        ("Alt 2", "〈예시〉 표준화 용이", "〈예시〉 초기 투자", "채택"),
                    ],
                    "col_w": [3.5, 7.0, 7.0, 6.3],
                    "row_h": 0.95,
                },
            },
            {
                "label": "3.\n결론",
                "session_message": "〈세션 메시지〉 권고안과 의사결정 요청",
                "content": [
                    (" • 〈예시〉 Alt 2 채택 권고 — 근거 요약",
                     {"size": SZ_BODY, "bold": True, "color": BLACK}),
                    (" • 〈예시〉 승인 요청 사항 / 다음 보고 일정",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
        ],
        footnote=FOOT,
    )


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    build_four(prs)
    build_three(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
