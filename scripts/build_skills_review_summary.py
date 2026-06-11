"""인기 Claude Skills 도입 권고 1장 요약 슬라이드 빌더.

본 빌더는 outputs/popular_skills_review_2026-05-20.md 의 결론을
LGES 표준 PPT 양식 (templates/사내양식/LGES_PPT_작업_가이드.md) 에 맞춰
임원 보고용 1장 + 보조 1장으로 정리한다.

기준 구현: scripts/build_flexibility_summary.py 헬퍼·색·폰트 패턴 그대로 채용.

출력: outputs/추천_스킬_요약_2026-05-20.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 팔레트 ─────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)

BLUE = RGBColor(0x00, 0x00, 0xFF)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

# ─── 폰트 사이즈 6단 상수 ─────────────────────────────
SZ_TITLE = 16
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

OUT_PATH = Path("outputs/추천_스킬_요약_2026-05-20.pptx")


def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    """run 폰트 설정 (한·영 자동 분리). build_flexibility_summary.py 동일 패턴."""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH
    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                color=BLACK, align=None, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        run = para.add_run()
        set_run(
            run,
            text,
            font=opts.get("font", font),
            size=opts.get("size", size),
            bold=opts.get("bold", bold),
            color=opts.get("color", color),
        )


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.1)
    tb.text_frame.margin_right = Cm(0.1)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, fill=DIM_GRAY, color=WHITE):
    shp = add_rect(slide, left, top, width, 0.7, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=SZ_BAND, bold=True, color=color)
    return shp


# ────────────────────────────────────────────────────
# 메인 슬라이드 데이터
# ────────────────────────────────────────────────────

TOP5_CARDS = [
    {
        "rank": "① xlsx",
        "tag": "공식·Document",
        "summary": "Excel 직접 처리 (수식·피벗·차트)",
        "bullets": [
            "KPI·Capex·MRM xlsx → MD/PPT 환승 비용 ↓",
            "data-teammate + sample_data/ 직결",
        ],
        "score": "5/5",
        "score_emph": True,
    },
    {
        "rank": "② pdf",
        "tag": "공식·Document",
        "summary": "PDF 본문/표/메타 추출 + 생성",
        "bullets": [
            "협력사 사양서·도면 검색 가능화",
            "ops-teammate 검색 커버리지 大 확장",
        ],
        "score": "5/5",
        "score_emph": True,
    },
    {
        "rank": "③ skill-creator",
        "tag": "공식·Meta",
        "summary": "사내 도메인 스킬 라이브러리 표준화",
        "bullets": [
            "협력사평가·안전점검·KPI진척 등 신규 스킬 출발점",
            "기존 3개 스킬 리팩토링 기준",
        ],
        "score": "5/5",
        "score_emph": False,
    },
    {
        "rank": "④ brand-guidelines",
        "tag": "공식·Design",
        "summary": "LGES 색·폰트·풋노트 룰 캡슐화",
        "bullets": [
            "ppt-writer 가 자동 적용",
            "기존 LGES_PPT_가이드.md → SKILL.md 변환",
        ],
        "score": "5/5",
        "score_emph": False,
    },
    {
        "rank": "⑤ mcp-builder",
        "tag": "공식·Engineering",
        "summary": "사내 ERP/MES/PLM Read-only MCP",
        "bullets": [
            "data·ops-teammate 실시간 조회로 확장",
            "보안/IT 검토 후 진행 (난이도 M~L)",
        ],
        "score": "4/5",
        "score_emph": False,
    },
]

REVIEW6 = [
    ("⑥ docx",            "meeting-merger 표준화·정책 SOP 문서화"),
    ("⑦ doc-coauthoring", "임원 보이스 보존, 주간/월간 톤 일관"),
    ("⑧ internal-comms",  "사내 공지·주간메일 양식 표준화"),
    ("⑨ webapp-testing",  "장기 — 사내 대시보드 자동 테스트"),
    ("⑩ Firecrawl",       "JS 사이트 안정 스크래핑 (tech-research)"),
    ("⑪ Sentry CodeRev",  "scripts/ 신규 빌더 품질 게이트"),
]

HOLD12 = [
    "canvas-design", "theme-factory", "algorithmic-art",
    "slack-gif-creator", "frontend-design", "web-artifacts-builder",
    "webapp-testing(현재)", "claude-api(이미 활용)", "Superpowers(SOP 중복)",
    "Trail of Bits", "Vercel Guidelines", "HuggingFace/Tinybird/Nimble",
]

# ────────────────────────────────────────────────────


def build_main_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 상단 제목 ───────────────────────────────────
    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [("인기 Claude Skills 도입 권고 — FA 자동화 프로젝트",
          {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.8, 0.35, 6.4, 0.8,
        [("FA기술혁신파트 · 2026-05-20",
          {"size": SZ_BODY, "bold": True, "color": BLUE,
           "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)

    # ── 핵심 메시지 박스 ───────────────────────────
    add_rect(slide, 0.6, 1.55, 26.3, 2.05,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.60, 26.0, 1.95,
        [
            ("[핵심 메시지]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 23개 후보 (공식 17 + 커뮤니티 6) 중 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("    즉시 도입 추천 5개 / 검토 추천 6개 / 보류 12개 ",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("    선별",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Top 3 (xlsx · pdf · skill-creator) 는 난이도 S, 1주 내 scaffold 가능 — KPI/사양서/스킬표준화 즉시 효과",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • brand-guidelines + mcp-builder 는 LGES PPT 일관화·사내시스템 연계의 장기 인프라 가치",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── 좌측 : 즉시 도입 Top 5 ──────────────────────
    LEFT_X = 0.6
    LEFT_W = 16.4
    TOP_Y = 3.85
    add_label_band(slide, LEFT_X, TOP_Y, LEFT_W, "① 즉시 도입 추천 Top 5")
    add_rect(slide, LEFT_X, TOP_Y + 0.7, LEFT_W, 13.55,
             fill=WHITE, line=MID_GRAY, line_w=0.75)

    card_y = TOP_Y + 0.95
    card_h = 2.55
    for i, c in enumerate(TOP5_CARDS):
        cy = card_y + i * (card_h + 0.10)
        # 카드 외곽
        add_rect(slide, LEFT_X + 0.25, cy, LEFT_W - 0.5, card_h,
                 fill=LIGHT_GRAY, line=LINE_GRAY, line_w=0.5)
        # 좌측 점수 띠
        score_color = BLUE if c["score_emph"] else DIM_GRAY
        add_rect(slide, LEFT_X + 0.25, cy, 1.6, card_h,
                 fill=score_color, line=None)
        add_text(
            slide, LEFT_X + 0.25, cy, 1.6, card_h,
            [(c["score"], {"size": SZ_SECTION, "bold": True, "color": WHITE,
                           "align": PP_ALIGN.CENTER})],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # 우측 콘텐츠
        add_text(
            slide, LEFT_X + 2.05, cy + 0.10, LEFT_W - 2.4, card_h - 0.20,
            [
                (f"{c['rank']}  ",
                 {"size": SZ_SECTION, "bold": True, "color": BLACK}),
                (f"  ({c['tag']})    {c['summary']}",
                 {"size": SZ_BODY, "bold": True, "color": BLACK,
                  "align": PP_ALIGN.LEFT}),
                (f" • {c['bullets'][0]}",
                 {"size": SZ_SUB, "color": BLACK}),
                (f" • {c['bullets'][1]}",
                 {"size": SZ_SUB, "color": BLACK}),
            ],
        )

    # ── 우측 상: 검토 추천 6개 ───────────────────────
    R_X = 17.20
    R_W = 9.75
    add_label_band(slide, R_X, TOP_Y, R_W, "② 검토 추천 6개")
    add_rect(slide, R_X, TOP_Y + 0.7, R_W, 7.55,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, R_X + 0.25, TOP_Y + 0.90, R_W - 0.5, 7.30,
        [
            (f" {label}", {"size": SZ_BODY, "bold": True, "color": BLACK}) if i % 2 == 0
            else (f"    {label}", {"size": SZ_SUB, "color": MID_GRAY})
            for pair in [[(name, ""), ("", desc)] for name, desc in REVIEW6]
            for i, (label, _) in enumerate(pair)
        ],
    )
    # 위 표현이 가독성 떨어지므로 더 명시적으로 재작성
    # (위 add_text 는 placeholder; 아래로 다시 정확히 그림)
    # placeholder textbox 위에 덮어쓰기 위해 다시 그림.
    review_lines = []
    for name, desc in REVIEW6:
        review_lines.append(
            (f" {name}", {"size": SZ_BODY, "bold": True, "color": BLACK})
        )
        review_lines.append(
            (f"     {desc}", {"size": SZ_SUB, "color": MID_GRAY})
        )
        review_lines.append(("", {"size": 4}))
    # 다시 한 번 add_text 로 덮어씌움
    add_rect(slide, R_X + 0.05, TOP_Y + 0.75, R_W - 0.1, 7.45,
             fill=WHITE, line=None)
    add_text(
        slide, R_X + 0.25, TOP_Y + 0.90, R_W - 0.5, 7.30,
        review_lines,
    )

    # ── 우측 하: 도입 효과 3줄 ────────────────────────
    EFF_Y = TOP_Y + 8.55
    add_label_band(slide, R_X, EFF_Y, R_W, "③ 기대 효과")
    add_rect(slide, R_X, EFF_Y + 0.7, R_W, 5.0,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, R_X + 0.25, EFF_Y + 0.85, R_W - 0.5, 4.75,
        [
            ("[Top 3 도입 시]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            ("", {"size": 3}),
            (" • KPI/Capex 엑셀 환승 비용",
             {"size": SZ_BODY, "color": BLACK}),
            ("     50% 이상 절감",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("     (xlsx + skill-creator)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 3}),
            (" • 협력사 PDF 검색 커버리지",
             {"size": SZ_BODY, "color": BLACK}),
            ("     0건 → 全 사양서 검색",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("     (pdf + ops-teammate)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 3}),
            (" • 모든 산출물 LGES 색/폰트 일관 (brand-guidelines)",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── 풋노트 ───────────────────────────────────────
    add_text(
        slide, 0.6, 18.30, 26.3, 0.55,
        [("* 출처 : outputs/popular_skills_review_2026-05-20.md  │  "
          "anthropic/skills 리포  │  본 세션 Explore 보고 (2026-05-20, FA기술혁신파트)",
          {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )


def build_aux_slide(prs):
    """보조 1장: 23개 후보 매트릭스 축약표."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [("[보조 1/1] 23개 후보 스킬 매트릭스 (공식 17 + 커뮤니티 6)",
          {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.8, 0.35, 6.4, 0.8,
        [("FA기술혁신파트 · 2026-05-20",
          {"size": SZ_BODY, "bold": True, "color": DIM_GRAY,
           "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)

    # 표 구성
    rows_data = [
        # (스킬, 카테고리, FA적합도, 결합 자산, 분류)
        ("xlsx", "공식·Doc", "상", "data-teammate, sample_data/", "Top5"),
        ("pdf", "공식·Doc", "상", "ops-teammate, references/기술자료/", "Top5"),
        ("skill-creator", "공식·Meta", "상", "신규 사내 스킬 라이브러리", "Top5"),
        ("brand-guidelines", "공식·Design", "상", "ppt-writer, LGES PPT 가이드", "Top5"),
        ("mcp-builder", "공식·Eng", "상", "data/ops-teammate, 사내시스템", "Top5"),
        ("docx", "공식·Doc", "중", "meeting-merger, document-writer", "검토"),
        ("doc-coauthoring", "공식·Comm", "중", "document-writer", "검토"),
        ("internal-comms", "공식·Comm", "중", "document-writer", "검토"),
        ("webapp-testing", "공식·Eng", "중", "(장기) 사내 대시보드", "검토"),
        ("Firecrawl", "커뮤니티", "중", "tech-research-teammate", "검토"),
        ("Sentry CodeReview", "커뮤니티", "중", "scripts/ 빌더 품질 게이트", "검토"),
        ("pptx", "공식·Doc", "중", "ppt-writer (이미 python-pptx)", "검토"),
        ("canvas-design", "공식·Design", "하", "—", "보류"),
        ("theme-factory", "공식·Design", "하", "—", "보류"),
        ("algorithmic-art", "공식·Design", "하", "—", "보류"),
        ("frontend-design", "공식·Eng", "하", "—", "보류"),
        ("web-artifacts-builder", "공식·Eng", "하", "—", "보류"),
        ("slack-gif-creator", "공식·Comm", "하", "— (Slack 미사용)", "보류"),
        ("claude-api", "공식·Eng", "중", "이미 활용 중", "보류"),
        ("Superpowers", "커뮤니티", "하", "SOP 중복", "보류"),
        ("Trail of Bits", "커뮤니티", "하", "보안 감사 우선순위 낮음", "보류"),
        ("Vercel Guidelines", "커뮤니티", "하", "UI 산출 없음", "보류"),
        ("HF/Tinybird/Nimble 等", "커뮤니티", "하", "도메인 무관", "보류"),
    ]

    tbl_left = 0.6
    tbl_top = 1.65
    tbl_w = 26.3
    tbl_h = 15.7
    n_rows = len(rows_data) + 1  # +1 header
    tbl_shape = slide.shapes.add_table(
        n_rows, 5, Cm(tbl_left), Cm(tbl_top), Cm(tbl_w), Cm(tbl_h)
    )
    tbl = tbl_shape.table
    col_widths = [4.8, 3.2, 2.0, 12.0, 4.3]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Cm(w)
    headers = ["스킬", "카테고리", "FA 적합도", "결합 자산", "분류"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

    for ri, row in enumerate(rows_data, start=1):
        skill, cat, fit, asset, cls = row
        # 분류별 배경 (보조 슬라이드는 회색·검정 위주, 강조 파랑 최소화)
        if cls == "Top5":
            bg = SOFT_GRAY
            txt_color = BLACK
            bold_first = True
        elif cls == "검토":
            bg = WHITE
            txt_color = BLACK
            bold_first = True
        else:
            bg = WHITE
            txt_color = MID_GRAY
            bold_first = False

        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            align = PP_ALIGN.CENTER if ci in (1, 2, 4) else PP_ALIGN.LEFT
            bold = bold_first if ci == 0 else False
            # FA 적합도 = 상 인 행의 적합도 셀만 검정 Bold (파랑 강조는 메인 한정)
            color = txt_color
            if ci == 2 and val == "상" and cls == "Top5":
                color = BLACK
                bold = True
            write_lines(cell.text_frame, [val], size=SZ_SUB, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    # 풋노트
    add_text(
        slide, 0.6, 18.30, 26.3, 0.55,
        [("* 출처 : outputs/popular_skills_review_2026-05-20.md (Top 10 6필드 상세)  │  "
          "anthropic/skills, firecrawl.dev, nimbleway.com  │  본 세션 조사 2026-05-20",
          {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )


def build():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)
    build_main_slide(prs)
    build_aux_slide(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
