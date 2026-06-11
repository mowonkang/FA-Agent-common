"""FA 기술혁신 Agent 팀 소개자료 PPT 양식 생성기.

`templates/소개자료.pptx` 를 생성한다. 14 본 슬라이드 + 부록 2 슬라이드 = 총 16장.
각 슬라이드는 빈칸 placeholder(`{{ key }}`) 와 미리 그려둔 도형으로 구성된다.

색상 팔레트는 `.claude/skills/meeting-merger/SKILL.md` 의 기준을 따른다.
- 메인: 딥 네이비 1C2B3A
- 강조: 웜 골드 B8883A
- 텍스트: 차콜 2D2D2D
- 배경 보조: 라이트 그레이 F5F5F5
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

NAVY = RGBColor(0x1C, 0x2B, 0x3A)
GOLD = RGBColor(0xB8, 0x88, 0x3A)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BLUE = RGBColor(0xEB, 0xF0, 0xF8)
SOFT_GREEN = RGBColor(0xEA, 0xF4, 0xEC)
SOFT_PEACH = RGBColor(0xFA, 0xEE, 0xE9)


def _set_run(run, text, size_pt, bold=False, color=CHARCOAL):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(slide, x, y, w, h, text, size_pt=14, bold=False, color=CHARCOAL,
              align=None, name=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    _set_run(run, text, size_pt, bold, color)
    return box


def _add_band(slide, x, y, w, h, fill=NAVY):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.text_frame.text = ""
    return box


def _add_title_bar(slide, title_text, slide_no=None):
    """슬라이드 상단 네이비 바 + 흰 제목."""
    _add_band(slide, 0, 0, 33.867, 1.6, fill=NAVY)
    box = slide.shapes.add_textbox(Cm(0.8), Cm(0.25), Cm(28), Cm(1.1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, title_text, 22, bold=True, color=WHITE)
    if slide_no is not None:
        no_box = slide.shapes.add_textbox(Cm(31.5), Cm(0.4), Cm(2), Cm(0.9))
        nrun = no_box.text_frame.paragraphs[0].add_run()
        _set_run(nrun, f"{slide_no:02d}", 14, bold=True, color=GOLD)


def _add_card(slide, x, y, w, h, header, body_placeholder,
              header_color=NAVY, body_color=CHARCOAL,
              fill=LIGHT_GRAY, header_size=14, body_size=12, name_prefix=None):
    """카드: 상단 헤더 띠(컬러) + 본문 placeholder."""
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Cm(x), Cm(y), Cm(w), Cm(h))
    base.fill.solid()
    base.fill.fore_color.rgb = fill
    base.line.color.rgb = header_color
    base.line.width = Pt(0.75)
    base.text_frame.text = ""

    head_h = 1.0
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(x), Cm(y), Cm(w), Cm(head_h))
    head.fill.solid()
    head.fill.fore_color.rgb = header_color
    head.line.fill.background()
    htf = head.text_frame
    htf.margin_left = Cm(0.3)
    htf.margin_top = Cm(0.1)
    hrun = htf.paragraphs[0].add_run()
    _set_run(hrun, header, header_size, bold=True, color=WHITE)
    if name_prefix:
        head.name = f"{name_prefix}_header"

    body = slide.shapes.add_textbox(Cm(x + 0.2), Cm(y + head_h + 0.2),
                                    Cm(w - 0.4), Cm(h - head_h - 0.4))
    body.text_frame.word_wrap = True
    brun = body.text_frame.paragraphs[0].add_run()
    _set_run(brun, body_placeholder, body_size, bold=False, color=body_color)
    if name_prefix:
        body.name = f"{name_prefix}_body"


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Cm(33.867)   # 16:9
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # ───── Slide 1: 표지 ─────
    s = prs.slides.add_slide(blank)
    _add_band(s, 0, 0, 33.867, 19.05, fill=NAVY)
    _add_band(s, 0, 8.5, 33.867, 0.15, fill=GOLD)
    title = _add_text(s, 2, 5.5, 30, 2.5, "{{ deck_title }}",
                      size_pt=40, bold=True, color=WHITE, name="deck_title")
    sub = _add_text(s, 2, 9.0, 30, 1.5, "{{ subtitle }}",
                    size_pt=20, bold=False, color=GOLD, name="subtitle")
    meta = _add_text(s, 2, 15.5, 30, 1.0,
                     "작성자 {{ author }}  |  작성일 {{ created_at }}  |  보고 대상 {{ audience }}",
                     size_pt=14, bold=False, color=WHITE, name="cover_meta")

    # ───── Slide 2: Executive Summary ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Executive Summary — 한 장 요약", slide_no=2)
    _add_text(s, 1.0, 2.2, 30, 0.8,
              "FA 기술혁신 Agent 팀, 4가지 핵심 가치",
              size_pt=14, bold=True, color=NAVY)
    cards_y = 3.4
    card_w, card_h, gap = 7.6, 6.5, 0.6
    for i, (title_, key) in enumerate([
        ("① 5인 AI 협업조", "summary_pt1"),
        ("② 출처 자동검증", "summary_pt2"),
        ("③ 도메인 스킬 3개", "summary_pt3"),
        ("④ 즉시 적용 ROI", "summary_pt4"),
    ]):
        x = 1.0 + i * (card_w + gap)
        _add_card(s, x, cards_y, card_w, card_h, title_,
                  "{{ " + key + " }}",
                  header_color=NAVY, fill=LIGHT_GRAY,
                  header_size=14, body_size=12, name_prefix=f"summary_card_{i+1}")
    _add_text(s, 1.0, 17.0, 30, 1.0,
              "출처: {{ summary_source }}",
              size_pt=10, bold=False, color=GOLD, name="summary_src")

    # ───── Slide 3: 문제 정의 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "문제 정의 — 보고 업무의 비효율", slide_no=3)
    # Before
    _add_band(s, 1.0, 2.5, 15, 1.0, fill=SOFT_PEACH)
    _add_text(s, 1.2, 2.65, 14, 0.8, "Before — 기존 방식",
              size_pt=14, bold=True, color=NAVY)
    bbox = s.shapes.add_textbox(Cm(1.0), Cm(3.7), Cm(15), Cm(12))
    bbox.text_frame.word_wrap = True
    brun = bbox.text_frame.paragraphs[0].add_run()
    _set_run(brun, "{{ problem_before }}", 13, color=CHARCOAL)
    bbox.name = "problem_before"
    # After
    _add_band(s, 17.5, 2.5, 15, 1.0, fill=SOFT_GREEN)
    _add_text(s, 17.7, 2.65, 14, 0.8, "After — Agent 팀 도입 후",
              size_pt=14, bold=True, color=NAVY)
    abox = s.shapes.add_textbox(Cm(17.5), Cm(3.7), Cm(15), Cm(12))
    abox.text_frame.word_wrap = True
    arun = abox.text_frame.paragraphs[0].add_run()
    _set_run(arun, "{{ problem_after }}", 13, color=CHARCOAL)
    abox.name = "problem_after"
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ problem_source }}",
              size_pt=10, color=GOLD, name="problem_src")

    # ───── Slide 4: 단일 vs Agent Teams ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "기존 단일 에이전트 vs FA Agent Teams", slide_no=4)
    table_shape = s.shapes.add_table(5, 3, Cm(1.0), Cm(2.5), Cm(31.8), Cm(13))
    table = table_shape.table
    headers = ["관점", "기존 Subagents (단일 호출)", "FA Agent Teams (본 팀)"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        _set_run(run, h, 14, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    rows_def = [
        ("협업 구조", "{{ comp_left_1 }}", "{{ comp_right_1 }}"),
        ("통신 방식", "{{ comp_left_2 }}", "{{ comp_right_2 }}"),
        ("실행 모드", "{{ comp_left_3 }}", "{{ comp_right_3 }}"),
        ("품질 보장", "{{ comp_left_4 }}", "{{ comp_right_4 }}"),
    ]
    for r, (label, left, right) in enumerate(rows_def, start=1):
        for c, t in enumerate([label, left, right]):
            cell = table.cell(r, c)
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            _set_run(run, t, 12, bold=(c == 0), color=CHARCOAL)
            if c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ comparison_source }}",
              size_pt=10, color=GOLD, name="comp_src")

    # ───── Slide 5: 시스템 구성도 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "시스템 구성도 — 5 에이전트 + 3 스킬 + 2 Hook", slide_no=5)
    # 사용자 박스
    _add_band(s, 14, 2.3, 6, 1.2, fill=GOLD)
    _add_text(s, 14, 2.45, 6, 1.0, "사용자 / 팀장",
              size_pt=14, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, name="arch_user")
    # 5명 에이전트 카드 (1행 5열)
    agents = [
        ("data-teammate", "KPI·Capex·MRM·DST"),
        ("ops-teammate", "협력사·HR·결재·문서"),
        ("tech-research-teammate", "외부 동향·경쟁사·평가"),
        ("document-writer", "Markdown 양식 채움"),
        ("ppt-writer", "PPT 양식 채움"),
    ]
    card_w = 6.1
    gap = 0.4
    start_x = (33.867 - (card_w * 5 + gap * 4)) / 2  # 가운데 정렬
    for i, (a, role) in enumerate(agents):
        x = start_x + i * (card_w + gap)
        _add_card(s, x, 4.5, card_w, 3.8, a, role,
                  header_color=NAVY, fill=LIGHT_GRAY,
                  header_size=11, body_size=10,
                  name_prefix=f"arch_agent_{i+1}")
    # 3개 스킬
    skills = [
        ("meeting-minutes", "녹취 → 회의록 PPTX"),
        ("meeting-merger", "녹취+OneNote → Word"),
        ("fa-task-discovery", "FA 6분야 → 과제 docx"),
    ]
    for i, (sk, desc) in enumerate(skills):
        x = 1.0 + i * 10.7
        _add_card(s, x, 9.5, 10.0, 3.0, sk, desc,
                  header_color=GOLD, fill=LIGHT_GRAY,
                  header_size=13, body_size=11,
                  name_prefix=f"arch_skill_{i+1}")
    # 2 Hook 띠
    _add_band(s, 1.0, 13.5, 31.8, 1.6, fill=SOFT_BLUE)
    _add_text(s, 1.3, 13.7, 31, 1.2,
              "[Hook 1] SessionStart → references/ 자동 인덱싱     "
              "|     [Hook 2] TaskCompleted → outputs/*.md 출처 자동 검증",
              size_pt=12, bold=True, color=NAVY, name="arch_hooks")
    _add_text(s, 1.0, 16.0, 31, 1.2,
              "{{ architecture_caption }}",
              size_pt=11, color=CHARCOAL, name="arch_caption")
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ architecture_source }}",
              size_pt=10, color=GOLD, name="arch_src")

    # ───── Slide 6: 팀원 5명 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "팀원 5명 — 사내 수치·운영 + 사외 리서치 + 작성", slide_no=6)
    # 1행: 상단 강조 카드 (tech-research-teammate — 외부 시야)
    _add_card(s, 1.0, 2.2, 31.8, 3.6,
              "tech-research-teammate  ✦ 외부 시야 담당",
              "{{ team_card_5 }}",
              header_color=GOLD, fill=LIGHT_GRAY,
              header_size=15, body_size=11,
              name_prefix="team_5")
    # 2행: 2×2 기존 4명
    positions = [(1.0, 6.2), (17.5, 6.2), (1.0, 12.3), (17.5, 12.3)]
    titles = [
        "data-teammate",
        "ops-teammate",
        "document-writer",
        "ppt-writer",
    ]
    for i, ((x, y), t) in enumerate(zip(positions, titles), start=1):
        _add_card(s, x, y, 15.3, 5.7, t,
                  "{{ team_card_" + str(i) + " }}",
                  header_color=NAVY, fill=LIGHT_GRAY,
                  header_size=14, body_size=11,
                  name_prefix=f"team_{i}")
    _add_text(s, 1.0, 18.2, 31, 0.7,
              "출처: {{ team_source }}",
              size_pt=10, color=GOLD, name="team_src")

    # ───── Slide 7: 도메인 스킬 3개 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "도메인 스킬 3개 — 회의·기술 동향 자동화", slide_no=7)
    skill_titles = [
        ("meeting-minutes", "녹취 → 회의록 PPTX"),
        ("meeting-merger", "녹취+OneNote → Word"),
        ("fa-task-discovery", "FA 6분야 → 과제 docx + 카톡"),
    ]
    for i, (sk, sub_) in enumerate(skill_titles, start=1):
        x = 1.0 + (i - 1) * 10.7
        _add_card(s, x, 2.5, 10.0, 13.5,
                  f"{sk}",
                  "{{ skill_card_" + str(i) + " }}",
                  header_color=GOLD, fill=LIGHT_GRAY,
                  header_size=14, body_size=11,
                  name_prefix=f"skill_{i}")
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ skills_source }}",
              size_pt=10, color=GOLD, name="skills_src")

    # ───── Slide 8: SOP ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "어떻게 만들었는가 — Agent Teams SOP 4단계", slide_no=8)
    steps = [
        ("Step 1\n분해", "{{ sop_step_1 }}"),
        ("Step 2\n병렬 수집", "{{ sop_step_2 }}"),
        ("Step 3\n위임 작성", "{{ sop_step_3 }}"),
        ("Step 4\n품질 게이트", "{{ sop_step_4 }}"),
    ]
    for i, (head, body_) in enumerate(steps):
        x = 1.0 + i * 8.0
        _add_card(s, x, 3.5, 7.4, 11, head, body_,
                  header_color=NAVY, fill=LIGHT_GRAY,
                  header_size=14, body_size=11,
                  name_prefix=f"sop_step_{i+1}")
    # 화살표 띠
    _add_band(s, 8.4, 7.5, 0.6, 0.6, fill=GOLD)
    _add_band(s, 16.4, 7.5, 0.6, 0.6, fill=GOLD)
    _add_band(s, 24.4, 7.5, 0.6, 0.6, fill=GOLD)
    _add_text(s, 1.0, 16.0, 31, 1.2,
              "{{ sop_caption }}",
              size_pt=11, color=CHARCOAL, name="sop_caption")
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ sop_source }}",
              size_pt=10, color=GOLD, name="sop_src")

    # ───── Slide 9: 특징 ① 격리 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "특징 ① 병렬·컨텍스트 격리", slide_no=9)
    _add_text(s, 1.0, 2.5, 31, 1.0,
              "팀원은 리더의 대화 기록을 받지 않고, 호출 프롬프트만 본다",
              size_pt=15, bold=True, color=NAVY)
    # 격리 박스 3개
    boxes = [
        ("토큰 절감", "{{ iso_card_1 }}"),
        ("환각 차단", "{{ iso_card_2 }}"),
        ("스코프 명확", "{{ iso_card_3 }}"),
    ]
    for i, (h, b) in enumerate(boxes):
        x = 1.0 + i * 10.7
        _add_card(s, x, 4.0, 10.0, 11.5, h, b,
                  header_color=GOLD, fill=LIGHT_GRAY,
                  header_size=14, body_size=12,
                  name_prefix=f"iso_card_{i+1}")
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ iso_source }}",
              size_pt=10, color=GOLD, name="iso_src")

    # ───── Slide 10: 특징 ② 검증·결정론 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "특징 ② 출처 자동검증 + 결정론 PPT 파이프라인", slide_no=10)
    # 좌측: hook 코드 박스
    _add_band(s, 1.0, 2.5, 15, 1.0, fill=NAVY)
    _add_text(s, 1.2, 2.65, 14, 0.8,
              "TaskCompleted Hook (출처 검증)",
              size_pt=13, bold=True, color=WHITE)
    code_box = s.shapes.add_textbox(Cm(1.0), Cm(3.7), Cm(15), Cm(8))
    code_box.text_frame.word_wrap = True
    crun = code_box.text_frame.paragraphs[0].add_run()
    _set_run(crun, "{{ hook_snippet }}", 11, color=CHARCOAL)
    code_box.name = "hook_snippet"
    # 우측: 파이프라인 박스
    _add_band(s, 17.5, 2.5, 15.3, 1.0, fill=NAVY)
    _add_text(s, 17.7, 2.65, 14, 0.8,
              "결정론 PPT 파이프라인",
              size_pt=13, bold=True, color=WHITE)
    pipe_box = s.shapes.add_textbox(Cm(17.5), Cm(3.7), Cm(15.3), Cm(8))
    pipe_box.text_frame.word_wrap = True
    prun = pipe_box.text_frame.paragraphs[0].add_run()
    _set_run(prun, "{{ pipeline_snippet }}", 11, color=CHARCOAL)
    pipe_box.name = "pipeline_snippet"
    _add_text(s, 1.0, 12.5, 31, 4.5,
              "{{ verification_caption }}",
              size_pt=12, color=CHARCOAL, name="ver_caption")
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ verification_source }}",
              size_pt=10, color=GOLD, name="ver_src")

    # ───── Slide 11: 할 수 있는 일 6가지 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "할 수 있는 일 — 6가지 산출 케이스", slide_no=11)
    items = [
        "주간보고 (.md)",
        "월간보고 (.md)",
        "회의록 (.md)",
        "회의록 (.pptx, 녹취 기반)",
        "회의록 (.docx, 녹취+OneNote)",
        "FA 기술과제 발굴 (.docx + 카톡)",
    ]
    for i, t in enumerate(items):
        col = i % 3
        row = i // 3
        x = 1.0 + col * 10.7
        y = 2.5 + row * 6.8
        _add_card(s, x, y, 10.0, 6.0, t,
                  "{{ cap_card_" + str(i + 1) + " }}",
                  header_color=NAVY, fill=LIGHT_GRAY,
                  header_size=13, body_size=11,
                  name_prefix=f"cap_card_{i+1}")
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ capabilities_source }}",
              size_pt=10, color=GOLD, name="cap_src")

    # ───── Slide 12: 협업 시나리오 — 우편함 develop 흐름 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "협업 시나리오 — 팀원 간 우편함 develop 흐름", slide_no=12)
    _add_text(s, 1.0, 2.1, 31.8, 1.0,
              "팀원들이 독립적이 아닌, 우편함으로 직접 소통하며 결과를 develop 한다 "
              "— 같은 데이터를 여러 라운드에 걸쳐 정제·심화한다",
              size_pt=13, bold=True, color=NAVY)
    # 좌측 카드: 시나리오 ①
    _add_card(s, 1.0, 3.3, 15.5, 12.7,
              "시나리오 ① 과제 진척률 심층 보고",
              "{{ collab_scenario_1 }}",
              header_color=NAVY, fill=LIGHT_GRAY,
              header_size=14, body_size=11,
              name_prefix="collab_card_1")
    # 우측 카드: 시나리오 ②
    _add_card(s, 17.3, 3.3, 15.5, 12.7,
              "시나리오 ② 외부 동향 → KPI 재설정 제안",
              "{{ collab_scenario_2 }}",
              header_color=GOLD, fill=LIGHT_GRAY,
              header_size=14, body_size=11,
              name_prefix="collab_card_2")
    _add_text(s, 1.0, 16.4, 31, 1.0,
              "{{ collab_caption }}",
              size_pt=11, color=CHARCOAL, name="collab_caption")
    _add_text(s, 1.0, 17.8, 31, 0.7,
              "출처: {{ collab_source }}",
              size_pt=10, color=GOLD, name="collab_src")

    # ───── Slide 13: 데모 ① 주간보고 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "데모 ① 주간보고 — outputs/주간보고_2026-05-12.md", slide_no=13)
    # 좌측: 이미지 슬롯 (사용자가 캡처 삽입)
    _add_band(s, 1.0, 2.5, 18, 14, fill=LIGHT_GRAY)
    _add_text(s, 1.5, 8.5, 17, 1.5,
              "[ 이미지 슬롯 ] outputs/주간보고_2026-05-12.md 캡처 삽입",
              size_pt=14, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, name="demo_weekly_image_slot")
    _add_text(s, 1.5, 10.0, 17, 1.0,
              "(빨간 박스로 (출처: ...) 부분 강조)",
              size_pt=11, color=CHARCOAL,
              align=PP_ALIGN.CENTER, name="demo_weekly_image_caption")
    # 우측: 설명
    _add_band(s, 20.5, 2.5, 12.3, 1.0, fill=NAVY)
    _add_text(s, 20.7, 2.65, 12, 0.8, "데모 포인트",
              size_pt=13, bold=True, color=WHITE)
    desc = s.shapes.add_textbox(Cm(20.5), Cm(3.8), Cm(12.3), Cm(12.5))
    desc.text_frame.word_wrap = True
    drun = desc.text_frame.paragraphs[0].add_run()
    _set_run(drun, "{{ demo_weekly_points }}", 12, color=CHARCOAL)
    desc.name = "demo_weekly_points"
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ demo_weekly_source }}",
              size_pt=10, color=GOLD, name="demo_weekly_src")

    # ───── Slide 14: 데모 ② 로드맵 + MMF ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "데모 ② 로드맵 v2 & MMF 컨셉", slide_no=14)
    # 좌측: 로드맵 캡처 슬롯
    _add_band(s, 1.0, 2.5, 15, 13.5, fill=LIGHT_GRAY)
    _add_text(s, 1.5, 8.0, 14, 1.0,
              "[ 이미지 슬롯 ] roadmap_v2.pptx 4축 KPI 슬라이드",
              size_pt=13, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, name="demo_roadmap_slot1")
    # 우측: MMF SVG 슬롯
    _add_band(s, 17.5, 2.5, 15.3, 13.5, fill=LIGHT_GRAY)
    _add_text(s, 18.0, 8.0, 14.5, 1.0,
              "[ 이미지 슬롯 ] MMF Pop-up Factory SVG",
              size_pt=13, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, name="demo_roadmap_slot2")
    _add_text(s, 1.0, 16.5, 31, 1.0,
              "{{ demo_roadmap_caption }}",
              size_pt=11, color=CHARCOAL, name="demo_roadmap_cap")
    _add_text(s, 1.0, 17.8, 31, 0.7,
              "출처: {{ demo_roadmap_source }}",
              size_pt=10, color=GOLD, name="demo_roadmap_src")

    # ───── Slide 15: ROI · 다음 단계 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "ROI · 적용 효과 · 다음 단계", slide_no=15)
    # 상단 KPI 박스 3개
    kpis = [
        ("작성 시간", "{{ roi_kpi_1 }}"),
        ("출처 검증률", "{{ roi_kpi_2 }}"),
        ("산출물 표준화", "{{ roi_kpi_3 }}"),
    ]
    for i, (h, v) in enumerate(kpis):
        x = 1.0 + i * 10.7
        _add_card(s, x, 2.5, 10.0, 5.0, h, v,
                  header_color=GOLD, fill=LIGHT_GRAY,
                  header_size=14, body_size=14,
                  name_prefix=f"roi_kpi_{i+1}")
    # 다음 단계 4개 체크리스트
    _add_band(s, 1.0, 8.0, 31.8, 1.0, fill=NAVY)
    _add_text(s, 1.2, 8.15, 31, 0.8,
              "다음 단계 (Next 4)",
              size_pt=14, bold=True, color=WHITE)
    for i in range(4):
        y = 9.5 + i * 1.7
        _add_band(s, 1.0, y, 0.4, 0.4, fill=GOLD)
        nb = s.shapes.add_textbox(Cm(1.7), Cm(y - 0.15), Cm(31), Cm(1.3))
        nb.text_frame.word_wrap = True
        nrun = nb.text_frame.paragraphs[0].add_run()
        _set_run(nrun, "{{ next_step_" + str(i + 1) + " }}", 13, color=CHARCOAL)
        nb.name = f"next_step_{i+1}"

    # ───── Slide 16 (부록 A): 회의록 데모 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "[부록 A] 회의록 자동화 데모", slide_no=16)
    _add_band(s, 1.0, 2.5, 18, 14, fill=LIGHT_GRAY)
    _add_text(s, 1.5, 8.5, 17, 1.5,
              "[ 이미지 슬롯 ] outputs/회의록_신규R&D과제착수_2026-05-12.md 캡처",
              size_pt=14, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, name="appendix_minutes_slot")
    desc = s.shapes.add_textbox(Cm(20.5), Cm(2.5), Cm(12.3), Cm(13.5))
    desc.text_frame.word_wrap = True
    drun = desc.text_frame.paragraphs[0].add_run()
    _set_run(drun, "{{ appendix_minutes_points }}", 12, color=CHARCOAL)
    desc.name = "appendix_minutes_points"
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ appendix_minutes_source }}",
              size_pt=10, color=GOLD, name="appendix_minutes_src")

    # ───── Slide 17 (부록 B): FA 과제 데모 ─────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "[부록 B] FA 기술과제 발굴 데모", slide_no=17)
    _add_band(s, 1.0, 2.5, 18, 14, fill=LIGHT_GRAY)
    _add_text(s, 1.5, 8.5, 17, 1.5,
              "[ 이미지 슬롯 ] outputs/FA_기술과제후보_20260512.docx 1페이지",
              size_pt=14, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, name="appendix_fa_slot")
    desc = s.shapes.add_textbox(Cm(20.5), Cm(2.5), Cm(12.3), Cm(13.5))
    desc.text_frame.word_wrap = True
    drun = desc.text_frame.paragraphs[0].add_run()
    _set_run(drun, "{{ appendix_fa_points }}", 12, color=CHARCOAL)
    desc.name = "appendix_fa_points"
    _add_text(s, 1.0, 17.5, 31, 0.8,
              "출처: {{ appendix_fa_source }}",
              size_pt=10, color=GOLD, name="appendix_fa_src")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build FA Agent intro PPT template.")
    parser.add_argument("--out", default="templates/소개자료.pptx")
    parser.add_argument("--force", action="store_true", help="기존 파일 덮어쓰기")
    args = parser.parse_args(argv)
    out = Path(args.out)
    if out.exists() and not args.force:
        sys.stderr.write(f"이미 존재합니다: {out}. 덮어쓰려면 --force 를 추가하세요.\n")
        return 1
    build(out)
    print(str(out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
