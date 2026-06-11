"""인프라그룹 AI경진대회 과제 출품 자료 빌더 (2026-05-19 갱신판).

`scripts/build_contest.py` (2026-05-13 판) 의 후속.
메인 양식 `templates/사내양식/260511_인프라그룹 AI경진대회 과제_양식.pptx`
의 Slide 1·2 placeholder 를 FA 기술혁신 Agent 팀 콘텐츠로 채우고,
2026-05-13 → 05-19 사이 추가 작업(LGES PPT 작업 가이드, 한·영 폰트
자동분리, 임원 보고 양식 카탈로그, 유연성지표·경쟁사 동향·전고체 검토
산출물, fa-task-discovery v2 등)을 반영한 보조 슬라이드 6장을 덧붙인다.

보조 슬라이드는 LGES_PPT_작업_가이드.md 기준(흰/회/검 + 파랑 강조 최소 +
초록곰팡이 풋노트 별도 박스 + 한·영 run 분리 + SZ_* 6단 사이즈)을 따른다.

출력:
  outputs/AI경진대회_FA기술혁신Agent팀_2026-05-19.pptx
  outputs/AI경진대회_FA기술혁신Agent팀_2026-05-19.md  (슬라이드 1:1 사이드카)
"""
from __future__ import annotations

import glob
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

# ─── 색 팔레트 (LGES_PPT_작업_가이드 §1) ──────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)

BLUE = RGBColor(0x00, 0x00, 0xFF)        # 강조 (포인트만, 슬라이드당 3~5)
LIGHT_BLUE = SOFT_GRAY                   # (사내 룰: 박스 배경=흰·회색만, alias)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)  # 풋노트 전용 초록곰팡이 (G102)

# 사내 룰: 본문=검정 / 강조=파랑 / 풋노트=진초록 / 박스=흰·회색
FORM_GREEN = BLUE                        # (구 폼 초록 → 강조 파랑으로 alias)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

# ─── 폰트 사이즈 스케일 (가이드 §2-2, 6단) ────────────────────────
SZ_TITLE = 16
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

TEMPLATE = glob.glob("templates/사내양식/260511*.pptx")
if not TEMPLATE:
    print("ERROR: template not found", file=sys.stderr)
    sys.exit(1)
TEMPLATE_PATH = TEMPLATE[0]
OUT_PPTX = Path("outputs/AI경진대회_FA기술혁신Agent팀_2026-05-19.pptx")
OUT_MD = Path("outputs/AI경진대회_FA기술혁신Agent팀_2026-05-19.md")


# ─── run/텍스트 헬퍼 (한·영 자동 분리, 가이드 §2-1) ───────────────
def set_run(run, text, *, font=FONT_BODY, size=SZ_BODY, bold=False,
            color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH
    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        ex = rPr.find(qn(f"a:{tag}"))
        if ex is not None:
            rPr.remove(ex)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=SZ_BODY, bold=False,
                color=BLACK, align=None, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, opts = item if isinstance(item, tuple) else (item, {})
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        run = para.add_run()
        set_run(run, text,
                font=opts.get("font", font),
                size=opts.get("size", size),
                bold=opts.get("bold", bold),
                color=opts.get("color", color))


def add_rect(slide, left, top, width, height, *, fill=WHITE, line=BLACK,
             line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(left), Cm(top), Cm(width), Cm(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    write_lines(tf, lines, **kw)
    return tb


def add_band(slide, left, top, width, text, *, h=0.72, fill=DIM_GRAY,
             color=WHITE, size=SZ_BAND):
    shp = add_rect(slide, left, top, width, h, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=size, bold=True, color=color)
    return shp


def add_card(slide, x, y, w, h, header, body_lines, *,
             head_fill=DIM_GRAY, head_color=WHITE, body_size=SZ_SUB,
             head_size=SZ_SECTION, base_fill=WHITE):
    add_rect(slide, x, y, w, h, fill=base_fill, line=MID_GRAY, line_w=0.75)
    head_h = 0.68
    bnd = add_rect(slide, x, y, w, head_h, fill=head_fill, line=None)
    bnd.text_frame.margin_left = Cm(0.15)
    bnd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bnd.text_frame, [header], size=head_size, bold=True,
                color=head_color)
    add_text(slide, x + 0.18, y + head_h + 0.08, w - 0.36,
             h - head_h - 0.18, body_lines, size=body_size, color=BLACK)


def add_footnote(slide, text):
    """풋노트 — 별도 텍스트박스, 초록곰팡이 #006600 Bold 8pt, y≈18.30."""
    add_text(slide, 0.6, 18.32, 26.3, 0.52,
             [(text, {"size": SZ_FOOT, "bold": True, "color": MOLD_GREEN})])


def _arrow_down(slide, cx_cm, y_top_cm, y_bot_cm, color=BLACK):
    """수직 화살표 (PowerPoint 네이티브 connector) — 편집 가능."""
    from pptx.util import Cm
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Cm(cx_cm), Cm(y_top_cm),
                                      Cm(cx_cm), Cm(y_bot_cm))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.0)
    # 화살표 끝 (end arrow)
    ln = conn.line._get_or_add_ln()
    from lxml import etree
    from pptx.oxml.ns import qn
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


def new_blank_slide(prs):
    """양식 레이아웃으로 슬라이드 추가 후 placeholder 제거 + 흰 배경."""
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for sh in list(s.shapes):
        if sh.is_placeholder:
            sh._element.getparent().remove(sh._element)
    add_rect(s, 0, 0, 27.52, 19.05, fill=WHITE, line=None)
    return s


def header(slide, label):
    """보조 슬라이드 상단 라벨 띠 + 우상단 책임자."""
    add_band(slide, 0.6, 0.5, 20.4, label, h=0.82, size=SZ_BAND)
    add_text(slide, 21.2, 0.62, 5.7, 0.6,
             [("FA기술혁신파트 강모원책임",
               {"size": SZ_BODY, "bold": True, "color": DIM_GRAY,
                "align": PP_ALIGN.RIGHT})])


# ─── 메인 양식 슬라이드 (Slide 1·2) ───────────────────────────────
TITLE = "FA 기술혁신 Agent 팀"
SUBTITLE = ("5 AI 에이전트 팀이 FA 9개 업무 영역을 자동화 — "
            "보고 완전 자동 · 외부 동향 매일 발굴 · 출처 100% 자동 검증")
RESPONSIBLE = "FA기술혁신파트 강모원책임"

ISSUE = [
    " • 주간/월간/회의록 작성 수시간 소요",
    " • 출처 누락·양식 제각각·재사용 난해",
    " • 회의록·녹취 정리 지연, 1인 부담",
    " • 외부 동향·경쟁사 KPI 반영 지연",
    " • 보고 표준 부재로 자산 재사용 불가",
    " • Excel/PDF/DOCX 산출 양식 분산",
    " • 사내 도구·MCP 별도 연결 미비",
]
# 개선 사항 — 1개 bullet 통째 BLUE Bold 강조 (Agent Teams 차별점)
IMPROVE = [
    (" • 5 AI 에이전트 우편함 협업", {}),
    (" • 에이전트끼리 직접 소통 (차별점)",
     {"bold": True, "color": BLUE}),
    (" • 도메인 스킬 14종 자동 발동",
     {"bold": True}),
    (" • MCP 16+ 자동 (카카오·Google·GitHub)",
     {"bold": True}),
    (" • 결정론 PPT — 양식 무결성 100%", {}),
    (" • LGES 가이드 + brand-guidelines 자동", {}),
    (" • Hook 2 + sentry-code-review 게이트", {}),
]
# 기대 효과 — 정량 5 + 정성 4, [정량]/[정성] 헤더만 BLUE Bold 강조
EFFECT = [
    ("[정량 효과]", {"bold": True, "color": BLUE}),
    (" • 작성 시간 수시간 → 분 단위", {}),
    (" • 출처 100% 자동 검증 (미기입 차단)", {}),
    (" • 검증 산출물 10+ 종 누적", {}),
    (" • Capex 자동 (예: △806.1억 21.3%)", {}),
    (" • 신규 스킬 11종 add — 산출 양식 확장", {}),
    ("[정성 효과]", {"bold": True, "color": BLUE}),
    (" • 신뢰성·재현성·감사가능성 확보", {}),
    (" • references/ 누적 → 보고 품질 ↑", {}),
    (" • brand-guidelines 자동 일관성", {}),
    (" • 워크플로우 동일 → 타 영역 확장", {}),
]


def fill_form(prs):
    """Slide 1: 폼 placeholder 채움 (사용자 피드백 반영판).

    피드백 반영:
    - 메인 타이틀·책임자 색 BLUE → BLACK (강조는 본문 내부에만)
    - SmartArt 흐름 띠 제거 (개선 전/후 chip 과 오버랩)
    - 템플릿 노트 "* 개인 참여시 1명 만 기입" 제거 (스킬 카드와 오버랩)
    - 참여 인원 표 H 압축, 우하 카드를 "FA 9업무 자동화 수준 요약" 으로 교체
      → 보조 ⑧ 내용을 한 장에서 즉시 볼 수 있도록.
    """
    s1 = prs.slides[0]

    # [5] 과제 명 — 모든 텍스트 BLACK (피드백: 제목은 검정)
    tf = s1.shapes[5].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    for seg in (("과제 명 ", BLACK), (": ", BLACK), (TITLE, BLACK)):
        r = p.add_run()
        set_run(r, seg[0], font=FONT_BODY, size=SZ_TITLE, bold=True,
                color=seg[1])
    # [6] 부제 — BLACK
    write_lines(s1.shapes[6].text_frame, [
        (SUBTITLE, {"font": FONT_BODY, "size": 12, "bold": True,
                    "color": BLACK})])

    def fit_box(sh, x, y, w, h):
        sh.left, sh.top, sh.width, sh.height = Cm(x), Cm(y), Cm(w), Cm(h)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

    # [7] 이슈 및 배경 (좌상)
    write_lines(s1.shapes[7].text_frame, ISSUE, size=SZ_BODY, color=BLACK)
    fit_box(s1.shapes[7], 0.76, 4.40, 12.5, 6.6)
    # [8] 우상단 책임자 — BLACK (피드백: 제목은 검정)
    write_lines(s1.shapes[8].text_frame, [RESPONSIBLE], font=FONT_EMPH,
                size=SZ_BODY, bold=True, color=BLACK,
                align=PP_ALIGN.RIGHT)
    # [10] 기대 효과 (좌하)
    write_lines(s1.shapes[10].text_frame, EFFECT, size=SZ_BODY, color=BLACK)
    fit_box(s1.shapes[10], 0.76, 12.16, 12.5, 6.0)
    # [12] 개선 사항 (우상)
    write_lines(s1.shapes[12].text_frame, IMPROVE, size=SZ_BODY, color=BLACK)
    fit_box(s1.shapes[12], 14.50, 4.40, 12.4, 6.6)

    # [14] 참여 인원 표 — 1행 (개인) · 표 H 압축
    tbl = s1.shapes[14].table

    def cell(r, c, txt, *, align=PP_ALIGN.CENTER):
        write_lines(tbl.cell(r, c).text_frame, [txt], size=SZ_BODY,
                    color=BLACK, align=align)

    cell(1, 0, "1")
    cell(1, 1, "FA기술혁신파트")
    cell(1, 2, "강모원")
    cell(1, 3, "FA 기술혁신 Agent 팀 설계·구축·운영 "
               "(5 에이전트 + 스킬 14종 + Hook 2 + MCP 16+)",
         align=PP_ALIGN.LEFT)
    try:
        s1.shapes[14].height = Cm(1.40)
    except Exception:
        pass

    # [15] 템플릿 노트 "* 개인 참여시 1명 만 기입" 제거
    # (스킬 카드 영역과 오버랩 → 불필요한 양식 안내문)
    try:
        note_el = s1.shapes[15]._element
        note_el.getparent().remove(note_el)
    except Exception:
        pass

    # ── 우하 신규 카드 : FA 9업무 자동화 수준 요약 (보조 ⑧ 핵심) ──
    # 사용자 피드백 #4: "현재 FA기술혁신파트 업무 + 자동화 수준" 메인장 노출
    # 위치: 참여인원표 아래 (T14.05), 풋노트 위 (T18.20) 까지 = H 3.20cm
    card_y = 14.05
    card_h = 3.20
    add_rect(s1, 14.22, card_y, 12.70, card_h, fill=WHITE, line=MID_GRAY,
             line_w=0.6)
    bn = add_rect(s1, 14.22, card_y, 12.70, 0.45, fill=DIM_GRAY, line=None)
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    bn.text_frame.margin_left = Cm(0.18)
    write_lines(bn.text_frame,
                ["FA 9업무 → Agent 자동화 수준 (보조 ⑧ 요약)"],
                size=SZ_SUB, bold=True, color=WHITE)

    # 4 행 자동화 수준 미니 표 (★★★/★★/★/△)
    levels = [
        ("★★★", "완전 자동 (1)",
         "보고 — CPO·센터장·생산성·사업부 (writer + meeting-minutes)",
         BLUE),
        ("★★", "자동 (4)",
         "중점추진·Project·MRM·미래기술 — KPI/Capex/과제/DST",
         BLACK),
        ("★", "부분 (2)",
         "협력사·경비 — 데이터 자동 · 등록/결재 사람",
         BLACK),
        ("△", "보조 (2)",
         "인사·기타 — 회의록/요약 자동 · 결정·진행 사람",
         BLACK),
    ]
    row_y = card_y + 0.55
    row_h = (card_h - 0.65) / 4
    for star, lvl, desc, c in levels:
        add_text(s1, 14.32, row_y, 1.30, row_h - 0.04,
                 [(star, {"size": SZ_SUB, "bold": True, "color": c,
                          "align": PP_ALIGN.CENTER})])
        add_text(s1, 15.65, row_y, 3.20, row_h - 0.04,
                 [(lvl, {"size": SZ_FOOT, "bold": True, "color": BLACK})])
        add_text(s1, 18.88, row_y, 7.95, row_h - 0.04,
                 [(desc, {"size": SZ_FOOT, "color": BLACK})])
        row_y += row_h

    # ── 우하 풋라인 : 한 줄 요약 (메인 핵심 메시지) ──
    # 단일 paragraph 다중 run — 자연스럽게 2 lines 로 wrap
    add_rect(s1, 14.22, 17.40, 12.70, 0.85, fill=LIGHT_GRAY, line=MID_GRAY,
             line_w=0.5)
    tb = s1.shapes.add_textbox(Cm(14.32), Cm(17.46), Cm(12.50), Cm(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.02); tf.margin_bottom = Cm(0.02)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    set_run(r1, "핵심 — ", size=SZ_FOOT, bold=True, color=BLACK)
    r2 = p.add_run()
    set_run(r2, "스킬 14 + MCP 16+ 로 9업무 자동화",
            size=SZ_FOOT, bold=True, color=BLUE)
    r3 = p.add_run()
    set_run(r3, " · 출처 100% 자동 · 산출물 10+종 검증",
            size=SZ_FOOT, color=BLACK)


def fill_form_aux(prs):
    """Slide 2: 시스템 구성 — 네이티브 아키텍처 다이어그램.

    피드백 반영: PNG 임베드 → PowerPoint 네이티브 도형으로 전환.
    각 사각형·화살표·텍스트가 PowerPoint 에서 직접 편집 가능.
    SVG/PNG 원본은 outputs/_assets/ 에 별도 보존 (외부 참고용).
    """
    s = prs.slides[1]

    # 1) 양식 잔존 도형 제거
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    add_rect(s, 0, 0, 27.52, 19.05, fill=WHITE, line=None)

    # 2) 보조 헤더 (3–12 와 동일)
    header(s, "[보조 ①/11] 시스템 구성 — 5 에이전트 + 스킬 14 + Hook 2 + MCP 16+")

    # ── Layer 1: 사용자 (상단 중앙) ──
    ux, uw = 11.86, 3.80  # center
    add_rect(s, ux, 1.55, uw, 0.70, fill=WHITE, line=BLACK, line_w=1.0)
    bn = s.shapes[-1]
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame, ["사용자 요청"], size=SZ_SUB,
                bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Arrow ↓
    _arrow_down(s, ux + uw / 2, 2.25, 2.45)

    # ── Layer 2: 팀장 ──
    tx, tw = 8.76, 10.00
    add_rect(s, tx, 2.50, tw, 0.95, fill=LIGHT_GRAY, line=BLACK, line_w=1.0)
    bn = s.shapes[-1]
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame, [
        ("팀장 (메인 Claude Code 세션)",
         {"size": SZ_SECTION, "bold": True, "color": BLACK,
          "align": PP_ALIGN.CENTER}),
        ("분해 · 우편함 라우팅 · 품질 게이트",
         {"size": SZ_FOOT, "color": DIM_GRAY, "align": PP_ALIGN.CENTER}),
    ])

    # Arrow ↓
    _arrow_down(s, tx + tw / 2, 3.45, 3.70)

    # ── Layer 3: 5 에이전트 카드 (가로 5칸) ──
    agents = [
        ("data-teammate", "사내 정량",
         "KPI · Capex · MRM · DST", "+ xlsx 트래커"),
        ("ops-teammate", "사내 운영",
         "협력사 · HR · 결재 · 문서", "+ pdf · docx"),
        ("tech-research", "사외 리서치",
         "동향 · 경쟁사 · 평가", "+ firecrawl"),
        ("document-writer", "작성 (.md)",
         "Markdown 양식 채움", "+ doc-coauthoring"),
        ("ppt-writer", "작성 (.pptx)",
         "결정론 파이프라인", "+ brand-guidelines"),
    ]
    cw, gap = 5.06, 0.18
    sx = (27.52 - (cw * 5 + gap * 4)) / 2
    agent_top = 3.80
    agent_h = 2.55
    for i, (nm, tag, role, plus) in enumerate(agents):
        x = sx + i * (cw + gap)
        add_rect(s, x, agent_top, cw, agent_h, fill=WHITE, line=MID_GRAY,
                 line_w=0.75)
        bnd = add_rect(s, x, agent_top, cw, 0.55, fill=DIM_GRAY, line=None)
        bnd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bnd.text_frame, [nm], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, agent_top + 0.60, cw, 0.45,
                 [(tag, {"size": SZ_FOOT, "bold": True, "color": BLACK,
                         "align": PP_ALIGN.CENTER})])
        add_text(s, x + 0.10, agent_top + 1.05, cw - 0.20, 0.55,
                 [(role, {"size": SZ_FOOT, "color": BLACK,
                          "align": PP_ALIGN.CENTER})])
        add_text(s, x + 0.10, agent_top + 1.65, cw - 0.20, 0.50,
                 [(plus, {"size": 7, "color": DIM_GRAY,
                          "align": PP_ALIGN.CENTER})])
        # Arrow ↓ for each card
        _arrow_down(s, x + cw / 2, agent_top + agent_h, agent_top + agent_h + 0.55)

    # 우편함 통신 안내 (5 카드 위 1줄)
    add_text(s, 0.6, 6.45, 26.32, 0.45,
             [("↕ 우편함(Mailbox) 통신 — 에이전트끼리 직접 소통 (팀장 거치지 않음)",
               {"size": SZ_FOOT, "italic": False, "color": DIM_GRAY,
                "align": PP_ALIGN.CENTER, "bold": True})])

    # ── Layer 4: 스킬 14 + MCP 16+ (전체 폭 카드) ──
    # 5 행 × (label bold + body) — 각 행 = 1 paragraph 2 runs (wrap 방지)
    sk_y, sk_h = 7.05, 2.65
    add_rect(s, 0.6, sk_y, 26.32, sk_h, fill=WHITE, line=MID_GRAY,
             line_w=0.75)
    bn = add_rect(s, 0.6, sk_y, 26.32, 0.55, fill=DIM_GRAY, line=None)
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame,
                ["스킬 14종 (도메인 3 + 신규 11) · 자동 발동  ·  "
                 "MCP 서버 16+ · 외부 자료 인-라인 호출"],
                size=SZ_SUB, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 5 행 — 각 행은 단일 paragraph 에 (label, body) 2 runs
    sk_rows = [
        ("도메인 3 :  ",
         "meeting-minutes · meeting-merger · fa-task-discovery"),
        ("산출 5 :  ",
         "docx · xlsx · pdf · doc-coauthoring · internal-comms"),
        ("메타 4 :  ",
         "brand-guidelines · skill-creator · mcp-builder · sentry-code-review"),
        ("장기 2 :  ",
         "webapp-testing · firecrawl"),
        ("MCP :  ",
         "카카오 · Google(Gmail/Calendar) · GitHub · Figma · Autodesk · 사내 PoC"),
    ]
    sk_tb = s.shapes.add_textbox(Cm(0.8), Cm(sk_y + 0.62),
                                 Cm(25.9), Cm(sk_h - 0.72))
    tf = sk_tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Cm(0.08); tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.02); tf.margin_bottom = Cm(0.02)
    tf.clear()
    for i, (lab, body) in enumerate(sk_rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        set_run(r1, lab, size=SZ_FOOT, bold=True, color=BLACK)
        r2 = p.add_run()
        set_run(r2, body, size=SZ_FOOT, color=BLACK)

    # Arrow ↓
    _arrow_down(s, 13.76, sk_y + sk_h, sk_y + sk_h + 0.55)

    # ── Layer 5: outputs/ ── (sk_h 확대(2.30→2.65) 만큼 0.35cm 아래로 이동)
    out_y, out_h = 10.30, 0.95
    add_rect(s, 6.76, out_y, 14.00, out_h, fill=WHITE, line=BLACK,
             line_w=1.0)
    bn = s.shapes[-1]
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame, [
        ("outputs/  —  보고서 · 회의록 · 분석 (10+ 종 검증)",
         {"size": SZ_SECTION, "bold": True, "color": BLACK,
          "align": PP_ALIGN.CENTER}),
        ("주간보고 · R&D 회의록 · 유연성지표 · 경쟁사 동향 · 임원보고 ·"
         " 전고체 검토 · FA 과제후보",
         {"size": SZ_FOOT, "color": DIM_GRAY, "align": PP_ALIGN.CENTER}),
    ])

    # Arrow ↓
    _arrow_down(s, 13.76, out_y + out_h, out_y + out_h + 0.55)

    # ── Layer 6: Hook + sentry (3 박스 가로) ── (0.35cm 아래로 이동)
    hk_y, hk_h = 11.85, 1.05
    add_rect(s, 0.6, hk_y, 8.40, hk_h, fill=LIGHT_GRAY, line=MID_GRAY,
             line_w=0.6)
    add_text(s, 0.80, hk_y + 0.10, 8.00, 0.40,
             [("SessionStart Hook",
               {"size": SZ_FOOT, "bold": True, "color": BLACK})])
    add_text(s, 0.80, hk_y + 0.50, 8.00, 0.50,
             [("build_reference_index.py · references/ 자동 인덱싱",
               {"size": 7, "color": DIM_GRAY})])

    add_rect(s, 9.20, hk_y, 9.12, hk_h, fill=LIGHT_GRAY, line=BLACK,
             line_w=1.0)
    add_text(s, 9.40, hk_y + 0.10, 8.72, 0.40,
             [("✓ TaskCompleted Hook — 출처 100% 자동 검증",
               {"size": SZ_FOOT, "bold": True, "color": BLACK})])
    add_text(s, 9.40, hk_y + 0.50, 8.72, 0.50,
             [("check_output_citations.py · 미기입 시 작업 차단",
               {"size": 7, "color": DIM_GRAY})])

    add_rect(s, 18.52, hk_y, 8.40, hk_h, fill=LIGHT_GRAY, line=MID_GRAY,
             line_w=0.6)
    add_text(s, 18.72, hk_y + 0.10, 8.00, 0.40,
             [("sentry-code-review",
               {"size": SZ_FOOT, "bold": True, "color": BLACK})])
    add_text(s, 18.72, hk_y + 0.50, 8.00, 0.50,
             [("scripts/build_*.py 빌더 운영 안전성 게이트",
               {"size": 7, "color": DIM_GRAY})])

    # ── SOP 요약 띠 ── (0.35cm 아래로 이동)
    add_rect(s, 0.6, 13.20, 26.32, 0.75, fill=LIGHT_BLUE, line=MID_GRAY,
             line_w=0.6)
    add_text(s, 0.80, 13.28, 25.92, 0.62, [
        ("SOP — 분해 → 병렬 수집(data·ops·tech) → 위임 작성"
         "(writer + 스킬 자동 발동) → 품질 게이트(Hook + sentry)",
         {"size": SZ_SUB, "bold": True, "color": BLUE,
          "align": PP_ALIGN.CENTER}),
    ])

    # ── 핵심 효과 3 박스 ──
    eff = [
        ("재현성", "같은 입력 = 같은 출력 · 양식 무결성 100%"),
        ("자동 검증", "출처 100% 자동 · 미기입 작업 차단"),
        ("확장성", "스킬 14 + MCP 16+ · 타 영역 즉시 이식"),
    ]
    ey, eh = 14.10, 3.30
    ew, eg = (26.32 - 0.4) / 3, 0.20
    for i, (lab, dsc) in enumerate(eff):
        ex = 0.6 + i * (ew + eg)
        add_rect(s, ex, ey, ew, eh, fill=WHITE, line=MID_GRAY,
                 line_w=0.6)
        bn = add_rect(s, ex, ey, ew, 0.55, fill=DIM_GRAY, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bn.text_frame, [lab], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, ex + 0.20, ey + 0.70, ew - 0.40, eh - 0.85,
                 [(dsc, {"size": SZ_FOOT, "color": BLACK,
                         "align": PP_ALIGN.CENTER})])

    add_footnote(s, "* 출처 : .claude/agents/*.md · .claude/skills/*/SKILL.md"
                    " (14종) · CLAUDE.md [67-84행] · scripts/hooks/"
                    " · outputs/_assets/agent_architecture.svg (참고용)")
    return  # 이하 unused — 이전 카드 레이아웃 보존용으로 남김

    # (이하 unused — 이전 버전 카드 레이아웃 보존용)

    # 3) 사용자→팀장 라우팅 라인
    add_text(s, 0.8, 1.5, 25.9, 0.5,
             [("사용자 요청 → 팀장(메인 세션, 분해·우편함 라우팅·품질 게이트)",
               {"size": SZ_SUB, "bold": True, "color": BLACK,
                "align": PP_ALIGN.CENTER})])

    # 4) 5 에이전트 카드 띠 (T2.0, H3.4)
    agents = [
        ("data-teammate", "사내 정량", "KPI·Capex·MRM·DST"),
        ("ops-teammate", "사내 운영", "협력사·HR·결재·문서"),
        ("tech-research", "사외 리서치", "동향·경쟁사·평가"),
        ("document-writer", "작성(.md)", "Markdown 양식"),
        ("ppt-writer", "작성(.pptx)", "결정론 파이프라인"),
    ]
    cw, gap = 5.06, 0.18
    sx = (27.52 - (cw * 5 + gap * 4)) / 2
    for i, (nm, tag, role) in enumerate(agents):
        x = sx + i * (cw + gap)
        add_rect(s, x, 2.0, cw, 3.40, fill=WHITE, line=MID_GRAY,
                 line_w=0.75)
        bnd = add_rect(s, x, 2.0, cw, 0.60, fill=DIM_GRAY, line=None)
        bnd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bnd.text_frame, [nm], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, 2.65, cw, 0.50,
                 [(tag, {"size": SZ_FOOT, "bold": True, "color": BLACK,
                         "align": PP_ALIGN.CENTER})])
        add_text(s, x + 0.12, 3.20, cw - 0.24, 2.10,
                 [(role, {"size": SZ_FOOT, "color": BLACK,
                          "align": PP_ALIGN.CENTER})])

    # 5) 도메인 스킬 3 (T5.55, H1.55)
    add_text(s, 0.8, 5.50, 25.9, 0.40,
             [("도메인 스킬 3 — 자동 발동 워크플로우",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])
    skills_dom = [
        ("meeting-minutes", "녹취 → LGES 표준 회의록 PPTX"),
        ("meeting-merger", "녹취 + OneNote → 통합 Word(.docx)"),
        ("fa-task-discovery", "FA 6분야 동향 → 2부 docx + 카톡 (v2)"),
    ]
    sw, sg = 8.46, 0.26
    ssx = (27.52 - (sw * 3 + sg * 2)) / 2
    for i, (nm, ds) in enumerate(skills_dom):
        x = ssx + i * (sw + sg)
        add_rect(s, x, 5.95, sw, 1.30, fill=WHITE, line=MID_GRAY,
                 line_w=0.6)
        bnd = add_rect(s, x, 5.95, sw, 0.45, fill=DIM_GRAY, line=None)
        bnd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bnd.text_frame, [nm], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.15, 6.45, sw - 0.30, 0.75,
                 [(ds, {"size": SZ_FOOT, "color": BLACK,
                        "align": PP_ALIGN.CENTER})])

    # 6) 신규 스킬 11종 — 카테고리 3 그룹 (T7.45, H3.85)
    add_text(s, 0.8, 7.40, 25.9, 0.40,
             [("신규 스킬 11 — 산출 양식 5 · 메타·운영 4 · 장기 검토 2",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])
    # 그룹 헤더 + 카드 리스트 (가로 3 칼럼)
    groups = [
        ("산출 양식 (5)", [
            ("docx", "Word 정책/SOP/협력사 평가서"),
            ("xlsx", "KPI/Capex/MRM 트래커"),
            ("pdf", "협력사 사양·도면 OCR"),
            ("doc-coauthoring", "임원 보이스 장문 협업"),
            ("internal-comms", "사내 공지·주간·3P·FAQ"),
        ]),
        ("메타·운영 (4)", [
            ("brand-guidelines", "색·폰트·풋노트 자동 적용"),
            ("skill-creator", "신규 스킬 scaffold·eval"),
            ("mcp-builder", "사내 ERP/MES/PLM MCP"),
            ("sentry-code-review", "빌더 품질 게이트"),
        ]),
        ("장기 검토 (2)", [
            ("webapp-testing", "Playwright (사내 대시보드)"),
            ("firecrawl", "JS 사이트 안정 스크래핑 PoC"),
        ]),
    ]
    gw, gg = 8.62, 0.13
    gsx = (27.52 - (gw * 3 + gg * 2)) / 2
    for gi, (gname, items) in enumerate(groups):
        gx = gsx + gi * (gw + gg)
        add_rect(s, gx, 7.85, gw, 3.40, fill=WHITE, line=MID_GRAY,
                 line_w=0.6)
        gh = add_rect(s, gx, 7.85, gw, 0.40, fill=DIM_GRAY, line=None)
        gh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        gh.text_frame.margin_left = Cm(0.12)
        write_lines(gh.text_frame, [gname], size=SZ_FOOT, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        # 각 카드 1행
        row_h = (3.40 - 0.45) / max(len(items), 5)
        for ri, (sk, dsc) in enumerate(items):
            ry = 8.30 + ri * row_h
            add_text(s, gx + 0.12, ry, gw - 0.24, row_h - 0.04, [
                (sk, {"size": SZ_FOOT, "bold": True, "color": BLACK}),
                (f"  {dsc}", {"size": 7, "color": DIM_GRAY}),
            ])

    # 7) Hook 2 + MCP 그룹 (T11.45, H2.60)
    add_text(s, 0.8, 11.40, 25.9, 0.40,
             [("자동 검증 Hook 2 + MCP 서버 16+ 자동 연결",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])

    # 좌: Hook 2 (W12.55)
    hx, hw, hg = 0.60, 5.95, 0.20
    for i, (nm, ds) in enumerate([
        ("Hook 1 — SessionStart",
         "build_reference_index.py — references/ 자동 인덱싱"),
        ("Hook 2 — TaskCompleted",
         "check_output_citations.py — 출처 미기입 시 차단"),
    ]):
        x = hx + i * (hw + hg)
        add_rect(s, x, 11.85, hw, 2.20, fill=WHITE, line=BLACK, line_w=0.6)
        add_rect(s, x, 11.85, 0.18, 2.20, fill=DIM_GRAY, line=None)
        add_text(s, x + 0.30, 11.95, hw - 0.40, 0.50,
                 [(nm, {"size": SZ_SUB, "bold": True, "color": BLACK})])
        add_text(s, x + 0.30, 12.40, hw - 0.40, 1.60,
                 [(ds, {"size": SZ_FOOT, "color": BLACK})])

    # 우: MCP 16+ (W13.80)
    mx, mw = 13.10, 13.82
    add_rect(s, mx, 11.85, mw, 2.20, fill=WHITE, line=MID_GRAY, line_w=0.75)
    mbnd = add_rect(s, mx, 11.85, mw, 0.42, fill=DIM_GRAY, line=None)
    mbnd.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    mbnd.text_frame.margin_left = Cm(0.15)
    write_lines(mbnd.text_frame, ["MCP 서버 16+ — 외부 자료 인-라인 호출"],
                size=SZ_SUB, bold=True, color=WHITE)
    add_text(s, mx + 0.18, 12.32, mw - 0.36, 1.65, [
        (" • 카카오 공식 — KakaoMap · KakaotalkChat · NaverSearch · OpenDart · YouTubeData",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • Google — Gmail · Calendar  |  GitHub — PR · Issue · Code Search",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • Figma — 디자인 컨텍스트  |  Autodesk Knowledge — 사양 카탈로그",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • 호텔·항공·도서·금융 등 외부 API + 사내 ERP/MES/PLM PoC (mcp-builder)",
         {"size": SZ_FOOT, "color": BLACK}),
    ])

    # 8) SOP + 빌더 띠 (T14.20, H3.60)
    add_rect(s, 0.6, 14.20, 26.32, 0.85, fill=LIGHT_BLUE, line=MID_GRAY,
             line_w=0.75)
    add_text(s, 0.80, 14.32, 25.92, 0.65,
             [("SOP : 분해 → 병렬 수집(data·ops·tech) → 위임 작성"
               "(writer + 스킬 14 자동발동) → 품질 게이트(Hook + sentry)",
               {"size": SZ_SUB, "bold": True, "color": BLUE,
                "align": PP_ALIGN.CENTER})])
    add_text(s, 0.6, 15.20, 26.4, 2.95, [
        ("빌더 계층 — 양식별 결정론 스크립트 (scripts/build_*.py)",
         {"size": SZ_SUB, "bold": True, "color": BLACK}),
        (" • ppt_extract.py / ppt_fill.py — 범용 placeholder 추출·채움",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • build_flexibility_summary.py — LGES PPT 가이드 기준 구현",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • build_report_catalog.py — 임원 보고 양식 12종 → 8블록",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • build_fa_exec_report.py / build_solidstate_transport_report.py"
         " · build_contest_2026_05_19.py — 카탈로그·가이드·본 deck",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • 신규 빌더 추가 시 sentry-code-review 스킬 자동 게이트",
         {"size": SZ_FOOT, "color": BLACK}),
    ])

    add_footnote(s, "* 출처 : .claude/agents/*.md · .claude/skills/*/SKILL.md"
                    " (14종) · CLAUDE.md [스킬 섹션 67-84행] · scripts/"
                    " (build_*.py, hooks/) · templates/사내양식/LGES_PPT_작업_가이드.md")


# ─── 보조 슬라이드 ────────────────────────────────────────────────
def aux_members(prs):
    s = new_blank_slide(prs)
    header(s, "[보조 ③/11] 5 팀원 · 표준 작업 절차(SOP) · 자동 검증 Hook")

    members = [
        ("data-teammate", ["[Read·Grep·Glob·Bash]", "사내 정량 데이터",
                           "KPI·Capex·MRM(양산)", "DST(R&D)",
                           "추측 금지·항상 출처", "외부는 tech 위임"]),
        ("ops-teammate", ["[Read·Grep·Glob·Bash]", "사내 운영",
                          "협력사·HR·결재 큐", "사내 PDF/PPT 검색",
                          "승인 작업은 제안만", "KPI·외부는 위임"]),
        ("tech-research", ["[+WebSearch·Fetch+MCP]", "사외 리서치",
                           "동향·경쟁사·평가", "AMR/로봇/트윈/AI",
                           "references/ 자동", "아카이빙"]),
        ("document-writer", ["[Read·Write·Edit]", "Markdown 작성",
                             "templates/*.md", "placeholder 채움",
                             "팀원에 직접 요청", "plan-mode 저장"]),
        ("ppt-writer", ["[+Bash 결정론]", "PPT 작성",
                        "extract→fill", "양식 무결성 보존",
                        "5단계 추론", "원본 불변"]),
    ]
    cw, gap = 5.06, 0.18
    sx = (27.52 - (cw * 5 + gap * 4)) / 2
    for i, (nm, lines) in enumerate(members):
        add_card(s, sx + i * (cw + gap), 1.75, cw, 6.35, nm,
                 [(t, {"size": SZ_SUB}) for t in lines],
                 head_size=SZ_SUB, body_size=SZ_SUB)

    # SOP 좌측
    add_text(s, 0.6, 8.45, 12.9, 0.62,
             [("표준 작업 절차(SOP) — Agent Teams 운영 패턴",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    sop = [
        ("Step 1 분해", "리더가 요청을 도메인별로 분해 (사내수치↔운영↔리서치, 작성↔검토)"),
        ("Step 2 병렬수집", "수집 팀원(data·ops·tech) 동시 spawn, 출처와 함께 (대화 미상속)"),
        ("Step 3 위임작성",
         "writer 가 placeholder 만 치환 · 스킬 14 자동 발동(brand-guidelines"
         "·xlsx·pdf·docx 등) · 부족분 우편함 라운드"),
        ("Step 4 품질게이트",
         "TaskCompleted hook 출처 검증 + sentry-code-review (빌더 운영 안전성)"),
    ]
    y = 9.05
    for t, d in sop:
        add_rect(s, 0.6, y, 3.0, 1.78, fill=DIM_GRAY, line=None)
        bn = s.shapes[-1]
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bn.text_frame, [t], size=SZ_SUB, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_rect(s, 3.7, y, 9.8, 1.78, fill=SOFT_GRAY, line=LINE_GRAY,
                 line_w=0.4)
        add_text(s, 3.9, y + 0.2, 9.45, 1.45,
                 [(d, {"size": SZ_SUB, "color": BLACK})])
        y += 1.86

    # Hook 우측
    add_text(s, 13.9, 8.45, 13.0, 0.62,
             [("Hook 2 — Claude Code 자동 검증",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    hooks = [
        ("SessionStart Hook", "scripts/build_reference_index.py",
         ["세션 시작 시 references/ 전체 스캔",
          "제목·카테고리·날짜·요약 휴리스틱 추출",
          "→ references/INDEX.md 자동 갱신",
          "tech-research 신규 자료 자동 등록"]),
        ("TaskCompleted Hook", "scripts/hooks/check_output_citations.py",
         ["outputs/*.md 생성 후 자동 실행",
          "(출처:…)/[N행]/자료 미확인 1개 이상 필수",
          "미기입 시 exit≠0 → 작업 차단",
          "URL 인용도 통과 처리"]),
    ]
    y = 9.05
    for nm, sc, lines in hooks:
        add_rect(s, 13.9, y, 13.0, 3.72, fill=LIGHT_GRAY, line=LINE_GRAY,
                 line_w=0.5)
        add_rect(s, 13.9, y, 0.2, 3.72, fill=DIM_GRAY, line=None)
        add_text(s, 14.35, y + 0.15, 12.4, 0.6,
                 [(nm, {"size": SZ_SECTION, "bold": True, "color": BLACK})])
        add_text(s, 14.35, y + 0.68, 12.4, 0.42,
                 [(sc, {"size": SZ_FOOT, "color": DIM_GRAY,
                        "font": FONT_EMPH})])
        add_text(s, 14.35, y + 1.12, 12.4, 2.5,
                 [(f" • {t}", {"size": SZ_SUB, "color": BLACK})
                  for t in lines])
        y += 3.88

    add_footnote(s, "* 출처 : .claude/agents/*.md · .claude/settings.json ·"
                    " scripts/hooks/check_output_citations.py ·"
                    " scripts/build_reference_index.py · CLAUDE.md")


def aux_skills(prs):
    s = new_blank_slide(prs)
    header(s, "[보조 ④/11] 스킬 라이브러리 14종 — 도메인 3 + 신규 11 (자동 발동)")

    # ── 1) 도메인 스킬 3종 (상단 띠) ──
    add_text(s, 0.6, 1.45, 26.3, 0.40,
             [("도메인 스킬 3 — 회의·과제 발굴 워크플로우 (기존)",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])
    cols = [
        ("meeting-minutes",
         "녹취 → LGES 표준 회의록 PPTX 자동 · 참석자·안건·결정·액션 UI"),
        ("meeting-merger",
         "녹취 + OneNote → 통합 Word(.docx) · 한·영 지원"),
        ("fa-task-discovery  (v2)",
         "FA 6분야 한/영 동향 → 2부 docx + 카톡 · references 자동 아카이빙"),
    ]
    cw, gap = 8.66, 0.22
    sx = 0.6
    for i, (nm, ds) in enumerate(cols):
        x = sx + i * (cw + gap)
        hf = LIGHT_BLUE if i == 2 else DIM_GRAY
        hc = BLUE if i == 2 else WHITE
        add_rect(s, x, 1.90, cw, 1.95, fill=WHITE, line=MID_GRAY, line_w=0.75)
        bn = add_rect(s, x, 1.90, cw, 0.55, fill=hf, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bn.text_frame, [nm], size=SZ_SUB, bold=True,
                    color=hc, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.18, 2.50, cw - 0.36, 1.30,
                 [(ds, {"size": SZ_FOOT, "color": BLACK})])

    # ── 2) 신규 스킬 11종 — 카테고리 3 그룹 ──
    add_text(s, 0.6, 4.05, 26.3, 0.40,
             [("신규 스킬 11 — 산출 양식 5 · 메타·운영 4 · 장기 검토 2 (2026-05 추가)",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])

    groups = [
        ("산출 양식 (5)", [
            ("docx", "Word 정책·SOP·협력사 평가서 (트래킹 체인지)"),
            ("xlsx", "Excel KPI·Capex·MRM 트래커 (수식·피벗)"),
            ("pdf", "협력사 사양·도면 추출·병합·OCR"),
            ("doc-coauthoring", "임원 보이스 장문 협업 워크플로"),
            ("internal-comms", "사내 공지·주간 메일·3P·FAQ·이슈 리포트"),
        ]),
        ("메타·운영 (4)", [
            ("brand-guidelines", "PPT/DOCX/MD 색·폰트·풋노트 자동 적용"),
            ("skill-creator", "신규 스킬 scaffold·eval·description 최적화"),
            ("mcp-builder", "사내 ERP/MES/PLM read-only MCP 설계·PoC"),
            ("sentry-code-review", "Python 빌더 운영 인식 코드 리뷰 게이트"),
        ]),
        ("장기 검토 (2)", [
            ("webapp-testing", "Playwright — 사내 대시보드 도입 시 활성화"),
            ("firecrawl", "JS 사이트 안정 스크래핑 — 보안/IT 검토 후 PoC"),
        ]),
    ]
    gw, gg = 8.66, 0.16
    gsx = (27.52 - (gw * 3 + gg * 2)) / 2
    grow_y, grow_h = 4.55, 4.30
    for gi, (gname, items) in enumerate(groups):
        gx = gsx + gi * (gw + gg)
        add_rect(s, gx, grow_y, gw, grow_h, fill=WHITE, line=MID_GRAY,
                 line_w=0.6)
        bn = add_rect(s, gx, grow_y, gw, 0.50, fill=DIM_GRAY, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.15)
        write_lines(bn.text_frame, [gname], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        # 각 카드 1행 (5/4/2 항목)
        inner_y = grow_y + 0.60
        inner_h = grow_h - 0.70
        per_row = inner_h / max(len(items), 1)
        for ri, (sk, dsc) in enumerate(items):
            ry = inner_y + ri * per_row
            add_text(s, gx + 0.18, ry, gw - 0.36, per_row - 0.04, [
                (sk, {"size": SZ_FOOT, "bold": True, "color": BLACK}),
                (f"  {dsc}", {"size": 7, "color": DIM_GRAY}),
            ])

    # ── 3) fa-task-discovery v2 최근 리팩터 핵심 ──
    add_rect(s, 0.6, 9.10, 26.32, 0.62, fill=DIM_GRAY, line=None)
    bn = s.shapes[-1]
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    bn.text_frame.margin_left = Cm(0.2)
    write_lines(bn.text_frame,
                ["fa-task-discovery v2 — 2026-05 리팩터 핵심 (최근 작업)"],
                size=SZ_SUB, bold=True, color=WHITE)

    v2 = [
        ("2세션 자동 분기",
         "전용 폴더의 오늘자 news_{날짜}.json 존재 여부로 세션 1(뉴스 수집) / "
         "세션 2(과제 도출) 자동 판정 — 수동 선택 없음"),
        ("2부 구성 docx",
         "제1부 최근 1개월(30일) 뉴스 브리핑 + 제2부 기술과제 후보 5~10건 "
         "(세션 1 은 뉴스 단일부, 세션 2 는 2부)"),
        ("전용 폴더 분리",
         "산출물은 outputs/fa-task-discovery/ 에 격리 — news/tasks JSON + "
         "FA_뉴스브리핑·FA_기술과제후보 docx"),
        ("날짜 브랜치 + PR 자동 생성",
         "기본 브랜치 직접 push 금지 — 날짜 브랜치 커밋·푸시 후 main 대상 PR "
         "자동 생성(중복 PR 회피) → 카카오톡 200자+ 요약"),
    ]
    y = 9.90
    row_h = 1.85
    for t, d in v2:
        add_rect(s, 0.6, y, 5.6, row_h, fill=SOFT_GRAY, line=LINE_GRAY,
                 line_w=0.4)
        add_text(s, 0.78, y + 0.18, 5.25, row_h - 0.30,
                 [(t, {"size": SZ_SUB, "bold": True, "color": BLACK})])
        add_rect(s, 6.3, y, 20.62, row_h, fill=WHITE, line=LINE_GRAY,
                 line_w=0.4)
        add_text(s, 6.5, y + 0.2, 20.25, row_h - 0.30,
                 [(d, {"size": SZ_SUB, "color": BLACK})])
        y += row_h + 0.08

    add_footnote(s, "* 출처 : .claude/skills/*/SKILL.md (14종 — meeting-minutes/"
                    "meeting-merger/fa-task-discovery + docx/xlsx/pdf/"
                    "doc-coauthoring/internal-comms/brand-guidelines/"
                    "skill-creator/mcp-builder/sentry-code-review/"
                    "webapp-testing/firecrawl) · CLAUDE.md [67-84행]")


def aux_pipeline(prs):
    s = new_blank_slide(prs)
    header(s, "[보조 ⑤/11] 결정론 PPT 파이프라인 · LGES PPT 가이드 · 임원 보고 카탈로그")

    # 좌측 — 결정론 PPT 파이프라인
    add_text(s, 0.6, 1.7, 12.9, 0.62,
             [("① 결정론 PPT 파이프라인 (원본 불변)",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    steps = [
        "templates/<양식>.pptx (원본 — 절대 미수정)",
        "↓ scripts/ppt_extract.py",
        "outputs/.cache/<양식>.structure.json",
        "↓ Claude 추론 + 팀원 데이터 요청",
        "outputs/.cache/<양식>.values.json",
        "↓ scripts/ppt_fill.py",
        "outputs/<양식>_<날짜>.pptx (최종)",
    ]
    y = 2.3
    for i, st in enumerate(steps):
        is_io = i % 2 == 0
        add_rect(s, 0.6, y, 12.9, 0.92,
                 fill=(LIGHT_GRAY if is_io else WHITE),
                 line=(MID_GRAY if is_io else None),
                 line_w=0.5)
        add_text(s, 0.85, y + 0.16, 12.4, 0.6,
                 [(st, {"size": SZ_SUB,
                        "bold": is_io,
                        "color": (BLACK if is_io else DIM_GRAY)})])
        y += 0.99
    add_rect(s, 0.6, y + 0.05, 12.9, 1.5, fill=LIGHT_BLUE, line=MID_GRAY,
             line_w=0.75)
    add_text(s, 0.8, y + 0.22, 12.5, 1.2,
             [("Claude 는 PPTX 바이너리를 직접 만지지 않음 — 같은 입력=같은 출력, "
               "양식 도형/색/폰트 100% 보존",
               {"size": SZ_SUB, "bold": True, "color": BLUE})])

    # 우상 — LGES PPT 작업 가이드
    add_text(s, 13.9, 1.7, 13.0, 0.62,
             [("② LGES PPT 작업 가이드 (에이전트 기본 룰)",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    guide = [
        " • 구조: 메인 1~2장(임원 요약) + 보조 N장(상세)",
        " • 색: 흰/검/회 기본 + 파랑 #0000FF 강조(3~5/장)",
        " • 풋노트: 별도 박스 초록곰팡이 #006600 Bold 8pt",
        " • 폰트: 한글 LG스마트체 + 영문 Arial Narrow run 분리",
        " • 사이즈: SZ 6단(16/12/11/10/9/8)만 사용",
        " • 무오버플로 필수 — 넘치면 보조로 분리",
    ]
    add_rect(s, 13.9, 2.3, 13.0, 4.05, fill=WHITE, line=MID_GRAY,
             line_w=0.75)
    add_text(s, 14.15, 2.5, 12.5, 3.7,
             [(g, {"size": SZ_SUB, "color": BLACK}) for g in guide])

    # 우하 — 임원 보고 양식 카탈로그 8블록
    add_text(s, 13.9, 6.55, 13.0, 0.62,
             [("③ 임원 보고 양식 카탈로그 (12 슬라이드 → 8 블록)",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    blocks = [
        ("A 보고 프레임", "메인 한 장 요약 레이아웃"),
        ("B Master Plan", "분기 Gantt·추진 일정"),
        ("C 투자비·KPI", "Capex As-is/To-be·KPI"),
        ("D 추진체계", "ES/PRI/CNS 조직 매트릭스"),
        ("E 전략·컨셉", "지향점·중점과제·컨셉"),
        ("F 사양·구성", "구성 방안·사양·경량화"),
        ("G As-Is/To-Be", "프로세스·시스템 전환"),
        ("H 향후계획", "보조 과제·PoC 검증"),
    ]
    y = 7.15
    for i in range(0, 8, 2):
        for j in range(2):
            b, d = blocks[i + j]
            x = 13.9 + j * 6.6
            add_rect(s, x, y, 2.0, 1.05, fill=DIM_GRAY, line=None)
            bn = s.shapes[-1]
            bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            bn.text_frame.margin_left = Cm(0.1)
            write_lines(bn.text_frame, [b], size=SZ_FOOT, bold=True,
                        color=WHITE)
            add_rect(s, x + 2.05, y, 4.45, 1.05, fill=WHITE, line=LINE_GRAY,
                     line_w=0.4)
            add_text(s, x + 2.2, y + 0.16, 4.2, 0.78,
                     [(d, {"size": SZ_FOOT, "color": BLACK})])
        y += 1.13

    add_rect(s, 13.9, y + 0.05, 13.0, 1.4, fill=SOFT_GRAY, line=LINE_GRAY,
             line_w=0.4)
    add_text(s, 14.1, y + 0.18, 12.6, 1.2, [
        ("상위 보고 = 메인(블록 A) + 보조(B~H 주제별 선택) 조립.",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • brand-guidelines 스킬 자동 적용 — 색·폰트·풋노트 위반 자가 점검",
         {"size": SZ_FOOT, "color": BLACK}),
        (" • sentry-code-review — scripts/build_*.py 신규 빌더 품질 게이트",
         {"size": SZ_FOOT, "color": BLACK}),
    ])

    add_footnote(s, "* 출처 : templates/사내양식/LGES_PPT_작업_가이드.md ·"
                    " 보고양식_카탈로그.md(.json) · README.md ·"
                    " scripts/ppt_extract.py · scripts/ppt_fill.py ·"
                    " .claude/skills/brand-guidelines/ ·"
                    " .claude/skills/sentry-code-review/")


def aux_outputs(prs):
    s = new_blank_slide(prs)
    header(s, "[보조 ⑥/11] 검증된 산출물 — 실제 outputs/ (2026-05)")

    rows = [
        ("주간보고 (출처 자동검증 모범)",
         "outputs/주간보고_2026-05-12.md",
         "모든 KPI/Capex 값에 (출처:…)·[N행] 자동 부착 · data+ops 병렬 수집 · "
         "document-writer 가 17개 placeholder 만 치환 · 매주 재현"),
        ("회의록 (meeting-minutes)",
         "outputs/회의록_신규R&D과제착수_2026-05-12.md",
         "녹취 텍스트 → LGES 표준 회의록 PPTX · 참석자·안건·결정·액션 "
         "인터랙티브 확정 · 회의 직후 즉시 정리"),
        ("유연성지표 요약 (LGES 가이드 기준구현)",
         "outputs/유연성지표_요약_2026-05-14.pptx",
         "Flexibility 지표 정의·산식 + ESMI Lansing 사례 임원 1장 · "
         "build_flexibility_summary.py = PPT 가이드 레퍼런스"),
        ("경쟁사 동향 통합보고서 (tech-research×4)",
         "outputs/FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14",
         "FA/AI/AX/Physical AI/경쟁사 4축 외부 자료 병렬 수집 · 26FA "
         "로드맵·KPI 재설정 근거 · md+docx+pptx 3종"),
        ("FA자동화 임원보고 예시 (카탈로그 적용)",
         "outputs/FA자동화_임원보고_예시_2026-05-18.pptx",
         "메인 1 + 보조 4, 스토리 6단계(현황→문제→방안→투자→일정→효과) · "
         "보고 블록 A·C·G·B·E·H 조합 · references 출처 인용"),
        ("전고체 밀폐이송장치 검토 (EFEM+FOUP)",
         "outputs/전고체라인_밀폐이송장치_검토_2026-05-18.pptx",
         "Robostar EFEM + 3S P-FOUP 기술자료 기반 전고체 라인 밀폐이송 "
         "도입 검토 임원 보고 (메인 1 + 보조 2)"),
        ("FA 기술과제후보 2부 docx (fa-task-discovery v2)",
         "outputs/fa-task-discovery/FA_기술과제후보_20260518.docx",
         "제1부 최근 1개월 뉴스 + 제2부 과제 후보 · 날짜 브랜치 PR 자동 "
         "생성 · 카카오톡 요약 동시 전송"),
    ]
    y = 1.72
    rh = 2.16
    for title_, path, desc in rows:
        add_rect(s, 0.6, y, 8.2, rh, fill=DIM_GRAY, line=None)
        add_text(s, 0.78, y + 0.16, 7.85, 1.05,
                 [(title_, {"size": SZ_SUB, "bold": True, "color": WHITE})])
        add_text(s, 0.78, y + 1.2, 7.85, 0.85,
                 [(path, {"size": SZ_FOOT, "color": WHITE,
                          "font": FONT_EMPH})])
        add_rect(s, 8.85, y, 18.07, rh, fill=WHITE, line=LINE_GRAY,
                 line_w=0.4)
        add_text(s, 9.05, y + 0.2, 17.7, rh - 0.36,
                 [(desc, {"size": SZ_SUB, "color": BLACK})])
        y += rh + 0.12

    add_footnote(s, "* 출처 : outputs/ (주간보고·회의록·유연성지표 요약·경쟁사"
                    " 동향보고서·FA자동화 임원보고 예시·전고체 검토·"
                    "fa-task-discovery/) — 모두 (출처:…) 자동 검증 통과")


def aux_evolution(prs):
    s = new_blank_slide(prs)
    header(s, "[보조 ⑨/11] 레포 진화 타임라인 · 도입 ROI")

    add_text(s, 0.6, 1.7, 26.32, 0.62,
             [("① 진화 타임라인 (2026-05-10 → 05-20)",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    timeline = [
        ("05-10", "초기 팀 셋업 — data/ops/document-writer 3 에이전트 + md 양식"),
        ("05-11", "ppt-writer + 결정론 PPT 파이프라인(ppt_extract/fill) 추가"),
        ("05-12", "tech-research-teammate + 주간보고/회의록/소개자료 산출"),
        ("05-13", "3 도메인 스킬 + 2 Hook 정착 · AI경진대회 출품자료(5장) 1차"),
        ("05-14", "LGES PPT 작업 가이드 신설 · 한·영 폰트 자동분리 · 유연성지표"
                  "·경쟁사 동향보고서"),
        ("05-18", "임원 보고 양식 카탈로그(8블록) · 전고체 검토 · "
                  "fa-task-discovery v2(2부 docx·PR 자동)"),
        ("05-19", "AI경진대회 출품자료 12장 1차 (5에이전트+스킬3+Hook2)"),
        ("05-20", "스킬 라이브러리 11종 확장(docx·xlsx·pdf·doc-coauthoring·"
                  "internal-comms·brand-guidelines·skill-creator·mcp-builder·"
                  "sentry-code-review·webapp-testing·firecrawl) + MCP 16+ 연결"),
    ]
    y = 2.28
    row_h = 0.88
    step = 0.95
    for d, t in timeline:
        add_rect(s, 0.6, y, 2.3, row_h, fill=DIM_GRAY, line=None)
        bn = s.shapes[-1]
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bn.text_frame, [d], size=SZ_SUB, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        last = d == "05-20"
        add_rect(s, 2.95, y, 23.97, row_h,
                 fill=(LIGHT_BLUE if last else WHITE),
                 line=(MID_GRAY if last else LINE_GRAY),
                 line_w=(0.75 if last else 0.4))
        add_text(s, 3.2, y + 0.10, 23.5, row_h - 0.20,
                 [(t, {"size": SZ_FOOT, "bold": last,
                       "color": (BLUE if last else BLACK)})])
        y += step

    add_text(s, 0.6, 10.05, 26.32, 0.62,
             [("② 도입 ROI",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    roi = [
        ("작성 시간", "보고서당 수시간 → 분 단위 (검증 산출물 10여 종 누적)"),
        ("출처 신뢰도", "100% 자동 검증 — TaskCompleted hook 이 미기입 차단"),
        ("자료 누적", "외부 동향·경쟁사 references/ 누적 → 보고 품질 비례 상승"),
        ("재현성", "결정론 파이프라인 — 같은 입력=같은 출력, 양식 무결성 100%"),
        ("표준화", "5인 팀+SOP+PPT 가이드+카탈로그 — 영역·작성자 무관 동일 품질"),
        ("그룹 확장성", "사람 워크플로우와 동일 — 생산·품질·구매에 즉시 이식"),
    ]
    y = 10.6
    for lab, d in roi:
        add_rect(s, 0.6, y, 4.3, 1.13, fill=DIM_GRAY, line=None)
        bn = s.shapes[-1]
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.15)
        write_lines(bn.text_frame, [lab], size=SZ_SUB, bold=True,
                    color=WHITE)
        add_rect(s, 4.95, y, 21.97, 1.13, fill=LIGHT_GRAY, line=LINE_GRAY,
                 line_w=0.4)
        add_text(s, 5.2, y + 0.2, 21.5, 0.8,
                 [(d, {"size": SZ_SUB, "color": BLACK})])
        y += 1.2

    add_footnote(s, "* 출처 : git log · CLAUDE.md · README.md ·"
                    " outputs/주간보고_2026-05-12_PR요약.md ·"
                    " templates/사내양식/LGES_PPT_작업_가이드.md")


def aux_next(prs):
    s = new_blank_slide(prs)
    header(s, "[보조 ⑪/11] 인프라 그룹 확장 시나리오 · 핵심 메시지")

    add_text(s, 0.6, 1.75, 26.32, 0.62,
             [("인프라 그룹 4 부서 (FA기술담당·에너지기술담당·건설담당·인프라 기획팀) "
               "— 사람 워크플로우 동일 → 즉시 이식",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    # 인프라 그룹 4 부서 (사용자 제공 조직 구성)
    scen = [
        ("FA기술담당\n(본 deck 실증)",
         "현재 활용 중 — FA 기술혁신파트 보고/회의록/과제 발굴/경쟁사 동향 자동 "
         "(검증 산출물 10+종, 스킬 14 + MCP 16+ 가동)"),
        ("에너지기술담당",
         "ESS·재생에너지·전력관리 KPI 자동 보고, 외부 동향(IEA·BNEF) 누적, "
         "에너지 R&D 회의록 자동  (스킬: xlsx · tech-research · meeting-minutes)"),
        ("건설담당",
         "도면·시방서 PDF 추출·검색·OCR, 공정 진척률 트래킹, 협력사(시공/감리) "
         "평가서 자동  (스킬: pdf · xlsx · docx)"),
        ("인프라 기획팀",
         "그룹 전략·정책 장문 문서, 임원 결재 PPT(메인+보조), 사내 공지·3P 메일 "
         "자동  (스킬: doc-coauthoring · ppt-writer + 카탈로그 · internal-comms)"),
    ]
    y = 2.5
    rh = 2.05
    for org, d in scen:
        add_rect(s, 0.6, y, 4.4, rh, fill=DIM_GRAY, line=None)
        bn = s.shapes[-1]
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.15)
        write_lines(bn.text_frame, [org], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, 5.1, y, 21.82, rh, fill=SOFT_GRAY, line=LINE_GRAY,
                 line_w=0.4)
        add_text(s, 5.35, y + 0.22, 21.4, rh - 0.4,
                 [(d, {"size": SZ_SUB, "color": BLACK})])
        y += rh + 0.18

    add_rect(s, 0.6, 11.5, 26.32, 2.0, fill=SOFT_GRAY, line=LINE_GRAY,
             line_w=0.5)
    add_text(s, 0.85, 11.7, 25.8, 1.7, [
        ("이식 비용이 낮은 이유",
         {"size": SZ_SUB, "bold": True, "color": BLACK}),
        (" • 코드/DB 없음 — .md 양식·sample_data·references 만 교체하면 "
         "동일 5인 팀이 그대로 동작", {"size": SZ_SUB, "color": BLACK}),
        (" • 양식은 templates/ 에 두면 자동 인식, 출처 검증 Hook 이 품질 "
         "자동 보증", {"size": SZ_SUB, "color": BLACK}),
    ])

    add_rect(s, 0.6, 13.75, 26.32, 3.4, fill=LIGHT_BLUE, line=MID_GRAY,
             line_w=0.75)
    add_text(s, 1.0, 14.1, 25.5, 2.9, [
        ("핵심 메시지",
         {"size": SZ_SECTION, "bold": True, "color": BLUE}),
        ("별도 서버·DB·코드 없이 Claude Code + Agent Teams + 마크다운 "
         "양식만으로 보고·리서치·기술과제 발굴 전 과정을 재현 가능하게 "
         "자동화한다.", {"size": SZ_SECTION, "bold": True, "color": BLUE}),
        ("단순 LLM 호출이 아니라 출처·재현성·감사가능성을 갖춘 팀 협업이며, "
         "이미 10여 종 실산출물로 검증되었다.",
         {"size": SZ_SUB, "color": BLACK}),
        ("스킬 14종 + MCP 16+ 로 양식·도구 확장 비용 최소화 — 도메인이 달라도 "
         "스킬 자동 발동만으로 적용.",
         {"size": SZ_SUB, "bold": True, "color": BLACK}),
    ])

    add_footnote(s, "* 출처 : CLAUDE.md · README.md · .claude/agents/*.md ·"
                    " references/roadmap/ · templates/사내양식/"
                    "보고양식_카탈로그.md")


def aux_examples(prs):
    """보조 ⑥ — 현재 에이전트가 실제 생성한 산출물의 구체 발췌 (강조)."""
    s = new_blank_slide(prs)
    header(s, "[보조 ⑦/11] 실제 산출 예시 발췌 — 에이전트가 만든 결과")

    add_text(s, 0.6, 1.58, 26.32, 0.55,
             [("아래는 가공 예시가 아니라 본 5인 팀이 실제 생성한 outputs/ "
               "원문 발췌 (수치·결정·과제는 실데이터)",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])

    cards = [
        ("① 주간보고  ·  data + ops + document-writer",
         "outputs/주간보고_2026-05-12.md",
         ["• KPI 실측 표: 양산가동률 87/90%(96.7%), 물류 원가절감 "
          "21.3%(목표 16.2% → 131.5%)",
          "• Capex: 사업계획 3,791.4억 → 실적예상 2,985.3억, "
          "△806.1억(21.3%) 절감 추진",
          "• 모든 수치에 (출처: references/26FA KPI.md [259-289행]) "
          "자동 부착",
          "• data·ops 병렬 수집 → document-writer 가 17 placeholder "
          "만 치환"]),
        ("② 신규 R&D 회의록  ·  meeting / document-writer",
         "outputs/회의록_신규R&D과제착수_2026-05-12.md",
         ["• 결정: 1순위 Isaac SIM 사전검증 2026 H2 PoC / 2순위 "
          "MMF 2027 초",
          "• 선행예산 약 2~3억, 착수조건 'Vendor 풀 2개사+ 신규발굴' "
          "KPI 연동",
          "• 액션아이템 6건(담당·내용·기한) 표 자동 정리",
          "• 자료 없는 항목은 '자료 미보유' 명시 — 추측 금지 원칙"]),
        ("③ FA 기술과제후보  ·  fa-task-discovery v2 + tech-research",
         "outputs/fa-task-discovery/FA_기술과제후보_20260518.docx",
         ["• ① 마커리스 비주얼 SLAM 고적재 AMR (ABB Flexley P603 · "
          "Seegrid 20M mi)",
          "• ③ Omniverse/Isaac Sim Sim2Real 가상 시운전 (유디엠텍 "
          "공급계약)",
          "• ⑦ 휴머노이드 자재창고·검사 보조 PoC (LG CNS·컬리 MOU)",
          "• 6분야 8건 도출 → 제1부 뉴스 + 제2부 과제 2부 docx + "
          "카톡 요약"]),
        ("④ 경쟁사 동향 통합보고서  ·  tech-research × 4 병렬",
         "outputs/FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.md",
         ["• CATL Lighthouse 생산성 +17% · 결함 -99% / BYD 시안 "
          "자율도 ~97%",
          "• 삼성SDI 헝가리 디지털트윈 무투자 capa↑ / SK온×NVIDIA "
          "5만+ GPU",
          "• 'LGES 즉시 액션 후보 5선' 도출 → 26FA 로드맵·KPI 재설정 "
          "근거",
          "• md+docx+pptx 3종 동시 생성, 수치마다 URL 출처 인용"]),
    ]
    cw, ch = 12.96, 7.35
    gx, gy = 0.6, 2.25
    for i, (hdr, path, lines) in enumerate(cards):
        col, row = i % 2, i // 2
        x = gx + col * (cw + 0.4)
        y = gy + row * (ch + 0.35)
        add_rect(s, x, y, cw, ch, fill=WHITE, line=MID_GRAY, line_w=0.75)
        bn = add_rect(s, x, y, cw, 0.7, fill=DIM_GRAY, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.18)
        write_lines(bn.text_frame, [hdr], size=SZ_SUB, bold=True,
                    color=WHITE)
        add_text(s, x + 0.2, y + 0.78, cw - 0.4, 0.42,
                 [(path, {"size": SZ_FOOT, "color": DIM_GRAY,
                          "font": FONT_EMPH})])
        add_text(s, x + 0.2, y + 1.24, cw - 0.4, ch - 1.4,
                 [(t, {"size": SZ_SUB, "color": BLACK}) for t in lines])

    add_footnote(s, "* 출처 : outputs/ 실데이터 원문 (주간보고_2026-05-12 ·"
                    " 회의록_신규R&D과제착수_2026-05-12 ·"
                    " fa-task-discovery/tasks_20260518 ·"
                    " 경쟁사_동향보고서_2026-05-14)")


def aux_roadmap(prs):
    """보조 ⑩ — 에이전트 발전 계획 (개발 항목 + 데이터 보충 항목, 구체화).

    사용자 피드백: 실제로 더 활용되기 위해 (a) 무엇을 개발하고 (b) 어떤
    데이터를 채워야 하는지 구체적 계획을 단/중/장기로 나누어 표기.
    """
    s = new_blank_slide(prs)
    header(s, "[보조 ⑩/11] 에이전트 발전 계획 — 개발 항목 + 데이터 보충 (구체화)")

    # 상단 안내
    add_text(s, 0.6, 1.50, 26.32, 0.45,
             [("에이전트가 더 활용되려면 ①신규 기능을 개발하고 "
               "②축적 데이터를 보충해야 한다. 단/중/장기로 분리.",
               {"size": SZ_SUB, "color": BLACK})])

    # 3 phase × 2 column (개발 / 데이터) 매트릭스
    # 열: phase 라벨 (W4.4) · 개발 (W11.0) · 데이터 (W11.0) · gap 0.16
    label_w, dev_w, data_w, gap = 4.40, 10.86, 10.86, 0.20

    # 헤더 행
    hy = 2.05
    add_rect(s, 0.6, hy, label_w, 0.50, fill=DIM_GRAY, line=None)
    bn = s.shapes[-1]; bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame, ["단계"], size=SZ_SUB, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 0.6 + label_w + gap, hy, dev_w, 0.50,
             fill=DIM_GRAY, line=None)
    bn = s.shapes[-1]; bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame, ["① 개발 항목 (무엇을 만들/확장할 것인가)"],
                size=SZ_SUB, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 0.6 + label_w + dev_w + 2 * gap, hy, data_w, 0.50,
             fill=DIM_GRAY, line=None)
    bn = s.shapes[-1]; bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(bn.text_frame, ["② 데이터 보충 (무엇을 채울 것인가)"],
                size=SZ_SUB, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 3 phase 행 — 각 행 H2.45
    phases = [
        ("단기  ·  Q3'26\n(3개월)", True,
         [" • 사내 OneNote·SharePoint MCP PoC (mcp-builder)",
          " • Outlook/Teams 알림 hook — 보고서 자동 발송",
          " • 빌더 추가 시 sentry-code-review 자동 게이트",
          " • 임원 보고 카탈로그 8 블록 templates/ 완성"],
         [" • 협력사 Pool 50+ (현재 일부) — sample_data/협력사.csv 보충",
          " • KPI 베이스라인 12개월치 — 양산/물류/Capex 시계열",
          " • DST 진척 (SMC/MMF/셋업/밀폐이송) 주간 업데이트",
          " • 과거 회의록 1년치 references/과거회의록/ 적재"]),
        ("중기  ·  Q4'26\n(3-6개월)", False,
         [" • 사내 ERP/MES/PLM read-only MCP (mcp-builder)",
          " • SSO·사내 인증 통과 — IT/보안 검토",
          " • firecrawl 활성화 — 경쟁사 IR/특허 안정 수집",
          " • webapp-testing PoC — 사내 대시보드 회귀 검증"],
         [" • FA 6분야 외부 동향 1년치 누적 (fa-task-discovery 매일)",
          " • 경쟁사 4사 IR/특허 12개월치 references/경쟁사/",
          " • 임원 보고 13종 양식 → references 매핑 완료",
          " • 그룹 4부서 KPI 양식 통합 (인프라기획팀 주관)"]),
        ("장기  ·  2027~\n(6개월+)", False,
         [" • 자기평가 Hook — 에이전트 출력 회귀 검증",
          " • 공유 Task List — 멀티 영역 동시 운영 (Agent Teams)",
          " • skill-creator 로 사내 도메인 스킬 50+ 표준화",
          " • 사내 표준 보고 플랫폼화 + 거버넌스 연계"],
         [" • 사내 SOP 50+개 · 협력사 평가 분기별 데이터",
          " • 그룹 전영역 KPI 통합 데이터 마트",
          " • 사내 용어집·정책 references/policy/ 완성",
          " • 산출물 카탈로그 100+종 (재사용 자산 풀)"]),
    ]
    ry = hy + 0.55
    row_h = 4.55
    for i, (ph_label, accent, dev_items, data_items) in enumerate(phases):
        # 좌: phase 라벨
        add_rect(s, 0.6, ry, label_w, row_h,
                 fill=(LIGHT_BLUE if accent else DIM_GRAY), line=None)
        bn = s.shapes[-1]; bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.2)
        write_lines(bn.text_frame, [ph_label],
                    size=SZ_SECTION, bold=True,
                    color=(BLUE if accent else WHITE),
                    align=PP_ALIGN.CENTER)
        # 중: 개발 항목 (SZ_SUB 9pt, 4 bullets + 행간 여유)
        add_rect(s, 0.6 + label_w + gap, ry, dev_w, row_h,
                 fill=WHITE, line=MID_GRAY, line_w=0.5)
        add_text(s, 0.6 + label_w + gap + 0.20, ry + 0.25,
                 dev_w - 0.40, row_h - 0.45,
                 [(t, {"size": SZ_SUB, "color": BLACK}) for t in dev_items])
        # 우: 데이터 보충
        add_rect(s, 0.6 + label_w + dev_w + 2 * gap, ry, data_w, row_h,
                 fill=LIGHT_GRAY, line=MID_GRAY, line_w=0.5)
        add_text(s, 0.6 + label_w + dev_w + 2 * gap + 0.20, ry + 0.25,
                 data_w - 0.40, row_h - 0.45,
                 [(t, {"size": SZ_SUB, "color": BLACK}) for t in data_items])
        ry += row_h + 0.20

    # 하단 운영 원칙 띠 (1줄)
    add_rect(s, 0.6, 17.40, 26.32, 0.80, fill=LIGHT_BLUE, line=MID_GRAY,
             line_w=0.5)
    add_text(s, 0.80, 17.50, 25.92, 0.60, [
        ("불변 원칙 (확장해도 유지) — ",
         {"size": SZ_FOOT, "bold": True, "color": BLUE}),
        ("추측 금지 · 출처 자동 검증(Hook 강제) · 원본 불변 · 같은 입력=같은 출력 · "
         "사람 승인 작업은 '제안'만 — 통제권 유지",
         {"size": SZ_FOOT, "color": BLACK}),
    ])

    add_footnote(s, "* 출처 : CLAUDE.md · README.md ·"
                    " references/roadmap/2026_FA기술담당_중장기로드맵_v2.md ·"
                    " .claude/agents/*.md · .claude/skills/mcp-builder/SKILL.md")


def aux_agent_teams(prs):
    """보조 ② — Agent Teams 운영 패턴: 서브에이전트와 차이 + 직접 소통."""
    s = new_blank_slide(prs)
    header(s, "[보조 ②/11] Agent Teams 운영 — 일반 서브에이전트와 차이 + "
              "에이전트끼리 직접 소통")

    add_text(s, 0.6, 1.58, 26.32, 0.55,
             [("본 레포는 단순 서브에이전트 호출이 아닌 'Agent Teams' "
               "(독립 세션 + 우편함 통신) 패턴 — Claude Code v2.1.32+ "
               "공식 기능 (2026.02, 실험적)",
               {"size": SZ_SUB, "bold": True, "color": BLACK})])

    # ── 좌우 비교 카드 ────────────────────────────────────────────
    cards = [
        ("일반 서브에이전트  (참고)", DIM_GRAY, WHITE, [
            "• 부모 세션이 매번 명시 호출 (부모-자식 1:1)",
            "• 결과는 부모로만 반환 — 자식끼리 통신 X",
            "• 부모 컨텍스트 일부 상속 (오염 가능)",
            "• 순차 호출 위주 — 병렬은 부모가 직접 조립",
            "• 공유 상태 없음, 표준 통신 프로토콜 없음",
        ]),
        ("Agent Teams  (본 레포 적용)", LIGHT_BLUE, BLUE, [
            "• 자동 spawn + 독립 세션 (in-process, 대화 미상속)",
            "• 우편함(Mailbox) 표준 통신:",
            "  요청 to:/subject:/payload:  ↔  응답 key:/value:/source:",
            "• **에이전트끼리 직접 요청** — writer 가 data·ops 에게 우편함",
            "• 같은 메시지에서 다수 팀원 동시 spawn (병렬 수집)",
            "• 환경: CLAUDE_CODE_EXPERIMENTAL_AGENT_TEAMS=1",
        ]),
    ]
    cw = 13.0
    cy = 2.32
    ch = 6.0
    for i, (hdr, hfill, hcolor, lines) in enumerate(cards):
        x = 0.6 + i * (cw + 0.32)
        add_rect(s, x, cy, cw, ch, fill=WHITE, line=MID_GRAY, line_w=0.75)
        bn = add_rect(s, x, cy, cw, 0.74, fill=hfill, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.18)
        write_lines(bn.text_frame, [hdr], size=SZ_SECTION, bold=True,
                    color=hcolor)
        add_text(s, x + 0.22, cy + 0.85, cw - 0.44, ch - 1.0,
                 [(t, {"size": SZ_SUB, "color": BLACK}) for t in lines])

    # ── 하단 — 우편함 통신 다이어그램 + 본 레포 실제 적용 ─────────
    dy = cy + ch + 0.32
    add_text(s, 0.6, dy, 26.32, 0.55,
             [("우편함 통신 — 팀장 거치지 않는 에이전트 직접 소통 "
               "(본 레포 실제 라우팅)",
               {"size": SZ_SECTION, "bold": True, "color": BLACK})])

    # 좌측 미니 다이어그램 (텍스트 박스로 표현)
    dgy = dy + 0.65
    add_rect(s, 0.6, dgy, 13.0, 4.55, fill=LIGHT_GRAY, line=LINE_GRAY,
             line_w=0.5)
    add_text(s, 0.85, dgy + 0.15, 12.5, 4.3, [
        ("사용자",
         {"size": SZ_BODY, "bold": True, "color": BLACK,
          "align": PP_ALIGN.CENTER}),
        ("↓",
         {"size": SZ_BODY, "color": DIM_GRAY,
          "align": PP_ALIGN.CENTER}),
        ("팀장 (메인 Claude Code 세션) — 분해·dispatch",
         {"size": SZ_SUB, "bold": True, "color": BLACK,
          "align": PP_ALIGN.CENTER}),
        ("↓",
         {"size": SZ_BODY, "color": DIM_GRAY,
          "align": PP_ALIGN.CENTER}),
        ("document-writer  /  ppt-writer    (작성 팀원)",
         {"size": SZ_SUB, "bold": True, "color": BLUE,
          "align": PP_ALIGN.CENTER}),
        ("↓  ←  직접 요청 (팀장 거치지 않음)  →  ↓",
         {"size": SZ_FOOT, "color": BLUE,
          "align": PP_ALIGN.CENTER}),
        ("data-teammate     ops-teammate     tech-research",
         {"size": SZ_SUB, "bold": True, "color": BLACK,
          "align": PP_ALIGN.CENTER}),
        ("↓",
         {"size": SZ_BODY, "color": DIM_GRAY,
          "align": PP_ALIGN.CENTER}),
        ("sample_data/  ·  references/  ·  외부 MCP",
         {"size": SZ_FOOT, "color": DIM_GRAY, "font": FONT_EMPH,
          "align": PP_ALIGN.CENTER}),
    ])

    # 우측 — 본 레포 적용 매핑
    add_rect(s, 13.92, dgy, 13.0, 4.55, fill=WHITE, line=MID_GRAY,
             line_w=0.75)
    add_text(s, 14.17, dgy + 0.18, 12.5, 4.2, [
        ("Agent Teams 핵심 효과 (본 레포 검증)",
         {"size": SZ_SUB, "bold": True, "color": BLACK}),
        (" • 독립 세션 → 팀원 컨텍스트 깨끗 — 추측·오염 차단",
         {"size": SZ_SUB, "color": BLACK}),
        (" • 우편함 통신 → 작성자가 데이터 팀원에게 직접 요청, "
         "팀장 토큰 낭비 없음",
         {"size": SZ_SUB, "color": BLACK}),
        (" • 병렬 spawn → 주간보고: data + ops 같은 라운드에 "
         "동시 수집 (실증)",
         {"size": SZ_SUB, "color": BLACK}),
        (" • 부족 데이터 라운드 — writer 가 우편함으로 추가 요청 "
         "→ 팀원 응답 → 재합성",
         {"size": SZ_SUB, "color": BLACK}),
        (" • 의도적 제외: 공유 Task List — 보고서 작성 흐름엔 불필요"
         " (코드 리뷰 등 도입 시 검토)",
         {"size": SZ_FOOT, "color": DIM_GRAY}),
    ])

    add_footnote(s, "* 출처 : README.md [27-41행] (Agent Teams ↔ 본 레포 "
                    "매핑 표) · CLAUDE.md [17-55행] · 짐코딩 'Claude "
                    "Agent Teams 완벽 정리'(2026-02-28)")


def aux_fa_mapping(prs):
    """보조 ⑧ — FA기술혁신파트 업무 → Agent 매핑 (도움/대체)."""
    s = new_blank_slide(prs)
    header(s, "[보조 ⑧/11] FA기술혁신파트 업무 → Agent 매핑 "
              "(현재 도움 / 대체 가능)")

    # 표 헤더 (4열, 총 26.32cm)
    cols = [
        ("FA 업무 (사용자 정의)", 5.0, DIM_GRAY),
        ("담당 에이전트 · 스킬", 5.2, DIM_GRAY),
        ("현재 도움 (검증된 사례)", 10.62, DIM_GRAY),
        ("자동화 수준", 5.5, DIM_GRAY),
    ]
    hx = 0.6
    hy = 1.7
    hh = 0.75
    cx = hx
    for label, w, fill in cols:
        bn = add_rect(s, cx, hy, w, hh, fill=fill, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.12)
        bn.text_frame.margin_right = Cm(0.12)
        write_lines(bn.text_frame, [label], size=SZ_SUB, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        cx += w

    rows = [
        ("1. 중점추진업무\n   KPI·M&R·로드맵",
         "data + writer + xlsx",
         "주간보고 KPI 표 자동(양산가동률 87/90%, 물류 -21.3%), "
         "로드맵 v2 4축 KPI, xlsx 트래커",
         "★★ 자동", BLUE),
        ("2. Project\n   투자비·면적·구매",
         "data + ops + writer + xlsx",
         "Capex 3,791.4억→2,985.3억 △806.1억(21.3%) 절감, 면적/구매 "
         "표 자동, xlsx 시뮬레이션",
         "★★ 자동", BLUE),
        ("3. MRM\n   과제 도출·진척·로드맵",
         "tech-research + fa-task-discovery + ppt-writer",
         "FA 기술과제후보 8건 docx + 카톡 (매일 자동), 회의록 결정 "
         "1순위 Isaac SIM",
         "★★ 자동 (매일)", BLUE),
        ("4. 협력사\n   Pool·신규·설비군",
         "ops + pdf + docx",
         "엡실론·델타로봇 변경, 물류 Vendor Pool 재편안 · 협력사 "
         "사양서 pdf 추출 + docx 평가서",
         "★ 부분 (등록·승인 사람)", BLACK),
        ("5. 미래기술 개발\n   DST: SMC/MMF/셋업/밀폐이송",
         "tech-research + data + ppt-writer + pdf",
         "전고체 밀폐이송장치 검토(EFEM+FOUP), 유연성지표 요약, "
         "SMC Demo 절감, 도면 pdf 추출",
         "★★ 보고 자동 (PoC 사람)", BLUE),
        ("6. 경비·투자\n   Project/경상/일반",
         "data + xlsx",
         "Capex 집행률 65%, SMC Demo 1.466→1.172억(20%) 절감 표 + "
         "xlsx 시뮬레이션",
         "★ 부분 (결재 사람)", BLACK),
        ("7. 보고\n   CPO·센터장·생산성·사업부",
         "writer 2종 + meeting-minutes + brand-guidelines + internal-comms",
         "주간보고, R&D 회의록, FA자동화 임원보고 예시(메인1+보조4), "
         "사내 공지·3P 메일 자동",
         "★★★ 완전 자동", BLUE),
        ("8. 인사\n   포상·채용·진급·교육·출장",
         "ops + internal-comms",
         "HR 데이터 입력 시 출장/교육/포상 정리 자동, 사내 공지·"
         "주간 메일 자동 (인사평가는 사람)",
         "△ 보조 (사람 결정)", BLACK),
        ("9. 기타\n   실사·타운홀·워크샵",
         "writer + meeting-merger + docx",
         "녹취+OneNote → 통합 Word(.docx) 자동, 회의록 PPTX 자동, "
         "트래킹 체인지",
         "△ 보조 (행사 진행 사람)", BLACK),
    ]
    ry = hy + hh
    rh = 1.13
    for ridx, (work, agent, help_, lvl, lcol) in enumerate(rows):
        zebra = LIGHT_GRAY if ridx % 2 == 0 else WHITE
        cx = hx
        widths = [c[1] for c in cols]
        for ci, w in enumerate(widths):
            add_rect(s, cx, ry, w, rh, fill=zebra, line=LINE_GRAY,
                     line_w=0.3)
            cx += w
        # 1열: 업무명 (Bold)
        add_text(s, hx + 0.12, ry + 0.12, widths[0] - 0.24, rh - 0.2,
                 [(work, {"size": SZ_FOOT, "bold": True, "color": BLACK})])
        # 2열: 담당
        add_text(s, hx + widths[0] + 0.12, ry + 0.18,
                 widths[1] - 0.24, rh - 0.3,
                 [(agent, {"size": SZ_FOOT, "color": BLACK,
                           "font": FONT_EMPH})])
        # 3열: 현재 도움
        add_text(s, hx + widths[0] + widths[1] + 0.12, ry + 0.12,
                 widths[2] - 0.24, rh - 0.2,
                 [(help_, {"size": SZ_FOOT, "color": BLACK})])
        # 4열: 수준 (강조)
        x4 = hx + widths[0] + widths[1] + widths[2] + 0.12
        add_text(s, x4, ry + 0.32, widths[3] - 0.24, rh - 0.5,
                 [(lvl, {"size": SZ_SUB, "bold": True, "color": lcol,
                         "align": PP_ALIGN.CENTER})])
        ry += rh

    # 하단 요약 박스
    sy = ry + 0.18
    add_rect(s, 0.6, sy, 26.32, 1.85, fill=LIGHT_BLUE, line=MID_GRAY,
             line_w=0.75)
    add_text(s, 0.9, sy + 0.18, 25.7, 1.5, [
        ("요약 — 9개 업무 영역 자동화 수준",
         {"size": SZ_SUB, "bold": True, "color": BLUE}),
        (" ★★★ 완전 자동(1): 보고 — 가장 큰 시간 절감.  "
         "★★ 자동(4): 중점추진·Project·MRM·미래기술 — 사람 검토 후 적용.",
         {"size": SZ_SUB, "color": BLACK}),
        (" ★ 부분(2): 협력사·경비 — 데이터 정리는 자동, 승인은 사람.  "
         "△ 보조(2): 인사·기타 — 회의록·요약은 자동, 결정·진행은 사람.",
         {"size": SZ_SUB, "color": BLACK}),
        (" → 스킬 14종(도메인 3 + 신규 11) 자동 발동으로 각 업무에 양식 맞춤 적용",
         {"size": SZ_SUB, "bold": True, "color": BLACK}),
    ])

    add_footnote(s, "* 출처 : 사용자 인풋 업무 리스트(2026-05-19) ·"
                    " outputs/(주간보고·회의록·fa-task-discovery) ·"
                    " references/roadmap/v2 · .claude/skills/ (14종)")


def aux_guide_setup(prs):
    """[가이드 ①/2] 사용 가이드 — 환경 셋업·첫 명령 (초보자용)."""
    s = new_blank_slide(prs)
    header(s, "[가이드 ①/2] 사용 가이드 — 환경 셋업·첫 명령 (초보자용)")

    # 4 step 박스 + 좌측 step 번호 띠
    add_text(s, 0.6, 1.55, 26.32, 0.45,
             [("Claude Code + 본 레포만 있으면 별도 서버·DB·설치 없이 "
               "10분 안에 첫 보고서를 자동 생성할 수 있다.",
               {"size": SZ_SUB, "color": BLACK})])

    steps = [
        ("STEP 1", "환경 준비 (1회만)", [
            (" • Claude Code 설치 — CLI / Desktop / Web / IDE 중 택일",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • 환경 변수: ", {"size": SZ_FOOT, "color": BLACK}),
            ("export CLAUDE_CODE_EXPERIMENTAL_AGENT_TEAMS=1",
             {"size": SZ_FOOT, "color": BLACK, "font": FONT_EMPH,
              "bold": True}),
            (" • 리포 clone: ", {"size": SZ_FOOT, "color": BLACK}),
            ("git clone <repo>/claude_FA-Agent_1  &&  cd claude_FA-Agent_1",
             {"size": SZ_FOOT, "color": BLACK, "font": FONT_EMPH}),
            (" • 세션 시작: ", {"size": SZ_FOOT, "color": BLACK}),
            ("claude",
             {"size": SZ_FOOT, "color": BLACK, "font": FONT_EMPH,
              "bold": True}),
            (" 입력 — CLAUDE.md 자동 로드, 인덱스 자동 갱신(SessionStart Hook)",
             {"size": SZ_FOOT, "color": BLACK}),
        ]),
        ("STEP 2", "폴더 구조 이해 (5분)", [
            (" • templates/ : 보고서 양식 (placeholder MD / PPTX) — 절대 수정 X",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • sample_data/ : 정량 데이터 (KPI · Capex · MRM · 협력사)",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • references/ : 참조 자료 (자동 인덱싱) — 추가만 하면 자동 인식",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • outputs/ : 결과물 폴더 (자동 생성) — 결과만 여기서 확인",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • .claude/agents/ : 5 에이전트 정의  /  .claude/skills/ : 14 스킬 (자동 발동)",
             {"size": SZ_FOOT, "color": BLACK}),
        ]),
        ("STEP 3", "첫 명령 (한 줄로 자동 위임)", [
            (" 사용자 입력 →  팀장이 분해 →  팀원 spawn →  스킬 자동 발동 →  outputs/ 생성",
             {"size": SZ_FOOT, "bold": True, "color": BLACK}),
            (" 예시 ① : ", {"size": SZ_FOOT, "color": BLACK}),
            ("\"이번 주 주간보고 작성해줘\"",
             {"size": SZ_FOOT, "color": BLUE, "bold": True}),
            (" → data + ops + document-writer 자동 호출",
             {"size": SZ_FOOT, "color": BLACK}),
            (" 예시 ② : ", {"size": SZ_FOOT, "color": BLACK}),
            ("\"녹취록으로 회의록 만들어줘\"",
             {"size": SZ_FOOT, "color": BLUE, "bold": True}),
            (" → meeting-minutes 스킬 발동, 인터랙티브 UI",
             {"size": SZ_FOOT, "color": BLACK}),
            (" 예시 ③ : ", {"size": SZ_FOOT, "color": BLACK}),
            ("\"FA 기술 동향 조사해줘\"",
             {"size": SZ_FOOT, "color": BLUE, "bold": True}),
            (" → fa-task-discovery v2 (6분야 자동 수집·과제 도출)",
             {"size": SZ_FOOT, "color": BLACK}),
        ]),
        ("STEP 4", "결과 확인 + 자동 검증", [
            (" • outputs/<날짜>.md (· .pptx) 자동 생성 — 모든 수치에 (출처:…) 자동 부착",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • TaskCompleted Hook 이 출처 미기입 자동 차단 (인용 ≥1 필수)",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • 자료 없는 항목은 자동으로 ", {"size": SZ_FOOT, "color": BLACK}),
            ("\"자료 미확인\"",
             {"size": SZ_FOOT, "color": BLUE, "bold": True}),
            (" 으로 명시 (추측 금지 원칙)",
             {"size": SZ_FOOT, "color": BLACK}),
            (" • brand-guidelines 자동 적용 — 색·폰트·풋노트 위반 자가 점검",
             {"size": SZ_FOOT, "color": BLACK}),
        ]),
    ]

    y = 2.05
    step_w = 3.2
    body_w = 26.32 - step_w - 0.1
    row_h = 3.65
    for sn, sname, lines in steps:
        # 좌: STEP n 번호 띠
        add_rect(s, 0.6, y, step_w, row_h, fill=DIM_GRAY, line=None)
        bn = s.shapes[-1]
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_lines(bn.text_frame, [
            (sn, {"size": SZ_BAND, "bold": True, "color": WHITE,
                  "align": PP_ALIGN.CENTER}),
            (sname, {"size": SZ_SUB, "bold": True, "color": WHITE,
                     "align": PP_ALIGN.CENTER}),
        ])
        # 우: 본문 박스
        add_rect(s, 0.6 + step_w + 0.1, y, body_w, row_h, fill=WHITE,
                 line=MID_GRAY, line_w=0.5)
        # 본문 lines — paragraph 별 옵션 적용
        tb = s.shapes.add_textbox(Cm(0.6 + step_w + 0.30), Cm(y + 0.15),
                                  Cm(body_w - 0.40), Cm(row_h - 0.30))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.clear()
        # 라인을 paragraph 단위로 그룹 — 각 라인 시작은 새 paragraph,
        # 같은 라인 내 (텍스트, 옵션) 튜플은 같은 paragraph 의 다음 run
        # 본 구현은 "한 라인 = 한 paragraph" 모델로, 마커 (예: "  ▸")
        # 가 시작인 항목만 새 paragraph 로 처리. 여기서는 단순화: 각
        # 원소를 한 paragraph 로 처리하되 인접 run 은 같은 paragraph 로.
        # → 본문은 다양한 색·폰트가 섞이므로, 각 원소가 한 paragraph 가
        # 되어도 시각적으로 줄이 자연스럽게 분리됨 (가독성 OK).
        first = True
        for txt, opts in lines:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            r = p.add_run()
            set_run(r, txt,
                    font=opts.get("font", FONT_BODY),
                    size=opts.get("size", SZ_FOOT),
                    bold=opts.get("bold", False),
                    color=opts.get("color", BLACK))
        y += row_h + 0.10

    add_footnote(s, "* 출처 : CLAUDE.md · README.md · .claude/settings.json"
                    " (SessionStart Hook) · scripts/hooks/"
                    "check_output_citations.py")


def aux_guide_scenarios(prs):
    """[가이드 ②/2] 사용 가이드 — 시나리오별 명령 예시·FAQ (초보자용)."""
    s = new_blank_slide(prs)
    header(s, "[가이드 ②/2] 사용 가이드 — 시나리오별 명령 예시·FAQ (초보자용)")

    add_text(s, 0.6, 1.55, 26.32, 0.42,
             [("아래 '명령 예시' 한 줄을 그대로 입력하면 해당 시나리오의 "
               "에이전트·스킬이 자동 발동된다.",
               {"size": SZ_SUB, "color": BLACK})])

    # 시나리오 6개 — 2 col × 3 row
    scenarios = [
        ("① 주간보고",
         "\"이번 주 주간보고 작성해줘\"",
         "data + ops + document-writer + xlsx",
         "outputs/주간보고_YYYY-MM-DD.md"),
        ("② 회의록",
         "\"녹취록 정리해서 회의록 만들어줘\"",
         "meeting-minutes 스킬 + document-writer",
         "outputs/회의록_*.pptx"),
        ("③ FA 과제 발굴 (매일)",
         "\"FA 기술 동향 조사 / FA 신규 과제 발굴\"",
         "fa-task-discovery v2 + tech-research",
         "outputs/fa-task-discovery/*.docx + 카톡"),
        ("④ 경쟁사 동향 분석",
         "\"경쟁사(CATL·BYD·삼성SDI·SK온) 동향 분석\"",
         "tech-research × 4 병렬 + writer 2종",
         "outputs/경쟁사_*.md/.docx/.pptx"),
        ("⑤ 임원 보고 PPT",
         "\"임원 보고 PPT 만들어줘 (블록 A+C+G+B+E+H)\"",
         "ppt-writer + brand-guidelines + 카탈로그",
         "outputs/*_임원보고_*.pptx (메인+보조)"),
        ("⑥ 협력사 사양 추출",
         "\"이 PDF 사양서에서 표·이미지 뽑아줘\"",
         "ops + pdf 스킬 + docx",
         "구조화 데이터 + 협력사 평가서 docx"),
    ]
    cw, ch = 13.06, 2.85
    gx, gy = 0.6, 2.10
    for i, (title, prompt, agents, output) in enumerate(scenarios):
        col = i % 2
        row = i // 2
        x = gx + col * (cw + 0.20)
        y = gy + row * (ch + 0.18)
        # 카드 박스
        add_rect(s, x, y, cw, ch, fill=WHITE, line=MID_GRAY, line_w=0.6)
        bn = add_rect(s, x, y, cw, 0.50, fill=DIM_GRAY, line=None)
        bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bn.text_frame.margin_left = Cm(0.18)
        write_lines(bn.text_frame, [title], size=SZ_SUB, bold=True,
                    color=WHITE)
        # 명령 예시 (BLUE Bold)
        add_text(s, x + 0.20, y + 0.58, cw - 0.40, 0.55,
                 [("명령 예시 : ", {"size": SZ_FOOT, "bold": True,
                                "color": BLACK}),
                  (prompt, {"size": SZ_FOOT, "bold": True, "color": BLUE})])
        # 발동
        add_text(s, x + 0.20, y + 1.18, cw - 0.40, 0.55,
                 [("발동 : ", {"size": SZ_FOOT, "bold": True,
                              "color": BLACK}),
                  (agents, {"size": SZ_FOOT, "color": BLACK,
                            "font": FONT_EMPH})])
        # 산출물
        add_text(s, x + 0.20, y + 1.78, cw - 0.40, 0.55,
                 [("산출 : ", {"size": SZ_FOOT, "bold": True,
                              "color": BLACK}),
                  (output, {"size": SZ_FOOT, "color": BLACK,
                            "font": FONT_EMPH})])
        # 추가 팁 (TIP)
        tips = {
            "① 주간보고": "TIP : 매주 같은 명령 반복 → 출처 자동 재검증 (재현성)",
            "② 회의록": "TIP : 텍스트 붙여넣기 OK. 참석자·결정·액션 인터랙티브 확정",
            "③ FA 과제 발굴 (매일)": "TIP : 매일 자동 실행 가능. 날짜 브랜치 + PR 자동",
            "④ 경쟁사 동향 분석": "TIP : 4 팀원 동시 spawn. 결과는 references/경쟁사/ 자동",
            "⑤ 임원 보고 PPT": "TIP : LGES 가이드 자동 적용 — 풋노트·색·폰트 자가 검증",
            "⑥ 협력사 사양 추출": "TIP : 스캔 PDF 도 OCR 자동. references/기술자료/ 권장",
        }
        add_text(s, x + 0.20, y + 2.32, cw - 0.40, 0.50,
                 [(tips[title], {"size": 7, "color": MOLD_GREEN,
                                 "bold": True})])

    # FAQ 카드 (5개)
    faq_y = 2.10 + 3 * (ch + 0.18) + 0.05
    add_rect(s, 0.6, faq_y, 26.32, 0.50, fill=DIM_GRAY, line=None)
    bn = s.shapes[-1]
    bn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    bn.text_frame.margin_left = Cm(0.18)
    write_lines(bn.text_frame, ["자주 묻는 질문 (FAQ)"],
                size=SZ_SUB, bold=True, color=WHITE)

    faqs = [
        ("Q1. 자료가 없는 항목은?",
         "→ 자동으로 \"자료 미확인\" 으로 명시. 추측 금지 원칙 (CLAUDE.md)"),
        ("Q2. 신규 양식을 추가하려면?",
         "→ templates/ 또는 templates/사내양식/ 에 .md/.pptx 두면 자동 인식"),
        ("Q3. 출처 검증을 끄고 싶다면?",
         "→ 불가 (Hook 이 강제). \"자료 미확인\" 명시로 우회 가능"),
        ("Q4. 에이전트가 응답 안 하면?",
         "→ /resume 후 \"팀원 다시 생성해줘\" 입력 (in-process 팀원 미복원 제약)"),
        ("Q5. 한·영 폰트 자동 분리는?",
         "→ ppt-writer 가 한글 LG스마트체 / 영문 Arial Narrow 자동 run 분리"),
    ]
    fy = faq_y + 0.60
    row_h_faq = 0.62
    for q, a in faqs:
        add_text(s, 0.8, fy, 10.5, row_h_faq,
                 [(q, {"size": SZ_FOOT, "bold": True, "color": BLACK})])
        add_text(s, 11.4, fy, 15.5, row_h_faq,
                 [(a, {"size": SZ_FOOT, "color": BLACK})])
        fy += row_h_faq

    add_footnote(s, "* 출처 : CLAUDE.md (추측 금지·출처 명시 원칙) ·"
                    " .claude/agents/*.md (에이전트 트리거) ·"
                    " .claude/skills/*/SKILL.md (스킬 트리거 키워드) ·"
                    " scripts/hooks/check_output_citations.py")


# ─── MD 사이드카 ──────────────────────────────────────────────────
MD = """# 인프라그룹 AI 경진대회 과제 — FA 기술혁신 Agent 팀 (2026-05-19 갱신판)

- **양식**: `templates/사내양식/260511_인프라그룹 AI경진대회 과제_양식.pptx` (LGES 사내 표준)
- **작성일**: 2026-05-19  (이전판: `outputs/AI경진대회_FA기술혁신Agent팀_2026-05-13.pptx`)
- **책임자**: FA기술혁신파트 강모원책임 · **참여 형태**: 개인 참여 (1명)
- **분량**: 메인 2장(양식 폼) + 보조 10장 + 가이드 2장 = 14 슬라이드
- **빌더**: `scripts/build_contest_2026_05_19.py` (결정론 — 같은 입력=같은 출력)

---

## Slide 1 — 메인 (양식 4사분면 + 참여인원 + FA 9업무 자동화 카드)

- **과제 명**: FA 기술혁신 Agent 팀 (**제목 BLACK** — 사용자 피드백 #3 반영)
- **요지(부제)**: 5 AI 에이전트 팀이 FA 9개 업무 영역을 자동화 — 보고 완전
  자동 · 외부 동향 매일 발굴 · 출처 100% 자동 검증
- **이슈·배경(7)**: 보고서당 수시간 · 출처 누락·양식 제각각 · 회의록 지연 ·
  외부 동향 반영 지연 · 보고 표준 부재 · 산출 양식 분산 · 사내 MCP 미연결
- **개선 사항(7)**: 5 에이전트 우편함 · **에이전트 직접 소통(파랑)** ·
  도메인 스킬 14종 자동 발동 · MCP 서버 16+ 자동 연결 ·
  결정론 PPT · LGES 가이드+brand-guidelines · Hook 2 + sentry-code-review
- **기대 효과(9)**: 정량 5 / 정성 4 (신규: 신규 스킬 11 추가, 빌더 18+ 품질
  게이트, brand-guidelines 일관)
- **참여 인원**: 강모원 1명 (개인 참여) — 표 H 4.0→1.4 cm 축소
- **신규 카드 : FA 9업무 → Agent 자동화 수준 (보조 ⑧ 요약)** —
  사용자 피드백 #4 반영. 메인 한 장에서 "결론적으로 무엇이 가능한지"
  4행 미니 표로 즉시 노출:
  - ★★★ 완전 자동 (1): 보고 — CPO·센터장·생산성·사업부
  - ★★ 자동 (4): 중점추진·Project·MRM·미래기술
  - ★ 부분 (2): 협력사·경비 — 데이터 자동, 등록/결재 사람
  - △ 보조 (2): 인사·기타 — 회의록·요약 자동, 결정·진행 사람
- **핵심 메시지 띠**: "5 에이전트 + 스킬 14 + Hook 2 + MCP 16+ 로 9업무
  자동화. 작성 시간 수시간→분, 출처 100% 자동, 산출물 10+종"
- **레이아웃 수정**: SmartArt 흐름 띠 제거(개선 전/후 chip 오버랩),
  템플릿 노트 "* 개인 참여시" 제거(스킬 카드 오버랩)

(출처: `CLAUDE.md`, `README.md [17-55행]`, `.claude/agents/*.md`,
`.claude/skills/*/SKILL.md` (14종), `outputs/주간보고_2026-05-12.md`)

## Slide 2 — 보조 ① 시스템 구성 (아키텍처 다이어그램)

**아키텍처 PNG 이미지 1장 (outputs/_assets/agent_architecture.png/.svg)**
— 사용자 피드백 #2 반영. 텍스트 카드 위주의 이전 레이아웃을 SVG/PNG
다이어그램으로 교체해 한눈에 시스템 흐름 파악 가능:

```
사용자 요청 → 팀장(메인 세션) → 5 에이전트 카드 (data·ops·tech·writer 2종)
                              ↕ 우편함 통신 (에이전트끼리 직접)
스킬 14 (도메인 3 + 산출 5 + 메타·운영 4 + 장기 2) + MCP 16+
→ outputs/ (10+종 산출물) → ✓ TaskCompleted Hook 출처 자동 검증
+ SessionStart Hook (references 인덱싱) + sentry-code-review (빌더 게이트)
```

하단 띠 : SOP 분해 → 병렬 수집 → 위임 작성(스킬 자동 발동) → 품질 게이트.
같은 입력 = 같은 출력 · 출처 100% 자동 검증 · 양식 무결성 100%.

(출처: `outputs/_assets/agent_architecture.svg (·png)`,
`.claude/agents/*.md`, `.claude/skills/*/SKILL.md` (14종),
`CLAUDE.md [67-84행]`)

## Slide 3 — 보조 ② Agent Teams 운영 (서브에이전트와 차이 + 직접 소통)

본 레포는 단순 서브에이전트 호출이 아니라 **Agent Teams**(독립 세션 +
우편함 통신) 패턴을 적용:

- **일반 서브에이전트**: 부모-자식 1:1 호출, 결과는 부모로만, 자식끼리 통신 X
- **Agent Teams (본 레포)**:
  - 자동 spawn + 독립 세션 (대화 미상속)
  - 우편함(Mailbox) 표준 통신 `to:/subject:/payload:` ↔ `key:/value:/source:`
  - **에이전트끼리 직접 요청** — writer 가 data·ops 에게 우편함으로
  - 같은 메시지에서 다수 팀원 동시 spawn (병렬 수집)
  - 환경: `CLAUDE_CODE_EXPERIMENTAL_AGENT_TEAMS=1` (Claude Code v2.1.32+)

검증: 주간보고 작성 시 document-writer 가 data·ops 에게 같은 라운드에
직접 요청 → 합쳐서 출력 (팀장 토큰 낭비 없음).

(출처: `README.md [27-41행]`, `CLAUDE.md [17-55행]`, 짐코딩 "Claude
Agent Teams 완벽 정리" 2026-02-28)

## Slide 4 — 보조 ③ 5 팀원 · SOP · Hook

팀원 5명 카드 + SOP 4단계(분해→병렬수집→위임작성→품질게이트) + Hook 2개
(SessionStart 자동 인덱싱 / TaskCompleted 출처 자동 검증).

(출처: `.claude/agents/*.md`, `.claude/settings.json`,
`scripts/hooks/check_output_citations.py`, `scripts/build_reference_index.py`)

## Slide 5 — 보조 ④ 스킬 라이브러리 14종 (도메인 3 + 신규 11)

- **도메인 3종 (회의·과제 발굴)**: meeting-minutes(녹취→회의록 PPTX) ·
  meeting-merger(녹취+OneNote→통합 Word) · fa-task-discovery v2.
  v2 핵심: 2세션 자동 분기 · 2부 구성 docx · 전용 폴더 · 날짜 브랜치 +
  메인 대상 PR 자동 생성.
- **신규 11종 (2026-05 추가)** — 카테고리 3 그룹:
  - **산출 양식 (5)**: `docx`(Word 정책·SOP·평가서) · `xlsx`(KPI/Capex/MRM
    트래커) · `pdf`(사양·도면 OCR) · `doc-coauthoring`(임원 보이스 장문) ·
    `internal-comms`(공지·주간·3P·FAQ)
  - **메타·운영 (4)**: `brand-guidelines`(PPT/DOCX 색·폰트·풋노트 자동) ·
    `skill-creator`(신규 스킬 scaffold·eval) · `mcp-builder`(사내 ERP/MES/
    PLM MCP) · `sentry-code-review`(빌더 품질 게이트)
  - **장기 검토 (2)**: `webapp-testing`(Playwright — 사내 대시보드 도입 시) ·
    `firecrawl`(JS 사이트 안정 스크래핑 — 보안/IT 검토 후 PoC)

(출처: `.claude/skills/*/SKILL.md` (14종 — `meeting-minutes/`,
`meeting-merger/`, `fa-task-discovery/`, `docx/`, `xlsx/`, `pdf/`,
`doc-coauthoring/`, `internal-comms/`, `brand-guidelines/`, `skill-creator/`,
`mcp-builder/`, `sentry-code-review/`, `webapp-testing/`, `firecrawl/`),
`CLAUDE.md [67-84행]`)

## Slide 6 — 보조 ⑤ 결정론 파이프라인 · PPT 가이드 · 보고 카탈로그

① 결정론 PPT 파이프라인(extract→Claude 추론→fill, 원본 불변) ②
LGES PPT 작업 가이드(메인 1~2+보조 N, 흰/검/회+파랑 강조, 초록 풋노트,
한·영 폰트 분리, SZ 6단) ③ 임원 보고 양식 카탈로그 8블록(A~H).
**브랜드 자동화**: `brand-guidelines` 스킬이 색·폰트·풋노트 위반을 자가
점검 · `sentry-code-review` 가 신규 `scripts/build_*.py` 빌더의 운영
안전성 게이트 역할.

(출처: `templates/사내양식/LGES_PPT_작업_가이드.md`,
`templates/사내양식/보고양식_카탈로그.md`, `scripts/ppt_extract.py`,
`scripts/ppt_fill.py`, `.claude/skills/brand-guidelines/SKILL.md`,
`.claude/skills/sentry-code-review/SKILL.md`)

## Slide 7 — 보조 ⑥ 검증된 산출물 (2026-05)

주간보고 · 회의록 · 유연성지표 요약 · 경쟁사 동향 통합보고서 ·
FA자동화 임원보고 예시 · 전고체 밀폐이송장치 검토 · FA 기술과제후보 2부 docx.
모두 `(출처:…)` 자동 검증 통과.

(출처: `outputs/주간보고_2026-05-12.md`,
`outputs/회의록_신규R&D과제착수_2026-05-12.md`,
`outputs/유연성지표_요약_2026-05-14.pptx`,
`outputs/FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.md`,
`outputs/FA자동화_임원보고_예시_2026-05-18.pptx`,
`outputs/전고체라인_밀폐이송장치_검토_2026-05-18.pptx`,
`outputs/fa-task-discovery/FA_기술과제후보_20260518.docx`)

## Slide 8 — 보조 ⑦ 실제 산출 예시 발췌 (강조)

본 5인 팀이 실제 생성한 `outputs/` 원문 발췌 4종:
- **주간보고** (data+ops+document-writer): KPI 실측(물류 원가절감 21.3%,
  목표 16.2% → 달성률 131.5%), Capex 3,791.4억 → 2,985.3억 △806.1억(21.3%)
  절감, 모든 수치 `(출처: …[N행])` 자동 부착
- **신규 R&D 회의록** (meeting/document-writer): 결정 1순위 Isaac SIM 2026
  H2 PoC / 2순위 MMF 2027 초, 선행예산 2~3억, 액션아이템 6건 표
- **FA 기술과제후보** (fa-task-discovery v2 + tech-research): 마커리스
  비주얼 SLAM AMR, Omniverse Sim2Real, 휴머노이드 PoC 등 6분야 8건
- **경쟁사 동향 통합보고서** (tech-research×4): CATL +17%/-99%, BYD ~97%,
  삼성SDI 디지털트윈, SK온×NVIDIA 5만+ GPU → 26FA KPI 재설정 근거

(출처: `outputs/주간보고_2026-05-12.md`,
`outputs/회의록_신규R&D과제착수_2026-05-12.md`,
`outputs/fa-task-discovery/tasks_20260518.json`,
`outputs/FA_AI_AX_PhysicalAI_경쟁사_동향보고서_2026-05-14.md`)

## Slide 9 — 보조 ⑧ FA기술혁신파트 업무 → Agent 매핑 (★ 강조)

사용자가 정의한 **FA기술혁신파트 현재 업무 9개 카테고리**와 본 5인 팀이
어떻게 돕고/대체할 수 있는지 매핑:

| # | 업무 | 담당 에이전트·스킬 | 자동화 수준 |
|---|---|---|---|
| 1 | 중점추진업무(KPI·M&R·로드맵) | data + writer + **xlsx** | ★★ 자동 |
| 2 | Project(투자비·면적·구매) | data + ops + writer + **xlsx** | ★★ 자동 |
| 3 | MRM(과제 도출·진척·로드맵) | tech-research + fa-task-discovery + ppt-writer | ★★ 자동(매일) |
| 4 | 협력사(Pool·신규·설비군) | ops + **pdf + docx** | ★ 부분(승인 사람) |
| 5 | 미래기술 DST(SMC/MMF/셋업/밀폐이송) | tech-research + data + ppt-writer + **pdf** | ★★ 보고 자동(PoC 사람) |
| 6 | 경비·투자(Project/경상/일반) | data + **xlsx** | ★ 부분(결재 사람) |
| 7 | 보고(CPO·센터장·생산성·사업부) | writer 2종 + meeting-minutes + **brand-guidelines + internal-comms** | ★★★ **완전 자동** |
| 8 | 인사(포상·채용·진급·교육·출장) | ops + **internal-comms** | △ 보조(인사평가 사람) |
| 9 | 기타(실사·타운홀·워크샵) | writer + meeting-merger + **docx** | △ 보조(행사 사람) |

- ★★★ 완전 자동(1): **보고** — 가장 큰 시간 절감
- ★★ 자동(4): 중점추진·Project·MRM·미래기술 — 사람 검토 후 적용
- ★ 부분(2): 협력사·경비 — 데이터 자동, 승인 사람
- △ 보조(2): 인사·기타 — 회의록/요약은 자동, 결정·진행은 사람
- → 스킬 14종 자동 발동으로 각 업무에 양식 맞춤 적용

(출처: 사용자 인풋 업무 리스트(2026-05-19),
`outputs/주간보고_2026-05-12.md`,
`outputs/회의록_신규R&D과제착수_2026-05-12.md`,
`outputs/fa-task-discovery/`,
`references/roadmap/2026_FA기술담당_중장기로드맵_v2.md`)

## Slide 10 — 보조 ⑨ 진화 타임라인 · ROI

2026-05-10 초기 셋업 → 05-19 deck 1차 → **05-20 스킬 라이브러리 11종 확장
+ MCP 16+ 연결** 까지 진화 타임라인 + 도입 ROI 6항목 (작성 시간·출처
신뢰도·자료 누적·재현성·표준화·그룹 확장성).

(출처: `git log`, `CLAUDE.md`, `README.md`,
`outputs/주간보고_2026-05-12_PR요약.md`,
`.claude/skills/*/SKILL.md` (14종))

## Slide 11 — 보조 ⑩ 에이전트 발전 계획 (개발 + 데이터, 구체화)

3 phase × 2 column 매트릭스 — **무엇을 개발할지 / 무엇을 데이터로 채울지**:

**단기 Q3'26 (3개월)**
- 개발: 사내 OneNote·SharePoint MCP PoC · Outlook/Teams 알림 hook ·
  sentry-code-review 자동 게이트 · 임원 보고 카탈로그 8블록 완성
- 데이터: 협력사 Pool 50+ 보충 · KPI 베이스라인 12개월치 · DST 진척
  (SMC/MMF/셋업/밀폐이송) · 과거 회의록 1년치 적재

**중기 Q4'26 (3-6개월)**
- 개발: 사내 ERP/MES/PLM read-only MCP (mcp-builder) · SSO·사내 인증
  통과 · firecrawl 활성화 (경쟁사 IR 수집) · webapp-testing 사내 대시보드
- 데이터: FA 6분야 외부 동향 1년치 · 경쟁사 4사 IR/특허 12개월치 · 임원
  보고 13종 양식 → references 매핑 · 인프라 그룹 4부서 KPI 양식 통합

**장기 2027~ (6개월+)**
- 개발: 자기평가 Hook (출력 회귀 검증) · 공유 Task List 멀티 영역 동시
  운영 · skill-creator 로 사내 도메인 스킬 50+ 표준화 · 사내 표준 보고
  플랫폼화 + 거버넌스 연계
- 데이터: 사내 SOP 50+개 · 협력사 평가 분기별 · 그룹 전영역 KPI 데이터
  마트 · 사내 용어집·정책 완성 · 산출물 카탈로그 100+종

**불변 원칙**: 추측 금지 · 출처 자동 검증(Hook 강제) · 원본 불변 · 같은
입력=같은 출력 · 사람 승인 작업은 '제안'만.

(출처: `CLAUDE.md`, `README.md`,
`references/roadmap/2026_FA기술담당_중장기로드맵_v2.md`,
`.claude/agents/*.md`, `.claude/skills/mcp-builder/SKILL.md`,
`.claude/skills/skill-creator/SKILL.md`, `git log`)

## Slide 12 — 보조 ⑪ 인프라 그룹 확장 시나리오 · 핵심 메시지

**인프라 그룹 = FA기술담당 + 에너지기술담당 + 건설담당 + 인프라 기획팀**.
사람 워크플로우 동일 패턴 → 즉시 이식 (사용자 정의 조직 기준):

- **FA기술담당 (본 deck 실증)** — 현재 가동, 검증 산출물 10+종
- **에너지기술담당** — ESS·재생에너지·전력관리 KPI 자동 보고, 외부 동향
  (IEA·BNEF), R&D 회의록  (스킬: xlsx · tech-research · meeting-minutes)
- **건설담당** — 도면·시방서 PDF OCR, 공정 진척률 트래킹, 협력사 평가서
  (스킬: pdf · xlsx · docx)
- **인프라 기획팀** — 그룹 전략·정책 장문, 임원 결재 PPT, 사내 공지·3P
  (스킬: doc-coauthoring · ppt-writer · internal-comms)

이식 비용이 낮은 이유 + 핵심 메시지: 서버·DB·코드 없이 Claude Code +
Agent Teams + 마크다운 양식만으로 보고·리서치·기술과제 발굴 자동화 +
**스킬 14종 + MCP 16+ 로 양식·도구 확장 비용 최소화**.

(출처: `CLAUDE.md`, `README.md`, `.claude/agents/*.md`,
`.claude/skills/*/SKILL.md` (14종),
`references/roadmap/`, `templates/사내양식/보고양식_카탈로그.md`)

## Slide 13 — 가이드 ①/2 사용 가이드 — 환경 셋업·첫 명령 (초보자용)

초보자가 10분 안에 첫 보고서를 자동 생성할 수 있도록 4 단계로 안내:

- **STEP 1 환경 준비**: Claude Code 설치(CLI/Desktop/Web/IDE 택일) ·
  환경 변수 `CLAUDE_CODE_EXPERIMENTAL_AGENT_TEAMS=1` 설정 ·
  `git clone .../claude_FA-Agent_1 && cd claude_FA-Agent_1` · `claude` 명령
- **STEP 2 폴더 구조 이해**: `templates/`(양식, 수정 X) · `sample_data/`(정량) ·
  `references/`(자동 인덱싱) · `outputs/`(결과만) · `.claude/agents/` · `.claude/skills/`
- **STEP 3 첫 명령 (한 줄로 자동 위임)**:
  - `"이번 주 주간보고 작성해줘"` → data + ops + document-writer 자동
  - `"녹취록으로 회의록 만들어줘"` → meeting-minutes 스킬 발동
  - `"FA 기술 동향 조사해줘"` → fa-task-discovery v2 (6분야)
- **STEP 4 결과 확인 + 자동 검증**: outputs/ 자동 생성 · TaskCompleted Hook
  출처 미기입 차단 · 자료 미확인 자동 명시 · brand-guidelines 자동 적용

(출처: `CLAUDE.md`, `README.md`, `.claude/settings.json`,
`scripts/hooks/check_output_citations.py`)

## Slide 14 — 가이드 ②/2 사용 가이드 — 시나리오별 명령 예시·FAQ (초보자용)

**6 시나리오** (각각 `명령 예시 → 발동 → 산출 → TIP`):

1. **주간보고** — `"이번 주 주간보고 작성해줘"` → data + ops + writer + xlsx →
   `outputs/주간보고_YYYY-MM-DD.md`. TIP: 매주 같은 명령 반복 (재현성)
2. **회의록** — `"녹취록 정리해서 회의록 만들어줘"` → meeting-minutes + writer →
   `outputs/회의록_*.pptx`. TIP: 텍스트 붙여넣기 OK, 인터랙티브 확정
3. **FA 과제 발굴 (매일)** — `"FA 기술 동향 조사"` → fa-task-discovery v2 +
   tech-research → `outputs/fa-task-discovery/*.docx + 카톡`. TIP: 매일 자동 + PR
4. **경쟁사 동향 분석** — `"경쟁사(CATL·BYD·삼성SDI·SK온) 동향 분석"` →
   tech-research × 4 병렬 + writer 2종 → `outputs/경쟁사_*`. TIP: 4팀원 동시
5. **임원 보고 PPT** — `"임원 보고 PPT 만들어줘 (블록 A+C+G+B+E+H)"` →
   ppt-writer + brand-guidelines → `outputs/*_임원보고_*.pptx`. TIP: LGES 가이드 자동
6. **협력사 사양 추출** — `"이 PDF 사양서에서 표·이미지 뽑아줘"` →
   ops + pdf 스킬 + docx → 구조화 데이터 + 평가서. TIP: 스캔 PDF OCR 자동

**FAQ 5문항**:
- Q1 자료 없는 항목 → "자료 미확인" 자동 명시 (추측 금지)
- Q2 신규 양식 추가 → `templates/` 폴더에 두면 자동 인식
- Q3 출처 검증 우회 → 불가 (Hook 강제). "자료 미확인" 명시로만 가능
- Q4 에이전트 무응답 → `/resume` 후 "팀원 다시 생성해줘" (in-process 미복원 제약)
- Q5 한·영 폰트 분리 → ppt-writer 가 LG스마트체/Arial Narrow 자동 run 분리

(출처: `CLAUDE.md`, `.claude/agents/*.md`, `.claude/skills/*/SKILL.md` (14종),
`scripts/hooks/check_output_citations.py`)

---

## 부록 — 보조 자료

- `outputs/소개자료_2026-05-12.pptx` / `.md` — 17 슬라이드 본격 발표자료
- `outputs/AI경진대회_FA기술혁신Agent팀_2026-05-13.pptx` / `.md` — 직전판(5장)

본 갱신판은 직전판 대비 ① LGES PPT 작업 가이드 ② 한·영 폰트 자동분리
③ 임원 보고 양식 카탈로그 ④ 유연성지표·경쟁사 동향·전고체 검토 산출물
⑤ fa-task-discovery v2 ⑥ Agent Teams 차이점·FA 9업무 매핑 슬라이드
⑦ **메인 1장 정량/정성 효과 분리 + 보조 11장 핵심 요약 띠**
⑧ **신규 스킬 11종 추가 + MCP 16+ 자동 연결 반영**
⑨ **메인장 SmartArt 흐름도 + 스킬·MCP 라이브러리 카드 신설**
⑩ **슬라이드 2 양식을 3–12 보조 양식으로 통일** 반영.

**색 규칙 통일 (사내 룰)**: 본문=검정 / 강조=파랑 #0000FF (B255) /
풋노트·참고=진초록 #006600 / 박스·표=흰색·회색만. 12장 전체에 일괄 적용.

(출처: `outputs/소개자료_2026-05-12.md`,
`outputs/AI경진대회_FA기술혁신Agent팀_2026-05-13.md`, `git log`)
"""


def main():
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(TEMPLATE_PATH)
    print(f"slide_size: {prs.slide_width/360000:.2f} x "
          f"{prs.slide_height/360000:.2f} cm")
    print(f"입력 슬라이드 수: {len(prs.slides)}")

    fill_form(prs)
    fill_form_aux(prs)
    aux_agent_teams(prs)
    aux_members(prs)
    aux_skills(prs)
    aux_pipeline(prs)
    aux_outputs(prs)
    aux_examples(prs)
    aux_fa_mapping(prs)
    aux_evolution(prs)
    aux_roadmap(prs)
    aux_next(prs)
    aux_guide_setup(prs)
    aux_guide_scenarios(prs)

    print(f"최종 슬라이드 수: {len(prs.slides)}")
    prs.save(OUT_PPTX)
    OUT_MD.write_text(MD, encoding="utf-8")
    print(f"saved: {OUT_PPTX}")
    print(f"saved: {OUT_MD}")


if __name__ == "__main__":
    main()
