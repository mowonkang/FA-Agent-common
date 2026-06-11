"""검증용 샘플 PPT 양식 생성기.

`templates/주간보고.pptx` 를 생성한다. 이미 존재하면 덮어쓰지 않는다 (--force 시 덮어씀).
실제 회사 양식이 아닌 검증·시연용이다.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    blank = prs.slide_layouts[6]

    # Slide 1 — 표지
    s1 = prs.slides.add_slide(blank)
    title = s1.shapes.add_textbox(Cm(2), Cm(3), Cm(21), Cm(3))
    title.name = "제목"
    title.text_frame.text = "{{ report_title }}"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    sub = s1.shapes.add_textbox(Cm(2), Cm(7), Cm(21), Cm(2))
    sub.name = "작성자/주차"
    sub.text_frame.text = "{{ author }} / {{ week_range }}"
    sub.text_frame.paragraphs[0].runs[0].font.size = Pt(20)

    # Slide 2 — KPI 표 (헤더만 있고 데이터행 비어있음)
    s2 = prs.slides.add_slide(blank)
    s2_title = s2.shapes.add_textbox(Cm(1), Cm(0.5), Cm(20), Cm(1.5))
    s2_title.name = "섹션 제목"
    s2_title.text_frame.text = "1. KPI 진척"
    s2_title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    s2_title.text_frame.paragraphs[0].runs[0].font.bold = True

    rows, cols = 4, 4
    table_shape = s2.shapes.add_table(rows, cols, Cm(1), Cm(2.5), Cm(23), Cm(8))
    table = table_shape.table
    headers = ["지표", "목표", "실적", "달성률"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    indicators = ["양산 가동률", "Capex 집행률", "신규 협력사 등록"]
    for r, name in enumerate(indicators, start=1):
        table.cell(r, 0).text = name
        # 데이터 셀(목표/실적/달성률)은 의도적으로 비워둠

    # Slide 3 — 협력사 라벨 + 빈 텍스트박스
    s3 = prs.slides.add_slide(blank)
    s3_title = s3.shapes.add_textbox(Cm(1), Cm(0.5), Cm(20), Cm(1.5))
    s3_title.name = "섹션 제목"
    s3_title.text_frame.text = "2. 협력사 신규/변경"
    s3_title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    s3_title.text_frame.paragraphs[0].runs[0].font.bold = True

    label = s3.shapes.add_textbox(Cm(1), Cm(2.5), Cm(6), Cm(1))
    label.name = "협력사 라벨"
    label.text_frame.text = "협력사 업데이트:"
    label.text_frame.paragraphs[0].runs[0].font.size = Pt(16)

    body = s3.shapes.add_textbox(Cm(1), Cm(3.8), Cm(23), Cm(8))
    body.name = "협력사 본문"
    # 빈 상태로 둠 (Claude 가 "라벨 인접 빈 도형" 으로 추론해야 함)
    body.text_frame.text = ""

    # Slide 4 — 명시적 placeholder 가 본문 안에 섞여 있는 케이스
    s4 = prs.slides.add_slide(blank)
    s4_title = s4.shapes.add_textbox(Cm(1), Cm(0.5), Cm(20), Cm(1.5))
    s4_title.name = "섹션 제목"
    s4_title.text_frame.text = "3. CPO 메시지"
    s4_title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    s4_title.text_frame.paragraphs[0].runs[0].font.bold = True

    msg = s4.shapes.add_textbox(Cm(1), Cm(2.5), Cm(23), Cm(10))
    msg.name = "CPO 메시지 본문"
    msg.text_frame.text = "{{ cpo_message }}"
    msg.text_frame.paragraphs[0].runs[0].font.size = Pt(16)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build sample PPT template.")
    parser.add_argument(
        "--out",
        default="templates/주간보고.pptx",
        help="출력 PPTX 경로 (기본: templates/주간보고.pptx)",
    )
    parser.add_argument("--force", action="store_true", help="기존 파일 덮어쓰기")
    args = parser.parse_args(argv)
    out = Path(args.out)
    if out.exists() and not args.force:
        sys.stderr.write(
            f"이미 존재합니다: {out}. 덮어쓰려면 --force 를 추가하세요.\n"
        )
        return 1
    build(out)
    print(str(out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
