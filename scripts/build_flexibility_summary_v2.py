"""유연성(Flexibility) 지표 v2 — Samsung SDI · CATL 비교 분석 CEO 보고용 결정판.

본 빌더는 `outputs/유연성_경쟁사_분석_v2_2026-05-21.md` (449줄)의 핵심을
LGES 사내 양식(`templates/사내양식/LGES_PPT_작업_가이드.md`)에 맞춰 임원
즉시 활용 가능한 메인 1장 + 보조 2장 PPT 로 정리한다.

기준 구현: `scripts/build_flexibility_summary.py` 의 헬퍼·색·사이즈 상수를
그대로 계승하며, 데이터만 v2(12셀 매트릭스 · 격차 원인 · 임계 지표)로 교체.

출력: outputs/유연성_경쟁사_요약_2026-05-21.pptx

슬라이드 구성:
  Slide 1 (메인, CEO 즉시 활용) — 12셀 매트릭스 + 3대 메시지 + CEO 액션
  Slide 2 (보조 1/2)             — 격차 원인 분석 (자동화 철학·재구성 비용)
  Slide 3 (보조 2/2)             — 임계 지표 3대 경계 신호 + 권고
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 팔레트 (회색·흰색 기조 / 파랑 강조 / 풋노트 진초록) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)

BLUE = RGBColor(0x00, 0x00, 0xFF)        # 강조 (R0 G0 B255 순수 파랑)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)  # 옅은 파랑 (핵심 메시지 박스)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)  # 풋노트 (G102)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

# ─── 폰트 사이즈 스케일 (6단 상수) ───────────────────
SZ_TITLE = 16
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

OUT_PATH = Path("outputs/유연성_경쟁사_요약_2026-05-21.pptx")


# ───────────────────────────── 헬퍼 ──────────────────────────────
def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    """한·영 폰트 자동 분리 run 설정."""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH

    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                color=BLACK, align=None, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        run = para.add_run()
        set_run(
            run,
            text,
            font=opts.get("font", font),
            size=opts.get("size", size),
            bold=opts.get("bold", bold),
            color=opts.get("color", color),
        )


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.1)
    tb.text_frame.margin_right = Cm(0.1)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, fill=DIM_GRAY, color=WHITE):
    shp = add_rect(slide, left, top, width, 0.7, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=SZ_BAND, bold=True, color=color)
    return shp


def add_footnote(slide, text):
    """풋노트(초록곰팡이 8pt Bold) 별도 텍스트박스 — 슬라이드 하단 y=18.30."""
    add_text(
        slide, 0.6, 18.30, 26.3, 0.55,
        [(text, {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )


def add_title_band(slide, title, sub_label):
    """슬라이드 공통 타이틀 + 우상단 라벨 띠."""
    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [(title, {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.4, 0.35, 6.8, 0.8,
        [(sub_label, {"size": SZ_BODY, "bold": True, "color": BLUE,
                      "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)


# ───────────────────────── Slide 1 (메인) ─────────────────────────
def build_slide_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 타이틀 + 라벨 띠
    add_title_band(
        slide,
        "Samsung SDI · CATL 유연성 지표 비교 분석 (CEO 보고용 결정판)",
        "FA기술담당 · 2026-05-21 · v2 하이브리드",
    )

    # ── 핵심 메시지 박스 (상단) ──
    add_rect(slide, 0.6, 1.55, 26.3, 1.75,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.60, 26.0, 1.65,
        [
            ("[CEO 핵심 결론]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • LGES Flexibility 정의(유연 설비 투자비 / 전체, 셋업비 ≤30% 컷오프) 환산 시 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("'28년 76% vs SDI 57% vs CATL 42% — 글로벌 선도 격차 확보 전망",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (' • CATL "자동화율 95%" 는 IoT 인터커넥션 비율 — Flexibility(유연 설비 비중) 와는 분자·분모가 다른 개념',
             {"size": SZ_BODY, "color": BLACK}),
            (" • SMC·MMF·AMR 확산이 LGES 유연성 견인축 — 신규 폼팩터 도입·다화학 운영 시 라인 재사용 1.5~2배 우위",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── 좌측: 12셀 매트릭스 표 ──
    LEFT = 0.6
    LEFT_W = 13.0
    add_label_band(slide, LEFT, 3.45, LEFT_W,
                   "① 3사 × 4개년 유연성 지표 매트릭스 (mid-point, %)")
    add_rect(slide, LEFT, 4.15, LEFT_W, 6.30,
             fill=WHITE, line=MID_GRAY, line_w=0.75)

    # 표 영역
    tbl_left = LEFT + 0.3
    tbl_top = 4.40
    tbl_w = LEFT_W - 0.6
    tbl_h = 4.30
    tbl_shape = slide.shapes.add_table(
        5, 5, Cm(tbl_left), Cm(tbl_top), Cm(tbl_w), Cm(tbl_h)
    )
    tbl = tbl_shape.table
    col_widths = [4.6, 1.85, 1.85, 1.85, 2.25]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Cm(w)

    headers = ["구분", "'25", "'26", "'27", "'28"]
    rows = [
        ("LGES MI Lansing",        "47%", "61%", "75%", "76%"),
        ("LGES ESGM2 각형",         "50%", "62%", "74%", "78%"),
        ("Samsung SDI (mid)",      "35%", "41%", "49%", "57%"),
        ("CATL (mid)",             "25%", "29%", "35%", "42%"),
    ]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        is_lges = row[0].startswith("LGES")
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            # LGES 행 옅은 강조 배경
            cell.fill.fore_color.rgb = SOFT_GRAY if is_lges else WHITE
            bold = (ci == 0) or (ci == 4 and is_lges)
            # '28 열 LGES 파랑 강조 (한 슬라이드 강조 포인트 中 1)
            if ci == 4 and is_lges:
                color = BLUE
            else:
                color = BLACK
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            size = SZ_SECTION if ci == 4 else SZ_BODY
            write_lines(cell.text_frame, [val], size=size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    # 표 하단 요약 텍스트
    add_text(
        slide, LEFT + 0.35, 8.85, LEFT_W - 0.7, 1.50,
        [
            ("[매트릭스 해석]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • LGES MI Lansing 47% → 76% ",
             {"size": SZ_BODY, "color": BLACK}),
            ("(+29%p, '25→'28)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • LGES vs SDI '28 +19%p, LGES vs CATL '28 +34%p",
             {"size": SZ_BODY, "color": BLACK}),
            (" • LGES 추정치 = F1 정의 ground truth, SDI/CATL = 4-Proxy 가중평균 (±5~10%p)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── 우측: 3대 핵심 메시지 ──
    RIGHT = 14.0
    RIGHT_W = 12.9
    add_label_band(slide, RIGHT, 3.45, RIGHT_W,
                   "② 3대 핵심 메시지 (CEO 의사결정 포인트)")
    add_rect(slide, RIGHT, 4.15, RIGHT_W, 6.30,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, RIGHT + 0.3, 4.30, RIGHT_W - 0.6, 6.10,
        [
            ("M1. LGES 글로벌 선도 격차 확보",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • MI Lansing 47% → 76% ",
             {"size": SZ_BODY, "color": BLACK}),
            ("(+29%p)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("    — 글로벌 각형 라인 중 명시적 유연성 지표 공개 유일",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 4}),
            ("M2. LGES vs CATL '28 +34%p",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 76% vs 42% — CATL 표준화·대형화·고정형 Stacker Crane 의존",
             {"size": SZ_BODY, "color": BLACK}),
            (' • "자동화율 95%" = IoT 인터커넥션 비율 (≠ 유연 설비 비중) ',
             {"size": SZ_BODY, "color": BLACK}),
            ("— 고정 설비 95% 인터커넥션은 라인 재구성과 무관",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 4}),
            ("M3. LGES vs SDI '28 +19%p",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 76% vs 57% — SDI 는 한국 EPC + 동일 공급사(SFA·코윈테크) 활용",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 헝가리 괴드 P5→P6 전환·천안 46파이 신규 라인이 격차 좁힘 변수",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — 코윈테크 AMR 수주 +465.2% YoY ('25, 헬로티)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── 하단: CEO 액션 박스 (3개, 단/중/장기) ──
    AC_TOP = 10.60
    AC_W = 8.65
    AC_H = 7.10
    # 단기
    add_label_band(slide, 0.6, AC_TOP, AC_W, "③ 단기 액션 (0~6개월)")
    add_rect(slide, 0.6, AC_TOP + 0.7, AC_W, AC_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, AC_TOP + 0.85, AC_W - 0.5, AC_H - 0.95,
        [
            ("[IR 메시지 · KPI 신설]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 유연성 지표를 차기 IR Day 신규 KPI 로 공식 도입",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — CATL 95% 對 LGES 76% 차별화 포지셔닝",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 유럽 OEM 다품종 대응 영업 메시지 강화",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — BMW·Mercedes·Stellantis 신규 셀 사양 협상",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • Tesla LFP Lansing $4.3억 후속 OEM 협상 시",
             {"size": SZ_BODY, "color": BLACK}),
            ("   SMC 컨셉 적극 활용",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )
    # 중기
    add_label_band(slide, 9.50, AC_TOP, AC_W, "④ 중기 액션 (6~18개월)")
    add_rect(slide, 9.50, AC_TOP + 0.7, AC_W, AC_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 9.75, AC_TOP + 0.85, AC_W - 0.5, AC_H - 0.95,
        [
            ("[KPI 가중치 · IP 강화]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • '26 KPI 8개 항목 모두 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("유연성 가중치 0% → 10~15% 신설 검토",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • SMC / MMF 기술 IP 강화",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — 코윈테크 AMR (+465.2% 수주) 발주 확대",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("   — 시너스텍 · SFA · 티로보틱스 파트너십 심화",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • CATL 'Xiaomo' 셀 제조 라인 확대 모니터링",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — 차세대 SMC R&D 가속화로 우위 사수",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )
    # 장기
    add_label_band(slide, 18.40, AC_TOP, AC_W, "⑤ 장기 액션 (18개월~)")
    add_rect(slide, 18.40, AC_TOP + 0.7, AC_W, AC_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 18.65, AC_TOP + 0.85, AC_W - 0.5, AC_H - 0.95,
        [
            ("[글로벌 표준화 · 정밀도]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • WEF Lighthouse flexibility index 표준화 대응",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — MI Lansing / ESGM2 신규 Lighthouse 추진",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • BloombergNEF · SNE Research 에 정의 제안",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — LGES 선도 지위 확보 (Flexibility 표준화)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • SDI / CATL 내부 분류 추가 정보 수집",
             {"size": SZ_BODY, "color": BLACK}),
            ("   — 설비 협력사 인터뷰 · LEAD H주 공시 분석",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("   — 추정 신뢰구간 ±5%p 이하 달성",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # 풋노트
    add_footnote(
        slide,
        "* 출처 : references/FA 유연성 지표 정의 (F1, 2026-05-21)  │  "
        "outputs/유연성_경쟁사_분석_v2_2026-05-21.md  │  "
        "references/경쟁사/* (D1~D5 외부 IR 아카이브 5건)",
    )


# ──────────────────────── Slide 2 (보조 1) ────────────────────────
def build_slide_aux1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_band(
        slide,
        "[보조 1/2] 격차 원인 분석 — 자동화 철학 · 재구성 비용",
        "FA기술담당 · 2026-05-21 · v2 하이브리드",
    )

    # ── 상단: 3사 자동화 철학 매트릭스 ──
    add_label_band(slide, 0.6, 1.55, 26.3,
                   "① 3사 자동화 철학 매트릭스 (8 차원 비교)")
    add_rect(slide, 0.6, 2.25, 26.3, 7.00,
             fill=WHITE, line=MID_GRAY, line_w=0.75)

    tbl_left = 0.9
    tbl_top = 2.50
    tbl_w = 25.7
    tbl_h = 6.50
    tbl_shape = slide.shapes.add_table(
        9, 4, Cm(tbl_left), Cm(tbl_top), Cm(tbl_w), Cm(tbl_h)
    )
    tbl = tbl_shape.table
    col_widths = [5.7, 6.7, 6.7, 6.6]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Cm(w)

    headers = ["차원", "LGES", "Samsung SDI", "CATL"]
    rows = [
        ("자동화 철학",
         "모듈·유연성 중심 (SMC/MMF/AMR)",
         "일본식 정밀 + 미국식 모듈 (P6)",
         "표준화·대형화·무인화 (PSL 8세대)"),
        ("자동화율 ('28)", "85~90%", "85%", "95%"),
        ("유연성 지표 ('28)", "76%", "57% (mid)", "42% (mid)"),
        ("다품종 대응",
         "高 (원통/각형/파우치, LFP/NMC)",
         "中 (각형 주력, 원통형 신규)",
         "中 (라인별 분리 전략)"),
        ("셀 사양 변경 비용",
         "低 (SMC 모듈 재배치)",
         "中 (라인 부분 개조)",
         "高 (라인 정지·재구성)"),
        ("셋업 변경 시간", "1.0x (기준)", "1.4x", "1.8x"),
        ("1초/셀 달성", "부분 (PPM 기준)", "미달성", "달성 (PSL 8세대)"),
        ("주요 협력사",
         "코윈테크·시너스텍·SFA·티로보틱스",
         "SFA·코윈테크·씨케이솔루션·한화모멘텀",
         "先导智能(LEAD)·Geek+·Hai Robotics"),
    ]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            bold = (ci == 0)
            color = CHARCOAL if ci == 0 else BLACK
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            # 유연성 지표 행만 LGES 셀 파랑 강조 (보조이므로 1포인트만)
            if ri == 3 and ci == 1:
                color = BLUE
                bold = True
            size = SZ_SUB if (ci != 0 and ri in (1, 4, 5, 8)) else SZ_BODY
            write_lines(cell.text_frame, [val], size=size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    # ── 하단 좌: 라인 재구성 비용 비교 표 ──
    BOT_TOP = 9.55
    LEFT_W = 13.0
    add_label_band(slide, 0.6, BOT_TOP, LEFT_W,
                   "② 라인 재구성 비용 비교 (CAPEX 대비 %, 정성 추정)")
    add_rect(slide, 0.6, BOT_TOP + 0.7, LEFT_W, 8.00,
             fill=WHITE, line=MID_GRAY, line_w=0.75)

    rb_left = 0.9
    rb_top = BOT_TOP + 0.95
    rb_w = LEFT_W - 0.6
    rb_h = 4.30
    rb_shape = slide.shapes.add_table(
        4, 4, Cm(rb_left), Cm(rb_top), Cm(rb_w), Cm(rb_h)
    )
    rb = rb_shape.table
    col_w2 = [4.6, 2.4, 2.4, 3.0]
    for i, w in enumerate(col_w2):
        rb.columns[i].width = Cm(w)

    rb_headers = ["시나리오", "LGES", "SDI", "CATL"]
    rb_rows = [
        ("동일 폼팩터 화학 변경 (NMC→LFP)", "5~10%", "12~18%", "20~30%"),
        ("폼팩터 부분 변경 (셀 크기)",      "8~15%", "18~25%", "25~35%"),
        ("신규 폼팩터 도입 (46파이 등)",    "15~25%", "25~35%", "35~50%"),
    ]
    for j, h in enumerate(rb_headers):
        cell = rb.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rb_rows, start=1):
        for ci, val in enumerate(row):
            cell = rb.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            bold = (ci == 0) or (ci == 1)  # LGES 열 Bold
            color = BLACK
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            write_lines(cell.text_frame, [val], size=SZ_BODY, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    add_text(
        slide, 0.9, BOT_TOP + 5.45, LEFT_W - 0.6, 3.10,
        [
            ("[재구성 비용 해석]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • LGES = SMC/MMF 재활용 가능 → 폐기 최소화",
             {"size": SZ_BODY, "color": BLACK}),
            (" • SDI = P5/P6 모듈 표준화, 라인 부분 개조 가능",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CATL = 라인별 분리 전략, 전용 공장 신축 일반",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CATL Liyang 사례: 3D printing fixtures 로 changeover 단축",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("    (catl.com / PRNewswire 2023-12-14)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── 하단 우: CATL 4대 격차 유연성 재해석 ──
    RIGHT = 14.0
    RIGHT_W = 12.9
    add_label_band(slide, RIGHT, BOT_TOP, RIGHT_W,
                   "③ CATL 4대 격차 vs 유연성 영향 재해석")
    add_rect(slide, RIGHT, BOT_TOP + 0.7, RIGHT_W, 8.00,
             fill=WHITE, line=MID_GRAY, line_w=0.75)

    cd_left = RIGHT + 0.3
    cd_top = BOT_TOP + 0.95
    cd_w = RIGHT_W - 0.6
    cd_h = 4.30
    cd_shape = slide.shapes.add_table(
        5, 3, Cm(cd_left), Cm(cd_top), Cm(cd_w), Cm(cd_h)
    )
    cd = cd_shape.table
    col_w3 = [4.0, 3.6, 4.7]
    for i, w in enumerate(col_w3):
        cd.columns[i].width = Cm(w)

    cd_headers = ["원인", "절대 CAPEX 영향", "유연성 지표 영향"]
    cd_rows = [
        ("1. 인건비",              "-15~20%",  "중립"),
        ("2. EPC·건축비",          "-30~40%",  "중립 (정의상 제외)"),
        ("3. 국산 설비 자가화",     "-15%",     "+영향 (LEAD 95% 표준)"),
        ("4. 표준화 (8세대 PSL)",   "-5~10%",   "+영향 (인력 -70%)"),
    ]
    for j, h in enumerate(cd_headers):
        cell = cd.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        write_lines(cell.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(cd_rows, start=1):
        for ci, val in enumerate(row):
            cell = cd.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            bold = (ci == 0)
            color = BLACK
            align = PP_ALIGN.LEFT if ci != 1 else PP_ALIGN.CENTER
            write_lines(cell.text_frame, [val], size=SZ_BODY, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    add_text(
        slide, RIGHT + 0.3, BOT_TOP + 5.45, RIGHT_W - 0.6, 3.10,
        [
            ("[재해석 시사점]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 절대 CAPEX 격차(인건비·EPC 합계 -45~60%pt) 중",
             {"size": SZ_BODY, "color": BLACK}),
            ("   절반 이상은 한국에서 통제 불가한 구조 격차",
             {"size": SZ_BODY, "color": BLACK}),
            (" • H/W 자체 단가 격차는 15~20% 수준",
             {"size": SZ_BODY, "color": BLACK}),
            (" • LEAD 모듈 95% / Pack 50% 표준이 CATL Flexibility",
             {"size": SZ_BODY, "color": BLACK}),
            ("   상승 견인 (출처: leadintelligent.com)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    add_footnote(
        slide,
        "* 출처 : outputs/유연성_경쟁사_분석_v2_2026-05-21.md §6 정성 분석  │  "
        "references/경쟁사/2026-05-21_CATL_PSL_8세대_SmartManufacturing.md  │  "
        "references/경쟁사/LGES_FA_Prismatic_Benchmarking_Integrated_v1-v5.md [1053~1087행]",
    )


# ──────────────────────── Slide 3 (보조 2) ────────────────────────
def build_slide_aux2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_band(
        slide,
        "[보조 2/2] 임계 지표 (경계 신호 3대 Trigger) + 권고 요약",
        "FA기술담당 · 2026-05-21 · v2 하이브리드",
    )

    # ── 상단 박스: 임계 지표 정의 ──
    add_rect(slide, 0.6, 1.55, 26.3, 1.30,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.60, 26.0, 1.20,
        [
            ("[임계 지표 (Benchmark / Threshold) 정의]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 본 슬라이드의 3대 경계 신호는 LGES 유연성 우위가 침식될 가능성을 사전 감지하기 위한 Trigger",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Trigger 발효 시 6개월 내 대응 액션 (KPI 재정의 · R&D 가속화 · Lighthouse 신청) 즉시 발동",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── 3개 경계 신호 카드 (좌·중·우) ──
    SIG_TOP = 3.10
    SIG_W = 8.65
    SIG_H = 9.70

    # 경계 신호 1
    add_label_band(slide, 0.6, SIG_TOP, SIG_W,
                   "① 경계 신호 1 — SDI 추격")
    add_rect(slide, 0.6, SIG_TOP + 0.7, SIG_W, SIG_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, SIG_TOP + 0.85, SIG_W - 0.5, SIG_H - 0.95,
        [
            ("[Trigger 조건]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • SDI 유연성 지표가 LGES 대비 5%p 이내로 좁혀짐",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 예: SDI '28년 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("71% 이상 도달",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 헝가리 괴드 P5→P6 전환 '27년 100% 달성",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 코윈테크 AMR 수주 비중 SDI향 > LGES향",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[모니터링 지표]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • SFA / 코윈테크 / 씨케이솔루션 분기 공시",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 천안 46파이 마더라인 가동 ramp-up",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 코코모 SPE 1·2공장 자동화 KPI",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[대응 액션]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • SMC IP 강화 + 코윈테크 발주 우선권 협상",
             {"size": SZ_BODY, "color": BLACK}),
            (" • LGES 측 ESGM2 각형 78% → 82%로 상향 검토",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # 경계 신호 2
    add_label_band(slide, 9.50, SIG_TOP, SIG_W,
                   "② 경계 신호 2 — CATL Xiaomo 셀 제조 확대")
    add_rect(slide, 9.50, SIG_TOP + 0.7, SIG_W, SIG_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 9.75, SIG_TOP + 0.85, SIG_W - 0.5, SIG_H - 0.95,
        [
            ("[Trigger 조건]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • CATL 휴머노이드 'Xiaomo' Pack → 셀 제조 라인 확대",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 자동화율 95% 유지 + ",
             {"size": SZ_BODY, "color": BLACK}),
            ("유연성 50% 돌파",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • '27년 내 공식 발표 (현재 뤄양 Pack 한정)",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[모니터링 지표]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • Spirit AI · Galbot · Figure AI 도입 사례",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CnEVPost · SCMP CATL 휴머노이드 보도",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 先导智能 智能物流 매출 비중 60% 돌파",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[대응 액션]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 차세대 SMC R&D 가속화",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 휴머노이드 생태계 협력 (티로보틱스 등)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • '28년 LGES 우위 침식 방지 액션 플랜 수립",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # 경계 신호 3
    add_label_band(slide, 18.40, SIG_TOP, SIG_W,
                   "③ 경계 신호 3 — 글로벌 표준화 (WEF Lighthouse)")
    add_rect(slide, 18.40, SIG_TOP + 0.7, SIG_W, SIG_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 18.65, SIG_TOP + 0.85, SIG_W - 0.5, SIG_H - 0.95,
        [
            ("[Trigger 조건]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • WEF Lighthouse Network 가 유연성 지표를",
             {"size": SZ_BODY, "color": BLACK}),
            (" • ",
             {"size": SZ_BODY, "color": BLACK}),
            ("'flexibility index' 공식 KPI 채택",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • Global Lighthouse Annual Report 등재 시점",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[모니터링 지표]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • WEF / McKinsey 공동 발표 (연 2회)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CATL Yibin · Liyang 등대공장 추가 발표",
             {"size": SZ_BODY, "color": BLACK}),
            (" • BloombergNEF · SNE Research 보고서",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[대응 액션]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 6개월 내 LGES Lighthouse 신청 준비",
             {"size": SZ_BODY, "color": BLACK}),
            (" • MI Lansing / ESGM2 각형 공장 신규 지정 추진",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 정의의 글로벌 표준화 선도 (LGES 자체 정의 채택)",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # ── 하단: 권고 요약 (단·중·장기 1줄씩) ──
    REC_TOP = 13.00
    add_label_band(slide, 0.6, REC_TOP, 26.3,
                   "④ 권고 요약 (단기·중기·장기 압축)")
    add_rect(slide, 0.6, REC_TOP + 0.7, 26.3, 4.45,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, REC_TOP + 0.85, 25.7, 4.25,
        [
            ("[단기 (0~6개월)] ",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            ("유연성 지표 IR 메시지 · 신규 KPI 도입 검토 · 유럽 OEM 다품종 대응 영업 강화",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[중기 (6~18개월)] ",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            ("'26 KPI 유연성 가중치 10~15% 신설 · SMC/MMF IP 강화 · 코윈테크 AMR (+465.2% 수주) 확대",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("[장기 (18개월~)] ",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            ("WEF Lighthouse flexibility index 표준화 대응 · 추정 정밀도 ±5%p 이하 달성",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            (" * 종합 시사점: LGES 의 SMC/MMF/AMR 기반 모듈·유연성 전략은 SDI/CATL 의 표준화·대형화 전략 대비 ",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("'28년 +19~34%p 우위",
             {"size": SZ_SUB, "bold": True, "color": MID_GRAY}),
            (" 확보 전망. 다만 SDI 추격 · CATL Xiaomo 확대 · WEF 표준화 등 3대 경계 신호를 분기 단위로 모니터링 필요.",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    add_footnote(
        slide,
        "* 출처 : outputs/유연성_경쟁사_분석_v2_2026-05-21.md §9 CEO 권고 · 임계 지표  │  "
        "references/FA 유연성 지표 정의 (F1)  │  references/26FA KPI.md  │  "
        "references/경쟁사/* (D1~D5)",
    )


# ─────────────────────────── 빌드 메인 ───────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)

    build_slide_main(prs)
    build_slide_aux1(prs)
    build_slide_aux2(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}  (slides: {len(prs.slides)})")


if __name__ == "__main__":
    build()
