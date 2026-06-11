"""분류형 레이아웃 PPT 공용 헬퍼 (사내 보고양식 슬라이드 1 기반).

사내 양식 ``templates/사내양식/(보고양식)YYMMDD_보고제목_v1.0.pptx`` 의
슬라이드 1 레이아웃을 재사용 가능한 형태로 제공한다.

레이아웃 (위 → 아래):
  - 제목(상단 좌측)
  - 배경 없는 검은 Head message
  - 좌측 세로 분류 레일(분류 1~N, 무채색 테두리, 높이 가변)
  - 우측: 세션 메시지 + 상세 내용(표/이미지 박스 선택)
  - 초록(#006600) 풋노트

핵심: 좌측 분류 박스의 높이는 대응하는 우측 내용 분량에 맞춰 자동 산정된다
(``render_classified_slide``). python-pptx 는 렌더 높이를 측정할 수 없으므로
줄 수·표·이미지 박스로부터 결정론적으로 추정한다.

드로잉 프리미티브(``set_run``/``write_lines``/``add_rect``/``add_text``/
``add_label_band``/``fill_table``/``title_block``/``add_footnote``)와 색·폰트·
사이즈 상수는 ``build_solidstate_transport_report.py`` 패턴을 이관·일반화한
것이다. 향후 빌더는 중복 정의 대신 본 모듈을 import 한다.

색·폰트·사이즈 규칙: ``templates/사내양식/LGES_PPT_작업_가이드.md``.
"""
from __future__ import annotations

import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 팔레트 (LGES_PPT_작업_가이드) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x00, 0xFF)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

SZ_TITLE = 16
SZ_HEAD = 15      # Head message (검정 Bold, 1~2줄)
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

# 슬라이드 규격 (cm) — 사내 양식과 동일 (10.83in × 7.5in)
SLIDE_W_CM = 27.52
SLIDE_H_CM = 19.05


# ─── 드로잉 프리미티브 ───
def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    """run 텍스트·폰트 설정 + 한/영 폰트를 run-level 로 분리."""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH
    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                color=BLACK, align=None, anchor=None):
    """text_frame 에 줄 목록을 쓴다. 항목은 str 또는 (str, opts) 튜플."""
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, opts = item if isinstance(item, tuple) else (item, {})
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        para.space_before = Pt(0)
        para.space_after = Pt(opts.get("space_after", 1))
        run = para.add_run()
        set_run(run, text,
                font=opts.get("font", font),
                size=opts.get("size", size),
                bold=opts.get("bold", bold),
                color=opts.get("color", color))


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    """사각형 도형. fill=None 이면 배경 투명(무채색)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.12)
    tb.text_frame.margin_right = Cm(0.12)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    tb.text_frame.word_wrap = True
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, height=0.7,
                   fill=DIM_GRAY, color=WHITE):
    shp = add_rect(slide, left, top, width, height, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.25)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=SZ_BAND, bold=True, color=color)
    return shp


def fill_table(tbl, headers, rows, *, col_w, body_size=SZ_SUB,
               head_size=SZ_BODY, emph_col0=True):
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Cm(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        cell.margin_left = Cm(0.1)
        cell.margin_right = Cm(0.1)
        cell.margin_top = Cm(0.03)
        cell.margin_bottom = Cm(0.03)
        write_lines(cell.text_frame, [h], size=head_size, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = Cm(0.12)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.03)
            cell.margin_bottom = Cm(0.03)
            bold = emph_col0 and ci == 0
            color = CHARCOAL if (emph_col0 and ci == 0) else BLACK
            write_lines(cell.text_frame, [val], size=body_size, bold=bold,
                        color=color, align=PP_ALIGN.LEFT,
                        anchor=MSO_ANCHOR.MIDDLE)


def title_block(slide, title, *, byline=None):
    """제목(좌측 상단) + 선택 byline(우측, 검정) + 구분선."""
    add_text(slide, 0.6, 0.25, 19.5, 1.0,
             [(title, {"size": SZ_TITLE, "bold": True, "color": BLACK})],
             anchor=MSO_ANCHOR.MIDDLE)
    if byline:
        lines = byline if isinstance(byline, list) else [byline]
        add_text(slide, 19.9, 0.30, 7.0, 0.95,
                 [(t, {"size": SZ_BODY, "bold": True, "color": BLACK,
                       "align": PP_ALIGN.RIGHT}) for t in lines],
                 anchor=MSO_ANCHOR.MIDDLE)
    ln = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    ln.line.color.rgb = MID_GRAY
    ln.line.width = Pt(0.75)


def add_footnote(slide, text, *, top=18.35):
    add_text(slide, 0.6, top, 26.3, 0.55,
             [(text, {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})])


# ─── 분류형 레이아웃 (자동 높이 맞춤) ───
RAIL_X = 0.60       # 좌측 분류 레일 x (cm) — 왼쪽으로 넓혀 라벨 칸 넘침 방지
RAIL_W = 2.30       # 분류 박스 폭 (라벨 자유 줄바꿈 수용)
RIGHT_X = 3.10      # 우측 내용 x
RIGHT_W = 23.80     # 우측 내용 폭 (→ 우측 여백 26.9)
GAP = 0.30          # 섹션 간 세로 간격
PAD_TOP = 0.20      # 섹션 내용 상단 패딩
PAD_BOT = 0.20      # 섹션 내용 하단 패딩
CM_PER_LINE = 0.52  # SZ_BODY 한 줄 높이 추정
CM_PER_HDR = 0.62   # 세션 메시지(헤더) 한 줄 높이
CPL = 60            # RIGHT_W 에서 한글 기준 줄당 글자 수 (영문 0.5 가중)


def _visual_len(text):
    """한글=1, 영문/숫자/공백=0.5 가중 길이."""
    total = 0.0
    for ch in text:
        total += 0.5 if ord(ch) < 0x1100 else 1.0
    return total


def _fit_label_size(label):
    """분류 라벨이 박스 폭(RAIL_W)을 넘지 않는 최대 폰트 사이즈를 고른다.

    가장 긴 단어(공백·줄바꿈 단위)가 박스 안쪽 폭에 들어가는 크기를 선택해
    칸 넘김(글자 잘림)을 방지한다.
    """
    inner_w = RAIL_W - 0.18
    longest = max((_visual_len(w) for w in label.replace("\n", " ").split()),
                  default=1.0)
    for size in (SZ_BAND, SZ_SECTION, SZ_BODY, SZ_SUB):
        if longest * (0.037 * size) <= inner_w:
            return size
    return SZ_SUB


def _line_for(item):
    text = item[0] if isinstance(item, tuple) else item
    return max(1, math.ceil(_visual_len(text) / CPL))


def _estimate_height(section):
    """섹션의 우측 내용 높이를 cm 로 추정한다."""
    hdr_h = CM_PER_HDR if section.get("session_message") else 0.0
    body_lines = sum(_line_for(it) for it in section.get("content", []))
    body_h = body_lines * CM_PER_LINE
    table_h = 0.0
    tbl = section.get("table")
    if tbl:
        rows = len(tbl["rows"]) + 1
        table_h = rows * tbl.get("row_h", 0.9) + 0.15
    img_h = float(section.get("image_box") or 0.0)
    return PAD_TOP + hdr_h + body_h + table_h + img_h + PAD_BOT


def render_classified_slide(slide, *, title, head_message, sections,
                            footnote=None, byline=None,
                            content_top=3.10, content_bottom=18.20,
                            fill_to_bottom=True):
    """사내 양식 슬라이드 1 레이아웃을 그린다.

    Parameters
    ----------
    slide : python-pptx slide (prs.slide_layouts[6] 권장)
    title : str — 상단 제목
    head_message : list — write_lines 형식. 배경 없는 검은 헤드 메시지.
    sections : list[dict] — 분류 1~N. 각 dict:
        label           : 좌측 박스 텍스트 (예 "분류 1.  개요")
        session_message  : 우측 헤더 한 줄 (생략 가능)
        content          : write_lines 항목 리스트 (생략 가능)
        table            : {"headers","rows","col_w","row_h"} 또는 None
        image_box        : 이미지 자리 박스 높이(cm) 또는 None
        height           : 명시 높이(cm) override 또는 None(자동 추정)
    footnote : str 또는 None — 초록 풋노트
    content_top / content_bottom : 분류 영역 상·하단 경계(cm)
    """
    title_block(slide, title, byline=byline)

    # Head message — 배경 없이 검은색 서술형 (1~2줄)
    add_text(slide, 0.6, 1.50, 26.3, 1.50,
             head_message, anchor=MSO_ANCHOR.TOP)

    avail_h = content_bottom - content_top

    # 1) 섹션별 높이 산정
    heights = []
    for sec in sections:
        h = sec.get("height")
        heights.append(float(h) if h is not None else _estimate_height(sec))

    n = len(sections)
    total = sum(heights) + GAP * (n - 1 if n else 0)

    # 2) 오버플로 시 스케일 다운 / 여백 시 하단까지 채움
    scale = 1.0
    if total > avail_h and total > 0:
        scale = avail_h / total
        heights = [h * scale for h in heights]
        if scale < 0.80:
            print(f"WARNING: 분류 내용이 슬라이드 높이를 초과 (scale={scale:.2f}) "
                  f"— 줄 수를 줄이거나 보조 슬라이드로 분리 권장")
    elif fill_to_bottom and total > 0 and total < avail_h:
        # 남는 세로 공간을 박스 높이에 비례 분배 (우측 내용은 상단 정렬 유지)
        slack = avail_h - total
        heights = [h + slack * (h / sum(heights)) for h in heights]

    space_after = 0 if scale < 1.0 else 1

    # 3) 그리기 루프
    y = content_top
    for sec, h in zip(sections, heights):
        # 좌측 분류 박스 — 무채색 + 테두리, 텍스트 세로 중앙
        box = add_rect(slide, RAIL_X, y, RAIL_W, h, fill=None,
                       line=BLACK, line_w=0.75)
        bt = box.text_frame
        bt.margin_left = Cm(0.06)
        bt.margin_right = Cm(0.06)
        bt.word_wrap = True
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        label_size = sec.get("label_size") or _fit_label_size(sec["label"])
        write_lines(bt, [sec["label"]], size=label_size, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER)

        # 우측 내용
        cy = y + PAD_TOP
        if sec.get("session_message"):
            add_text(slide, RIGHT_X, cy, RIGHT_W, CM_PER_HDR,
                     [(sec["session_message"],
                       {"size": SZ_SECTION, "bold": True, "color": BLACK})])
            cy += CM_PER_HDR

        if sec.get("content"):
            body_lines = sum(_line_for(it) for it in sec["content"])
            body_h = body_lines * CM_PER_LINE * scale
            content = sec["content"]
            if space_after == 0:
                content = [
                    (it[0], {**it[1], "space_after": 0}) if isinstance(it, tuple)
                    else (it, {"space_after": 0})
                    for it in content
                ]
            add_text(slide, RIGHT_X, cy, RIGHT_W, body_h, content)
            cy += body_h

        if sec.get("table"):
            tbl = sec["table"]
            nrows = len(tbl["rows"]) + 1
            row_h = tbl.get("row_h", 0.9) * scale
            tw = sum(tbl["col_w"])
            shape = slide.shapes.add_table(
                nrows, len(tbl["headers"]),
                Cm(RIGHT_X), Cm(cy), Cm(tw), Cm(nrows * row_h))
            fill_table(shape.table, tbl["headers"], tbl["rows"],
                       col_w=tbl["col_w"],
                       body_size=tbl.get("body_size", SZ_BODY),
                       head_size=tbl.get("head_size", SZ_BODY))
            for r in range(nrows):
                shape.table.rows[r].height = Cm(row_h)
            cy += nrows * row_h + 0.15

        if sec.get("image_box"):
            img_h = float(sec["image_box"]) * scale
            ph = add_rect(slide, RIGHT_X, cy, RIGHT_W, img_h,
                          fill=LIGHT_GRAY, line=MID_GRAY, line_w=0.75)
            pt = ph.text_frame
            pt.vertical_anchor = MSO_ANCHOR.MIDDLE
            write_lines(pt, [("〈이미지 첨부 영역〉",
                              {"size": SZ_SUB, "color": MID_GRAY})],
                        align=PP_ALIGN.CENTER)

        y += h + GAP

    if footnote:
        add_footnote(slide, footnote)
