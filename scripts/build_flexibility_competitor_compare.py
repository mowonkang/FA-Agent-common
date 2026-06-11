"""Samsung SDI · CATL 유연성 비교 분석 자료 — 메인(분류형) + 보조 3장.

메인은 사내 보고양식 슬라이드 1 분류형 레이아웃(좌측 분류 레일 + 우측 내용)
으로 3사 비교 요약을 1장에 담고, 보조 3장은 표 중심으로 심층 데이터를 싣는다.

  - 메인        : 비교 요약 / SDI 진단 / CATL 진단 / 시사점·권고 (분류형)
  - 보조 1/3    : Samsung SDI 심층 — 공장별 매트릭스·협력사 공시·시나리오
  - 보조 2/3    : CATL 심층 — 공장별 매트릭스·PSL 8세대/LEAD/Xiaomo·시나리오
  - 보조 3/3    : 자동화 철학 비교·라인 재구성 비용·임계지표(경계신호)·권고

데이터 출처(추측 없음): outputs/유연성_경쟁사_분석_v2_2026-05-21.md
(LGES F1 정의 기준 4-Proxy 가중 평균, SDI ±5~8%p · CATL ±5~10%p 추정).
출력: outputs/유연성_SDI_CATL_비교분석_2026-05-27.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_layout import (  # noqa: E402
    render_classified_slide, title_block, add_label_band, add_rect, add_text,
    fill_table, add_footnote, SLIDE_W_CM, SLIDE_H_CM,
    BLACK, BLUE, MID_GRAY, CHARCOAL, SZ_HEAD, SZ_SECTION, SZ_BODY, SZ_SUB,
)

OUT_PATH = Path("outputs/유연성_SDI_CATL_비교분석_2026-05-27.pptx")
BYLINE = ["FA기술담당", "(2026.05.27)"]


def band_table(slide, x, y, w, label, headers, rows, col_w, row_h, *,
               body=SZ_BODY, head=SZ_BODY, band_h=0.6, emph_col0=True):
    """라벨 띠 + 표 블록. 표 하단 y 를 반환."""
    add_label_band(slide, x, y, w, label, height=band_h)
    ty = y + band_h + 0.08
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(
        nrows, len(headers), Cm(x), Cm(ty), Cm(sum(col_w)), Cm(nrows * row_h))
    fill_table(shape.table, headers, rows, col_w=col_w,
               body_size=body, head_size=head, emph_col0=emph_col0)
    for r in range(nrows):
        shape.table.rows[r].height = Cm(row_h)
    return ty + nrows * row_h


def band_text(slide, x, y, w, label, lines, *, band_h=0.6, box_h=6.0):
    add_label_band(slide, x, y, w, label, height=band_h)
    add_text(slide, x + 0.05, y + band_h + 0.1, w - 0.1, box_h, lines)


# ───────────────────────── 메인 (분류형) ─────────────────────────
def build_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="Samsung SDI · CATL 유연성 비교 분석",
        byline=BYLINE,
        head_message=[
            ("LGES F1 정의(유연 설비 투자비 / 전체, 셋업비 ≤30% 컷오프)로 환산 시 "
             "'28년 유연성은 LGES 76% > SDI 57% > CATL 42% 로,",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
            ("표준화·고정형 중심의 경쟁사 대비 다품종·다화학 대응 유연성에서 글로벌 "
             "선도 격차를 확보할 전망입니다.",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "1\n비교 요약",
                "session_message": "3사 × 4개년 유연성 지표 (공장 가중 평균, mid-point)",
                "content": [
                    (" • LGES vs SDI '28 +19%p, LGES vs CATL '28 +34%p — '25→'28 "
                     "격차 유지·확대 / LGES = F1 ground truth, SDI·CATL = 4-Proxy "
                     "가중 환산 추정(±5~10%p)",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
                "table": {
                    "headers": ["유연성 지표", "LGES", "Samsung SDI", "CATL"],
                    "rows": [
                        ("'25년", "47%", "35%", "25%"),
                        ("'26년", "61%", "41%", "29%"),
                        ("'27년", "75%", "49%", "35%"),
                        ("'28년", "76%", "57%", "42%"),
                    ],
                    "col_w": [5.6, 5.4, 6.4, 5.4],
                    "row_h": 0.62,
                },
            },
            {
                "label": "2\nSamsung\nSDI",
                "session_message": "SDI '25 35% → '28 57% (mid) — LGES 대비 12~14%p 低",
                "content": [
                    (" • P5/P6 각형 라인의 표준화된 Conveyor 기반 설계로 라인 변경 "
                     "자유도 제약 (일본식 정밀 라인 전통)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 헝가리 괴드 wound→stacked 전환 = 기존 설비 개조 비중 高 → "
                     "셋업비 30% 초과로 고정형 분류",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 추격 변수 : 천안 46파이 신규 라인·코코모 SPE·뉴칼라일 GM JV·"
                     "SFA AGV/AS-RS 도입 (코윈테크 AMR 수주 +465.2% YoY)",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "3\nCATL",
                "session_message": "CATL '25 25% → '28 42% (mid) — 자동화율 95% 최고이나 "
                                   "유연성은 最低",
                "content": [
                    (" • PSL 8세대 : 1초/셀·인력 △70%·속도 +300%, 고정형 Conveyor + "
                     "Stacker Crane 의존 구조",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • \"95% 인터커넥션 ≠ 유연성\" — 설비 IoT 연결률이지 유연 설비 "
                     "투자 비중 아님 (정의 차이)",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 다화학(LFP/삼원계/Na-ion)은 전용 공장·라인 분리 전략 → 셀 사양 "
                     "변경 시 라인 재구성 비용 高 (Xiaomo 는 현재 Pack 한정)",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "4\n시사점\n·권고",
                "session_message": "유연성 = LGES 차별화 축, KPI·IR·R&D 에 전략 반영 필요",
                "content": [
                    (" • 단기 : 유연성 지표 IR 신규 KPI 도입 + 유럽 OEM 다품종 대응 "
                     "영업 메시지 강화",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 중기 : '26 KPI 유연성 가중치 0% → 하반기 재정의 시 10~15% 신설 "
                     "+ SMC/MMF 기술 IP 강화",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 장기 : WEF Lighthouse flexibility index 표준화 대응 + 추정 "
                     "정밀도 ±5%p 이하  (경계신호 3대 → 보조 3 참조)",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
        ],
        footnote=("* 출처 : outputs/유연성_경쟁사_분석_v2_2026-05-21.md "
                  "[22~30·166~169·182·237~239·252·262~270행]  ※ LGES F1 정의 환산 추정, "
                  "SDI ±5~8%p·CATL ±5~10%p"),
    )


# ─────────────────────── 보조 1/3 — SDI 심층 ───────────────────────
def build_aux_sdi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(
        slide, "[보조 1/3] Samsung SDI 유연성 심층 — 공장 매트릭스·공시·시나리오",
        byline=BYLINE)

    band_table(
        slide, 0.6, 1.55, 26.3,
        "① 공장별 유연성 매트릭스 (7개 공장, 유연성 추정 = '25년 기준)",
        ["공장", "CAPA('25)", "폼팩터", "자동화 특성", "유연성('25)"],
        [
            ("천안 (46파이 마더라인)", "1~2→확장", "46파이 원통",
             "신규 라인·국산 장비(코엠/필옵틱스/하나/이노메트리)", "40~48%"),
            ("울산 (각형)", "~30 GWh", "P5/P6 각형",
             "기존 Conveyor 中心, 일부 P6 전환", "28~34%"),
            ("헝가리 괴드 1·2", "40 GWh", "P5/P6 각형",
             "1공장 wound→stacked 전환 진행 中", "30~36%"),
            ("미국 코코모 1 (SPE)", "33 GWh", "P6 각형 NCA",
             "신규 라인, IRA 대응, 일부 LFP 전환", "38~44%"),
            ("미국 코코모 2 (SPE)", "34 GWh('27)", "P6 각형", "신규 라인", "45~52%"),
            ("미국 뉴칼라일 (GM JV)", "27→36('27)", "각형+원통",
             "신규 라인, 다폼팩터 대응", "48~55%"),
            ("말레이시아 SDIEM", "21700 원통", "21700 원통",
             "기존 Conveyor + 신규 BBU 라인", "35~40%"),
        ],
        col_w=[4.4, 3.0, 3.0, 12.5, 3.4], row_h=0.74, body=SZ_BODY)

    band_text(
        slide, 0.6, 8.40, 13.0,
        "② 설비 협력사 공시 기반 역추정",
        [
            (" • SFA(056190) : '26 1Q 신규 수주 2,179억(+18.3% YoY), 로봇물류 견인",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — 괴드 공장에 AI AGV·스마트 컨베이어·AS-RS 공급",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 코윈테크(282880) : AMR 상용화 후 로봇 수주 +465.2% YoY",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — 4-Way Shuttle·하이브리드 스토커 / 북미 전극 자동화 턴키 320억",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 씨케이솔루션 : 괴드 스택 전환 €16.8m('25.11)·€13.93m('24.5)",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — P5(권취)→P6(적층) 전환 본격화 신호",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 한화모멘텀·피엔티 : 전극 공정 장비, 헝가리 증설 협력사",
             {"size": SZ_BODY, "color": BLACK}),
        ], box_h=4.2)

    bot = band_table(
        slide, 14.0, 8.40, 12.9,
        "③ 시나리오별 유연성 추정 ('25~'28, %)",
        ["시나리오", "'25", "'26", "'27", "'28"],
        [
            ("보수 (기존 라인 比 高)", "32%", "36%", "42%", "50%"),
            ("기준 (Base)", "35%", "41%", "49%", "57%"),
            ("낙관 (신규 라인 比 高)", "38%", "44%", "52%", "62%"),
        ],
        col_w=[5.3, 1.9, 1.9, 1.9, 1.9], row_h=0.72, body=SZ_BODY)
    add_text(
        slide, 14.05, bot + 0.25, 12.8, 4.0,
        [
            ("[SDI 종합 판단]", {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • LGES 대비 보수 12~14%p·낙관 8~10%p 低 — P5/P6 표준 Conveyor 설계, "
             "일본식 정밀 라인, 괴드 개조 비중 高가 원인",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 검증 한계 : SDI 는 자동화 KPI 미공개, DOE LPO 발표에 AGV/AMR 미확인 "
             "→ opendart 분기보고서·코코모 환경영향평가 추가 페치 필요",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])

    add_footnote(
        slide,
        "* 출처 : outputs/유연성_경쟁사_분석_v2_2026-05-21.md §4 (SDI) "
        "[135~184행]  │  references/경쟁사/2026-05-21_SamsungSDI_P6_Hungary_"
        "InterBattery2025.md  │  The Elec·헬로티·joseilbo (accessed 2026-05-21)")


# ─────────────────────── 보조 2/3 — CATL 심층 ───────────────────────
def build_aux_catl(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(
        slide, "[보조 2/3] CATL 유연성 심층 — 공장 매트릭스·PSL 8세대·시나리오",
        byline=BYLINE)

    band_table(
        slide, 0.6, 1.55, 26.3,
        "① 공장별 유연성 매트릭스 (7개 공장, 유연성 추정 = '25년 기준)",
        ["공장", "CAPA", "폼팩터", "자동화 특성", "유연성('25)"],
        [
            ("푸젠 닝더 본사", "다수", "각형",
             "PSL 8세대·1.7초/셀, 고정형 Conveyor+Crane 의존", "18~24%"),
            ("푸젠 푸딩 No.5 (다층)", "25 GWh('25.8)", "각형",
             "230,000㎡·4라인 다층 → Stacker Crane 의존 高", "15~20%"),
            ("쓰촨 이빈 (CATL-SC)", "16 GWh", "각형",
             "WEF 등대, Pack 80% 자동화, gluing 인력 △70%", "25~30%"),
            ("장쑤 리양", "126 GWh", "각형",
             "WEF 등대, 3D프린팅 changeover 단축·출력 +320%", "28~34%"),
            ("독일 에르푸르트 (CATT)", "8→24 GWh", "모듈+셀",
             "6라인 가동, 1,700명 고용", "22~28%"),
            ("헝가리 데브레첸 P1", "~40 GWh('26)", "각형",
             "'26초 가동, BMW/Mercedes/Stellantis/VW 공급", "25~32%"),
            ("허난 뤄양 (Pack)", "—", "Pack",
             "휴머노이드 'Xiaomo' EOL/DCR 도입('25.12)", "20~26%"),
        ],
        col_w=[4.6, 3.0, 2.4, 12.9, 3.4], row_h=0.74, body=SZ_BODY)

    band_text(
        slide, 0.6, 8.40, 13.0,
        "② PSL 8세대 · LEAD 공시 · Xiaomo",
        [
            (" • PSL 8세대 : 1초/셀·모듈 20초/셀, 인력 △70%·속도 +300%",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — 7,000+ QC 포인트 실시간 모니터링 (고속 Conveyor+Crane)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • \"95% 인터커넥션\" = 설비 IoT 연결률 (≠ 유연 설비 비중)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 先导智能(LEAD, 300450) : 모듈 95%/Pack 50% 자동화 표준",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — CATL 매출비중 28.32%('18), 智能物流 글로벌 23.8%",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • Xiaomo(Spirit AI) : 뤄양 EOL/DCR 99% 성공률·인력 3배",
             {"size": SZ_BODY, "color": BLACK}),
            ("    — 현재 Pack 후공정 한정, 셀 라인 확대 시 +5~10%p 잠재",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ], box_h=4.2)

    bot = band_table(
        slide, 14.0, 8.40, 12.9,
        "③ 시나리오별 유연성 추정 ('25~'28, %)",
        ["시나리오", "'25", "'26", "'27", "'28"],
        [
            ("보수 (PSL 고정형 中心)", "22%", "24%", "28%", "32%"),
            ("기준 (Base)", "25%", "29%", "35%", "42%"),
            ("낙관 (휴머노이드/AMR 확산)", "28%", "33%", "40%", "46%"),
        ],
        col_w=[5.3, 1.9, 1.9, 1.9, 1.9], row_h=0.72, body=SZ_BODY)
    add_text(
        slide, 14.05, bot + 0.25, 12.8, 4.0,
        [
            ("[CATL 종합 판단]", {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • LGES 대비 '28 30~38%p 低, 단 자동화율 95%로 5~10%p 우위 — "
             "'효율성 vs 유연성' trade-off 명확",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 다화학은 전용 공장·라인 분리(Liyang P4·Yibin LFP·Fuding No.5) → "
             "'동일 라인 사양 변경' 접근과 본질적 상이",
             {"size": SZ_BODY, "color": BLACK}),
            (" • '25 실적 : 매출 RMB 423.7B(+17%)·순익 72.2B(+42%)·판매 661GWh·"
             "글로벌 PB 39.2%(9년 연속 1위)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])

    add_footnote(
        slide,
        "* 출처 : outputs/유연성_경쟁사_분석_v2_2026-05-21.md §5 (CATL) "
        "[192~254행]  │  references/경쟁사/2026-05-21_CATL_PSL_8세대_Smart"
        "Manufacturing.md [40~65행]  │  catl.com·leadintelligent.com·CnEVPost")


# ──────────────── 보조 3/3 — 자동화 철학·임계지표·권고 ────────────────
def build_aux_strategy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(
        slide, "[보조 3/3] 자동화 철학 비교 · 라인 재구성 비용 · 임계지표 · 권고",
        byline=BYLINE)

    band_table(
        slide, 0.6, 1.55, 26.3,
        "① 3사 자동화 철학 매트릭스 (8 차원)",
        ["차원", "LGES", "Samsung SDI", "CATL"],
        [
            ("자동화 철학", "모듈·유연성 (SMC/MMF/AMR)",
             "일본식 정밀 + 미국식 모듈(P6)", "표준화·대형화·무인화 (PSL 8세대)"),
            ("자동화율 ('28)", "85~90%", "85%", "95%"),
            ("유연성 지표 ('28)", "76%", "57% (mid)", "42% (mid)"),
            ("다품종 대응", "高 (원통/각형/파우치)", "中 (각형 주력)",
             "中 (라인별 분리)"),
            ("셀 사양 변경 비용", "低 (SMC 재배치)", "中 (라인 부분 개조)",
             "高 (라인 정지·재구성)"),
            ("셋업 변경 시간", "1.0x (기준)", "1.4x", "1.8x"),
            ("1초/셀 달성", "부분 (PPM 기준)", "미달성", "달성 (PSL 8세대)"),
            ("주요 협력사", "코윈테크·시너스텍·SFA",
             "SFA·코윈테크·씨케이솔루션", "先导智能·Geek+·Hai Robotics"),
        ],
        col_w=[4.6, 6.6, 6.6, 8.5], row_h=0.60, body=SZ_BODY)

    bot = band_table(
        slide, 0.6, 8.05, 13.0,
        "② 라인 재구성 비용 비교 (CAPEX 대비 %, 정성 추정)",
        ["시나리오", "LGES", "SDI", "CATL"],
        [
            ("동일 폼팩터 화학변경(NMC→LFP)", "5~10%", "12~18%", "20~30%"),
            ("폼팩터 부분변경(셀 크기)", "8~15%", "18~25%", "25~35%"),
            ("신규 폼팩터(46파이 등)", "15~25%", "25~35%", "35~50%"),
        ],
        col_w=[6.4, 2.2, 2.2, 2.2], row_h=0.74, body=SZ_BODY)
    add_text(
        slide, 0.65, bot + 0.25, 12.8, 4.0,
        [
            ("[재구성 비용 해석]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • LGES = SMC/MMF 재활용 → 폐기 최소화 / SDI = 모듈 부분 개조 / "
             "CATL = 라인별 분리·전용 공장 신축 일반",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CATL Liyang : 3D 프린팅 fixtures 로 changeover 단축 (본원적 설비 "
             "유연성과는 구분, catl.com 2023-12-14)",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])

    band_text(
        slide, 14.0, 8.05, 12.9,
        "③ 임계지표 3대 경계신호 + 권고",
        [
            (" ① SDI 추격 : '28 71%↑(LGES 5%p 이내) + 괴드 P5→P6 '27 100%",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("     → SMC IP 강화·코윈테크 발주 우선권, ESGM2 78%→82% 상향 검토",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" ② CATL Xiaomo : 셀 제조 라인 확대 + 95% 유지하며 유연성 50% 돌파",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("     → '27 내 공식 발표 시 차세대 SMC R&D 가속·휴머노이드 협력",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" ③ WEF 표준화 : Lighthouse 'flexibility index' 공식 KPI 채택",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("     → 6개월 내 MI Lansing/ESGM2 신규 Lighthouse 신청 준비",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 권고 : 단기 IR KPI·OEM 영업 / 중기 '26 KPI 유연성 10~15% 신설·IP "
             "강화 / 장기 표준화·정밀도 ±5%p",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ], box_h=4.6)

    add_footnote(
        slide,
        "* 출처 : outputs/유연성_경쟁사_분석_v2_2026-05-21.md §6·§9 "
        "[262~307·359~379행]  │  references/FA 유연성 지표 정의  │  "
        "references/26FA KPI.md (Slide 3 [62~70행])")


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    build_main(prs)
    build_aux_sdi(prs)
    build_aux_catl(prs)
    build_aux_strategy(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}  (slides: {len(prs.slides)})")


if __name__ == "__main__":
    build()
