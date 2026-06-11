"""전고체 라인 밀폐형 이송장치 보고서 빌더 (EFEM + FOUP).

references/기술자료/EFEM_로보스타.md (Robostar FoPLP EFEM, 원본 13 slide) 와
references/기술자료/FOUP(510x515)_3S.md (3S P-FOUP 510x515 8-slot, 원본 16 slide)
를 토대로, 전고체(all-solid-state) 배터리 라인용 밀폐형 이송장치 도입 검토를
FA기술담당 임원 보고용으로 정리한다.

사내 보고양식 `(보고양식)YYMMDD_보고제목_v1.0.pptx` 의 "분류 1~4 + Head
message" 구조 정체성을 유지하되, 양식의 작은 placeholder 박스로는 본문이
넘치므로 LGES_PPT_작업_가이드 기준(build_flexibility_summary.py 패턴)으로
박스·폰트를 정확히 맞춘 형태로 재구성한다.

구성 (3장, 전부 박스 내 수렴):
  - 슬라이드 1 : 메인 4사분면 (분류 1~4) + Head message + 풋노트
  - 슬라이드 2 : 보조 1/2 — EFEM vs FOUP 비교 및 전고체 적용 시사점 표
  - 슬라이드 3 : 보조 2/2 — 업체 공유 원본 PPT 슬라이드 삽입 가이드 표

전고체 배터리 특화 수치는 원문(반도체/FoPLP 패널용)에 없으므로 추측하지 않고
"적용 검토 / 가정" 으로 회색 명시한다. 유첨 본문은 새로 만들지 않고 업체
원본 PPT 슬라이드를 그대로 복사해 붙인다 (슬라이드 3 표로 안내).

출력: outputs/전고체라인_밀폐이송장치_검토_2026-05-18.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 팔레트 (LGES_PPT_작업_가이드) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x00, 0xFF)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

SZ_TITLE = 16
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

OUT_PATH = Path("outputs/전고체라인_밀폐이송장치_검토_2026-05-18.pptx")
SRC = ("* 출처 : references/기술자료/EFEM_로보스타.md │ "
       "references/기술자료/FOUP(510x515)_3S.md   "
       "※ 전고체 특화 제원은 원자료(반도체/FoPLP) 미수록 → '적용 검토/가정'")


def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH
    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                 color=BLACK, align=None, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, opts = item if isinstance(item, tuple) else (item, {})
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        para.space_before = Pt(0)
        para.space_after = Pt(opts.get("space_after", 1))
        run = para.add_run()
        set_run(run, text,
                font=opts.get("font", font),
                size=opts.get("size", size),
                bold=opts.get("bold", bold),
                color=opts.get("color", color))


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.12)
    tb.text_frame.margin_right = Cm(0.12)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    tb.text_frame.word_wrap = True
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, height=0.7,
                    fill=DIM_GRAY, color=WHITE):
    shp = add_rect(slide, left, top, width, height, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.25)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=SZ_BAND, bold=True, color=color)
    return shp


def add_footnote(slide):
    add_text(slide, 0.6, 18.35, 26.3, 0.55,
             [(SRC, {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})])


def fill_table(tbl, headers, rows, *, col_w, body_size=SZ_SUB,
               head_size=SZ_BODY, emph_col0=True):
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Cm(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        cell.margin_left = Cm(0.1)
        cell.margin_right = Cm(0.1)
        cell.margin_top = Cm(0.03)
        cell.margin_bottom = Cm(0.03)
        write_lines(cell.text_frame, [h], size=head_size, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = Cm(0.12)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.03)
            cell.margin_bottom = Cm(0.03)
            bold = emph_col0 and ci == 0
            color = CHARCOAL if (emph_col0 and ci == 0) else BLACK
            align = PP_ALIGN.LEFT if ci != 0 else PP_ALIGN.LEFT
            write_lines(cell.text_frame, [val], size=body_size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)


def title_block(slide, title):
    add_text(slide, 0.6, 0.25, 20.0, 1.0,
             [(title, {"size": SZ_TITLE, "bold": True, "color": BLACK})],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 20.7, 0.35, 6.5, 0.8,
             [("FA기술혁신파트 강모원책임",
               {"size": SZ_BODY, "bold": True, "color": BLUE,
                "align": PP_ALIGN.RIGHT})],
             anchor=MSO_ANCHOR.MIDDLE)
    ln = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    ln.line.color.rgb = MID_GRAY
    ln.line.width = Pt(0.75)


def quadrant(slide, left, top, w, h_band, h_card, label, lines):
    add_label_band(slide, left, top, w, label, height=h_band)
    add_rect(slide, left, top + h_band, w, h_card,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(slide, left + 0.3, top + h_band + 0.18, w - 0.6, h_card - 0.32,
             lines)


def build_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(
        slide, "전고체 라인 밀폐형 이송장치 도입 검토 (EFEM + FOUP 기반)")

    # Head message
    add_rect(slide, 0.6, 1.50, 26.3, 2.20,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.56, 25.8, 2.10,
        [
            ("[Head Message]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 전고체 공정은 0.3% 수준 Dry Room(DR) 조건 상시 유지가 "
             "필수 → 밀폐 이송장치(EFEM·FOUP) 검증 진행",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • EFEM : LPM 만 DR-less Test 부스 검증 필요(투자비 이슈) / "
             "FOUP : 1차 테스트 결과 내부 DR 유지 실패",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 다음 단계 → LPM Test · FOUP 개선 방안 도출 · Pilot "
             "재검증  (상세 → 보조 1 · 2)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ])

    Q_TOP, Q_MID = 3.75, 11.10
    Q_L, Q_R = 0.6, 14.0
    Q_W = 13.05
    H_BAND = 0.7
    H_TOP = 6.55
    H_BOT = 6.10

    quadrant(
        slide, Q_L, Q_TOP, Q_W, H_BAND, H_TOP,
        "분류 1.  도입 배경 · 요구사항",
        [
            ("[전고체 = 0.3% DR 조건 상시 유지 필수]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 전고체 셀·전극은 수분(H₂O)에 극히 민감",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 공정 전 구간 0.3% 수준 Dry Room(DR, 저습) 조건 유지 필요",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 공정 간 이송 시에도 DR 분위기 단절 없이 밀폐 유지가 관건",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 밀폐 이송장치(EFEM 전단 자동 이송 + FOUP 밀폐 포드) 검증",
             {"size": SZ_BODY, "color": BLACK}),
            (" · '0.3%' 기준값은 사용자 제시 — 측정 단위 정의는 검증 시 확정",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])

    quadrant(
        slide, Q_R, Q_TOP, Q_W, H_BAND, H_TOP,
        "분류 2.  EFEM — LPM DR-less Test 필요",
        [
            ("[LPM 만 별도 DR-less 검증 필요]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • EFEM 구성 : Transfer Robot · LPM · Aligner · Controller",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 이 중 LPM(Load Port Module)만 DR-less Test 부스에서",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("   별도 검증 필요 (DR 미적용 환경 거동 확인)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 전용 Test 부스 구축에 따른 투자비 이슈 존재",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 기타 모듈은 기존 EFEM 사양 활용 가능",
             {"size": SZ_BODY, "color": BLACK}),
            (" · EFEM 일반 사양 출처 : EFEM_로보스타.md 32~70행",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])

    quadrant(
        slide, Q_L, Q_MID, Q_W, H_BAND, H_BOT,
        "분류 3.  FOUP — 1차 테스트 내부 DR 유지 실패",
        [
            ("[1차 테스트 결과 : 내부 DR 유지 실패]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 3S P-FOUP 510×515 8-slot 대상 1차 테스트 진행 (에너지)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 결과 — 포드 내부 DR(저습) 유지 실패",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 밀폐/Purge 구조가 0.3% DR 조건 미충족 (개선 필요)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 점검 대상 : Door Gasket · Support Pin/O-ring · "
             "Purge Filter",
             {"size": SZ_BODY, "color": BLACK}),
            (" · FOUP 기본 사양 출처 : FOUP(510x515)_3S.md 30~160행",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])

    quadrant(
        slide, Q_R, Q_MID, Q_W, H_BAND, H_BOT,
        "분류 4.  다음 단계",
        [
            ("[Action Plan]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" ① LPM Test",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    - DR-less Test 부스에서 LPM 거동 검증 (투자비 협의 병행)",
             {"size": SZ_BODY, "color": BLACK}),
            (" ② FOUP 개선 방안 도출",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    - 1차 실패 원인 분석 → 밀폐/Purge 구조 개선안 수립",
             {"size": SZ_BODY, "color": BLACK}),
            (" ③ Pilot 재검증",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("    - 개선 FOUP·LPM 통합 Pilot 로 0.3% DR 유지 재확인",
             {"size": SZ_BODY, "color": BLACK}),
        ])

    add_footnote(slide)


def build_compare(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide, "전고체 밀폐 이송장치 — EFEM vs FOUP 비교")
    add_label_band(
        slide, 0.6, 1.55, 26.3,
        "보조 1/2.  EFEM vs FOUP 비교 및 전고체 적용 시사점", height=0.7)

    headers = ["구분", "EFEM (Robostar)", "FOUP (3S P-FOUP)",
               "전고체 적용 시사점", "비고"]
    rows = [
        ("역할", "장비 전단 자동 적출·Align", "공정 간 밀폐 이송 포드",
         "EFEM = 장비 IF / FOUP = 라인 이송", "-"),
        ("DR 유지", "LPM 구간 DR-less 거동 미확인",
         "1차 테스트 — 내부 DR 유지 실패",
         "0.3% DR 조건 미충족", "이슈"),
        ("핵심 과제", "LPM DR-less Test 부스 검증",
         "밀폐/Purge 구조 개선 방안 도출",
         "투자비 이슈(부스 구축)", "진행"),
        ("다음 단계", "LPM Test", "FOUP 개선 → 통합 Pilot",
         "Pilot 로 0.3% DR 재검증", "예정"),
    ]
    tbl_shape = slide.shapes.add_table(
        5, 5, Cm(0.6), Cm(2.55), Cm(26.3), Cm(8.6))
    fill_table(tbl_shape.table, headers, rows,
               col_w=[2.4, 6.0, 6.0, 7.5, 4.4],
               body_size=SZ_BODY, head_size=SZ_BODY)
    for r in range(5):
        tbl_shape.table.rows[r].height = Cm(1.72)

    add_rect(slide, 0.6, 11.7, 26.3, 5.8,
             fill=LIGHT_GRAY, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, 11.9, 25.7, 5.5,
        [
            ("[해석 — 1차 테스트 현황 및 다음 단계]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 전고체는 0.3% DR 조건 상시 유지가 전제 → 이송 구간에서도 "
             "DR 단절이 없어야 함",
             {"size": SZ_BODY, "color": BLACK}),
            (" • EFEM : LPM 만 DR-less Test 부스 검증 필요 — 전용 부스 "
             "구축에 따른 투자비 이슈 존재",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • FOUP : 1차 테스트 결과 내부 DR 유지 실패 → 밀폐/Purge "
             "구조 개선 방안 도출 후 Pilot 재검증 필요",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" · 1차 테스트 결과·투자비 이슈는 사용자(현업) 제시 현황 — "
             "원자료(공급사 md)는 일반 사양 출처로만 인용",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])
    add_footnote(slide)


def build_appendix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide, "유첨 — 업체 공유 원본 PPT 슬라이드 삽입 가이드")
    add_label_band(
        slide, 0.6, 1.55, 26.3,
        "보조 2/2.  업체 원본 PPT 에서 아래 slide 를 순서대로 복사 → 본 "
        "보고서 뒤 삽입", height=0.7)

    add_text(
        slide, 0.6, 2.45, 26.3, 0.7,
        [(" ※ 유첨 본문은 신규 작성하지 않음. 업체(Robostar / 3S) 공유 "
          "원본 PPT 의 해당 slide 를 그대로 복사·삽입한다.",
          {"size": SZ_SUB, "color": MID_GRAY})])

    headers = ["구분", "삽입할 원본 slide", "주요 내용 / 비고"]
    rows = [
        ("EFEM 핵심", "slide 3 · 4 · 5 · 8",
         "FoPLP EFEM 구성 · 시스템 구성요소 · Robostar 강점 · Aligner"),
        ("EFEM 기타", "slide 6 (또는 7 택1) · 10~12",
         "EFEM Robot · 2P/4P 3D · 외형치수"),
        ("FOUP 핵심", "slide 3 · 4 · 12 · 13 · 15",
         "FOUP 개요 · 설계 CONCEPT · 변형방지 설계 · Door Gasket · "
         "N₂ Purge Filter"),
        ("FOUP 기타", "slide 5~6 · 7 · 8~11 · 14",
         "기본/외형 제원 · SLOT 내부치수 · BOM(선택) · ESD"),
        ("제외 / 안내", "EFEM 1·2·9·13 / FOUP 1·2·16",
         "표지 · 회사소개 · 목차 · 감사 슬라이드는 보고서 불요 (제외)"),
    ]
    tbl_shape = slide.shapes.add_table(
        6, 3, Cm(0.6), Cm(3.35), Cm(26.3), Cm(11.6))
    fill_table(tbl_shape.table, headers, rows,
               col_w=[3.2, 7.3, 15.8],
               body_size=SZ_BODY, head_size=SZ_BODY)
    tbl_shape.table.rows[0].height = Cm(1.0)
    for r in range(1, 6):
        tbl_shape.table.rows[r].height = Cm(2.12)

    add_rect(slide, 0.6, 15.3, 26.3, 2.2,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, 15.45, 25.7, 1.95,
        [
            ("[삽입 순서]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" 본 보고서(메인 + 보조 1) 뒤에 → ① EFEM 핵심 → ② EFEM 기타 "
             "→ ③ FOUP 핵심 → ④ FOUP 기타 순으로 원본 slide 복사·삽입",
             {"size": SZ_BODY, "color": BLACK}),
            (" 핵심 표기 slide 우선 삽입, 기타·BOM 은 보고 분량에 따라 "
             "취사 선택",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ])
    add_footnote(slide)


def build():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)
    build_main(prs)
    build_compare(prs)
    build_appendix(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
