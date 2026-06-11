"""유연성(Flexibility) 지표 정의 자료 — 분류형 레이아웃 (메인 1장 보강).

`scripts/_ppt_layout.py` 의 ``render_classified_slide`` 로 사내 보고양식
슬라이드 1 레이아웃(제목 → 검정 서술형 Head message → 좌측 분류 레일 →
우측 세션메시지/내용)을 적용해 유연성 지표 정의를 1장으로 정리한다.

데이터 출처(추측 없음): references/FA 유연성 지표 정의
출력: outputs/유연성지표_정의_2026-05-27.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_layout import (  # noqa: E402
    render_classified_slide, SLIDE_W_CM, SLIDE_H_CM,
    BLACK, BLUE, MID_GRAY, SZ_HEAD, SZ_BODY, SZ_SUB,
)

OUT_PATH = Path("outputs/유연성지표_정의_2026-05-27.pptx")
FOOT = ("* 출처 : references/FA 유연성 지표 정의 "
        "[7~10·14~16·18~21·26~27·32~44·48~56·58·64~73행]  (작성 : FA기술담당, 26.05.19)")


def build_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="유연성(Flexibility) 지표 정의",
        byline=["FA기술담당", "(2026.05.27)"],
        head_message=[
            ("제품·폼팩터 변화에 유연하게 대응 가능한 설비를 제작하고, 이를 정량적으로 "
             "평가·관리하기 위해",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
            ("'유연성 지표(= 유연 설비 투자비 / 전체 설비 투자비)'를 정의하고 全 Site로 "
             "확산하겠습니다.",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "1\n정의·산식",
                "session_message": "Flexibility = 유연 설비 투자비 / 전체 설비 투자비 "
                                   "(분자·분모 모두 설비비 + 셋업비 합계)",
                "content": [
                    (" • 설비비 단독 산정안은 부결(지표 과대평가 우려) — 셋업비까지 합산해 "
                     "라인 구축 실투자 기준으로 평가",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 기존 Reusable(재활용성) 산식의 명칭·정의를 재정립한 확정안(F1)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • MI Lansing 투자비 1,322.3억 = 설비비 1,015.7 + 셋업비 306.6 "
                     "(전극 275.3 / 조립 611.0 / 활성화 436.0억)",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "2\n분류 기준",
                "session_message": "유연 설비 선정 = 셋업비 비중 30% 이하 (초과 시 고정형)",
                "content": [
                    (" • 예외 : 전용 설비·유니언·개조 高비용 설비(Stacker Crane, Conveyor, "
                     "Lift, 포장기, 세척기, MSS Rack)는 강제 고정형 분류 / 개조 셋업비 高 "
                     "Case 는 산정 제외",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["구분", "Sub-Group", "제작 : 설치", "주요 설비"],
                    "rows": [
                        ("고정 설비", "구조물", "50 : 50", "Rack, Stage 등"),
                        ("고정 설비", "기계장치", "70 : 30",
                         "Crane, Conveyor, Lift, 세척기 등"),
                        ("유연 설비", "Vehicle (AMR·Shuttle)", "90 : 10",
                         "AGV, AMR, OHT, Shuttle"),
                        ("유연 설비", "제작품", "100 : 0", "Skid, 대차"),
                    ],
                    "col_w": [3.6, 5.2, 3.4, 11.6],
                    "row_h": 0.66,
                },
            },
            {
                "label": "3\n연도별 지표",
                "session_message": "MI Lansing / ESGM2 각형 4개년 추이 (+29%p, '25→'28)",
                "content": [
                    (" • '25→'26 +14%p, '26→'27 +13%p (전극 SMC·MMF 효과로 큰 폭 상승) / "
                     "'27→'28 +1%p (유연 설비 확산 안정화 구간)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 견인 과제 : '25 혁신물류 → '26 전극 SMC → '27 MMF(활성화 SMC) → "
                     "'28 밀폐형 AMR",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["구분", "'25년", "'26년", "'27년", "'28년"],
                    "rows": [
                        ("MI Lansing (ESMI)", "47%", "61%", "75%", "76%"),
                        ("ESGM2 각형", "50%", "62%", "74%", "78%"),
                    ],
                    "col_w": [7.4, 3.6, 3.6, 3.6, 3.6],
                    "row_h": 0.66,
                },
            },
            {
                "label": "4\nKPI·로드맵\n연계",
                "session_message": "산식·MI Lansing 지표 확정 → 全 Site KPI 로 확산 (ESGM2)",
                "content": [
                    (" • SMC·MMF·AMR 등 '26 FA KPI 핵심 과제가 유연성 지표를 직접 견인",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • '26 KPI 모니터링 체계 구축 — 분기별 유연성 지표 트래킹",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 단, 현 '26 KPI 8개 항목 모두 유연성 가중치 0% → 하반기 재정의 시 "
                     "10~15% 신설 검토 필요",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                ],
            },
        ],
        footnote=FOOT,
    )


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    build_main(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
