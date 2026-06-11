"""PPT 처리 공통 헬퍼."""
from __future__ import annotations

import re
from typing import Iterable, Iterator

try:
    from pptx.shapes.autoshape import Shape
except Exception:
    Shape = object  # type: ignore

PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*([a-zA-Z0-9_\-]+)\s*\}\}")


def shape_id(slide_index: int, shape_index: int) -> str:
    return f"s{slide_index}_sh{shape_index}"


def parse_shape_id(value: str) -> tuple[int, int]:
    m = re.fullmatch(r"s(\d+)_sh(\d+)", value)
    if not m:
        raise ValueError(f"invalid shape_id: {value!r}")
    return int(m.group(1)), int(m.group(2))


def find_explicit_placeholders(text: str | None) -> list[str]:
    if not text:
        return []
    return PLACEHOLDER_PATTERN.findall(text)


def normalize_text(text: str | None) -> str:
    return (text or "").replace("\r\n", "\n")


def is_blank(text: str | None) -> bool:
    return not (text or "").strip()


def iter_shapes_recursive(shapes: Iterable) -> Iterator:
    """Yield all shapes including those inside group shapes."""
    for shp in shapes:
        yield shp
        if getattr(shp, "shape_type", None) == 6 and hasattr(shp, "shapes"):
            yield from iter_shapes_recursive(shp.shapes)
