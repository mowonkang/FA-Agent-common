"""roadmap_v2.md → docx + pptx 빌더.

그룹장 코멘트(설비 유연성·재활용성 지표 / 모듈화·셋업 자동화 강조 /
CNS Physical Works × AX × Physical AI 연계)를 반영한 v2 로드맵 콘텐츠를
Word(docx)와 PowerPoint(pptx)로 함께 출력한다.

사용:
    python3 scripts/build_roadmap_v2.py
"""
from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PPRGBColor
from pptx.enum.shapes import MSO_SHAPE


ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "outputs"
OUT_DIR.mkdir(exist_ok=True)

DOCX_PATH = OUT_DIR / "roadmap_v2.docx"
PPTX_PATH = OUT_DIR / "roadmap_v2.pptx"

VISION_TITLE = "FA기술담당 중장기 기술 로드맵 v2 (재정리)"
VISION_SUB = (
    "그룹장 코멘트 반영 — 설비 유연성·재활용성 지표 신설 / "
    "모듈화·셋업 자동화의 유연성 강조 / "
    "CNS Physical Works × AX × Physical AI 연계 명시"
)

VISION_BODY = (
    "시장·폼팩터 변화에 유연하게 대응 가능한 스마트 혁신 공장 구축으로 "
    "매몰 Loss 최소화 + CAPEX/OPEX 절감을 동시 달성."
)

VISION_BULLETS = [
    "유연성: Cell 단위 모듈화 + Phasing 투자 + 셋업 자동화로 폼팩터·Capa 변화에 즉응",
    "재활용성: Cube/Rack 등 핵심 자산을 철거 후 타 라인·공장으로 재활용 가능한 구조",
    "AX × Physical AI: CNS Physical Works(Omniverse, Isaac SIM)와 연계, AI 설계 자동화 + 가상 사전 검증",
]

# 4축 KPI 표 데이터 (header + rows)
KPI_HEADER = ["축", "지표", "2026", "2027", "2028"]
KPI_ROWS = [
    ["CAPEX (기존)", "각형 투자비 (억/GWh)", "88.4", "79.6 (-10%)", "70.7 (-20%)"],
    ["면적 (기존)", "각형 면적 (평/GWh)", "1,129", "1,073 (-15%)", "1,019 (-20%)"],
    ["유연성 (신규)", "폼팩터/Capa 전환 LT, 셋업 LT", "베이스라인 측정", "-30%", "-50%"],
    ["유연성 (신규)", "Phasing 투자 가능 비율", "Pilot 도입", "50% 적용", "80% 적용"],
    ["재활용성 (신규)", "설비·Rack 재활용률 (%)", "베이스라인 측정", "30%", "50%"],
    ["재활용성 (신규)", "매몰 Loss 절감액 (억)", "베이스라인 측정", "1차 산정", "표준화 적용"],
]

# CAPEX 모듈러 트랙
CAPEX_MOD_HEADER = ["연도", "과제", "키워드", "내용"]
CAPEX_MOD_ROWS = [
    ["2026", "SMC (Smart Modular Cube)", "[모듈화·재활용]",
     "다품종 대응 모듈형 Cube Rack, Cell 단위 Phasing 셋업, 철거 후 타 라인 재활용"],
    ["2027", "MMF (Mobile Micro Factory)", "[셋업 자동화·모듈화]",
     "이동 가능한 Rack 자동 조립 로봇, Modular 라인 이동성 극대화"],
    ["2028", "밀폐형 Cube 물류 설비", "[모듈화·재활용]",
     "DR Less 공정 내 밀폐형 Cube Rack, 폼팩터 변경 대응 표준 구조"],
]

# 셋업 자동화 트랙
SETUP_HEADER = ["연도", "과제", "내용"]
SETUP_ROWS = [
    ["2026", "SMA 설치 효율화", "셋업 관리 플랫폼 + 모듈 단위 사전 조립으로 현장 공수 최소화"],
    ["2027", "SMS 자동 관리", "3D 디지털 트윈 기반 셋업 진척 관리, 셋업·반입 시나리오 자동 생성"],
    ["2028", "SMS 셋업 자동화", "셋업 자동 제어 + 이상 사전 감지·리스크 대응"],
]

# AX 트랙
AX_HEADER = ["연도", "과제", "AX 연계 포지셔닝"]
AX_ROWS = [
    ["2026", "AI 레이아웃 자동 설계",
     "AX 기반 설계 자동화 — 도면 분석→설계안 생성→검토 전 주기 자동화"],
    ["2027", "제품 연동 시뮬레이션", "AX 기반 동기화 — 제품 기반 선행 레이아웃 실시간 동기화"],
    ["2028", "3D 디지털 트윈 공장 검증 통합",
     "AX + Physical AI 통합 — Omniverse 디지털 트윈 기반 통합 시뮬레이션"],
]

# Physical AI 트랙
PAI_HEADER = ["연도", "과제", "Physical AI 연계"]
PAI_ROWS = [
    ["2026", "신규 PC 제어기 도입", "(기반 기술) 복합 동작용 PC 제어기 확대"],
    ["2027", "Isaac SIM Physical 사전 검증",
     "Physical AI (NVIDIA Isaac SIM) — 가상 검증, CNS Physical Works 연계"],
    ["2028", "지능형 설비 제작 플랫폼",
     "AX + Physical AI 통합 플랫폼 — 디지털 트윈 + AI 설비 설계 통합 관리"],
]

# OPEX 트랙
OPEX_HEADER = ["구분", "2026", "2027", "2028"]
OPEX_ROWS = [
    ["공장 자동화", "MOMA 공코어 회수/투입", "정형 포장 해체 Robot", "비정형 포장 해체 Robot"],
    ["운영 효율화", "AMR 2D Navigation", "AMR 3D Vision Navigation",
     "AMR Vision + Edge AI 자율 운영 (구 '자율 운영' 흡수)"],
]

# 그룹장 코멘트 매핑
MAPPING_HEADER = ["그룹장 코멘트", "본 문서 반영 위치"]
MAPPING_ROWS = [
    ["① 설비 유연성·재활용성 지표로 고민",
     "§2 4축 KPI — 유연성·재활용성 2축 신설"],
    ["② 모듈화·셋업 자동화로 시장·폼팩터 변화 대응 강조",
     "§3 CAPEX 트랙 — 리드 톤 전환 + 키워드 명시"],
    ["③ CNS Physical Works, AX·Physical AI 연계",
     "§4 AX × Physical AI 트랙 — 트랙 재명명 + Isaac SIM/Omniverse 키워드"],
]


# ────────────────────────────────────────────────────────────────
# DOCX 빌드
# ────────────────────────────────────────────────────────────────
def _set_font(run, size=10, bold=False, color=None):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.bold = bold
    if color is not None:
        run.font.color.rgb = RGBColor(*color)


def _add_heading(doc, text, level=1):
    p = doc.add_heading(level=level)
    run = p.add_run(text)
    _set_font(run, size=14 if level == 1 else 12, bold=True,
              color=(0x1F, 0x3A, 0x5F))


def _add_para(doc, text, *, bold=False, bullet=False):
    style = "List Bullet" if bullet else None
    p = doc.add_paragraph(style=style)
    run = p.add_run(text)
    _set_font(run, size=10, bold=bold)


def _add_table(doc, header, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(header))
    table.style = "Light Grid Accent 1"
    for j, h in enumerate(header):
        cell = table.rows[0].cells[j]
        cell.text = ""
        run = cell.paragraphs[0].add_run(h)
        _set_font(run, size=10, bold=True, color=(0xFF, 0xFF, 0xFF))
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.rows[i].cells[j]
            cell.text = ""
            run = cell.paragraphs[0].add_run(val)
            _set_font(run, size=9)


def build_docx() -> None:
    doc = Document()

    # 표지 / 제목
    title = doc.add_paragraph()
    tr = title.add_run(VISION_TITLE)
    _set_font(tr, size=18, bold=True, color=(0x1F, 0x3A, 0x5F))

    sub = doc.add_paragraph()
    sr = sub.add_run(VISION_SUB)
    _set_font(sr, size=10, color=(0x55, 0x55, 0x55))
    doc.add_paragraph()

    # 1. 지향점
    _add_heading(doc, "1. 지향점 (Vision)", level=1)
    _add_para(doc, VISION_BODY, bold=True)
    for b in VISION_BULLETS:
        _add_para(doc, b, bullet=True)
    _add_para(doc,
              "[시장 변화 배경] 북미 숙련 작업자 수배 어려움 및 인건비 증가, "
              "생산 계획 변경에 따른 매몰 Loss 절감을 위한 유연성 극대화된 "
              "설비 개발 필요. (원본 슬라이드 2 인용)")

    # 2. 4축 KPI
    _add_heading(doc, "2. 4축 KPI (목표 및 효과, 2026 → 2028)", level=1)
    _add_para(doc,
              "기존 CAPEX·면적 2축에 유연성·재활용성 2축을 추가하여 그룹장 "
              "코멘트 ①을 반영.")
    _add_table(doc, KPI_HEADER, KPI_ROWS)

    # 3. CAPEX 트랙
    _add_heading(doc, "3. CAPEX 트랙 — 유연성·재활용성 극대화", level=1)
    _add_para(doc,
              "시장 수요·폼팩터 변화에 유연하게 대응하기 위해 Cell 단위 "
              "모듈화 + 셋업 자동화 + Phasing 투자 + 철거·재활용 가능 "
              "구조를 통합 적용. 매몰 Loss 최소화 + 셋업 공수 절감 동시 달성.",
              bold=True)

    _add_heading(doc, "3.1 모듈러 설비 — 폼팩터·Capa 변화 대응", level=2)
    _add_table(doc, CAPEX_MOD_HEADER, CAPEX_MOD_ROWS)
    _add_para(doc, "SMC 기준 효과: 면적 25% / 공수 33% / 투자비 15% 개선 + "
                   "유연성(Phasing·재활용) 신규 확보")

    _add_heading(doc, "3.2 셋업 자동화 — 셋업 LT 단축으로 시장 변화 즉응",
                 level=2)
    _add_table(doc, SETUP_HEADER, SETUP_ROWS)

    # 4. AX × Physical AI 트랙
    _add_heading(doc,
                 "4. AX × Physical AI 트랙 — CNS Physical Works 연계",
                 level=1)
    _add_para(doc,
              "기존 '설비 사전 검증' + '백엔드 레이아웃' 트랙을 AX 기반 설계 "
              "자동화 + Physical AI(NVIDIA Isaac SIM / Omniverse) 기반 가상 "
              "검증으로 재정의. CNS Physical Works 디지털 트윈 자산과 연계.")

    _add_heading(doc, "4.1 AX 기반 레이아웃 설계 자동화", level=2)
    _add_table(doc, AX_HEADER, AX_ROWS)

    _add_heading(doc,
                 "4.2 Physical AI 기반 설비 사전 검증 (CNS Physical Works 연계)",
                 level=2)
    _add_table(doc, PAI_HEADER, PAI_ROWS)

    # 5. OPEX 트랙
    _add_heading(doc, "5. OPEX 트랙", level=1)
    _add_table(doc, OPEX_HEADER, OPEX_ROWS)

    # 6. 매핑
    _add_heading(doc, "6. 그룹장 코멘트 ↔ 본 재정리 매핑", level=1)
    _add_table(doc, MAPPING_HEADER, MAPPING_ROWS)

    doc.save(DOCX_PATH)
    print(f"[OK] Word saved: {DOCX_PATH}")


# ────────────────────────────────────────────────────────────────
# PPTX 빌드
# ────────────────────────────────────────────────────────────────
def _add_textbox(slide, left, top, width, height, text, *,
                 size=14, bold=False, color=(0x1F, 0x3A, 0x5F)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = PPRGBColor(*color)
    return box


def _add_bullet_box(slide, left, top, width, height, items, *,
                    size=12, color=(0x33, 0x33, 0x33)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = "Malgun Gothic"
        run.font.size = PPt(size)
        run.font.color.rgb = PPRGBColor(*color)
    return box


def _add_pptx_table(slide, left, top, width, height, header, rows,
                    *, font_size=10):
    n_rows = 1 + len(rows)
    n_cols = len(header)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left),
                                   Inches(top), Inches(width),
                                   Inches(height))
    table = shape.table
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = "Malgun Gothic"
        run.font.size = PPt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = PPRGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPRGBColor(0x1F, 0x3A, 0x5F)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name = "Malgun Gothic"
            run.font.size = PPt(font_size)
            run.font.color.rgb = PPRGBColor(0x33, 0x33, 0x33)
    return shape


def _add_title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.33), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPRGBColor(0x1F, 0x3A, 0x5F)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.3)
    tf.text = ""
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = PPt(20)
    run.font.bold = True
    run.font.color.rgb = PPRGBColor(0xFF, 0xFF, 0xFF)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── 슬라이드 1: 표지 + 4축 KPI
    s1 = prs.slides.add_slide(blank)
    _add_textbox(s1, 0.5, 0.4, 12.3, 0.9, VISION_TITLE,
                 size=28, bold=True)
    _add_textbox(s1, 0.5, 1.3, 12.3, 0.7, VISION_SUB, size=12,
                 color=(0x55, 0x55, 0x55))
    _add_textbox(s1, 0.5, 2.0, 12.3, 0.6, "■ 지향점", size=16,
                 bold=True, color=(0x1F, 0x3A, 0x5F))
    _add_bullet_box(s1, 0.7, 2.55, 12.0, 1.3,
                    [VISION_BODY] + VISION_BULLETS, size=12)
    _add_textbox(s1, 0.5, 4.0, 12.3, 0.5, "■ 4축 KPI (2026 → 2028)",
                 size=16, bold=True)
    _add_pptx_table(s1, 0.5, 4.5, 12.3, 2.7, KPI_HEADER, KPI_ROWS,
                    font_size=10)

    # ── 슬라이드 2: CAPEX 트랙
    s2 = prs.slides.add_slide(blank)
    _add_title_bar(s2, "CAPEX 트랙 — 유연성·재활용성 극대화")
    _add_textbox(s2, 0.5, 0.9, 12.3, 0.9,
                 "시장 수요·폼팩터 변화에 유연하게 대응하기 위해 "
                 "Cell 단위 모듈화 + 셋업 자동화 + Phasing 투자 + "
                 "철거·재활용 가능 구조 통합 적용",
                 size=13, bold=True, color=(0x33, 0x33, 0x33))
    _add_textbox(s2, 0.5, 1.9, 12.3, 0.4,
                 "3.1 모듈러 설비 — 폼팩터·Capa 변화 대응",
                 size=14, bold=True)
    _add_pptx_table(s2, 0.5, 2.3, 12.3, 2.0, CAPEX_MOD_HEADER,
                    CAPEX_MOD_ROWS, font_size=10)
    _add_textbox(s2, 0.5, 4.5, 12.3, 0.4,
                 "3.2 셋업 자동화 — 셋업 LT 단축", size=14, bold=True)
    _add_pptx_table(s2, 0.5, 4.9, 12.3, 2.0, SETUP_HEADER, SETUP_ROWS,
                    font_size=10)

    # ── 슬라이드 3: AX × Physical AI 트랙
    s3 = prs.slides.add_slide(blank)
    _add_title_bar(s3, "AX × Physical AI 트랙 — CNS Physical Works 연계")
    _add_textbox(s3, 0.5, 0.9, 12.3, 0.9,
                 "기존 '설비 사전 검증' + '백엔드 레이아웃'을 "
                 "AX 기반 설계 자동화 + Physical AI 기반 가상 검증으로 재정의. "
                 "Isaac SIM / Omniverse / CNS Physical Works 디지털 트윈 연계.",
                 size=12, color=(0x33, 0x33, 0x33))
    _add_textbox(s3, 0.5, 1.9, 12.3, 0.4,
                 "4.1 AX 기반 레이아웃 설계 자동화", size=14, bold=True)
    _add_pptx_table(s3, 0.5, 2.3, 12.3, 2.0, AX_HEADER, AX_ROWS,
                    font_size=10)
    _add_textbox(s3, 0.5, 4.5, 12.3, 0.4,
                 "4.2 Physical AI 기반 설비 사전 검증", size=14, bold=True)
    _add_pptx_table(s3, 0.5, 4.9, 12.3, 2.0, PAI_HEADER, PAI_ROWS,
                    font_size=10)

    # ── 슬라이드 4: OPEX 트랙
    s4 = prs.slides.add_slide(blank)
    _add_title_bar(s4, "OPEX 트랙 — 공장 자동화 / 운영 효율화")
    _add_textbox(s4, 0.5, 1.0, 12.3, 0.4,
                 "AMR · MOMA · Edge AI 기반 운영 자율화", size=13,
                 color=(0x33, 0x33, 0x33))
    _add_pptx_table(s4, 0.5, 1.6, 12.3, 2.5, OPEX_HEADER, OPEX_ROWS,
                    font_size=11)
    _add_textbox(s4, 0.5, 4.3, 12.3, 0.4, "■ 핵심 메시지", size=14,
                 bold=True)
    _add_bullet_box(s4, 0.7, 4.8, 12.0, 2.0, [
        "2028년 AMR Edge AI 자율 운영 단계에서 구 '자율 운영' 트랙 흡수",
        "OPEX 트랙은 모듈러·셋업 자동화(CAPEX 트랙)와 유기적으로 연동 — "
        "재활용 가능한 Cube/Rack을 AMR이 자동 적재·이송",
        "Vision + Edge AI 기반 비정상 상황 실시간 대응으로 무인화 가속",
    ], size=12)

    # ── 슬라이드 5: 그룹장 코멘트 매핑 + Closing
    s5 = prs.slides.add_slide(blank)
    _add_title_bar(s5, "그룹장 코멘트 ↔ 본 재정리 매핑")
    _add_pptx_table(s5, 0.5, 1.2, 12.3, 2.3, MAPPING_HEADER,
                    MAPPING_ROWS, font_size=12)
    _add_textbox(s5, 0.5, 3.8, 12.3, 0.4, "■ 다음 단계", size=14,
                 bold=True)
    _add_bullet_box(s5, 0.7, 4.3, 12.0, 2.5, [
        "유연성·재활용성 베이스라인 측정 방법론 정의 (’26 3Q)",
        "Isaac SIM PoC 과제 정식 착수 + CNS Physical Works 디지털 트윈 자산 연계 협의",
        "SMC/MMF Phasing 투자 가이드라인 표준화 후, 타 라인 재활용 사례화",
    ], size=12)

    prs.save(PPTX_PATH)
    print(f"[OK] PPT saved: {PPTX_PATH}")


if __name__ == "__main__":
    build_docx()
    build_pptx()
