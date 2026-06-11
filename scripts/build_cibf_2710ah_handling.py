"""CIBF 2026 2700Ah급 ESS 셀과 고중량 셀 핸들링 자동화 — LGES FA 종합 보고 슬라이드 빌더.

유일 출처:
  references/기술자료/2026-05-26_CIBF2026_2700Ah셀_고중량핸들링자동화_종합보고서.md

관점: "셀 대형화로 셀 단중 10~25kg+ 시대 진입 → 신규 ESS 라인 핸들링 완전 무인화 필수".
셀 사실은 핸들링 요구사항을 도출하기 위한 입력으로 다룬다.

구조 (LGES_PPT_작업_가이드.md 준수):
  메인 1장 (TL;DR 4포인트 + 설계 기준선)
  보조 6장:
    [보조 1] 셀 대형화 경쟁 라인업 비교표 (280→2710Ah)
    [보조 2] BYD 2710Ah/Haohan 핵심 사실표 + 핸들링 시사점
    [보조 3] 전공정 핸들링 자동화 (전극·그리퍼·로봇·비전·ASRS)
    [보조 4] 후공정 물류 (2t급 AMR·ASRS·갠트리) + 벤더 매핑표
    [보조 5] FA 신규 자동화 과제 후보 10종 표
    [보조 6] 단계별 실행안(Phase 1~3) + 전환기준 + Caveats

출력: outputs/CIBF2026_2700Ah셀_고중량핸들링자동화_종합보고_2026-05-26.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt
from lxml import etree

# ─── 색상 팔레트 (회색·흰색 기조 / 파랑 강조 / 풋노트 진초록) ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SOFT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
GRAY = RGBColor(0xC0, 0xC0, 0xC0)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
DIM_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
CHARCOAL = RGBColor(0x2A, 0x2A, 0x2A)
BLACK = RGBColor(0x00, 0x00, 0x00)

BLUE = RGBColor(0x00, 0x00, 0xFF)        # 강조색 (R0 G0 B255 순수 파랑)
LIGHT_BLUE = RGBColor(0xE6, 0xEC, 0xFF)  # 하이라이트 배경 (옅은 파랑)
MOLD_GREEN = RGBColor(0x00, 0x66, 0x00)  # 풋노트 전용 초록곰팡이 (G102)

FONT_BODY = "LG스마트체 Regular"
FONT_EMPH = "Arial Narrow"

SZ_TITLE = 16
SZ_BAND = 12
SZ_SECTION = 11
SZ_BODY = 10
SZ_SUB = 9
SZ_FOOT = 8

OUT_PATH = Path(
    "outputs/CIBF2026_2700Ah셀_고중량핸들링자동화_종합보고_2026-05-26.pptx"
)
RESPONSIBLE = "FA기술담당"
SRC = "references/기술자료/2026-05-26_CIBF2026_2700Ah셀_고중량핸들링자동화_종합보고서.md"

SLIDE_W = 27.52
SLIDE_H = 19.05


# ───────────────────── 헬퍼 (build_flexibility_summary.py 재사용) ─────────────────────
def set_run(run, text, *, font=FONT_BODY, size=11, bold=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font == FONT_EMPH:
        ea_font, latin_font = FONT_EMPH, FONT_EMPH
    else:
        ea_font, latin_font = FONT_BODY, FONT_EMPH
    run.font.name = ea_font
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs", "latin"):
        existing = rPr.find(qn(f"a:{tag}"))
        if existing is not None:
            rPr.remove(existing)
    for tag, fnt in (("ea", ea_font), ("cs", ea_font), ("latin", latin_font)):
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", fnt)


def write_lines(tf, lines, *, font=FONT_BODY, size=11, bold=False,
                color=BLACK, align=None, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            para.alignment = align
        if "align" in opts:
            para.alignment = opts["align"]
        run = para.add_run()
        set_run(
            run, text,
            font=opts.get("font", font),
            size=opts.get("size", size),
            bold=opts.get("bold", bold),
            color=opts.get("color", color),
        )


def add_rect(slide, left, top, width, height, *, fill=WHITE,
             line=BLACK, line_w=0.5):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, lines, **kw):
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tb.text_frame.margin_left = Cm(0.1)
    tb.text_frame.margin_right = Cm(0.1)
    tb.text_frame.margin_top = Cm(0.05)
    tb.text_frame.margin_bottom = Cm(0.05)
    write_lines(tb.text_frame, lines, **kw)
    return tb


def add_label_band(slide, left, top, width, text, *, fill=DIM_GRAY, color=WHITE,
                   size=SZ_BAND):
    shp = add_rect(slide, left, top, width, 0.7, fill=fill, line=None)
    tf = shp.text_frame
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.02)
    tf.margin_bottom = Cm(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write_lines(tf, [text], size=size, bold=True, color=color)
    return shp


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, title, *, sub=None):
    add_text(
        slide, 0.6, 0.22, 21.0, 1.0,
        [(title, {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.8, 0.30, 6.4, 0.8,
        [("FA기술담당  │  종합보고",
          {"size": SZ_BODY, "bold": True, "color": BLUE,
           "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.30), Cm(26.9), Cm(1.30))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)


def add_footnote(slide, text):
    add_text(
        slide, 0.6, 18.35, 26.3, 0.55,
        [(text, {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )


def fill_table(tbl, headers, rows, col_widths, *, header_size=SZ_BODY,
               body_size=SZ_BODY, blue_cols=None, blue_rows=None,
               bold_col0=True, left_cols=None, sub_size=None,
               sub_cols=None):
    """범용 표 채움. blue_rows: 파랑 강조할 (행idx, 컬럼idx 또는 'all') 집합."""
    blue_cols = blue_cols or set()
    blue_rows = blue_rows or set()
    left_cols = left_cols if left_cols is not None else set()
    sub_cols = sub_cols or set()
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Cm(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GRAY
        cell.margin_left = Cm(0.08)
        cell.margin_right = Cm(0.08)
        cell.margin_top = Cm(0.02)
        cell.margin_bottom = Cm(0.02)
        write_lines(cell.text_frame, [h], size=header_size, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.02)
            cell.margin_bottom = Cm(0.02)
            color = BLACK
            if ci in blue_cols:
                color = BLUE
            if (ri, "all") in blue_rows or (ri, ci) in blue_rows:
                color = BLUE
            bold = (ci == 0 and bold_col0) or color == BLUE
            align = PP_ALIGN.LEFT if ci in left_cols else PP_ALIGN.CENTER
            size = sub_size if (sub_size and ci in sub_cols) else body_size
            write_lines(cell.text_frame, [val], size=size, bold=bold,
                        color=color, align=align, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════ 메인 슬라이드 ════════════════════════════
def build_main(prs):
    slide = add_blank(prs)
    add_title(slide,
              "CIBF 2026 2700Ah급 ESS 셀 · 고중량 셀 핸들링 자동화 — FA 종합 보고")

    # 핵심 메시지 박스
    add_rect(slide, 0.6, 1.50, 26.3, 2.30,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.55, 26.0, 2.25,
        [
            ("[핵심 메시지]  셀 대형화로 \"셀 단중 10~25kg+ 시대\" 진입 → 신규 ESS 라인 핸들링 완전 무인화는 선택이 아닌 필수",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 인간 단독 핸들링 한계(통상 23kg, ISO 11228-1)를 셀 한 개가 이미 초과 → 적재·이재·조립 무인화가 신규 라인 필수 조건",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • FA 우선순위 재정렬 : 신규 부지 대규모 투자(X) → 기존 라인 대형 셀 호환 retrofit + 자체 셀(JF1/JH4) 대형화 대비",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 설계 기준선 권고 : 확정 스펙 있는 587Ah(주류)~1175Ah(차세대) 우선, 2710Ah Blade는 교체형 EOAT로 별도 시나리오",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # 4사분면
    Q_TOP = 4.05
    Q_MID = 11.35
    Q_LEFT = 0.6
    Q_RIGHT = 14.0
    Q_W = 13.05
    Q_H_TOP = 7.05
    Q_H_BOT = 6.85

    # ── Q1: 2700Ah = BYD 단독 ──
    add_label_band(slide, Q_LEFT, Q_TOP, Q_W, "① 2700Ah = BYD 단독 라인업")
    add_rect(slide, Q_LEFT, Q_TOP + 0.7, Q_W, Q_H_TOP - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_LEFT + 0.35, Q_TOP + 0.85, Q_W - 0.7, Q_H_TOP - 1.0,
        [
            (" • BYD 2710Ah LFP Blade(공칭 3.2V, 약 8.67kWh/셀) = 세계 최대 용량",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    2025-09 IDEE 최초 공개 → 2026 CIBF 양산·다수 省 납입 강조",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 시스템 단위(14.5MWh Haohan)로만 판매, 셀 단품 치수·중량 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("    공식 미공개",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 1위권 경쟁사는 587/588Ah(주류) + 628~790Ah(차세대)로 수렴",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 1000Ah+ 양산품은 Hithium 1175Ah(20.4kg)·1300Ah(2026 Q4 예정)가",
             {"size": SZ_BODY, "color": BLACK}),
            ("    사실상 유일 — 셀=시스템(CTC/CTB) 모델은 BYD 단독 베팅",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── Q2: 핸들링 변곡점 ──
    add_label_band(slide, Q_RIGHT, Q_TOP, Q_W,
                   "② 핸들링 변곡점 = 셀 단중 10~25kg+ 시대")
    add_rect(slide, Q_RIGHT, Q_TOP + 0.7, Q_W, Q_H_TOP - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    # 단중 사다리 표
    t = slide.shapes.add_table(
        5, 3, Cm(Q_RIGHT + 0.3), Cm(Q_TOP + 0.9), Cm(Q_W - 0.6), Cm(3.25)
    ).table
    fill_table(
        t,
        ["셀", "용량", "셀 단중"],
        [
            ("EVE MB56", "628Ah", "10.7kg"),
            ("Sunwoda", "684Ah", "11.35kg"),
            ("Hithium ∞", "1175Ah", "20.4kg"),
            ("BYD Blade", "2710Ah", "35~40kg(추정)"),
        ],
        [4.3, 3.0, 5.15],
        blue_rows={(4, "all")},
        left_cols={0},
    )
    add_text(
        slide, Q_RIGHT + 0.35, Q_TOP + 4.40, Q_W - 0.7, Q_H_TOP - 4.55,
        [
            (" • 1300Ah은 22~25kg 추정, 2710Ah는 시스템 부피에너지밀도",
             {"size": SZ_BODY, "color": BLACK}),
            ("    233.8kWh/m³·Vcts 52.1% 역산 → 길이 1.5~2m·35~40kg 추정",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 셀 1개가 인력 한계 23kg 초과 → 모든 신규 ESS 라인 ",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    완전 무인화 필수",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ],
    )

    # ── Q3: 즉시 검토 5대 과제 ──
    add_label_band(slide, Q_LEFT, Q_MID, Q_W, "③ FA 즉시 검토 5대 신규 과제")
    add_rect(slide, Q_LEFT, Q_MID + 0.7, Q_W, Q_H_BOT - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_LEFT + 0.35, Q_MID + 0.85, Q_W - 0.7, Q_H_BOT - 1.0,
        [
            (" ① 대형 셀 적응형 그리퍼 (다규격 280~1000Ah+ 공용)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" ② 고중량(2t급) AMR/AGF 도입 (ESS 컨테이너 부분조립품 운반)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" ③ 화성/에이징 ASRS 모듈러화 (트레이·shuttle 1톤급 상향)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" ④ AI 비전 1000Ah+ 외관/얼라인먼트 검사 (0.3mm 정밀도)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" ⑤ 셀-팩 일체형(CTP/CTC) 라인 디지털 트윈",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("", {"size": 4}),
            (" * 상세 과제 후보 10종은 보조 5 참조",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # ── Q4: 우선순위 재정렬 + 결정 포인트 ──
    add_label_band(slide, Q_RIGHT, Q_MID, Q_W,
                   "④ 우선순위 재정렬 (LGES 2026 가이던스 연계)")
    add_rect(slide, Q_RIGHT, Q_MID + 0.7, Q_W, Q_H_BOT - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.35, Q_MID + 0.85, Q_W - 0.7, 3.10,
        [
            (" • 전사 capex 약 −40%, ESS 셀 출하 60GWh+ 목표, ESS 신규 수주 90GWh+",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 美 Holland(MI) EV→ESS retool 중, 1.4B USD 4개 신규 셀 라인",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 신규 대규모 투자 대신 기존 라인 retrofit + 자체 셀(JF1/JH4)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    대형화 대비로 FA 우선순위 재정렬",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )
    add_rect(slide, Q_RIGHT + 0.3, Q_MID + 4.30, Q_W - 0.6, 1.95,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.45, Q_MID + 4.35, Q_W - 0.9, 1.90,
        [
            ("[임원 결정 포인트]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • Phase 1(즉시) : 적응형 그리퍼+협동로봇 NG제거 PoC, 2t급 AMR 5대 시범",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 전환 Trigger : BYD 2710Ah 1GWh+ 인도 / 美 ESS 60GWh 출하 달성 시",
             {"size": SZ_SUB, "color": BLACK}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_footnote(
        slide,
        f"* 출처 : {SRC} (accessed 2026-05-26)  │  ess-news.com 2025/9/19 · "
        "EnergyTrend 2026/5/18 · Energy-Storage.News 2026/4  │  2710Ah 치수·중량은 공식 미공개·추정치")


# ════════════════════ 보조 1 : 경쟁 라인업 비교표 ════════════════════
def build_aux1(prs):
    slide = add_blank(prs)
    add_label_band(slide, 0.6, 0.30, 26.3,
                   "[보조 1/6] 셀 대형화 경쟁 라인업 비교 (CIBF 2026 시점, 280 → 2710Ah)",
                   size=SZ_BAND)
    add_text(
        slide, 0.6, 1.20, 26.3, 0.7,
        [(" • 587/588Ah로 수렴(컨테이너 표준화) vs. 2710Ah는 컨테이너-셀 직접 통합(CTC) 모델 — 핸들링 사양은 확정 스펙(587~1175Ah) 기준선 권고",
          {"size": SZ_SUB, "color": MID_GRAY})],
    )

    headers = ["셀", "제조사", "용량", "에너지", "치수 L×W×H (mm)", "중량",
               "Wh/kg", "Cycle", "상태"]
    rows = [
        ("280Ah", "다수", "280Ah", "0.896kWh", "173×72×207", "~5.4kg", "165", "6~8천", "점차 단종"),
        ("314Ah", "다수", "314Ah", "1.005kWh", "173×72×207", "~5.6kg", "173", "8~10천", "현주류"),
        ("587Ah", "CATL", "587Ah", "1.878kWh", "미공시", "미공시", "≥185", "12천+", "2025/6 양산"),
        ("587Ah ∞", "Hithium", "587Ah", "1.878kWh", "286×73.5×216.3", "~11kg", "~185", "11천+", "2025 양산"),
        ("588Ah", "REPT", "588Ah", "1.88kWh", "미공시", "미공시", "190", "10~12천", "2026 양산"),
        ("628Ah MB56", "EVE", "628Ah", "2.01kWh", "352.3×71.7×207.2", "10.7kg", "188", "8천+", "2024/12 양산"),
        ("684Ah", "Sunwoda", "684Ah", "2.19kWh", "501.5×55.5×175.5", "11.35kg", "193", "12~15천", "2025/9 양산"),
        ("755Ah 반고체", "Shuangdeng", "755Ah", "~2.5kWh", "비공개", "비공개", "-", "-", "CIBF 신제품"),
        ("1175Ah ∞", "Hithium", "1175Ah", "3.76kWh", "580.2×75.2×216.3", "20.4kg", "180", "11천+", "2025/6 양산"),
        ("1300Ah ∞", "Hithium", "1300Ah", "4.16kWh", "580.2×75.2×234.3", "~22kg추정", "≥190", "≥10천", "2026 Q4 예정"),
        ("2710Ah Blade", "BYD", "2710Ah", "~8.67kWh", "비공개(추정 1.5~2m)", "비공개(35~40kg)", "-", "≥10천", "2026 양산·납입"),
    ]
    t = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Cm(0.6), Cm(2.05), Cm(26.3), Cm(13.4)
    ).table
    col_w = [3.0, 2.55, 2.0, 2.2, 4.6, 3.1, 1.9, 2.4, 3.5]
    blue_rows = {(11, ci) for ci in range(9)}  # 2710Ah 행 강조
    fill_table(
        t, headers, rows, col_w,
        header_size=SZ_SUB, body_size=SZ_SUB,
        blue_rows=blue_rows, left_cols={0, 1},
    )

    add_text(
        slide, 0.6, 15.65, 26.3, 2.4,
        [
            ("[핵심 통찰]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 587/588Ah 수렴 이유 : ① 20ft 컨테이너 ② 1500V PCS ③ 50t 해상운송 한계 3대 제약 안에서 IRR 극대화",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 2710Ah는 IP66 + Vcts 52.1%로 1GWh 시스템당 컨테이너 수 −52%·부지 −33% 달성 (시스템 통합 단계에서 제약 돌파)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 중량(EVE MB56)은 자료별 10.7kg/11.5kg 혼재, 2710Ah·1300Ah는 비공개·추정 — 핸들링 사양화는 확정 스펙 우선",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )
    add_footnote(
        slide,
        f"* 출처 : {SRC} §3 (accessed 2026-05-26)  │  "
        "Energy-Storage.News 'Beyond 314Ah' · EnergyTrend 2026/5  │  일부 단중·치수는 제3자 소매상 자료(±20%)")


# ════════════════ 보조 2 : BYD 2710Ah 핵심 사실표 + 핸들링 시사점 ════════════════
def build_aux2(prs):
    slide = add_blank(prs)
    add_label_band(slide, 0.6, 0.30, 26.3,
                   "[보조 2/6] BYD 2710Ah Blade / Haohan 핵심 사실표 + 핸들링 시사점")

    # 좌: 사실표
    add_text(slide, 0.6, 1.20, 13.0, 0.6,
             [("[셀·시스템 핵심 사실]", {"size": SZ_SECTION, "bold": True, "color": BLACK})])
    fact_rows = [
        ("제조사", "BYD FinDreams Battery"),
        ("최초 공개", "2025/9 IDEE 2025 (선전), 曹虎 원장"),
        ("양산 단계", "2026 양산, 다수 省 Haohan 첫 납입"),
        ("폼팩터", "Prismatic Blade (스택 Z-folding)"),
        ("정격 용량 / 전압", "2,710Ah / 3.2V (LFP)"),
        ("셀 에너지", "약 8.67kWh (공식 미공시·산출치)"),
        ("부극", "실리콘-카본 + 掺硅补锂 (ESS급 최초)"),
        ("사이클 / 온도", "≥10,000 cycle / −30 ~ +50°C"),
        ("셀 치수·단중", "공식 미공개 (시스템 단위만 판매)"),
        ("시스템 Haohan", "14.5MWh, 233.8kWh/m³, Vcts 52.1%, IP66"),
        ("첫 글로벌 배포", "사우디 12.5GWh, 일본 BYD법인(6개월 납기)"),
    ]
    t = slide.shapes.add_table(
        len(fact_rows), 2, Cm(0.6), Cm(1.85), Cm(13.0), Cm(11.6)
    ).table
    for i, w in enumerate([4.3, 8.7]):
        t.columns[i].width = Cm(w)
    for ri, (k, v) in enumerate(fact_rows):
        kc, vc = t.cell(ri, 0), t.cell(ri, 1)
        kc.fill.solid(); kc.fill.fore_color.rgb = SOFT_GRAY
        vc.fill.solid(); vc.fill.fore_color.rgb = WHITE
        for c in (kc, vc):
            c.margin_left = Cm(0.12); c.margin_right = Cm(0.08)
            c.margin_top = Cm(0.02); c.margin_bottom = Cm(0.02)
        blue = "미공개" in v or "미공시" in v
        write_lines(kc.text_frame, [k], size=SZ_SUB, bold=True, color=BLACK,
                    anchor=MSO_ANCHOR.MIDDLE)
        write_lines(vc.text_frame, [v], size=SZ_SUB,
                    bold=blue, color=(BLUE if blue else BLACK),
                    anchor=MSO_ANCHOR.MIDDLE)

    # 우: 핸들링 시사점
    add_label_band(slide, 14.0, 1.20, 12.9, "핸들링 관점 시사점", fill=DIM_GRAY)
    add_rect(slide, 14.0, 1.90, 12.9, 11.55, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 14.3, 2.05, 12.3, 11.3,
        [
            ("[셀 단품 추정]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 시스템 부피에너지밀도 233.8kWh/m³ + Vcts 52.1% 역산",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 셀당 부피 약 35~40L, 길이 1.5~2m, 중량 35~40kg 추정",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • (EV용 Blade 138Ah가 96×9×1.4cm·2.63kg 기준 외삽)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 5}),
            ("[핸들링 요구]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 사람 1인 핸들링 한계 명백히 초과 → 단일 셀 운송도 2인 1조 또는 기계 보조 의무화 대상(산안법)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Blade 형 길쭉한 폼팩터 → 다점 흡착 + 처짐(deflection) 방지 강성 필수, 폼팩터별 교체형 그리퍼 필요",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 셀 단독 트레이 가능성 → ASRS shuttle 페이로드 1톤급 상향",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 5}),
            ("[안전]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 손상 시 8.67kWh 일시 방출 가능성 → 단일 셀 격리(NTP) 트레이 설계 필수",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 한국 KFI는 아직 2710Ah급 ESS 별도 시험 기준 없음",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )
    add_footnote(
        slide,
        f"* 출처 : {SRC} §2·§8 (accessed 2026-05-26)  │  "
        "ess-news.com 2025/9/19 · Energy-Storage.News 2026/4  │  셀 단품 치수·중량 공식 미공개, 추정치(±20%)")


# ════════════════ 보조 3 : 전공정 핸들링 자동화 ════════════════
def build_aux3(prs):
    slide = add_blank(prs)
    add_label_band(slide, 0.6, 0.30, 26.3,
                   "[보조 3/6] 전공정 핸들링 자동화 — 전극·그리퍼·로봇·비전·화성/에이징")

    Q_LEFT, Q_RIGHT = 0.6, 14.0
    Q_W = 13.05
    Q_TOP = 1.25
    Q_H = 7.55

    # Q1 전극 공정
    add_label_band(slide, Q_LEFT, Q_TOP, Q_W, "전극 공정 (Roll-to-Roll)", fill=DIM_GRAY)
    add_rect(slide, Q_LEFT, Q_TOP + 0.7, Q_W, Q_H - 0.7, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_LEFT + 0.35, Q_TOP + 0.85, Q_W - 0.7, Q_H - 1.0,
        [
            (" • 전극 폭(~1000mm)은 크게 안 늘고 길이 방향 확장(시트 1.5m+)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 권취 단위 무게 현 200~400kg → 잠재 600kg+ 가 자동화 이슈",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Hithium 587Ah Xiamen : 단일라인 +17.1%, Cpk>1.33, 약 2,000개 지능형 장비",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • PNT : 광폭 코터 + AI 비전 동시 공급, 1μm 압연 재현성",
             {"size": SZ_BODY, "color": BLACK}),
            ("    (PNT 재무 수치는 DART/IR 1차 검증 필요)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 건식 전극(Dry electrode) 전환 시 코터·압연 구조 단순화, 중량은 유사",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # Q2 그리퍼 + 로봇
    add_label_band(slide, Q_RIGHT, Q_TOP, Q_W, "그리퍼 · 중량 가반 로봇 선택", fill=DIM_GRAY)
    add_rect(slide, Q_RIGHT, Q_TOP + 0.7, Q_W, Q_H - 0.7, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.35, Q_TOP + 0.85, Q_W - 0.7, Q_H - 1.0,
        [
            (" • 587Ah↑ : 진공 흡착 안정성 ↑ + 측면 클램핑(자중 슬립 방지) + 다규격 호환 필수",
             {"size": SZ_BODY, "color": BLACK}),
            (" • EOAT : Schmalz 진공 리프터·크레인, Re:Build 커스텀 컴플라이언스",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • 협동로봇 : UR30(35kg·리치1300) → 1175/1300Ah 단품 마진, FANUC CRX(3~40kg)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 산업용 6축 : FANUC M-900iB(700kg), ABB IRB 7600/8700(500~800kg)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • KUKA 리니어 유닛 5,000kg → 화성/에이징 트랙 연결",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 587Ah(~11kg)=25~35kg 협동로봇 영역 / 2710Ah(추정 35~40kg)·모듈=6축+갠트리",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ],
    )

    # 하단: 비전 + 화성/에이징
    Q_MID = 9.15
    Q_HB = 8.0
    add_label_band(slide, Q_LEFT, Q_MID, Q_W, "AI 비전 정렬", fill=DIM_GRAY)
    add_rect(slide, Q_LEFT, Q_MID + 0.7, Q_W, Q_HB - 0.7, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_LEFT + 0.35, Q_MID + 0.85, Q_W - 0.7, Q_HB - 1.0,
        [
            (" • Hithium 587Ah : 390+ CCD·3D카메라 13 스테이션, 1,200+ 검사작업, 0.3mm 결함 검출 = 산업 표준",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 얼라인먼트 허용치 ±0.3mm 유지 → 셀 길이 1.5~2m 시 카메라 시야각·셀 단위 마진 재설계",
             {"size": SZ_BODY, "color": BLACK}),
            (" • EOAT 컴플라이언스 + 비전 정렬 피드백 루프 결합 권장",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • LGES : 코윈테크·PNT 비전 적용 중, 대형 셀 대응 재설계 필요",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )
    add_label_band(slide, Q_RIGHT, Q_MID, Q_W, "화성 / 에이징 공정", fill=DIM_GRAY)
    add_rect(slide, Q_RIGHT, Q_MID + 0.7, Q_W, Q_HB - 0.7, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, Q_RIGHT + 0.35, Q_MID + 0.85, Q_W - 0.7, Q_HB - 1.0,
        [
            (" • 트레이 단중 : 314Ah 24셀 ~150kg → 587Ah ~260kg → 1175Ah ~480kg",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 2710Ah는 셀당 단독 트레이 가능성 → 트레이 로봇·ASRS shuttle 1톤급 상향",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 셀 에너지 ↑ → 에이징 창고 면적 셀당 에너지 비례 축소(긍정), 단일 셀 화재 위험 ↑",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Cimcorp 턴키 : Soaking/Aging ASRS + Pre-charge/Formation ASRS + 갠트리 + AGV",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" • Holland 라인 retool 시 검토 후보",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )
    add_footnote(
        slide,
        f"* 출처 : {SRC} §4 (accessed 2026-05-26)  │  "
        "Schmalz · UR30 · FANUC CRX · Re:Build · Cimcorp 공식 페이지")


# ════════════════ 보조 4 : 후공정 물류 + 벤더 매핑 ════════════════
def build_aux4(prs):
    slide = add_blank(prs)
    add_label_band(slide, 0.6, 0.30, 26.3,
                   "[보조 4/6] 후공정 물류 (2t급 AMR · ASRS · 갠트리) + 글로벌/한국 벤더 매핑")

    # 좌: AMR/물류 텍스트
    add_label_band(slide, 0.6, 1.25, 12.9, "공장 내 물류 (Intralogistics)", fill=DIM_GRAY)
    add_rect(slide, 0.6, 1.95, 12.9, 11.5, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, 2.10, 12.3, 11.3,
        [
            ("[2t급 AMR 페이로드 현황]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • OTTO 1500 V2 (Rockwell) : 1,900kg, 2 m/s, CE·B56.5·ISO 인증",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • MiR 1350 : 1,350kg / ROEQ TMS-C1500 부착 시 1,500kg 확장",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Youibot L1000 : 회전 리프팅 1,000kg, 8h 연속 운용",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 코윈테크 : 고중량 ESS 이송 AMR + 조립로봇 통합, 글로벌 ESS 수주",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("    (2025/9·12 추가 계약, 한국발 솔루션 침투)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 5}),
            ("[중량물·창고]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • AGF(무인지게차) : BYD/Toyota/KION 1.5~3톤급",
             {"size": SZ_BODY, "color": BLACK}),
            (" • ESS 컨테이너(7~30t)는 갠트리 크레인·reach stacker 영역",
             {"size": SZ_BODY, "color": BLACK}),
            (" • ASRS/갠트리 : Cimcorp 한 SI 통합, 한국 SFA·신성이엔지 동등급",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CTB/CTP 보편화 → 모듈 단계 축소, 셀이 직접 팩/컨테이너로",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # 우: 벤더 매핑표
    add_label_band(slide, 14.0, 1.25, 12.9, "벤더 매핑 (배터리 manufacturing 한정)", fill=DIM_GRAY)
    headers = ["영역", "글로벌 리더", "한국"]
    rows = [
        ("전극 코터·압연", "Manz, Hirano, Bühler", "PNT, CIS, 디에이"),
        ("노칭·스택·와인딩", "Manz, Dürr, 利元亨", "엠플러스, 필옵틱스"),
        ("조립·용접", "KUKA, FANUC, ABB", "CIS, 한화모멘텀"),
        ("화성·에이징", "Chroma, Hioki, PNE", "하나기술, 원익피앤이"),
        ("AMR/AGV", "MiR, OTTO, Youibot", "코윈테크, 모비어스, 트위니"),
        ("ASRS·갠트리", "Cimcorp, Daifuku, SSI", "SFA, 신성이엔지"),
        ("비전 검사", "Cognex, Keyence, ISRA", "코그넥스 한국, 라온피플"),
        ("Gripper·진공", "Schmalz, FIPA, Piab", "OnRobot, 슈말츠 한국"),
        ("중량 협동로봇", "UR(UR20/30), FANUC CRX", "두산로보틱스(M2513)"),
    ]
    t = slide.shapes.add_table(
        len(rows) + 1, 3, Cm(14.0), Cm(1.95), Cm(12.9), Cm(9.6)
    ).table
    fill_table(
        t, headers, rows, [3.7, 5.0, 4.2],
        header_size=SZ_SUB, body_size=SZ_SUB, left_cols={1, 2},
        blue_rows={(5, 2)},  # AMR 한국=코윈테크 강조
    )
    add_text(
        slide, 14.0, 11.75, 12.9, 1.8,
        [
            (" • 한국발 강점 : 코윈테크(고중량 ESS 전문 AMR), SFA(ASRS), 두산로보틱스(중량 협동로봇)",
             {"size": SZ_SUB, "color": BLACK}),
            (" • EVE MB56 압축 사양 : 단자당 800N·6N·m 토크 → 셀 자중×수량이 토크 사양 직결",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )
    add_footnote(
        slide,
        f"* 출처 : {SRC} §5 (accessed 2026-05-26)  │  "
        "ottomotors.com · Youibot L-series · IFR case study · Pinpoint News")


# ════════════════ 보조 5 : FA 신규 자동화 과제 후보 10종 ════════════════
def build_aux5(prs):
    slide = add_blank(prs)
    add_label_band(slide, 0.6, 0.30, 26.3,
                   "[보조 5/6] LGES FA 신규 자동화 과제 후보 10종 (공정 · TRL · CAPEX · 우선순위)")

    headers = ["#", "과제명", "적용공정", "TRL", "CAPEX", "핵심 고려사항", "우선순위"]
    rows = [
        ("1", "다규격 호환 적응형 셀 그리퍼", "조립·모듈", "6-7", "중(10~30억)", "컴플라이언스, 진공+측면 클램핑", "High"),
        ("2", "2톤급 AMR 도입 (컨테이너 운반)", "후공정 물류", "8", "대(50~100억)", "충전 인프라, KCs 인증", "High"),
        ("3", "AI 비전 1000Ah+ 외관/얼라인먼트", "전극·조립", "7-8", "중(20~40억)", "클린룸, 카메라 시야각 재설계", "High"),
        ("4", "화성/에이징 ASRS 모듈러화", "화성·에이징", "6-7", "대(100억+)", "트레이/shuttle 1톤 상향", "Medium"),
        ("5", "셀-팩 일체형(CTP/CTB) 디지털 트윈", "전라인", "4-5", "중(15~30억)", "MES 통합, OT/IT 보안", "Medium"),
        ("6", "건식 전극(Dry) 핸들링 라인", "전극", "5-6", "대(200~300억)", "LGES 자체 IP, PNT 협업", "Medium"),
        ("7", "협동로봇 화성트레이 NG셀 자동 제거", "화성·에이징", "8", "소(3~5억)", "25kg↑ 작업자 안전, UR30 35kg", "High"),
        ("8", "20ft 컨테이너 CTC 통합 라인", "후공정 통합", "4-5", "초대(500억+)", "BYD·CATL 선점, 자체 셀 검토", "Low"),
        ("9", "고중량 셀 적재 갠트리+모노레일", "조립·후공정", "9", "중(20~40억)", "천장 구조, 고소 작업 안전", "Medium"),
        ("10", "셀 화재 조기 검출·비상 격리(IR+AI)", "화성·창고", "7", "소~중(5~15억)", "UL 9540A, 화재 에너지 ↑", "High"),
    ]
    t = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Cm(0.6), Cm(1.25), Cm(26.3), Cm(13.2)
    ).table
    col_w = [1.0, 6.6, 3.0, 1.6, 3.3, 8.0, 2.8]
    # High 우선순위 셀(마지막 컬럼)만 파랑 강조
    high_rows = {1, 2, 3, 7, 10}
    blue_rows = {(ri, 6) for ri in high_rows}
    fill_table(
        t, headers, rows, col_w,
        header_size=SZ_SUB, body_size=SZ_SUB,
        left_cols={1, 5}, bold_col0=False,
        blue_rows=blue_rows,
    )
    # 1번 컬럼(#) 가운데 정렬은 기본, 과제명 Bold 처리
    for ri in range(1, len(rows) + 1):
        c = t.cell(ri, 1)
        write_lines(c.text_frame, [rows[ri - 1][1]], size=SZ_SUB, bold=True,
                    color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(
        slide, 0.6, 14.65, 26.3, 3.2,
        [
            ("[해석]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • High 우선순위 5종(#1·2·3·7·10) = 적응형 그리퍼 · 2t AMR · AI 비전 · 협동로봇 NG제거 · 화재 조기검출 → 메인 즉시검토 과제와 정합",
             {"size": SZ_BODY, "color": BLACK}),
            (" • Medium 4종(#4·5·6·9)은 라인 면적·ramp-up·IP 확보 효과, Low 1종(#8 CTC)은 BYD·CATL 선점 영역으로 전략 검토용",
             {"size": SZ_BODY, "color": BLACK}),
            (" • CAPEX/OPEX 수치는 한국 시장 평균 추정이며 SI별 ±50% 차이 가능",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )
    add_footnote(
        slide,
        f"* 출처 : {SRC} §7 (accessed 2026-05-26)  │  "
        "CAPEX/OPEX 한국 시장 평균 추정치(±50%), TRL 자체 평가")


# ════════════════ 보조 6 : 단계별 실행안 + 전환기준 + Caveats ════════════════
def build_aux6(prs):
    slide = add_blank(prs)
    add_label_band(slide, 0.6, 0.30, 26.3,
                   "[보조 6/6] 단계별 실행안 (Phase 1~3) · 전환 기준(Trigger) · Caveats")

    # 좌상 Phase 1
    add_label_band(slide, 0.6, 1.25, 12.9, "Phase 1 (2026 Q3-Q4, 즉시 착수)", fill=DIM_GRAY)
    add_rect(slide, 0.6, 1.95, 12.9, 4.65, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, 2.10, 12.3, 4.45,
        [
            (" ① 적응형 그리퍼 + 협동로봇 NG제거 PoC (오창 또는 Holland 1라인)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    예산 5~10억, 외부 SI(코윈테크·두산로보틱스), 6개월 내 가동",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" ② OTTO 1500 V2 또는 코윈테크 1.5~2t AMR fleet 5대 시범 도입",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    ESS 모듈 운반 동선 우선, 9개월 ROI 검증",
             {"size": SZ_SUB, "color": MID_GRAY}),
            (" ③ 화재 IR 검출 + AI 격리 시스템 1개 라인 시험",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
        ],
    )

    # 좌하 Phase 2/3
    add_label_band(slide, 0.6, 6.95, 12.9, "Phase 2 (2027 H1) · Phase 3 (2027 H2~2028)", fill=DIM_GRAY)
    add_rect(slide, 0.6, 7.65, 12.9, 5.0, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.9, 7.80, 12.3, 4.8,
        [
            ("[Phase 2]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" ④ AI 비전 0.3mm 정밀도 전 라인 표준화",
             {"size": SZ_BODY, "color": BLACK}),
            (" ⑤ 화성/에이징 ASRS 모듈러 RFI (Cimcorp/Chroma/SFA)",
             {"size": SZ_BODY, "color": BLACK}),
            (" ⑥ 디지털 트윈 파일럿 1라인 (Siemens 또는 NVIDIA)",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 5}),
            ("[Phase 3]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" ⑦ CTC 통합 라인 전략 검토 (Holland 2nd, JF1/JH4 셀 기반)",
             {"size": SZ_BODY, "color": BLACK}),
            (" ⑧ 2톤급 AMR 전사 표준화",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # 우상 전환 기준
    add_label_band(slide, 14.0, 1.25, 12.9, "전환 기준 (Trigger)", fill=DIM_GRAY)
    add_rect(slide, 14.0, 1.95, 12.9, 4.65, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 14.3, 2.10, 12.3, 4.45,
        [
            (" • BYD 2710Ah 1GWh+ 글로벌 인도 실적 확인",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    → CTC 전략 가속",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 미국 ESS 60GWh 출하 달성",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    → Holland 2단계 retool 의사결정",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • KFI/UL 9540A에서 1000Ah+ 셀 별도 기준 마련",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            ("    → 화재 안전 자동화 표준 재정렬",
             {"size": SZ_SUB, "color": MID_GRAY}),
        ],
    )

    # 우하 Caveats
    add_label_band(slide, 14.0, 6.95, 12.9, "Caveats (전제·한계)", fill=DIM_GRAY)
    add_rect(slide, 14.0, 7.65, 12.9, 5.0, fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 14.3, 7.80, 12.3, 4.8,
        [
            (" • BYD 2710Ah 셀 단품 치수·중량 공식 미공개 — 추정치는 시스템 단위 역산값(±20%) → 핸들링 사양화는 587/1175Ah 기준선 우선",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 일부 셀 단중·치수는 제3자 소매상 자료, 제조사 데이터시트와 차이 가능(EVE MB56 10.7/11.5kg 혼재)",
             {"size": SZ_SUB, "color": BLACK}),
            (" • CIBF 일부 '신제품'은 2025 발표 후 양산 진입 강조 마케팅 성격, 관람객 35만→20만+ 보정 수치 존재",
             {"size": SZ_SUB, "color": BLACK}),
            (" • LGES 2026 가이던스(capex −40%·60GWh+·90GWh+)는 IR 기준, IRA·FEOC 정책 변화 시 재조정 가능",
             {"size": SZ_SUB, "color": BLACK}),
        ],
    )
    add_footnote(
        slide,
        f"* 출처 : {SRC} Recommendations·Caveats (accessed 2026-05-26)  │  "
        "Energy-Storage.News 2026/4 (2710Ah 정보 제한 명시)")


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    build_main(prs)
    build_aux1(prs)
    build_aux2(prs)
    build_aux3(prs)
    build_aux4(prs)
    build_aux5(prs)
    build_aux6(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}  (slides={len(prs.slides._sldIdLst)})")


if __name__ == "__main__":
    build()
