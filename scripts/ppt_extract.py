"""PPTX 양식의 구조를 JSON 으로 추출한다.

사용법:
    python scripts/ppt_extract.py <input.pptx> [--out <output.json>]

출력 JSON 스키마는 README "PPT 양식 자동 채움" 섹션 참조.
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import os
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.stderr.write(
        "python-pptx 가 설치되어 있지 않습니다. 먼저 다음을 실행하세요:\n"
        "    pip install -r requirements.txt\n"
    )
    sys.exit(2)

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_common import (  # noqa: E402
    find_explicit_placeholders,
    is_blank,
    normalize_text,
    shape_id,
)

SCHEMA_VERSION = "1.0"

EMPHASIS_COLORS = {
    "FFFF00",
    "FFFF99",
    "FFC000",
    "00B050",
    "92D050",
    "00B0F0",
    "FF0000",
}


def _emu(value) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _font_color_rgb(font) -> str | None:
    try:
        rgb = font.color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    try:
        return str(rgb).upper()
    except Exception:
        return None


def _detect_emphasis(runs: list[dict[str, Any]]) -> str | None:
    for r in runs:
        c = r.get("font_color_rgb")
        if c and c.upper() in EMPHASIS_COLORS:
            return c.upper()
    return None


def _extract_runs(text_frame) -> list[dict[str, Any]]:
    runs = []
    for p_idx, paragraph in enumerate(text_frame.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            font = run.font
            runs.append(
                {
                    "paragraph_index": p_idx,
                    "run_index": r_idx,
                    "text": run.text,
                    "font_name": font.name,
                    "font_size_pt": font.size.pt if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                    "font_color_rgb": _font_color_rgb(font),
                }
            )
    return runs


def _extract_table(shape) -> dict[str, Any]:
    table = shape.table
    rows = list(table.rows)
    cols = list(table.columns)
    n_rows = len(rows)
    n_cols = len(cols)
    cells: list[list[dict[str, Any]]] = []
    for r in range(n_rows):
        row_cells = []
        for c in range(n_cols):
            cell = table.cell(r, c)
            text = normalize_text(cell.text)
            cell_runs = _extract_runs(cell.text_frame)
            row_cells.append(
                {
                    "row": r,
                    "col": c,
                    "text": text,
                    "is_empty": is_blank(text),
                    "explicit_placeholders": find_explicit_placeholders(text),
                    "runs": cell_runs,
                    "highlight_color_rgb": _detect_emphasis(cell_runs),
                }
            )
        cells.append(row_cells)
    header = [c["text"] for c in cells[0]] if cells else []
    return {
        "rows": n_rows,
        "cols": n_cols,
        "header": header,
        "cells": cells,
    }


def _extract_shape(shape, slide_idx: int, shape_idx: int) -> dict[str, Any]:
    sid = shape_id(slide_idx, shape_idx)
    shape_type = getattr(shape, "shape_type", None)
    shape_type_name = (
        shape_type.name if hasattr(shape_type, "name") else str(shape_type)
    )

    is_placeholder = bool(getattr(shape, "is_placeholder", False))
    placeholder_type = None
    if is_placeholder:
        try:
            placeholder_type = shape.placeholder_format.type.name
        except Exception:
            placeholder_type = None

    is_table = shape.has_table if hasattr(shape, "has_table") else False
    is_chart = shape.has_chart if hasattr(shape, "has_chart") else False
    is_image = shape_type == MSO_SHAPE_TYPE.PICTURE
    is_group = shape_type == MSO_SHAPE_TYPE.GROUP
    is_smartart = False
    try:
        if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            pass
        if "smartArt" in shape._element.xml:
            is_smartart = True
    except Exception:
        pass

    text = ""
    runs: list[dict[str, Any]] = []
    if shape.has_text_frame:
        text = normalize_text(shape.text_frame.text)
        runs = _extract_runs(shape.text_frame)

    table_data = _extract_table(shape) if is_table else None
    explicit = find_explicit_placeholders(text)
    if table_data:
        for row in table_data["cells"]:
            for cell in row:
                explicit.extend(cell["explicit_placeholders"])

    position: dict[str, Any] = {}
    for attr in ("left", "top", "width", "height"):
        position[f"{attr}_emu"] = _emu(getattr(shape, attr, None))

    highlight = _detect_emphasis(runs)
    if not highlight and table_data:
        for row in table_data["cells"]:
            for cell in row:
                if cell.get("highlight_color_rgb"):
                    highlight = cell["highlight_color_rgb"]
                    break
            if highlight:
                break

    return {
        "shape_id": sid,
        "shape_index": shape_idx,
        "shape_type": shape_type_name,
        "name": getattr(shape, "name", None),
        "is_placeholder": is_placeholder,
        "placeholder_type": placeholder_type,
        "position": position,
        "text": text,
        "runs": runs,
        "explicit_placeholders": explicit,
        "is_empty": is_blank(text) and not is_table,
        "is_table": is_table,
        "is_chart": is_chart,
        "is_image": is_image,
        "is_group": is_group,
        "is_smartart": is_smartart,
        "highlight_color_rgb": highlight,
        "table": table_data,
    }


def extract(path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    slides_data: list[dict[str, Any]] = []
    explicit_keys: set[str] = set()
    shapes_total = 0
    shapes_empty = 0
    tables_total = 0
    images_total = 0
    charts_total = 0
    smartart_total = 0

    emphasis_runs_total = 0

    for s_idx, slide in enumerate(prs.slides):
        shapes_data = []
        for sh_idx, shape in enumerate(slide.shapes):
            data = _extract_shape(shape, s_idx, sh_idx)
            shapes_data.append(data)
            shapes_total += 1
            if data["is_empty"]:
                shapes_empty += 1
            if data["is_table"]:
                tables_total += 1
            if data["is_image"]:
                images_total += 1
            if data["is_chart"]:
                charts_total += 1
            if data["is_smartart"]:
                smartart_total += 1
            for key in data["explicit_placeholders"]:
                explicit_keys.add(key)
            for r in data["runs"]:
                c = r.get("font_color_rgb")
                if c and c in EMPHASIS_COLORS:
                    emphasis_runs_total += 1
            if data["table"]:
                for row in data["table"]["cells"]:
                    for cell in row:
                        for r in cell.get("runs", []):
                            c = r.get("font_color_rgb")
                            if c and c in EMPHASIS_COLORS:
                                emphasis_runs_total += 1

        layout_name = None
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        slides_data.append(
            {
                "slide_index": s_idx,
                "slide_id": getattr(slide, "slide_id", None),
                "layout_name": layout_name,
                "shapes": shapes_data,
            }
        )

    return {
        "schema_version": SCHEMA_VERSION,
        "source_file": str(path),
        "extracted_at": dt.datetime.now().astimezone().isoformat(timespec="seconds"),
        "slide_size": {
            "width_emu": _emu(prs.slide_width),
            "height_emu": _emu(prs.slide_height),
        },
        "slides": slides_data,
        "summary": {
            "slide_count": len(slides_data),
            "explicit_placeholder_keys": sorted(explicit_keys),
            "shapes_total": shapes_total,
            "shapes_empty": shapes_empty,
            "tables_total": tables_total,
            "images_total": images_total,
            "charts_total": charts_total,
            "smartart_total": smartart_total,
            "emphasis_runs_total": emphasis_runs_total,
        },
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX structure to JSON.")
    parser.add_argument("input", help="입력 PPTX 경로")
    parser.add_argument(
        "--out",
        help="출력 JSON 경로 (생략 시 outputs/.cache/<basename>.structure.json)",
    )
    args = parser.parse_args(argv)

    in_path = Path(args.input)
    if not in_path.exists():
        sys.stderr.write(f"입력 파일을 찾을 수 없습니다: {in_path}\n")
        return 1

    if args.out:
        out_path = Path(args.out)
    else:
        cache_dir = Path("outputs") / ".cache"
        cache_dir.mkdir(parents=True, exist_ok=True)
        out_path = cache_dir / f"{in_path.stem}.structure.json"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        data = extract(in_path)
    except Exception as e:
        sys.stderr.write(f"PPTX 파싱 실패: {e}\n")
        return 2

    out_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(str(out_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
