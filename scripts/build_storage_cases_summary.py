"""보관 영역 12 케이스 브레인스토밍 회의록 요약 PPT 빌더.

회의록(`outputs/회의록_보관영역브레인스토밍_2026-05-21.md`) 의 12개
케이스 토론 결과를 4장 슬라이드로 정리.

- Slide 1 (메인): 회의 요약 + 12 케이스 평가 매트릭스 + 권장안 3트랙
- Slide 2 (보조): 사용자 시드 3 케이스 (A/B/C) 상세
- Slide 3 (보조): 추가 9 케이스 (D~L) 3x3 카드
- Slide 4 (보조): 권장안 트랙 + 액션아이템 9건

기준 빌더 ``scripts/build_flexibility_summary.py`` 의 상수·헬퍼 재사용.

출력: outputs/보관영역_케이스정리_2026-05-21.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

from build_flexibility_summary import (
    BLACK,
    BLUE,
    CHARCOAL,
    DIM_GRAY,
    GRAY,
    LIGHT_BLUE,
    LIGHT_GRAY,
    LINE_GRAY,
    MID_GRAY,
    MOLD_GREEN,
    SOFT_GRAY,
    SZ_BAND,
    SZ_BODY,
    SZ_FOOT,
    SZ_SECTION,
    SZ_SUB,
    SZ_TITLE,
    WHITE,
    add_label_band,
    add_rect,
    add_text,
    set_run,
    write_lines,
)

OUT_PATH = Path("outputs/보관영역_케이스정리_2026-05-21.pptx")


def add_title(slide, text, subtitle=None):
    """공통 타이틀 + 책임자."""
    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [(text, {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.8, 0.35, 6.4, 0.8,
        [("FA기술혁신파트 강모원책임",
          {"size": SZ_BODY, "bold": True, "color": BLUE,
           "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.35), Cm(26.9), Cm(1.35))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)


def add_footnote(slide, text):
    add_text(
        slide, 0.6, 18.30, 26.3, 0.55,
        [(text, {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )


def build_slide_1_summary(prs):
    """Slide 1 — 메인 요약."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "보관영역 12 케이스 평가 및 권장안 — 브레인스토밍 회의록(2026-05-21) 요약")

    # 핵심 메시지 박스
    add_rect(slide, 0.6, 1.55, 26.3, 2.10,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, 1.60, 26.0, 2.00,
        [
            ("[핵심 메시지 — 5 에이전트 토론 결과]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • 보관 영역 12 케이스(사용자 시드 A·B·C + 추가 9건)를 7개 항목 ◎/○/△/× 평가",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 단기(~'26) Case C (일반 SMC + 미니 DR) — Vendor risk 최저, 기술 성숙도 ◎",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 중기('27~'28) Case L (Hybrid A+C) + 보조 Case I (모듈 컨테이너 DR)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            (" • 장기('29~) Case K (셀 표면 코팅) — R&D 트랙 분리",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 배제 Case H (진공 포장) — 인라인 공수 폭증",
             {"size": SZ_BODY, "color": DIM_GRAY}),
        ],
    )

    # ── 평가 매트릭스 (13행 × 8열) ────────────────────────
    add_label_band(slide, 0.6, 3.80, 17.0,
                   "12 케이스 평가 매트릭스 (◎ 매우 우수 / ○ 양호 / △ 보통 / × 미흡)")
    tbl_shape = slide.shapes.add_table(
        13, 8, Cm(0.6), Cm(4.55), Cm(17.0), Cm(11.10)
    )
    tbl = tbl_shape.table
    col_w = [3.8, 1.9, 1.9, 1.9, 1.9, 1.9, 1.9, 1.8]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Cm(w)

    headers = ["Case", "DR 달성", "CAPEX", "면적 절감",
               "매몰 Loss", "재활용성", "기술 성숙도", "도입 시점"]
    rows = [
        ("A. SMC ME화",         "◎", "△", "◎", "◎", "○", "△", "중기'27~"),
        ("B. FOUP 고사양",      "○", "△", "○", "○", "○", "△", "중기'27~"),
        ("C. 일반 SMC+미니DR",  "◎", "○", "○", "○", "○", "◎", "단기~'26"),
        ("D. EFEM 차용",        "○", "△", "△", "○", "△", "○", "중기"),
        ("E. ASRS 국소 DR",     "◎", "×", "◎", "◎", "○", "○", "중·장기"),
        ("F. 글러브박스",       "○", "△", "◎", "○", "○", "△", "중기"),
        ("G. 흡습제",           "△", "◎", "△", "△", "△", "○", "단기"),
        ("H. 진공 포장",        "◎", "△", "○", "○", "×", "○", "× 인라인"),
        ("I. 모듈 컨테이너 DR", "○", "△", "○", "○", "◎", "○", "중기"),
        ("J. 설비 내장 Buffer", "○", "○", "◎", "△", "△", "○", "단기(부분)"),
        ("K. 셀 코팅 완화",     "△", "×", "○", "○", "△", "×", "장기'29~"),
        ("L. Hybrid (A+C)",     "◎", "△", "◎", "◎", "○", "○", "단계 도입"),
    ]
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = GRAY
        write_lines(c.text_frame, [h], size=SZ_SUB, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    # 권장 케이스 (C, L, I, K) 행은 LIGHT_BLUE 배경
    highlight_rows = {2, 8, 10, 11}  # C(인덱스 2 행=row1+1...) 실제 row index
    # rows 인덱스: 0=A,1=B,2=C,3=D,4=E,5=F,6=G,7=H,8=I,9=J,10=K,11=L
    highlight_rows = {2, 8, 10, 11}
    for ri, row in enumerate(rows):
        bg = LIGHT_BLUE if ri in highlight_rows else WHITE
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            is_first = ci == 0
            color = CHARCOAL if is_first else BLACK
            bold = is_first
            align = PP_ALIGN.LEFT if is_first else PP_ALIGN.CENTER
            size = SZ_SUB
            write_lines(c.text_frame, [val], size=size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    # ── 우측 추진 트랙 3단 ────────────────────────────────
    R_LEFT = 18.05
    R_W = 9.40

    add_label_band(slide, R_LEFT, 3.80, R_W, "추진 트랙 (단·중·장기)")

    # 단기 박스
    add_rect(slide, R_LEFT, 4.55, R_W, 3.35,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, R_LEFT + 0.25, 4.65, R_W - 0.5, 3.20,
        [
            ("[단기 (~'26)]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" Case C — 일반 SMC + 미니 DR",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • Vendor risk 최저 (기존 5사 풀)",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 754평 + α 수준으로 압축 가능",
             {"size": SZ_SUB, "color": BLACK}),
            (" • '26 ESWA 신증설 시 우선 적용",
             {"size": SZ_SUB, "color": BLACK}),
        ],
    )

    # 중기 박스
    add_rect(slide, R_LEFT, 8.05, R_W, 3.85,
             fill=LIGHT_BLUE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, R_LEFT + 0.25, 8.15, R_W - 0.5, 3.70,
        [
            ("[중기 ('27~'28)]",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" Case L — Hybrid (A + C) (메인)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • Case C 위에 Case A(SMC ME화) 적층",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 면적 KPI 1,019평/GWh(-20%) 달성",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            (" Case I — 모듈 컨테이너 DR (보조)",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 재활용률 50% KPI 친화",
             {"size": SZ_SUB, "color": BLACK}),
        ],
    )

    # 장기 박스
    add_rect(slide, R_LEFT, 12.05, R_W, 3.05,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, R_LEFT + 0.25, 12.15, R_W - 0.5, 2.90,
        [
            ("[장기 ('29~)] — R&D 트랙",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" Case K — 셀 표면 보호 코팅",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 0.3% → 0.5% RH 완화",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 별도 R&D 과제로 분리",
             {"size": SZ_SUB, "color": DIM_GRAY}),
        ],
    )

    # 배제 박스
    add_rect(slide, R_LEFT, 15.25, R_W, 1.30,
             fill=LIGHT_GRAY, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, R_LEFT + 0.25, 15.30, R_W - 0.5, 1.20,
        [
            ("[배제] Case H — 진공 포장",
             {"size": SZ_BODY, "bold": True, "color": DIM_GRAY}),
            (" • 인라인 공수 폭증으로 배제",
             {"size": SZ_SUB, "color": DIM_GRAY}),
        ],
    )

    add_footnote(slide,
        "* 출처: outputs/회의록_보관영역브레인스토밍_2026-05-21.md │ "
        "references/26FA KPI.md [211-220, 354-432행] │ "
        "references/roadmap/2026_FA기술담당_중장기로드맵_v2.md [38-45, 67-69행] │ "
        "references/26FA Vendor [25-32행] │ "
        "외부 비교 일부는 tech-research-teammate 후속 보강 예정")


def build_slide_2_seed_cases(prs):
    """Slide 2 — 시드 3 케이스 (A/B/C) 상세."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "사용자 시드 3 케이스 상세 (A · B · C)")

    # 3 column
    C_TOP = 1.55
    C_H = 16.50
    C_W = 8.64
    cols = [
        (0.60, "Case A — SMC 자체 ME(Mini-Environment) DR화", [
            ("[정의]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" Cube Rack 내부 자체가 0.3% DR 자체 유지",
             {"size": SZ_BODY, "color": BLACK}),
            (" (N2 또는 CDA 순환, 외부는 일반 환경)",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            ("", {"size": 4}),
            ("[정량 근거 — data-teammate]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • SMC Demo 11.72억 (14.66→-20%)",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 활성화 Rack 4.97억 (-5%)",
             {"size": SZ_SUB, "color": BLACK}),
            (" • SMC 효과 면적+25%·공수-33%·-15%",
             {"size": SZ_SUB, "color": BLACK}),
            (" • DR 4,050평 Rack 內 흡수 시 공조",
             {"size": SZ_SUB, "color": BLACK}),
            ("   CAPEX·OPEX 절감 압도적",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            ("[구성 — ops-teammate]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 전극 SMC (LG PRI): Rack 10대 + AMR",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 활성화 SMC (LG CNS): Rack 96대 +",
             {"size": SZ_SUB, "color": BLACK}),
            ("   Lifter + Shuttle",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 추가 필요: Door 가스켓·N2 매니폴드·",
             {"size": SZ_SUB, "color": BLACK}),
            ("   차압 센서",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            ("[강점/약점]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" + 면적 절감 효과 최대",
             {"size": SZ_SUB, "color": BLACK}),
            (" + 매몰 Loss·SMC 효과 위 적층",
             {"size": SZ_SUB, "color": BLACK}),
            (" − Rack 단가 상승(미정량)",
             {"size": SZ_SUB, "color": BLACK}),
            (" − Door 개폐 RH 회복 시간 미검증",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            ("실현 가능성 : 중기('27~)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
        (9.44, "Case B — FOUP 고사양화 (2~3일 장기 보관)", [
            ("[정의]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" 1~3h 단기 → 2~3일 장기 보관 가능",
             {"size": SZ_BODY, "color": BLACK}),
            (" 고사양 FOUP 개발 (N2 Purge 상시·",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            ("  고밀폐 가스켓·내장 흡습)",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            ("", {"size": 4}),
            ("[정량 근거 — data-teammate]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 2~3일 보관 가능 시 활성화 Rack 수",
             {"size": SZ_SUB, "color": BLACK}),
            ("   (현행 96대) 자체 축소 가능",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 면적 KPI 1,019평/GWh (-20%)",
             {"size": SZ_SUB, "color": BLACK}),
            ("   에 직접 기여",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            ("[기준 사양 — ops-teammate]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 3S 510×515 8슬롯 P-FOUP",
             {"size": SZ_SUB, "color": BLACK}),
            (" • DOOR GASKET·AL 합금·N2 Purge",
             {"size": SZ_SUB, "color": BLACK}),
            ("   Filter Option 가능 공간 확보",
             {"size": SZ_SUB, "color": BLACK}),
            (" • MOQ/단가/납기: 자료 미확인",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            ("", {"size": 4}),
            ("[강점/약점]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" + 시간 분산으로 라인 동기화 부담 ↓",
             {"size": SZ_SUB, "color": BLACK}),
            (" + 캐리어 단독 교체로 라인 영향 최소",
             {"size": SZ_SUB, "color": BLACK}),
            (" − 셀 장기 N2 Purge 안정성 미검증",
             {"size": SZ_SUB, "color": BLACK}),
            (" − FOUP 단가 SMC Rack 단가 누름 가능",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            ("실현 가능성 : 중기('27~) R&D 동반",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
        (18.28, "Case C — 작은 공간에 일반 SMC + 미니 DR ★ 단기 권장", [
            ("[정의]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" Rack/캐리어 사양 그대로 + 컴팩트",
             {"size": SZ_BODY, "color": BLACK}),
            (" 미니 DR 룸 (수십~수백평) 조합",
             {"size": SZ_BODY, "color": BLACK}),
            (" 4,050평 → 수백평 압축",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            ("", {"size": 4}),
            ("[정량 근거 — data-teammate]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • ESWA 1동 DR 4,050평 내역:",
             {"size": SZ_SUB, "color": BLACK}),
            ("   0.5% 공조 +330평",
             {"size": SZ_SUB, "color": BLACK}),
            ("   점보롤/팬케이크 +754평",
             {"size": SZ_SUB, "color": BLACK}),
            ("   기타 (46120·UT Matrix 등) +120평",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 미니 DR 도입 시 754평 + α 수준",
             {"size": SZ_SUB, "bold": True, "color": BLUE}),
            (" • DR 평당 CAPEX 단가: 자료 미확인",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            ("", {"size": 4}),
            ("[Vendor — ops-teammate]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 기존 SMC 협력사 풀 그대로 활용",
             {"size": SZ_SUB, "color": BLACK}),
            ("   LG CNS·코윈테크·아바코·시너스텍·SFA",
             {"size": SZ_SUB, "color": DIM_GRAY}),
            (" • Vendor risk 최저",
             {"size": SZ_SUB, "bold": True, "color": BLUE}),
            ("", {"size": 4}),
            ("[강점/약점]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" + 기술 성숙도 ◎ (기존 SMC 그대로)",
             {"size": SZ_SUB, "color": BLACK}),
            (" + 단기 도입 용이",
             {"size": SZ_SUB, "color": BLACK}),
            (" − Case A 대비 면적 절감 ↓",
             {"size": SZ_SUB, "color": BLACK}),
            (" − 미니 DR 자체 CAPEX 상존",
             {"size": SZ_SUB, "color": BLACK}),
            ("", {"size": 4}),
            ("실현 가능성 : 단기 (~'26) — 즉시 적용",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
    ]
    for left, head, lines in cols:
        # 헤더 띠 — Case C 만 파랑 강조
        is_recommended = head.startswith("Case C")
        band_color = BLUE if is_recommended else DIM_GRAY
        add_label_band(slide, left, C_TOP, C_W, head, fill=band_color)
        add_rect(slide, left, C_TOP + 0.7, C_W, C_H - 0.7,
                 fill=WHITE, line=MID_GRAY, line_w=0.75)
        add_text(slide, left + 0.20, C_TOP + 0.85,
                 C_W - 0.40, C_H - 1.00, lines)

    add_footnote(slide,
        "* 출처: outputs/회의록_보관영역브레인스토밍_2026-05-21.md [37~88행] │ "
        "references/26FA KPI.md [211-220, 354-432행] │ "
        "references/roadmap/2026_FA기술담당_중장기로드맵_v2.md [38-45, 67-69행] │ "
        "references/기술자료/FOUP(510x515)_3S.md [38-46, 274, 307-308행] │ "
        "references/26FA Vendor [25-32행]")


def build_slide_3_extra_cases(prs):
    """Slide 3 — 추가 9 케이스 (D~L) 3x3 카드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "추가 9 케이스 (D~L) — 브레인스토밍 카드 요약")

    cards = [
        ("D. EFEM 미니환경 차용", "△", [
            "반도체 Load Port + Mini-env + FFU",
            "EFEM 차용 (로보스타 FoPLP)",
            "셀 사이즈/중량 매칭 변수",
            "+ 반도체 양산 검증",
            "− 셀 특성 비매칭 가능성",
        ], False),
        ("E. ASRS 국소 DR", "○", [
            "Stacker Crane/Shuttle ASRS",
            "챔버 자체를 DR 화",
            "활성화 SMC 진화형",
            "+ 면적·인력 동시 절감",
            "− 챔버 단위 CAPEX 큼",
        ], False),
        ("F. N2 충전 글러브박스", "△", [
            "진공/N2 충전 후",
            "일반 환경 보관 컨테이너",
            "재충전 사이클 관리 변수",
            "+ 면적 절감 ◎",
            "− 운영 공수 ↑",
        ], False),
        ("G. 데시컨트 흡습제", "△", [
            "캐리어 내 흡습제 카트리지",
            "수동/반자동 DR 유지",
            "0.3% 달성 △",
            "+ CAPEX ◎",
            "− OPEX 교체 공수",
        ], False),
        ("H. 셀 진공 포장", "× 배제", [
            "셀 단위 진공 Pouch 포장",
            "일반 환경 보관",
            "포장·재포장 공수 발생",
            "+ 환경 무관 ◎",
            "− 인라인 공수 폭증",
        ], False),
        ("I. 모듈 컨테이너 DR", "○ 중기 보조", [
            "표준 컨테이너 외관",
            "내부 DR 보관·이송",
            "재활용률 50% KPI 친화",
            "+ 재활용성·유연성 ◎",
            "− 컨테이너 단가 부담",
        ], True),
        ("J. 설비 내장 Buffer", "△", [
            "공정 설비 안 소형 buffer",
            "라인 외 보관 不요",
            "처리량 제한·백업 부재",
            "+ 면적 ◎",
            "− 처리량 ↓",
        ], False),
        ("K. 셀 표면 코팅", "○ 장기 R&D", [
            "셀 표면 보호 코팅",
            "DR 0.3% → 0.5% 완화",
            "R&D 트랙 분리",
            "+ 환경 부담 ↓",
            "− R&D 장기·재료비 ↑",
        ], True),
        ("L. Hybrid (A+C)", "◎ 중기 메인", [
            "ME화 SMC (A) + 미니 DR (C)",
            "이송 = AMR + 단기 캐리어",
            "단계 도입 자연스러움",
            "+ 위험 분산 ◎",
            "+ C→A→L 전이 시나리오",
        ], True),
    ]

    # 3x3 그리드
    G_TOP = 1.55
    G_LEFT = 0.60
    G_W = 8.64
    G_H = 5.30
    G_GAP_Y = 0.10

    for i, (head, grade, lines, is_recommended) in enumerate(cards):
        row = i // 3
        col = i % 3
        left = G_LEFT + col * (G_W + 0.12)
        top = G_TOP + row * (G_H + G_GAP_Y)

        # 권장안은 LIGHT_BLUE 배경, 배제는 LIGHT_GRAY
        if is_recommended:
            bg = LIGHT_BLUE
            border = BLUE
        elif "배제" in grade:
            bg = LIGHT_GRAY
            border = MID_GRAY
        else:
            bg = WHITE
            border = MID_GRAY
        add_rect(slide, left, top, G_W, G_H,
                 fill=bg, line=border, line_w=0.75)

        # 헤더 (제목 + 평가)
        add_text(
            slide, left + 0.2, top + 0.10, G_W - 0.4, 0.70,
            [
                (head, {"size": SZ_BODY, "bold": True, "color": BLACK,
                         "align": PP_ALIGN.LEFT}),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide, left + 0.2, top + 0.10, G_W - 0.4, 0.70,
            [
                (grade, {"size": SZ_SUB, "bold": True,
                          "color": BLUE if is_recommended else DIM_GRAY,
                          "align": PP_ALIGN.RIGHT}),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # 본문
        body_lines = [(" • " + ln, {"size": SZ_SUB, "color": BLACK})
                      for ln in lines]
        add_text(
            slide, left + 0.25, top + 0.85, G_W - 0.50, G_H - 1.00,
            body_lines,
        )

    add_footnote(slide,
        "* 출처: outputs/회의록_보관영역브레인스토밍_2026-05-21.md [91~136행] │ "
        "references/기술자료/EFEM_로보스타.md [32-69, 89-101행] │ "
        "references/roadmap/2026_FA기술담당_중장기로드맵_v2.md [38-45행] │ "
        "Case E·H 외부 사례는 tech-research-teammate 후속 보강 예정")


def build_slide_4_actions(prs):
    """Slide 4 — 권장안 정리 + 액션아이템 9건."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "권장안 3 트랙 정리 및 후속 액션아이템 (~'26.06)")

    # 좌측: 권장안 트랙 박스
    L_LEFT = 0.60
    L_W = 13.0
    add_label_band(slide, L_LEFT, 1.55, L_W, "권장안 3 트랙 + 배제 1")

    # 단기
    add_rect(slide, L_LEFT, 2.30, L_W, 3.20,
             fill=LIGHT_BLUE, line=BLUE, line_w=1.0)
    add_text(
        slide, L_LEFT + 0.25, 2.40, L_W - 0.5, 3.05,
        [
            ("[단기 (~'26)] ★ Case C — 일반 SMC + 미니 DR",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • Vendor risk 최저 (기존 SMC 협력사 5사 풀)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • DR 4,050평 → 754평 + α 수준 압축 가능",
             {"size": SZ_BODY, "color": BLACK}),
            (" • '26년 내 ESWA 신증설 시 우선 적용",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 미해결 : DR 평당 CAPEX 단가 → data-teammate 후속 산출",
             {"size": SZ_SUB, "color": DIM_GRAY}),
        ],
    )

    # 중기
    add_rect(slide, L_LEFT, 5.70, L_W, 4.10,
             fill=LIGHT_BLUE, line=BLUE, line_w=1.0)
    add_text(
        slide, L_LEFT + 0.25, 5.80, L_W - 0.5, 3.95,
        [
            ("[중기 ('27~'28)] ★ Case L — Hybrid (A+C) 메인 / Case I 보조",
             {"size": SZ_SECTION, "bold": True, "color": BLUE}),
            (" • Case C 위에 Case A (SMC ME화) 점진 적층",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 면적 1,019평/GWh (-20%) KPI 달성 경로와 정합",
             {"size": SZ_BODY, "color": BLACK}),
            (" • C → A → L 전이 시나리오 자연스러움",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 보조 Case I (모듈 컨테이너 DR) : 재활용률 50% KPI 친화",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 미해결 : Rack 내부 N2 Purge OPEX, Door 개폐 RH 회복 시간",
             {"size": SZ_SUB, "color": DIM_GRAY}),
        ],
    )

    # 장기
    add_rect(slide, L_LEFT, 10.00, L_W, 2.55,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, L_LEFT + 0.25, 10.10, L_W - 0.5, 2.40,
        [
            ("[장기 ('29~)] Case K — 셀 표면 보호 코팅 R&D",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 0.3% → 0.5% 완화로 환경 부담 자체 ↓ 근본 해법",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 셀 설계·재료 R&D 필요 → 별도 과제 등록 검토",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # 배제
    add_rect(slide, L_LEFT, 12.75, L_W, 1.30,
             fill=LIGHT_GRAY, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, L_LEFT + 0.25, 12.85, L_W - 0.5, 1.10,
        [
            ("[배제] Case H — 셀 진공 포장 (인라인 공수 폭증)",
             {"size": SZ_BODY, "bold": True, "color": DIM_GRAY}),
        ],
    )

    # 차기 회의
    add_rect(slide, L_LEFT, 14.30, L_W, 1.85,
             fill=SOFT_GRAY, line=MID_GRAY, line_w=0.5)
    add_text(
        slide, L_LEFT + 0.25, 14.40, L_W - 0.5, 1.70,
        [
            ("[차기 회의 '26.06 초]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • data-teammate DR 평당 CAPEX 단가 + Case C/L 정량 비교 리뷰",
             {"size": SZ_SUB, "color": BLACK}),
            (" • tech-research-teammate 외부 사례 보강 결과 리뷰",
             {"size": SZ_SUB, "color": BLACK}),
            (" • 권장안 3개 최종 확정 → MRM 후속 자료 반영",
             {"size": SZ_SUB, "color": BLACK}),
        ],
    )

    # 우측: 액션아이템 표
    R_LEFT = 14.00
    R_W = 13.0
    add_label_band(slide, R_LEFT, 1.55, R_W, "액션아이템 9건 (~'26.06.11)")

    tbl_shape = slide.shapes.add_table(
        10, 3, Cm(R_LEFT), Cm(2.30), Cm(R_W), Cm(13.85)
    )
    tbl = tbl_shape.table
    col_w = [3.4, 7.5, 2.1]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Cm(w)
    headers = ["담당", "내용", "기한"]
    rows = [
        ("data-teammate",
         "DR 평당 CAPEX 단가 + DR Less 시뮬레이션 (Case C/L 정량 비교)",
         "'26.06.04"),
        ("data-teammate",
         "Case A SMC ME화 시 Rack 단가 상승 + Door 개폐 RH 회복 시간 시뮬레이션",
         "'26.06.11"),
        ("ops-teammate",
         "3S P-FOUP MOQ/단가/납기 협력사 확인 (사양서 보강)",
         "'26.05.28"),
        ("ops-teammate",
         "SMC ME화 사양 RFI — Logistic 5사 (LG CNS·코윈테크·아바코·시너스텍·SFA)",
         "'26.06.04"),
        ("tech-research-teammate",
         "EFEM·ASRS·Stocker 외부 사례 (SEMI E47/E84/E87, TSMC/Samsung)",
         "'26.05.28"),
        ("tech-research-teammate",
         "CATL/BYD/SDI 셀 단계 DR 운영 데이터 (references/경쟁사/ 갱신)",
         "'26.06.04"),
        ("ppt-writer",
         "CEO 보조자료 PPTX 갱신 — Stock 영역 단기 C / 중기 L 2단 트랙",
         "'26.05.27"),
        ("document-writer",
         "본 회의록 배포 + 차기 회의록 양식 준비",
         "'26.05.22"),
        ("강모원책임 (의장)",
         "Case K (셀 코팅) R&D 트랙 분리 — 별도 과제 등록 검토",
         "'26.06.11"),
    ]
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = GRAY
        write_lines(c.text_frame, [h], size=SZ_BODY, bold=True,
                    color=BLACK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            bold = (ci == 0)
            color = CHARCOAL if ci == 0 else BLACK
            align = PP_ALIGN.LEFT
            size = SZ_SUB
            write_lines(c.text_frame, [val], size=size, bold=bold,
                        color=color, align=align,
                        anchor=MSO_ANCHOR.MIDDLE)

    add_footnote(slide,
        "* 출처: outputs/회의록_보관영역브레인스토밍_2026-05-21.md [186~217행] │ "
        "회의 결정 9개 항목 + 차기 회의 '26.06 초 확정")


def build():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)

    build_slide_1_summary(prs)
    build_slide_2_seed_cases(prs)
    build_slide_3_extra_cases(prs)
    build_slide_4_actions(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
