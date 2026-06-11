"""FA기술담당 로봇 소모임 커리큘럼·교부재 구매안 보고 PPT 빌더.

사내 분류형 레이아웃(scripts/_ppt_layout.py)을 사용한다.
- 슬라이드 1 : 메인(분류형 임원 요약) — 개요 / 24주 커리큘럼 / FA 캡스톤 / 교부재·예산
- 슬라이드 2 : 보조 — 24주 커리큘럼 로드맵(표)
- 슬라이드 3 : 보조 — FA 응용 캡스톤 3과제(표)
- 슬라이드 4 : 보조 — 교부재(HW) 구매안(표)
- 슬라이드 5 : 보조 — 학습연계 SW + 개략 예산(표)

베이스 강의: 패스트캠퍼스 ROS 2 & AI 자율주행 로봇 개발 올인원.
출력: outputs/로봇소모임_커리큘럼_2026-06-05.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _ppt_layout import (  # noqa: E402
    render_classified_slide, title_block, add_text, add_footnote, fill_table,
    SLIDE_W_CM, SLIDE_H_CM,
    BLACK, BLUE, MID_GRAY, MOLD_GREEN,
    SZ_HEAD, SZ_SECTION, SZ_BODY, SZ_SUB,
)
from pptx.enum.text import MSO_ANCHOR  # noqa: E402

OUT_PATH = Path("outputs/로봇소모임_커리큘럼_2026-06-05.pptx")
BYLINE = "FA기술담당 로봇 소모임 (2026.06.05)"
FOOT = ("* 베이스 강의: 패스트캠퍼스 ROS 2 & AI 자율주행 로봇 개발 올인원 / "
        "로드맵 연계: FA기술담당 중장기 로드맵 v2 / 단가는 시장 개략가(추정), 정식 견적 별도")


def table_slide(prs, *, title, head, headers, rows, col_w, row_h=1.05,
                foot=FOOT, body_size=SZ_SUB, top=3.30, second=None):
    """제목 + 헤드메시지 + 전폭 표(최대 2개) 보조 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_block(slide, title, byline=BYLINE)
    add_text(slide, 0.6, 1.50, 26.3, 1.4,
             [(head, {"size": SZ_HEAD, "bold": True, "color": BLACK})],
             anchor=MSO_ANCHOR.TOP)

    def _draw(headers, rows, col_w, t):
        nrows = len(rows) + 1
        tw = sum(col_w)
        shape = slide.shapes.add_table(
            nrows, len(headers), Cm(0.6), Cm(t), Cm(tw), Cm(nrows * row_h))
        fill_table(shape.table, headers, rows, col_w=col_w, body_size=body_size)
        for r in range(nrows):
            shape.table.rows[r].height = Cm(row_h)
        return t + nrows * row_h

    next_top = _draw(headers, rows, col_w, top)
    if second:
        add_text(slide, 0.6, next_top + 0.25, 26.3, 0.6,
                 [(second["head"], {"size": SZ_SECTION, "bold": True,
                                    "color": BLACK})])
        _draw(second["headers"], second["rows"], second["col_w"],
              next_top + 0.95)
    add_footnote(slide, foot)
    return slide


def build_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_classified_slide(
        slide,
        title="FA기술담당 로봇 소모임 커리큘럼·교부재 구매안",
        byline=BYLINE,
        head_message=[
            ("외부 강의로 ROS 2·자율주행·비전 AI 기본기를 다진 뒤, FA 실무 직결 "
             "응용 로직(랙 내 사람감지·전극롤 풀림·AMR 회피)을 캡스톤으로 자체 개발",
             {"size": SZ_HEAD, "bold": True, "color": BLACK}),
        ],
        sections=[
            {
                "label": "1.\n추진 개요",
                "session_message": "10명 · 6개월(주1회 2h·24주) · 2~3인 1조 팀 프로젝트",
                "content": [
                    (" • 목적 — ROS 2·자율주행·AI 비전 내부 인력 풀 양성, FA 과제 직결",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 베이스 — 패스트캠퍼스 ROS 2 & AI 올인원(6파트·104클립·약 29h)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • 산출 — 캡스톤 데모 → 로드맵 차기 과제 제안으로 연결",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "2.\n24주 커리큘럼",
                "session_message": "조립·구동 → ROS 2 → 자율주행 센서 → SLAM·Nav2 → AI/비전 → 캡스톤",
                "content": [
                    (" • Phase 1~2(W2~9) HW 조립·구동 + ROS 2 기초(Node/Topic/Action)",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • Phase 3~5(W10~19) LiDAR·카메라·센서융합 → SLAM·Nav2 → YOLO11·MediaPipe",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • Phase 6~7(W20~24) FA 응용 캡스톤 + 성과 공유 (게이트 4회 평가)",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
            },
            {
                "label": "3.\nFA 응용\n캡스톤",
                "session_message": "로드맵 트랙과 직접 연결되는 자체 응용 3과제",
                "content": [
                    (" Ⓐ Vision 랙 내부 사람 여부 확인 → 진입 인터락 (Physical AI·인터락)",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" Ⓑ 전극롤 풀림 감지 → 비전 이상감지 (Roll 대응·이상감지)",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" Ⓒ AMR 회피 기동 → Nav2 동적 장애물 회피 (AMR 3D Vision Nav)",
                     {"size": SZ_BODY, "bold": True, "color": BLUE}),
                    (" • 심화: 로봇암·드론·개로봇 연계 미니 트랙(Physical AI 체험)",
                     {"size": SZ_SUB, "color": MID_GRAY}),
                ],
            },
            {
                "label": "4.\n교부재·예산",
                "session_message": "교부재 HW 추가 구매 + 학습연계 SW(라이선스 검토) · 개략 예산(추정)",
                "content": [
                    (" • 로봇 추가 4대(현1→5대)·뎁스카메라·로봇암·드론·개로봇·Edge·워크스테이션",
                     {"size": SZ_BODY, "color": BLACK}),
                    (" • SW 대부분 무료(ROS2/Ubuntu/Docker), YOLO 상용·SaaS·Omniverse 검토",
                     {"size": SZ_BODY, "color": BLACK}),
                ],
                "table": {
                    "headers": ["묶음", "구성", "개략 합계(추정)"],
                    "rows": [
                        ("핵심 로봇", "모바일 로봇 추가·RPi·뎁스카메라·LiDAR", "약 350~450만원"),
                        ("확장 교부재", "로봇암·드론·개로봇·Edge·워크스테이션", "약 1,500~2,800만원"),
                        ("부자재·인프라", "테스트베드(랙·전극롤 목업)·네트워크", "약 200~400만원"),
                        ("총계(추정)", "SW 유상분 별도 산정", "약 2,050~3,650만원"),
                    ],
                    "col_w": [4.2, 13.6, 6.0],
                    "row_h": 0.82,
                    "body_size": SZ_SUB,
                },
            },
        ],
        footnote=FOOT,
    )


def build_curriculum(prs):
    table_slide(
        prs,
        title="[보조1] 24주 커리큘럼 로드맵",
        head="베이스 강의 6개 파트를 7개 Phase로 재구성 — 매주 이론 30분 + 실습 80분 + 회고 10분",
        headers=["Phase", "주차", "핵심 내용", "게이트/산출"],
        rows=[
            ("0. 온보딩", "W1", "OT·조 편성(4~5조)·개발환경(Ubuntu+ROS2 Humble+Docker)·Git", "환경 구축 완료"),
            ("1. HW 조립·구동", "W2~4", "로봇 조립, RPi5 세팅, 모터·엔코더·LiDAR·카메라 기초 구동", "① 조별 로봇 가동"),
            ("2. ROS 2 기초", "W5~9", "Package/Node/Launch, Topic·Service·Action, 디버깅(rqt·bag·tf)", "② 텔레옵+로깅 과제"),
            ("3. 자율주행 센서", "W10~13", "LiDAR 회피/Follower, 카메라 색상·QR·Line·Tracking, 센서융합(IMU·Odom)", "센서 실습 완료"),
            ("4. SLAM·Nav2", "W14~16", "URDF·TF, SLAM Toolbox 매핑, AMCL, Nav2 완전 자율주행", "③ 맵+자율주행"),
            ("5. AI·비전", "W17~19", "MediaPipe 제스처·포즈, RPi5+YOLO11(Det/Seg/Pose), 데이터셋 구축·라벨링", "비전 파이프라인"),
            ("6. FA 캡스톤", "W20~23", "랙 사람감지·전극롤 풀림·AMR 회피 3과제 + 확장 교부재 미니 트랙", "④ 캡스톤 데모"),
            ("7. 성과 공유", "W24", "조별 발표·정량 지표 리뷰, 로드맵 차기 과제 도출", "차기 과제 제안"),
        ],
        col_w=[4.0, 2.6, 13.9, 5.8],
        row_h=1.55,
        body_size=SZ_SUB,
    )


def build_capstone(prs):
    table_slide(
        prs,
        title="[보조2] FA 응용 캡스톤 3과제",
        head="외부 강의 비전·Nav2 역량을 FA 실무 과제로 전이 — 모두 중장기 로드맵 트랙과 직접 연결",
        headers=["과제", "문제 정의", "핵심 기술", "로드맵 연계"],
        rows=[
            ("Ⓐ 랙 내부\n사람 여부 확인",
             "랙·설비 진입 전 작업자 잔류 여부를 비전 판정 → AMR·크레인 기동 인터락",
             "YOLO11 person + 랙 ROI zone 로직 + 상태머신(점유/비점유/불확실) + 뎁스 거리 게이팅 + 인터락 토픽",
             "Physical AI 설비 사전검증 'Vision·인터락·이상동작 제거'"),
            ("Ⓑ 전극롤\n풀림 감지",
             "전극롤 권취 풀림·끝단 들뜸·텐션 이상을 비전으로 조기 감지·경보",
             "엣지/윤곽 + YOLO11-Seg/이상감지 모델, 롤 단면 ROI 추적, 임계 초과 경보 토픽",
             "다품종 Roll 대응 SMC / Vision 기반 이상 감지"),
            ("Ⓒ AMR\n회피 기동",
             "동적 장애물(사람·타 AMR) 출현 시 Nav2 기반 실시간 회피·재계획·경로 복귀",
             "Nav2 costmap 동적 레이어 + DWB/MPPI 튜닝 + LiDAR/뎁스 융합, 정지·우회·복귀 평가",
             "AMR 3D Vision Navigation '동적 장애물 회피'(2027)"),
            ("심화 미니 트랙",
             "확장 교부재로 Physical AI 체험 (선택)",
             "로봇암 픽플레이스(전극롤/트레이) · 드론 상공 비전 점검 · 개로봇 비정형 순찰·회피",
             "MOMA / Physical AI / Edge AI 자율운영"),
        ],
        col_w=[3.4, 7.0, 9.5, 6.4],
        row_h=2.4,
        body_size=SZ_SUB,
    )


def build_hw(prs):
    table_slide(
        prs,
        title="[보조3] 교부재(HW) 구매안 — 10명 기준",
        head="조별 전담 로봇 확보(현 1대 → 5대) + 캡스톤·심화용 확장 교부재 · 단가는 개략 추정",
        headers=["구분", "품목(모델 예시)", "수량", "용도", "개략단가(추정)"],
        rows=[
            ("핵심", "모바일 로봇 (MentorPi M1급)", "추가 4 (총5)", "2인1조×5조 전담, 전 과정", "약 40만원/대"),
            ("핵심", "입문/예비 로봇 (TurboPi급)", "2", "입문 실습·고장 대체", "약 14만원/대"),
            ("핵심", "Raspberry Pi 5 8GB + SD", "로봇당+예비2", "로봇 온보드 컴퓨팅", "약 15만원/세트"),
            ("핵심", "뎁스카메라 (RealSense D435i급)", "3", "캡스톤 A·B 거리/3D 비전", "약 40만원/대"),
            ("확장", "협동/데스크톱 로봇 암", "1~2", "전극롤·트레이 픽플레이스+비전", "약 150~700만원/대"),
            ("확장", "교육용 드론 (Tello EDU 류)", "2~3", "상공 비전·랙 상단/재고 점검", "약 20만원/대"),
            ("확장", "4족보행 개로봇 (Go2 EDU급)", "1", "Physical AI·비정형 순찰·회피", "약 400~700만원/대"),
            ("확장", "Edge AI 보드 (Jetson Orin Nano)", "2~3", "YOLO 엣지 추론", "약 50만원/대"),
            ("확장", "GPU 워크스테이션 (RTX 4090급)", "1", "YOLO 학습·Isaac Sim", "약 500~800만원"),
            ("인프라", "테스트베드·네트워크·작업도구", "1식", "랙/전극롤 목업·충전·보관", "약 200~400만원"),
        ],
        col_w=[2.4, 9.0, 3.4, 7.5, 4.0],
        row_h=1.18,
        body_size=SZ_SUB,
    )


def build_sw(prs):
    table_slide(
        prs,
        title="[보조4] 학습연계 SW + 라이선스·예산",
        head="ROS 2 생태계는 대부분 오픈소스(무료) — 단 YOLO 상용·외부 SaaS·디지털트윈은 보안/법무 검토 필요",
        headers=["SW", "유형", "비용/라이선스", "검토 포인트"],
        rows=[
            ("Ubuntu 22.04 / ROS 2 Humble", "OS·미들웨어", "무료(Apache-2.0)", "표준 개발환경 — 검토 불요"),
            ("Docker / Gazebo·Ignition", "컨테이너·시뮬", "무료", "Docker Desktop 상용 약관만 확인"),
            ("NVIDIA Isaac Sim / Isaac ROS", "Physical AI 시뮬", "무료(내부)·GPU 필수", "로드맵 Physical AI·CNS Physical Works 연계"),
            ("Ultralytics YOLO11", "비전 모델", "AGPL-3.0 / 상용 Enterprise 별도", "★ 사내 비공개·상용 시 라이선스 검토 필수"),
            ("Roboflow / CVAT", "데이터 라벨링", "SaaS 유상 / 오픈소스", "외부 SaaS 데이터 반출 보안검토 (CVAT 온프렘 대안)"),
            ("Foxglove · VS Code · Git(사내)", "시각화·IDE·형상", "무료 / 사내", "사내 계정 발급"),
        ],
        col_w=[7.2, 4.2, 7.0, 7.9],
        row_h=1.05,
        body_size=SZ_SUB,
        second={
            "head": "개략 예산 요약 (추정 · 정식 견적 별도)",
            "headers": ["묶음", "개략 합계(추정)", "우선순위"],
            "rows": [
                ("핵심 실습 로봇", "약 350~450만원", "1순위 (전 과정 필수)"),
                ("확장 교부재", "약 1,500~2,800만원", "2~3순위 (캡스톤·심화)"),
                ("부자재·인프라", "약 200~400만원", "1순위 (테스트베드)"),
                ("SW 유상분", "라이선스 검토 후 별도", "선검토 (YOLO·SaaS·Omniverse)"),
                ("총계(HW+인프라, 추정)", "약 2,050~3,650만원", "—"),
            ],
            "col_w": [6.0, 9.0, 11.3],
        },
    )


def build():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)
    build_main(prs)
    build_curriculum(prs)
    build_capstone(prs)
    build_hw(prs)
    build_sw(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    if OUT_PATH.exists():
        OUT_PATH.unlink()
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
