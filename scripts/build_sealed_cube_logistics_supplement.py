"""밀폐형 큐브 물류 CEO 보고용 보조자료 PPT 1장 빌더 (v2).

CEO 임원 보고서에 첨부할 보조자료. 본 슬라이드는 두 핵심 요소로 구성:
1) **FOUP (Front Opening Unified Pod) 개념 정리** — 반도체 표준 밀폐
   캐리어의 정의·5대 사용 이유·4대 설계 특성, FOUP 외형 이미지 자리 포함
2) **FA담당 밀폐형 큐브 물류 적용 컨셉** — Stock(SMC) / 이송(AMR Bridge)
   / 공정(국소 DR) 의 3-stack 구성 요약

보고양식 카탈로그 (`templates/사내양식/보고양식_카탈로그.md`) 의
블록 A(보고 프레임 표준) + 블록 E(전략·컨셉) + 블록 F(기술 사양) 조합.
케이스 비교(이전 v1) → 개념·컨셉 중심으로 재구성. 핵심 메시지 박스 제거.

기준 빌더 ``scripts/build_flexibility_summary.py`` 의 상수·헬퍼 재사용.

출력: outputs/밀폐형큐브물류_보조자료_2026-05-20.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

from build_flexibility_summary import (
    BLACK,
    BLUE,
    DIM_GRAY,
    LIGHT_BLUE,
    LIGHT_GRAY,
    LINE_GRAY,
    MID_GRAY,
    MOLD_GREEN,
    SOFT_GRAY,
    SZ_BAND,
    SZ_BODY,
    SZ_FOOT,
    SZ_SECTION,
    SZ_SUB,
    SZ_TITLE,
    WHITE,
    add_label_band,
    add_rect,
    add_text,
)

OUT_PATH = Path("outputs/밀폐형큐브물류_보조자료_2026-05-20.pptx")


def add_classification_band(slide, left, top, width, items):
    """블록 A '분류 1~4' 라벨 띠. DIM_GRAY 띠 + 흰색 텍스트, 셀 분리선."""
    add_rect(slide, left, top, width, 0.6, fill=DIM_GRAY, line=None)
    col_w = width / len(items)
    for i, (label, value) in enumerate(items):
        cell_left = left + i * col_w
        add_text(
            slide, cell_left + 0.1, top, col_w - 0.2, 0.6,
            [
                (label + " ", {"size": SZ_FOOT, "color": LIGHT_GRAY}),
                (value, {"size": SZ_BODY, "bold": True, "color": WHITE,
                          "align": PP_ALIGN.LEFT}),
            ],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # 셀 분리선 (좌측, 첫 셀 제외)
        if i > 0:
            sep = slide.shapes.add_connector(
                1, Cm(cell_left), Cm(top + 0.08),
                Cm(cell_left), Cm(top + 0.52),
            )
            sep.line.color.rgb = MID_GRAY
            sep.line.width = Pt(0.5)


def add_image_placeholder(slide, left, top, width, height, caption):
    """이미지 삽입 자리 placeholder — 회색 fill + 보더 + 안내 텍스트."""
    add_rect(slide, left, top, width, height,
             fill=LIGHT_GRAY, line=MID_GRAY, line_w=1.2)
    # 중앙 안내 텍스트
    add_text(
        slide, left, top, width, height,
        [
            ("[이미지 삽입 자리]",
             {"size": SZ_SECTION, "bold": True, "color": MID_GRAY,
              "align": PP_ALIGN.CENTER}),
            ("", {"size": 6}),
            (caption,
             {"size": SZ_BODY, "color": DIM_GRAY,
              "align": PP_ALIGN.CENTER}),
            ("", {"size": 6}),
            ("(PowerPoint 에서 이미지 우클릭 → 그림으로 채우기)",
             {"size": SZ_SUB, "color": MID_GRAY,
              "align": PP_ALIGN.CENTER}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def build_slide():
    prs = Presentation()
    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 상단 제목 ─────────────────────────────────────────
    add_text(
        slide, 0.6, 0.25, 20.0, 1.0,
        [("밀폐형 큐브 물류 설비 — FOUP 개념 및 적용 컨셉",
          {"size": SZ_TITLE, "bold": True, "color": BLACK})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, 20.8, 0.35, 6.4, 0.8,
        [("CEO 보조자료 / FA기술담당",
          {"size": SZ_BODY, "bold": True, "color": BLUE,
           "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_connector(1, Cm(0.6), Cm(1.30), Cm(26.9), Cm(1.30))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.75)

    # ── 분류 라벨 띠 (블록 A 표준) ─────────────────────────
    add_classification_band(
        slide, 0.6, 1.42, 26.3,
        [
            ("분류 1.", "FA기술담당"),
            ("분류 2.", "신규 과제 검토"),
            ("분류 3.", "밀폐형 Cube 물류"),
            ("분류 4.", "2028 로드맵 §3.1"),
        ],
    )

    # ── 상단 : FOUP 개념 (좌측 텍스트 + 우측 이미지) ───────
    TOP_Y = 2.20
    TOP_H = 6.70

    add_label_band(slide, 0.6, TOP_Y, 26.3, "① FOUP (Front Opening Unified Pod) — 개념 및 사용 이유")

    # 좌측 텍스트 박스
    add_rect(slide, 0.6, TOP_Y + 0.7, 13.0, TOP_H - 0.7,
             fill=WHITE, line=MID_GRAY, line_w=0.75)
    add_text(
        slide, 0.85, TOP_Y + 0.80, 12.5, TOP_H - 0.90,
        [
            ("[정의]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" • 반도체 표준 ",
             {"size": SZ_BODY, "color": BLACK}),
            ("밀폐형 캐리어",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
            ("   외부 파티클·수분·금속 오염으로부터",
             {"size": SZ_BODY, "color": BLACK}),
            ("   Wafer/Panel 표면을 보호",
             {"size": SZ_BODY, "color": BLACK}),
            (" • SEMI 표준 호환 (AMHS·OHT·RFID 자동 이송)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 3S P-FOUP 510×515mm 8-슬롯 (PLP용)",
             {"size": SZ_SUB, "color": MID_GRAY}),
            ("", {"size": 6}),
            ("[5대 사용 이유]",
             {"size": SZ_SECTION, "bold": True, "color": BLACK}),
            (" ① 청정도 확보 — 완전 밀폐 구조로 파티클·",
             {"size": SZ_BODY, "color": BLACK}),
            ("     수분·금속 오염 차단 (1차 목적)",
             {"size": SZ_BODY, "color": BLACK}),
            (" ② 자동화 대응 — AMHS·OHT 등 자동 이송",
             {"size": SZ_BODY, "color": BLACK}),
            ("     시스템과 SEMI 표준 호환",
             {"size": SZ_BODY, "color": BLACK}),
            (" ③ Wafer/Panel 보호 — SLOT 단위 고정으로",
             {"size": SZ_BODY, "color": BLACK}),
            ("     기계적 충격·진동·정전기 차단",
             {"size": SZ_BODY, "color": BLACK}),
            (" ④ 추적성 — RFID/바코드 태그로 이력 실시간",
             {"size": SZ_BODY, "color": BLACK}),
            ("     관리, 생산제어 시스템 연동",
             {"size": SZ_BODY, "color": BLACK}),
            (" ⑤ 효율 — 다수 매수 동시 이송, 가동률 향상",
             {"size": SZ_BODY, "color": BLACK}),
            ("     및 휴먼 에러 감소",
             {"size": SZ_BODY, "color": BLACK}),
        ],
    )

    # 우측 이미지 placeholder
    add_image_placeholder(
        slide,
        13.85, TOP_Y + 0.7, 13.65, TOP_H - 0.7,
        "FOUP 외형 / 도면 / 내부 구조 이미지",
    )

    # ── 중간 : FOUP 4대 설계 특성 (4-column) ──────────────
    MID_Y = 9.10
    MID_H = 3.40

    add_label_band(slide, 0.6, MID_Y, 26.3, "② FOUP 핵심 설계 특성 (4대 요소)")

    feature_cols = [
        ("① 강한 내구성", [
            "운송·재사용에 따른",
            "파손·변형 최소화",
            "고강도 원료 적용",
            "Top Flange 신뢰성 ↑",
        ]),
        ("② 밀폐성", [
            "DOOR GASKET 완전 밀착",
            "외부 Particle 차단",
            "내부 N2 압력 유지",
            "DOOR-REAR AL 합금",
        ]),
        ("③ ESD 방지", [
            "BODY·DOOR·BOTTOM·",
            "  O-RING·SIDE-SLOT 全 부품",
            "표면저항 ESD 처리",
            "정전기 손상 방지",
        ]),
        ("④ N2 Purge (옵션)", [
            "Purge Filter 적용",
            "공간 사전 확보",
            "3S 독자 개발 모듈",
            "저수분(DR) 환경 대응",
        ]),
    ]
    col_w = 26.3 / 4
    for i, (head, bullets) in enumerate(feature_cols):
        left = 0.6 + i * col_w
        add_rect(slide, left + 0.05, MID_Y + 0.70, col_w - 0.10, MID_H - 0.70,
                 fill=WHITE, line=MID_GRAY, line_w=0.75)
        # 헤더 박스 (상단 SOFT_GRAY 띠)
        add_rect(slide, left + 0.05, MID_Y + 0.70, col_w - 0.10, 0.65,
                 fill=SOFT_GRAY, line=LINE_GRAY, line_w=0.5)
        add_text(
            slide, left + 0.10, MID_Y + 0.70, col_w - 0.20, 0.65,
            [(head, {"size": SZ_SECTION, "bold": True, "color": BLACK,
                      "align": PP_ALIGN.CENTER})],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # 본문 bullets
        body_lines = [(" • " + b, {"size": SZ_SUB, "color": BLACK})
                      for b in bullets]
        add_text(
            slide, left + 0.15, MID_Y + 1.40, col_w - 0.30, MID_H - 1.45,
            body_lines,
        )

    # ── 하단 : FA담당 적용 컨셉 (3-column 스택) ────────────
    BOT_Y = 12.75
    BOT_H = 5.20

    add_label_band(
        slide, 0.6, BOT_Y, 26.3,
        "③ FA담당 밀폐형 큐브 물류 적용 컨셉 (Stock — 이송 — 공정)",
    )

    stack_cols = [
        ("Stock 영역 — SMC 일체형", [
            ("Cube Rack 자체가 밀폐 구조",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 분리·재투입 없이 라인 옆 운송",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 면적 +25% (908→1,149)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 공수 -33% (3,234→2,100 MD)",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 투자비 -15% (239.8→203.3억)",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("→ 2028 로드맵 표준 (DR Less)",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
        ("이송 영역 — AMR Bridge", [
            ("AMR + FOUP-like 단기 캐리어",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 일반 AMR 이 캐리어 적재 운반",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 이송 시에만 1~3h 밀폐",
             {"size": SZ_BODY, "color": BLACK}),
            (" • +N2 Purge 결합 시 0.3% DR",
             {"size": SZ_BODY, "color": BLACK}),
            ("   임시 유지 (조건부 가능)",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("→ Stock·공정 간 bridge 역할",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
        ("공정 영역 — 국소 DR", [
            ("공정 설비가 0.3% 자체 유지",
             {"size": SZ_BODY, "bold": True, "color": BLACK}),
            (" • 설비 단위 미니 환경 구축",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 전체 DR 룸 不요 → 면적 ↓",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 매몰 Loss · CAPEX 동시 절감",
             {"size": SZ_BODY, "color": BLACK}),
            (" • 2027~28 채택 검토",
             {"size": SZ_BODY, "color": BLACK}),
            ("", {"size": 4}),
            ("→ DR Less 공정 최종 목표",
             {"size": SZ_BODY, "bold": True, "color": BLUE}),
        ]),
    ]
    col_w3 = 26.3 / 3
    for i, (head, lines) in enumerate(stack_cols):
        left = 0.6 + i * col_w3
        add_rect(slide, left + 0.05, BOT_Y + 0.70, col_w3 - 0.10, BOT_H - 0.70,
                 fill=WHITE, line=MID_GRAY, line_w=0.75)
        # 헤더 띠
        add_rect(slide, left + 0.05, BOT_Y + 0.70, col_w3 - 0.10, 0.7,
                 fill=SOFT_GRAY, line=LINE_GRAY, line_w=0.5)
        add_text(
            slide, left + 0.10, BOT_Y + 0.70, col_w3 - 0.20, 0.7,
            [(head, {"size": SZ_SECTION, "bold": True, "color": BLACK,
                      "align": PP_ALIGN.CENTER})],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide, left + 0.20, BOT_Y + 1.45, col_w3 - 0.40, BOT_H - 1.55,
            lines,
        )

    # ── 풋노트 ─────────────────────────────────────────
    add_text(
        slide, 0.6, 18.30, 26.3, 0.55,
        [("* 출처 : references/기술자료/FOUP(510x515)_3S.md [L60-72, L92-108, L268-274, L286-300, L308-309]  │  "
          "references/roadmap/2026_FA기술담당_중장기로드맵_v2.md [L62-70]  │  "
          "references/26FA KPI.md [L354-376]  │  "
          "이송 1~3h·국소 DR 컨셉은 보고 작성자 정의 (사내 추정치)",
          {"size": SZ_FOOT, "color": MOLD_GREEN, "bold": True})],
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH}")


if __name__ == "__main__":
    build_slide()
